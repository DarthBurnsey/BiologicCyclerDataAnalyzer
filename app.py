import sys
from pathlib import Path

ROOT_DIR = Path(__file__).resolve().parent

from legacy_import_bootstrap import prefer_legacy_modules

prefer_legacy_modules(
    repo_root=ROOT_DIR,
    backend_root=ROOT_DIR / "Cellscope 2.0" / "backend",
    module_names=("database", "data_processing", "porosity_calculations"),
)

import streamlit as st
import pandas as pd
import io
from collections import Counter
from openpyxl import load_workbook
import matplotlib.pyplot as plt
import numpy as np
from datetime import datetime, date
import re
import sqlite3
import json
import os
import html
from io import StringIO

# Import our modular components
from database import (
    get_db_connection, init_database, migrate_database, get_project_components,
    get_user_projects, create_project, save_cell_experiment, update_cell_experiment,
    get_experiment_by_name_and_file, get_project_experiments, check_experiment_exists,
    get_experiment_data, delete_cell_experiment, delete_project, rename_project,
    rename_experiment, save_experiment, update_experiment, check_experiment_name_exists,
    get_experiment_by_name, get_all_project_experiments_data, TEST_USER_ID,
    update_project_type, get_project_by_id, duplicate_experiment,
    get_user_projects_with_counts, get_project_experiment_index, get_hydrated_experiment_payload,
    get_experiments_by_formulation_component, get_formulation_summary,
    get_experiments_grouped_by_formulation
)
from data_analysis import (
    calculate_cell_summary, calculate_experiment_average,
    calculate_cycle_life_80, get_qdis_series
)
from display_components import (
    display_experiment_summaries_table, display_individual_cells_table,
    display_best_performers_analysis
)
from draggable_tabs import (
    get_available_main_tab_labels,
    get_ordered_tab_labels,
    render_tab_settings_section,
)
from file_processing import extract_date_from_filename
from data_processing import load_and_preprocess_data, calculate_efficiency_based_on_project_type
from dialogs import confirm_delete_project, confirm_delete_experiment, show_delete_dialogs

# Initialize database
init_database()
migrate_database()
from ui_components import (
    render_hybrid_separator_input, render_autocomplete_input, render_hybrid_electrolyte_input,
    get_all_battery_materials, render_toggle_section, render_experiment_color_customization,
    render_comparison_plot_options, render_comparison_color_customization, render_comparison_name_customization,
    display_summary_stats, display_averages, render_cell_inputs, get_initial_areal_capacity,
    render_formulation_table, get_substrate_options, coerce_float_input, coerce_int_input
)
from plotting import plot_capacity_graph, plot_capacity_retention_graph, plot_comparison_capacity_graph, plot_combined_capacity_retention_graph
from llm_summary import generate_experiment_summary
from preference_components import render_preferences_sidebar, render_formulation_editor_modal, get_default_values_for_experiment, render_default_indicator
from formulation_analysis import (
    extract_formulation_component, extract_all_formulation_components,
    compare_formulations, create_formulation_comparison_dataframe,
    group_experiments_by_formulation_range, extract_formulation_component_from_experiment
)
from dashboard_analytics import (
    get_global_statistics, get_project_summaries, get_top_performers,
    get_recent_activity, get_stalled_projects, get_cells_with_cycle_data
)
from insights_engine import generate_insights
from dashboard_components import (
    render_dashboard_header, render_filter_controls, render_project_summary_grid,
    render_top_performers_table, render_insights_alerts
)
from dashboard_plots import (
    plot_multi_project_retention, plot_fade_rate_scatter,
    plot_project_comparison_bar, plot_activity_timeline
)
from interactive_plots import (
    plot_interactive_capacity, plot_interactive_retention, plot_interactive_comparison_capacity,
    plot_interactive_comparison_metrics
)
from cycler_tracking import (
    create_tracking_draft_experiment, get_tracking_dashboard_payload,
    set_tracking_status_override, summarize_tracking_rows, sync_tracking_rows_to_database
)
from batch_builder import (
    apply_batch_builder_cell_input_request,
    clear_batch_builder_cell_input_state,
    get_active_batch_builder_additional_data,
    get_active_batch_builder_template,
    queue_batch_builder_cell_inputs_request,
    render_batch_builder,
)
from batch_builder_service import build_batch_cell_inputs_template
from backend_analysis_client import get_cohort_snapshot_payload_preferred
from cohort_explorer import render_cohort_explorer
from cohort_tools import (
    build_cohort_snapshot_export_payload,
    format_cohort_snapshot_markdown,
)
from ontology_explorer import queue_lineage_explorer_focus, render_lineage_explorer
from ontology_workflow import normalize_ontology_context
from study_workspace_view import render_study_workspace

EDITOR_STATE_PREFIXES = (
    'loading_', 'active_', 'testnum_', 'formation_cycles_', 'electrolyte_', 'substrate_',
    'formulation_data_', 'formulation_saved_', 'component_dropdown_', 'component_text_',
    'component_', 'fraction_', 'add_row_', 'delete_row_', 'multi_file_upload_', 'assign_all_cells_'
)
EDITOR_STATE_SUFFIXES = (
    '_query', '_suggestions', '_selected', '_show_suggestions', '_input', '_clear'
)
EDITOR_STATE_KEYS = ('datasets', 'processed_data_cache', 'cache_key')


@st.cache_data(show_spinner=False, ttl=60)
def load_sidebar_projects(user_id):
    return get_user_projects_with_counts(user_id)


@st.cache_data(show_spinner=False, ttl=60)
def load_sidebar_experiments(project_id):
    return get_project_experiment_index(project_id)


@st.cache_data(show_spinner=False, ttl=60)
def load_sidebar_experiment_payload(experiment_id):
    return get_hydrated_experiment_payload(experiment_id)


def clear_navigation_caches():
    st.cache_data.clear()


def clear_experiment_editor_state(clear_loaded_experiment=False):
    keys_to_clear = []
    for key in list(st.session_state.keys()):
        if key.startswith(EDITOR_STATE_PREFIXES) or key.endswith(EDITOR_STATE_SUFFIXES) or key in EDITOR_STATE_KEYS:
            keys_to_clear.append(key)

    for key in keys_to_clear:
        st.session_state.pop(key, None)

    if clear_loaded_experiment:
        st.session_state.pop('loaded_experiment', None)


def set_active_project(project_id, project_name, start_new_experiment=False):
    clear_experiment_editor_state(clear_loaded_experiment=True)
    st.session_state['current_project_id'] = project_id
    st.session_state['current_project_name'] = project_name
    st.session_state['start_new_experiment'] = start_new_experiment
    if start_new_experiment:
        st.session_state['show_cell_inputs_prompt'] = True


def format_nav_date(value):
    if not value:
        return "No date"

    text_value = str(value)
    try:
        return datetime.fromisoformat(text_value.replace('Z', '+00:00')).strftime('%b %d, %Y')
    except Exception:
        pass

    for candidate in (text_value[:19], text_value[:10]):
        for fmt in ('%Y-%m-%d %H:%M:%S', '%Y-%m-%d'):
            try:
                return datetime.strptime(candidate, fmt).strftime('%b %d, %Y')
            except Exception:
                pass

    return text_value[:10]


def get_experiment_sort_date(raw_data_json, created_date):
    if raw_data_json:
        try:
            experiment_data = json.loads(raw_data_json)
            if isinstance(experiment_data, dict) and experiment_data.get('experiment_date'):
                return experiment_data['experiment_date']
        except Exception:
            pass
    return created_date or ""


def _format_ontology_summary_bits(ontology_metadata):
    ontology = normalize_ontology_context(ontology_metadata)
    if not ontology:
        return []

    batch_name = ontology.get('display_batch_name') or ontology.get('batch_name')
    root_batch_name = ontology.get('display_root_batch_name') or ontology.get('root_batch_name') or batch_name
    summary_bits = []
    if batch_name:
        summary_bits.append(f"Canonical batch: {batch_name}")
    if root_batch_name:
        summary_bits.append(f"Parent batch: {root_batch_name}")
    linked_build_count = ontology.get('linked_cell_build_count')
    if linked_build_count:
        summary_bits.append(f"{linked_build_count} linked cell build(s)")
    mapping_basis = ontology.get('mapping_basis')
    if mapping_basis == 'cell_build_edge':
        summary_bits.append("Mapped from ontology cell builds")
    elif mapping_basis == 'direct_experiment_edge':
        summary_bits.append("Mapped from ontology experiment lineage")
    elif mapping_basis == 'exact_batch_name_lookup':
        summary_bits.append("Matched by canonical batch name")
    return summary_bits


def render_ontology_context_banner(
    ontology_metadata,
    *,
    key_prefix,
    legacy_experiment_id=None,
    legacy_experiment_name=None,
):
    ontology = normalize_ontology_context(ontology_metadata)
    if not ontology:
        return

    summary_bits = _format_ontology_summary_bits(ontology)
    banner_cols = st.columns([0.78, 0.22])
    with banner_cols[0]:
        if summary_bits:
            st.caption(" | ".join(summary_bits))

    focus_button_key = f"{key_prefix}_focus_lineage"
    if banner_cols[1].button("Focus in Lineage Explorer", key=focus_button_key, use_container_width=True):
        queue_lineage_explorer_focus(
            ontology,
            legacy_experiment_id=legacy_experiment_id,
            legacy_experiment_name=legacy_experiment_name,
        )
        root_batch_name = ontology.get('display_root_batch_name') or ontology.get('root_batch_name') or ontology.get('display_batch_name') or ontology.get('batch_name')
        if root_batch_name:
            st.info(f"Lineage Explorer is focused on {root_batch_name}. Open the `🧬 Lineage Explorer` tab to inspect the graph.")
        else:
            st.info("Lineage Explorer focus updated. Open the `🧬 Lineage Explorer` tab to inspect the graph.")


def open_experiment_from_sidebar(experiment_id, project_id, project_name):
    payload = load_sidebar_experiment_payload(experiment_id)
    if not payload:
        st.error("Unable to load that experiment.")
        return

    _, experiment_name, hydrated_json = payload
    clear_experiment_editor_state(clear_loaded_experiment=False)
    st.session_state['current_project_id'] = project_id
    st.session_state['current_project_name'] = project_name
    st.session_state['loaded_experiment'] = {
        'experiment_id': experiment_id,
        'experiment_name': experiment_name,
        'project_id': project_id,
        'experiment_data': json.loads(hydrated_json)
    }
    st.session_state['_sidebar_pending_experiment_jump'] = (project_id, experiment_id)
    st.session_state['start_new_experiment'] = False


def open_tracking_row_from_dashboard(tracking_row):
    experiment_id = tracking_row.get('db_experiment_id')
    project_id = tracking_row.get('project_id')
    project_name = tracking_row.get('project_name')

    if not experiment_id:
        if not project_id or not project_name:
            st.error("Unable to determine which project should receive that tracking record.")
            return
        experiment_id = create_tracking_draft_experiment(tracking_row)
        clear_navigation_caches()

    if not experiment_id or not project_id or not project_name:
        st.error("Unable to open that tracking record in Cell Inputs.")
        return

    open_experiment_from_sidebar(experiment_id, project_id, project_name)
    st.session_state['show_cell_inputs_prompt'] = True


def open_batch_in_cell_inputs(batch_id):
    template = build_batch_cell_inputs_template(batch_id=batch_id)
    if not template:
        st.error("Unable to build a Cell Inputs template for that canonical batch.")
        return False

    project_id = template.get('project_id') or st.session_state.get('current_project_id')
    project_name = template.get('project_name') or st.session_state.get('current_project_name')
    if project_id and not project_name:
        project_payload = get_project_by_id(project_id)
        if project_payload:
            project_name = project_payload[1]

    if not project_id or not project_name:
        st.error(
            "This canonical batch does not have a default project yet. Select a project in the sidebar or add a default project in Batch Builder first."
        )
        return False

    queued_template = dict(template)
    queued_template['project_id'] = project_id
    queued_template['project_name'] = project_name
    queue_batch_builder_cell_inputs_request(queued_template)
    st.session_state['_sidebar_pending_project_selector'] = project_id
    set_active_project(project_id, project_name, start_new_experiment=True)
    return True


def render_tracking_lineage_batch_section(lineage_root_batches):
    if not lineage_root_batches:
        return

    with st.container(border=True):
        st.markdown("#### Canonical Batch Coverage")
        st.caption(
            "Root batches from the Lineage Explorer stay visible here even before a tracking sheet row exists, so parent batches like `N8` do not disappear between tabs."
        )

        lineage_frame = pd.DataFrame(
            [
                {
                    'Status': row.get('status') or '',
                    'Parent Batch': row.get('root_batch_name') or '',
                    'Default Project': row.get('default_project_name') or '',
                    'Legacy App Experiments': row.get('legacy_experiment_count', 0),
                    'Tracking Rows': row.get('tracking_row_count', 0),
                    'Active Rows': row.get('active_tracking_count', 0),
                    'Completed Rows': row.get('completed_tracking_count', 0),
                    'Study Focus': row.get('study_focus') or '',
                }
                for row in lineage_root_batches
            ]
        )
        st.dataframe(lineage_frame, use_container_width=True, hide_index=True)

        selected_root_batch_name = st.selectbox(
            "Selected canonical batch",
            options=[row.get('root_batch_name') for row in lineage_root_batches if row.get('root_batch_name')],
            key='tracking_lineage_root_select',
        )
        selected_root_batch = next(
            row for row in lineage_root_batches
            if row.get('root_batch_name') == selected_root_batch_name
        )

        selected_template = None
        selected_batch_id = selected_root_batch.get('batch_id')
        if selected_batch_id:
            selected_template = build_batch_cell_inputs_template(batch_id=int(selected_batch_id))

        preview_bits = [selected_root_batch.get('status')]
        if selected_template:
            preview_bits.extend(
                [
                    selected_template.get('project_name'),
                    selected_template.get('preferred_experiment_name'),
                    (
                        f"{len(selected_template.get('formulation') or [])} formulation component(s)"
                        if selected_template.get('formulation')
                        else "No saved formulation yet"
                    ),
                ]
            )
            defaults = selected_template.get('cell_inputs_defaults') or {}
            if defaults.get('electrolyte'):
                preview_bits.append(f"Electrolyte: {defaults['electrolyte']}")
        elif selected_root_batch.get('default_project_name'):
            preview_bits.append(selected_root_batch.get('default_project_name'))
        if selected_root_batch.get('study_focus'):
            preview_bits.append(f"Study focus: {selected_root_batch['study_focus']}")
        st.caption(" | ".join(bit for bit in preview_bits if bit))

        action_cols = st.columns(2)
        if action_cols[0].button(
            "Prepare in Cell Inputs",
            key='tracking_lineage_prepare',
            use_container_width=True,
            disabled=not selected_batch_id,
        ):
            if selected_batch_id and open_batch_in_cell_inputs(int(selected_batch_id)):
                st.rerun()

        if action_cols[1].button(
            "Focus in Lineage Explorer",
            key='tracking_lineage_focus',
            use_container_width=True,
        ):
            queue_lineage_explorer_focus(
                {
                    "root_batch_id": selected_batch_id,
                    "root_batch_name": selected_root_batch_name,
                    "display_root_batch_name": selected_root_batch_name,
                    "batch_id": selected_batch_id,
                    "batch_name": selected_root_batch_name,
                    "display_batch_name": selected_root_batch_name,
                }
            )
            st.info(f"Lineage Explorer is focused on {selected_root_batch_name}.")


def queue_cohort_comparison(project_id, project_name, experiment_names, cohort_name=None):
    canonical_names = sorted({str(name).strip() for name in experiment_names if str(name).strip()})
    if len(canonical_names) < 2:
        st.warning("A cohort needs at least two experiments from one project before it can be sent to Comparison.")
        return

    st.session_state['current_project_id'] = project_id
    st.session_state['current_project_name'] = project_name
    st.session_state['comparison_selected_experiments'] = canonical_names
    cohort_label = cohort_name or "Current cohort"
    st.session_state['comparison_prefill_notice'] = (
        f"{cohort_label} is preloaded for comparison in {project_name}."
    )


def queue_dashboard_cohort_snapshot(snapshot_id, snapshot_name=None):
    st.session_state['dashboard_cohort_snapshot_id'] = int(snapshot_id)
    if snapshot_name:
        st.session_state['dashboard_cohort_snapshot_notice'] = (
            f"{snapshot_name} is focused on the Dashboard."
        )


def queue_study_workspace_snapshot(snapshot_id, snapshot_name=None):
    st.session_state['study_workspace_pending_snapshot_id'] = int(snapshot_id)
    if snapshot_name:
        st.session_state['study_workspace_pending_snapshot_name'] = snapshot_name


def render_dashboard_cohort_snapshot_panel(queue_workspace_callback=None):
    snapshot_id = st.session_state.get('dashboard_cohort_snapshot_id')
    if not snapshot_id:
        return

    snapshot = get_cohort_snapshot_payload_preferred(snapshot_id)
    if not snapshot:
        st.session_state.pop('dashboard_cohort_snapshot_id', None)
        st.session_state.pop('dashboard_cohort_snapshot_notice', None)
        st.warning("The selected cohort snapshot is no longer available.")
        return

    notice = st.session_state.pop('dashboard_cohort_snapshot_notice', None)
    if notice:
        st.info(notice)

    summary = snapshot.get('summary') or {}
    root_rows = pd.DataFrame(snapshot.get('root_batch_summary') or [])
    export_payload = build_cohort_snapshot_export_payload(snapshot)

    with st.container(border=True):
        header_cols = st.columns([0.58, 0.14, 0.14, 0.14])
        with header_cols[0]:
            st.subheader(f"🗂️ Cohort Snapshot: {snapshot.get('name') or 'Unnamed Snapshot'}")
            meta_bits = []
            if snapshot.get('cohort_name'):
                meta_bits.append(f"Cohort: {snapshot['cohort_name']}")
            if snapshot.get('updated_date'):
                meta_bits.append(f"Updated: {snapshot['updated_date']}")
            if meta_bits:
                st.caption(" | ".join(meta_bits))
            if snapshot.get('description'):
                st.caption(snapshot['description'])
        if header_cols[1].download_button(
            "Snapshot JSON",
            data=json.dumps(export_payload, indent=2, default=str),
            file_name="dashboard_cohort_snapshot.json",
            mime="application/json",
            key="dashboard_snapshot_json_download",
            use_container_width=True,
        ):
            pass
        if header_cols[2].button(
            "Open Workspace",
            key="dashboard_snapshot_open_workspace",
            use_container_width=True,
            disabled=queue_workspace_callback is None,
        ):
            queue_workspace_callback(snapshot_id, snapshot.get('name'))
            st.info(f"{snapshot.get('name') or 'Snapshot'} is loaded into Study Workspace.")
        if header_cols[3].button("Clear Focus", key="dashboard_clear_snapshot_focus", use_container_width=True):
            st.session_state.pop('dashboard_cohort_snapshot_id', None)
            st.session_state.pop('dashboard_cohort_snapshot_notice', None)
            st.rerun()

        metric_cols = st.columns(5)
        metric_cols[0].metric("Experiments", summary.get('experiment_count', 0))
        metric_cols[1].metric("Cells", summary.get('cell_count', 0))
        metric_cols[2].metric("Metrics Ready", summary.get('metrics_ready_experiment_count', 0))
        metric_cols[3].metric(
            "Avg Retention",
            f"{summary['avg_retention_pct']:.2f}%"
            if summary.get('avg_retention_pct') is not None
            else "N/A",
        )
        metric_cols[4].metric(
            "Best Cycle Life",
            f"{summary['best_cycle_life_80']:.0f}"
            if summary.get('best_cycle_life_80') is not None
            else "N/A",
        )

        if snapshot.get('ai_summary_text'):
            with st.expander("AI-ready Cohort Brief", expanded=False):
                st.text_area(
                    "Snapshot Brief",
                    value=snapshot['ai_summary_text'],
                    height=220,
                    key="dashboard_snapshot_ai_brief",
                )
                st.download_button(
                    "Download Brief",
                    data=format_cohort_snapshot_markdown(snapshot),
                    file_name="dashboard_cohort_snapshot_brief.md",
                    mime="text/markdown",
                    key="dashboard_snapshot_brief_download",
                )

        if not root_rows.empty:
            st.markdown("#### Parent Batch Rollup")
            st.dataframe(root_rows, use_container_width=True, hide_index=True)


TRACKING_STATUS_COLORS = {
    'Completed': '#2e7d32',
    'Ongoing': '#1565c0',
    'Unknown / Missing': '#b58900',
}
TRACKING_STATUS_DISPLAY_LABELS = {
    'Active': 'Ongoing',
    'Completed': 'Completed',
    'Unknown': 'Unknown / Missing',
}
TRACKING_STATUS_FILTER_OPTIONS = ['All', 'Ongoing', 'Completed', 'Unknown / Missing']
TRACKING_FOCUS_FILTER_OPTIONS = ['Missing Cells', 'Issues']
TRACKING_STATUS_OVERRIDE_OPTIONS = [
    ('Auto (match sheet + database)', None),
    ('Ongoing', 'Active'),
    ('Completed', 'Completed'),
    ('Unknown / Missing', 'Unknown'),
]
TRACKING_STATUS_OVERRIDE_LABELS = [label for label, _ in TRACKING_STATUS_OVERRIDE_OPTIONS]
TRACKING_STATUS_OVERRIDE_VALUE_TO_LABEL = {
    value: label for label, value in TRACKING_STATUS_OVERRIDE_OPTIONS
}


def get_tracking_display_status(tracking_row):
    normalized_status = tracking_row.get('status')
    if normalized_status in TRACKING_STATUS_DISPLAY_LABELS:
        return TRACKING_STATUS_DISPLAY_LABELS[normalized_status]

    alerts = tracking_row.get('alerts', [])
    has_missing_channels = any('Missing cycler/channel assignments' in alert for alert in alerts)
    unresolved_duplicate = tracking_row.get('is_duplicate_name') and not tracking_row.get('db_experiment_id')
    unresolved_project = not tracking_row.get('project_name')

    if has_missing_channels or unresolved_duplicate or unresolved_project:
        return 'Unknown / Missing'
    if tracking_row.get('status') == 'Completed':
        return 'Completed'
    return 'Ongoing'


def get_tracking_override_label(manual_status):
    return TRACKING_STATUS_OVERRIDE_VALUE_TO_LABEL.get(
        manual_status,
        TRACKING_STATUS_OVERRIDE_OPTIONS[0][0]
    )


def style_tracking_table_row(row):
    color = TRACKING_STATUS_COLORS.get(row['Status'], '#111827')
    return [f'color: {color}; font-weight: 600;' for _ in row]


def reset_tracking_filters():
    st.session_state['tracking_search_query'] = ''
    st.session_state['tracking_status_scope'] = 'All'
    st.session_state['tracking_project_scope'] = []
    st.session_state['tracking_focus_scope'] = []


def render_tracking_filter_styles():
    st.markdown(
        """
        <style>
        .tracking-legend {
            display: flex;
            flex-wrap: wrap;
            gap: 0.55rem;
            margin: 0.35rem 0 1rem 0;
        }

        .tracking-legend-chip {
            display: inline-flex;
            align-items: center;
            gap: 0.45rem;
            padding: 0.42rem 0.8rem;
            border-radius: 999px;
            font-size: 0.84rem;
            font-weight: 600;
            border: 1px solid transparent;
        }

        .tracking-legend-dot {
            width: 0.58rem;
            height: 0.58rem;
            border-radius: 999px;
            display: inline-block;
        }

        .tracking-legend-completed {
            color: #1b5e20;
            background: rgba(46, 125, 50, 0.08);
            border-color: rgba(46, 125, 50, 0.18);
        }

        .tracking-legend-ongoing {
            color: #0d47a1;
            background: rgba(21, 101, 192, 0.08);
            border-color: rgba(21, 101, 192, 0.18);
        }

        .tracking-legend-unknown {
            color: #8a6d00;
            background: rgba(181, 137, 0, 0.12);
            border-color: rgba(181, 137, 0, 0.22);
        }

        .tracking-filter-hint {
            color: #6b7280;
            font-size: 0.83rem;
            margin-top: 0.2rem;
        }

        .tracking-active-filters {
            display: flex;
            flex-wrap: wrap;
            gap: 0.45rem;
            margin: 0.4rem 0 0.25rem 0;
        }

        .tracking-active-chip {
            display: inline-flex;
            align-items: center;
            gap: 0.35rem;
            padding: 0.35rem 0.65rem;
            border-radius: 999px;
            background: #f8fafc;
            border: 1px solid #dbe4f0;
            color: #334155;
            font-size: 0.82rem;
            font-weight: 600;
        }

        .tracking-active-label {
            color: #64748b;
            text-transform: uppercase;
            letter-spacing: 0.04em;
            font-size: 0.72rem;
            font-weight: 700;
        }

        .tracking-summary-line {
            color: #6b7280;
            font-size: 0.95rem;
            margin: 0.1rem 0 0.8rem 0;
        }
        </style>
        """,
        unsafe_allow_html=True
    )


def render_tracking_legend():
    st.markdown(
        """
        <div class="tracking-legend">
            <span class="tracking-legend-chip tracking-legend-completed">
                <span class="tracking-legend-dot" style="background:#2e7d32;"></span>
                Completed
            </span>
            <span class="tracking-legend-chip tracking-legend-ongoing">
                <span class="tracking-legend-dot" style="background:#1565c0;"></span>
                Ongoing
            </span>
            <span class="tracking-legend-chip tracking-legend-unknown">
                <span class="tracking-legend-dot" style="background:#b58900;"></span>
                Unknown / Missing
            </span>
        </div>
        """,
        unsafe_allow_html=True
    )


def render_tracking_active_filters(search_query, status_scope, selected_projects, focus_filters):
    active_filters = []
    if search_query:
        display_query = search_query if len(search_query) <= 28 else f"{search_query[:25]}..."
        active_filters.append(("Search", display_query))
    if status_scope != 'All':
        active_filters.append(("Status", status_scope))
    if selected_projects:
        if len(selected_projects) <= 2:
            project_label = ", ".join(selected_projects)
        else:
            project_label = f"{selected_projects[0]}, {selected_projects[1]} +{len(selected_projects) - 2}"
        active_filters.append(("Projects", project_label))
    for focus_filter in focus_filters:
        active_filters.append(("Focus", focus_filter))

    if not active_filters:
        st.caption("No filters applied. Showing all tracking rows.")
        return

    chips = "".join(
        (
            f'<span class="tracking-active-chip">'
            f'<span class="tracking-active-label">{html.escape(str(label or "Filter"))}</span>'
            f'{html.escape(str(value if value is not None else "Unspecified"))}</span>'
        )
        for label, value in active_filters
    )
    st.markdown(f'<div class="tracking-active-filters">{chips}</div>', unsafe_allow_html=True)


def render_cycler_tracking_tab():
    st.header("⚡ Cycler Tracking")

    tracking_payload = get_tracking_dashboard_payload()
    source_path = tracking_payload.get('source_path')
    lineage_root_batches = tracking_payload.get('lineage_root_batches') or []
    if source_path:
        st.caption(f"Tracking sheet: {source_path}")

    render_tracking_lineage_batch_section(lineage_root_batches)

    if not tracking_payload.get('available'):
        if tracking_payload.get('reason') == 'missing_file':
            st.info("The Cycler tracking comparison is ready, but the tracking sheet was not found.")
            searched_paths = tracking_payload.get('searched_paths') or []
            if searched_paths:
                st.caption("Checked: " + " | ".join(searched_paths))
        else:
            st.info("The tracking sheet was found, but it doesn't contain any rows to compare yet.")
        return

    synced_records = sync_tracking_rows_to_database(tracking_payload['rows'])
    if synced_records:
        tracking_payload = get_tracking_dashboard_payload()

    tracking_rows = [
        {
            **row,
            'display_status': get_tracking_display_status(row)
        }
        for row in tracking_payload['rows']
    ]
    tracking_summary = tracking_payload['summary']
    status_counts = Counter(row['display_status'] for row in tracking_rows)
    available_projects = sorted({
        row.get('project_name') or 'Unresolved'
        for row in tracking_rows
    })
    current_project_id = st.session_state.get('current_project_id')
    current_project_name = st.session_state.get('current_project_name')

    if st.session_state.get('tracking_status_scope') not in TRACKING_STATUS_FILTER_OPTIONS:
        st.session_state['tracking_status_scope'] = 'All'
    metric_scope_options = ['Current Project', 'All Projects'] if current_project_id else ['All Projects']
    preferred_metric_scope = 'Current Project' if current_project_id else 'All Projects'
    if st.session_state.get('tracking_metric_scope') not in metric_scope_options:
        st.session_state['tracking_metric_scope'] = preferred_metric_scope
    if not isinstance(st.session_state.get('tracking_search_query', ''), str):
        st.session_state['tracking_search_query'] = ''
    existing_project_scope = st.session_state.get('tracking_project_scope', [])
    if not isinstance(existing_project_scope, list):
        existing_project_scope = []
    st.session_state['tracking_project_scope'] = [
        project for project in existing_project_scope if project in available_projects
    ]
    existing_focus_scope = st.session_state.get('tracking_focus_scope', [])
    if not isinstance(existing_focus_scope, list):
        existing_focus_scope = []
    st.session_state['tracking_focus_scope'] = [
        item for item in existing_focus_scope if item in TRACKING_FOCUS_FILTER_OPTIONS
    ]

    metric_scope = st.segmented_control(
        "Metric Scope",
        options=metric_scope_options,
        key='tracking_metric_scope',
        help="Choose whether the top summary cards represent the current open project or every tracked project."
    )
    if metric_scope == 'Current Project' and current_project_id:
        metric_rows = [
            row for row in tracking_rows
            if row.get('project_id') == current_project_id
        ]
        metric_scope_label = current_project_name or 'Current Project'
    else:
        metric_rows = tracking_rows
        metric_scope_label = 'All Tracked Projects'
    tracking_summary = summarize_tracking_rows(metric_rows)

    st.caption(
        f"Top metrics currently represent: {metric_scope_label} ({len(metric_rows)} tracking row(s))."
    )
    tracking_metric_cols = st.columns(4)
    tracking_metric_cols[0].metric("Active Experiments", tracking_summary['active_experiments'])
    tracking_metric_cols[1].metric("Active Cells", tracking_summary['active_cells'])
    tracking_metric_cols[2].metric("Completed Experiments", tracking_summary['completed_experiments'])
    dud_value = tracking_summary['suspected_dud_cells']
    dud_delta = (
        f"{tracking_summary['defect_fraction']:.1%} dud rate"
        if tracking_summary.get('defect_fraction') is not None
        else None
    )
    tracking_metric_cols[3].metric("Suspected Dud Cells", dud_value, dud_delta)
    render_tracking_filter_styles()
    render_tracking_legend()

    with st.container(border=True):
        toolbar_cols = st.columns([0.82, 0.18])
        with toolbar_cols[0]:
            st.markdown("#### Find the rows that matter")
            st.caption("Search by experiment, project, cycler, notes, or flags, then narrow the list with one-click filters.")
        with toolbar_cols[1]:
            st.button(
                "Reset Filters",
                key='tracking_reset_filters',
                on_click=reset_tracking_filters,
                use_container_width=True
            )

        search_query = st.text_input(
            "Search rows",
            key='tracking_search_query',
            placeholder="Try T27, LIB Anodes, A6, SiCx, missing, or duplicate"
        ).strip()

        status_scope = st.segmented_control(
            "Status",
            options=TRACKING_STATUS_FILTER_OPTIONS,
            key='tracking_status_scope',
            format_func=lambda option: (
                f"All ({len(tracking_rows)})"
                if option == 'All'
                else f"{option} ({status_counts.get(option, 0)})"
            ),
            help="Switch between all rows and the main workflow states."
        ) or 'All'

        selected_tracking_projects = st.pills(
            "Projects",
            options=available_projects,
            selection_mode='multi',
            key='tracking_project_scope',
            help="Select one or more projects. Leave empty to keep every project visible."
        ) or []
        st.markdown(
            '<div class="tracking-filter-hint">Projects are optional. If nothing is selected, the table shows every project.</div>',
            unsafe_allow_html=True
        )

        focus_filters = st.pills(
            "Focus",
            options=TRACKING_FOCUS_FILTER_OPTIONS,
            selection_mode='multi',
            key='tracking_focus_scope',
            help="Use focus filters to jump straight to missing-cell mismatches or tracking anomalies."
        ) or []

    filtered_tracking_rows = []
    for row in tracking_rows:
        if status_scope != 'All' and row['display_status'] != status_scope:
            continue
        project_label = row.get('project_name') or 'Unresolved'
        if selected_tracking_projects and project_label not in selected_tracking_projects:
            continue
        if 'Missing Cells' in focus_filters and row.get('missing_cell_count', 0) <= 0:
            continue
        if 'Issues' in focus_filters and not row.get('alerts'):
            continue
        if search_query:
            search_blob = " ".join([
                row.get('experiment_name') or '',
                project_label,
                row.get('tracking_date_display') or '',
                row.get('cycler_channel_text') or '',
                row.get('loading_text') or '',
                row.get('ontology_batch_name') or '',
                row.get('ontology_root_batch_name') or '',
                row.get('notes') or '',
                " ".join(row.get('alerts', [])),
                row.get('display_status') or '',
            ]).casefold()
            if search_query.casefold() not in search_blob:
                continue
        filtered_tracking_rows.append(row)

    render_tracking_active_filters(
        search_query=search_query,
        status_scope=status_scope,
        selected_projects=selected_tracking_projects,
        focus_filters=focus_filters
    )
    st.markdown(
        f'<div class="tracking-summary-line"><strong>{len(filtered_tracking_rows)}</strong> of {len(tracking_rows)} tracking row(s) shown</div>',
        unsafe_allow_html=True
    )

    tracking_table = pd.DataFrame([
        {
            'Status': row['display_status'],
            'Experiment': row['experiment_name'],
            'Project': row.get('project_name') or 'Unresolved',
            'Canonical Batch': row.get('ontology_batch_name') or '',
            'Parent Batch': row.get('ontology_root_batch_name') or '',
            'Tracking Date': row.get('tracking_date_display') or '',
            'Tracked Cells': row.get('tracked_cell_count', 0),
            'DB Cells': row.get('database_cell_count', 0),
            'Missing Cells': row.get('missing_cell_count', 0),
            'Cyclers': row.get('cycler_channel_text') or 'Missing',
            'Notes': row.get('notes') or '',
            'Flags': ' | '.join(row.get('alerts', []))
        }
        for row in filtered_tracking_rows
    ])

    if not tracking_table.empty:
        styled_tracking_table = tracking_table.style.apply(style_tracking_table_row, axis=1)
        st.dataframe(styled_tracking_table, use_container_width=True, hide_index=True)
    else:
        st.info("No tracking rows match the current filters. Reset filters or broaden the search to bring rows back.")

    actionable_rows = [row for row in filtered_tracking_rows if row.get('can_open_in_editor')]
    if actionable_rows:
        selected_tracking_row_key = st.selectbox(
            "Selected tracking row",
            options=[row['row_key'] for row in actionable_rows],
            format_func=lambda key: next(
                (
                    f"{row['experiment_name']} | {row.get('project_name') or 'Unresolved'} | "
                    f"{row['display_status']} | {row.get('tracking_date_display') or 'No date'}"
                )
                for row in actionable_rows if row['row_key'] == key
            ),
            key='tracking_open_selector'
        )
        selected_tracking_row = next(
            row for row in actionable_rows if row['row_key'] == selected_tracking_row_key
        )
        selected_experiment_id = selected_tracking_row.get('db_experiment_id')
        current_manual_status = selected_tracking_row.get('manual_status')
        current_override_label = get_tracking_override_label(current_manual_status)

        selected_tracking_ontology = normalize_ontology_context(selected_tracking_row.get('ontology'))
        if selected_tracking_ontology:
            with st.container(border=True):
                st.markdown("#### Canonical Lineage")
                render_ontology_context_banner(
                    selected_tracking_ontology,
                    key_prefix=f"tracking_{selected_experiment_id or selected_tracking_row_key}",
                    legacy_experiment_id=selected_experiment_id,
                    legacy_experiment_name=selected_tracking_row.get('experiment_name'),
                )

        with st.container(border=True):
            st.markdown("#### Status Control")
            st.caption(
                "Use a manual status when uploaded data is still partial. `Auto` keeps the sheet/database-derived state."
            )
            status_cols = st.columns([0.7, 0.3])
            selected_override_label = status_cols[0].selectbox(
                "Manual status override",
                options=TRACKING_STATUS_OVERRIDE_LABELS,
                index=TRACKING_STATUS_OVERRIDE_LABELS.index(current_override_label),
                key=f"tracking_manual_status_{selected_experiment_id or selected_tracking_row_key}",
                disabled=not selected_experiment_id
            )
            selected_override_value = next(
                value
                for label, value in TRACKING_STATUS_OVERRIDE_OPTIONS
                if label == selected_override_label
            )
            if status_cols[1].button(
                "Save Status",
                key=f"tracking_save_status_{selected_experiment_id or selected_tracking_row_key}",
                use_container_width=True,
                disabled=not selected_experiment_id
            ):
                set_tracking_status_override(selected_experiment_id, selected_override_value)
                clear_navigation_caches()
                st.rerun()

            if selected_experiment_id:
                if current_manual_status:
                    auto_status_label = TRACKING_STATUS_DISPLAY_LABELS.get(
                        selected_tracking_row.get('auto_status'),
                        'Unknown / Missing'
                    )
                    st.caption(
                        f"Current status is manually pinned to {get_tracking_display_status(selected_tracking_row)}. "
                        f"Auto status would be {auto_status_label}."
                    )
                else:
                    st.caption(
                        f"Current status is automatic: {get_tracking_display_status(selected_tracking_row)}."
                    )
            else:
                st.caption(
                    "Create a draft in Cell Inputs first if you want to save a manual status override for this row."
                )

        action_label = (
            "Open Existing Experiment"
            if selected_experiment_id
            else "Create Draft in Cell Inputs"
        )
        if st.button(action_label, key='tracking_open_button', use_container_width=True):
            open_tracking_row_from_dashboard(selected_tracking_row)
            st.rerun()
    else:
        st.info("Rows with duplicate unresolved names or unknown projects cannot be opened automatically yet.")

# =============================
# Battery Data Gravimetric Capacity Calculator App
# =============================

st.markdown(
    """
    <style>
    .block-container {
        padding-top: 1.1rem;
        padding-bottom: 2rem;
    }

    .cellscope-app-shell {
        margin: 0 0 0.18rem 0;
    }

    .cellscope-app-shell-eyebrow {
        color: #64748b;
        font-size: 0.76rem;
        font-weight: 700;
        letter-spacing: 0.08em;
        text-transform: uppercase;
        margin-bottom: 0.25rem;
    }

    .cellscope-app-shell-title-row {
        display: flex;
        align-items: center;
        flex-wrap: wrap;
        gap: 0.7rem;
    }

    .cellscope-app-shell-title {
        margin: 0;
        color: #2b2d42;
        font-size: clamp(2.5rem, 3.2vw, 3.3rem);
        line-height: 0.95;
        font-weight: 800;
        letter-spacing: -0.04em;
    }

    .cellscope-app-shell-badge {
        display: inline-flex;
        align-items: center;
        gap: 0.45rem;
        width: 100%;
        min-height: 2.55rem;
        padding: 0.38rem 0.9rem;
        border-radius: 14px;
        background: linear-gradient(180deg, #eff6ff 0%, #dbeafe 100%);
        border: 1px solid #c7ddff;
        color: #0f4c8a;
        font-size: 0.85rem;
        font-weight: 700;
        box-sizing: border-box;
    }

    .cellscope-app-shell-badge-label {
        color: #1d4ed8;
        text-transform: uppercase;
        letter-spacing: 0.06em;
        font-size: 0.68rem;
        font-weight: 800;
        flex-shrink: 0;
    }

    .cellscope-app-shell-badge-text {
        overflow: hidden;
        text-overflow: ellipsis;
        white-space: nowrap;
    }

    .cellscope-app-shell-badge-muted {
        background: linear-gradient(180deg, #f8fafc 0%, #f1f5f9 100%);
        border-color: #dbe4f0;
        color: #475569;
    }

    .cellscope-app-shell-badge-muted .cellscope-app-shell-badge-label {
        color: #64748b;
    }

    .cellscope-app-shell-divider {
        height: 1px;
        background: linear-gradient(90deg, rgba(203, 213, 225, 0.95) 0%, rgba(226, 232, 240, 0.25) 100%);
        margin: 0.22rem 0 0.48rem 0;
    }

    div[class*="st-key-save_changes_btn"] button {
        min-height: 2.55rem !important;
        border-radius: 12px !important;
        font-weight: 700 !important;
        margin-top: 0.1rem !important;
    }

    div[data-testid="stTabs"] {
        margin-top: 0.1rem;
    }

    div[data-testid="stTabs"] [role="tablist"] {
        gap: 0.35rem;
    }

    div[data-testid="stTabs"] [role="tab"] {
        min-height: 2.35rem !important;
        padding: 0.35rem 0.9rem !important;
    }

    div[data-testid="stTabs"] [data-testid="stMarkdownContainer"] p {
        font-size: 0.9rem;
        font-weight: 600;
    }
    </style>
    """,
    unsafe_allow_html=True
)

loaded_experiment = st.session_state.get('loaded_experiment')
current_project_name = st.session_state.get('current_project_name')

header_badge_html = ""
if loaded_experiment:
    header_badge_html = (
        '<div class="cellscope-app-shell-badge">'
        '<span class="cellscope-app-shell-badge-label">Experiment</span>'
        f'<span class="cellscope-app-shell-badge-text">{html.escape(loaded_experiment["experiment_name"])}</span>'
        '</div>'
    )
elif current_project_name:
    header_badge_html = (
        '<div class="cellscope-app-shell-badge cellscope-app-shell-badge-muted">'
        '<span class="cellscope-app-shell-badge-label">Project</span>'
        f'<span class="cellscope-app-shell-badge-text">{html.escape(current_project_name)}</span>'
        '</div>'
    )

# --- Top Bar ---
with st.container():
    if loaded_experiment:
        header_title_col, header_badge_col, header_action_col = st.columns([0.58, 0.27, 0.15], gap="small")
    elif header_badge_html:
        header_title_col, header_badge_col = st.columns([0.72, 0.28], gap="small")
        header_action_col = None
    else:
        header_title_col = None
        header_badge_col = None
        header_action_col = None

    if header_title_col is None:
        st.markdown(
            f"""
            <div class="cellscope-app-shell">
                <div class="cellscope-app-shell-eyebrow">Battery Experiment Workspace</div>
                <div class="cellscope-app-shell-title-row">
                    <h1 class="cellscope-app-shell-title">CellScope</h1>
                </div>
            </div>
            """,
            unsafe_allow_html=True
        )
    else:
        with header_title_col:
            st.markdown(
                f"""
                <div class="cellscope-app-shell">
                    <div class="cellscope-app-shell-eyebrow">Battery Experiment Workspace</div>
                    <div class="cellscope-app-shell-title-row">
                        <h1 class="cellscope-app-shell-title">CellScope</h1>
                    </div>
                </div>
                """,
                unsafe_allow_html=True
            )

    if header_badge_col is not None:
        with header_badge_col:
            st.markdown(header_badge_html, unsafe_allow_html=True)

    if header_action_col is not None:
        with header_action_col:
            # Save button for loaded experiments
            if st.button("Save", key="save_changes_btn", use_container_width=True):
                # Get current experiment data
                experiment_data = loaded_experiment['experiment_data']
                experiment_id = loaded_experiment['experiment_id']
                project_id = loaded_experiment['project_id']

                # Get current values from session state or use loaded values
                current_experiment_date = st.session_state.get('current_experiment_date', experiment_data.get('experiment_date'))
                current_disc_diameter = st.session_state.get('current_disc_diameter_mm', experiment_data.get('disc_diameter_mm'))
                current_group_assignments = st.session_state.get('current_group_assignments', experiment_data.get('group_assignments'))
                current_group_names = st.session_state.get('current_group_names', experiment_data.get('group_names'))

                # Convert date string to date object if needed
                if isinstance(current_experiment_date, str):
                    try:
                        current_experiment_date = datetime.fromisoformat(current_experiment_date).date()
                    except:
                        current_experiment_date = date.today()

                # Get updated cells data from session state (includes exclude changes)
                current_datasets = st.session_state.get('datasets', [])
                pressed_thickness = st.session_state.get('pressed_thickness', experiment_data.get('pressed_thickness'))
                updated_cells_data = []
                recalculated_cells = []

                for i, dataset in enumerate(current_datasets):
                    # Get original cell data
                    original_cell = experiment_data['cells'][i] if i < len(experiment_data['cells']) else {}

                    # Read current input values from session state widgets
                    # These might be more recent than the dataset values
                    widget_loading = st.session_state.get(f'edit_loading_{i}')
                    widget_active = st.session_state.get(f'edit_active_{i}')
                    widget_formation = st.session_state.get(f'edit_formation_{i}')
                    widget_testnum = st.session_state.get(f'edit_testnum_{i}')

                    # Use widget values if available, otherwise use dataset values
                    new_loading = widget_loading if widget_loading is not None else dataset.get('loading', 0)
                    new_active = widget_active if widget_active is not None else dataset.get('active', 0)
                    new_formation = widget_formation if widget_formation is not None else dataset.get('formation_cycles', 4)
                    new_testnum = widget_testnum if widget_testnum is not None else dataset.get('testnum', f'Cell {i+1}')

                    # Check if loading or active material has changed
                    original_loading = original_cell.get('loading', 0)
                    original_active = original_cell.get('active_material', 0)

                    # Recalculate gravimetric capacities if loading or active material changed
                    updated_data_json = original_cell.get('data_json')
                    if (new_loading != original_loading or new_active != original_active) and updated_data_json:
                        try:
                            # Parse the original DataFrame
                            original_df = pd.read_json(StringIO(updated_data_json))

                            # Recalculate gravimetric capacities
                            updated_df = recalculate_gravimetric_capacities(original_df, new_loading, new_active)

                            # Update the data JSON with recalculated values
                            updated_data_json = updated_df.to_json()
                            recalculated_cells.append(new_testnum)
                        except Exception:
                            # If recalculation fails, keep original data
                            pass

                    # Recalculate porosity if loading changed and we have the required data
                    porosity = original_cell.get('porosity')
                    if (new_loading != original_loading and
                        pressed_thickness and pressed_thickness > 0 and
                        dataset.get('formulation') and
                        current_disc_diameter):
                        try:
                            from porosity_calculations import calculate_porosity_from_experiment_data
                            porosity_data = calculate_porosity_from_experiment_data(
                                disc_mass_mg=new_loading,
                                disc_diameter_mm=current_disc_diameter,
                                pressed_thickness_um=pressed_thickness,
                                formulation=dataset['formulation']
                            )
                            porosity = porosity_data['porosity']
                        except Exception:
                            pass

                    # Read other widget values too
                    widget_electrolyte = st.session_state.get(f'edit_electrolyte_{i}') or st.session_state.get(f'edit_single_electrolyte_{i}')
                    widget_substrate = st.session_state.get(f'edit_substrate_{i}') or st.session_state.get(f'edit_single_substrate_{i}')
                    widget_separator = st.session_state.get(f'edit_separator_{i}') or st.session_state.get(f'edit_single_separator_{i}')
                    new_cutoff_lower = dataset.get('cutoff_voltage_lower', original_cell.get('cutoff_voltage_lower'))
                    new_cutoff_upper = dataset.get('cutoff_voltage_upper', original_cell.get('cutoff_voltage_upper'))

                    # Convert session state dataset back to cells data format
                    updated_cell = original_cell.copy()
                    updated_cell.update({
                        'loading': new_loading,
                        'active_material': new_active,
                        'formation_cycles': new_formation,
                        'test_number': new_testnum,
                        'cell_name': new_testnum,
                        'electrolyte': widget_electrolyte if widget_electrolyte is not None else dataset.get('electrolyte', '1M LiPF6 1:1:1'),
                        'substrate': widget_substrate if widget_substrate is not None else dataset.get('substrate', 'Copper'),
                        'separator': widget_separator if widget_separator is not None else dataset.get('separator', '25um PP'),
                        'cutoff_voltage_lower': new_cutoff_lower,
                        'cutoff_voltage_upper': new_cutoff_upper,
                        'formulation': dataset.get('formulation', []),
                        'excluded': dataset.get('excluded', False),
                        'data_json': updated_data_json,
                        'porosity': porosity,
                        'file_name': original_cell.get('file_name'),
                        'cycler': dataset.get('cycler', original_cell.get('cycler')),
                        'channel': dataset.get('channel', original_cell.get('channel')),
                        'cycler_channel': dataset.get('cycler_channel', original_cell.get('cycler_channel')),
                        'tracking_placeholder': original_cell.get('tracking_placeholder', False)
                    })
                    updated_cells_data.append(updated_cell)

                # Get additional experiment data
                solids_content = st.session_state.get('solids_content', experiment_data.get('solids_content'))
                experiment_notes = st.session_state.get('experiment_notes', experiment_data.get('experiment_notes'))

                # Prepare cell format data if it's a Full Cell project
                project_type = "Full Cell"  # Default
                if project_id:
                    project_info = get_project_by_id(project_id)
                    if project_info:
                        project_type = project_info[3]

                cell_format_data = {}
                if project_type == "Full Cell":
                    cell_format = st.session_state.get('current_cell_format', experiment_data.get('cell_format', 'Coin'))
                    cell_format_data['cell_format'] = cell_format
                    if cell_format == "Pouch":
                        cell_format_data['cathode_length'] = st.session_state.get('current_cathode_length', experiment_data.get('cathode_length', 50.0))
                        cell_format_data['cathode_width'] = st.session_state.get('current_cathode_width', experiment_data.get('cathode_width', 50.0))
                        cell_format_data['num_stacked_cells'] = st.session_state.get('current_num_stacked_cells', experiment_data.get('num_stacked_cells', 1))

                # Update the experiment with current data including exclude changes
                update_experiment(
                    experiment_id=experiment_id,
                    project_id=project_id,
                    experiment_name=loaded_experiment['experiment_name'],
                    experiment_date=current_experiment_date,
                    disc_diameter_mm=current_disc_diameter,
                    group_assignments=current_group_assignments,
                    group_names=current_group_names,
                    cells_data=updated_cells_data,  # Use updated data with exclude changes
                    solids_content=solids_content,
                    pressed_thickness=pressed_thickness,
                    experiment_notes=experiment_notes,
                    cell_format_data=cell_format_data
                )
                clear_navigation_caches()

                # Update the loaded experiment in session state with all current changes
                st.session_state['loaded_experiment']['experiment_data'].update({
                    'experiment_date': current_experiment_date.isoformat(),
                    'disc_diameter_mm': current_disc_diameter,
                    'group_assignments': current_group_assignments,
                    'group_names': current_group_names,
                    'cells': updated_cells_data,
                    'solids_content': solids_content,
                    'pressed_thickness': pressed_thickness,
                    'experiment_notes': experiment_notes
                })

                # Add cell format data if applicable
                if cell_format_data:
                    st.session_state['loaded_experiment']['experiment_data'].update(cell_format_data)

                # Clear any cached processed data to force recalculation
                if 'processed_data_cache' in st.session_state:
                    del st.session_state['processed_data_cache']
                if 'cache_key' in st.session_state:
                    del st.session_state['cache_key']

                # Set flag to indicate calculations have been updated
                st.session_state['calculations_updated'] = True
                st.session_state['update_timestamp'] = datetime.now()

                st.success("Changes saved!")
                if recalculated_cells:
                    st.info(f"Recalculated specific capacity values for {len(recalculated_cells)} cell(s): {', '.join(recalculated_cells)}")
                st.rerun()

st.markdown('<div class="cellscope-app-shell-divider"></div>', unsafe_allow_html=True)
# Show delete confirmation dialogs when triggered
show_delete_dialogs()

# --- Sidebar ---
with st.sidebar:
    try:
        st.image("logo.png", width=118)
    except Exception:
        st.image("https://placehold.co/150x80?text=Logo", width=118)

    st.markdown(
        """
        <style>
        section[data-testid="stSidebar"] {
            background:
                radial-gradient(circle at top right, rgba(59, 130, 246, 0.14), transparent 34%),
                linear-gradient(180deg, #0f172a 0%, #162033 100%);
        }

        section[data-testid="stSidebar"] div[data-testid="stSidebarUserContent"] {
            padding-top: 0.45rem;
        }

        section[data-testid="stSidebar"] .block-container {
            padding-top: 0.4rem !important;
            padding-bottom: 1rem !important;
            padding-left: 0.85rem !important;
            padding-right: 0.85rem !important;
        }

        section[data-testid="stSidebar"] h1,
        section[data-testid="stSidebar"] h2,
        section[data-testid="stSidebar"] h3,
        section[data-testid="stSidebar"] h4,
        section[data-testid="stSidebar"] p,
        section[data-testid="stSidebar"] span,
        section[data-testid="stSidebar"] label,
        section[data-testid="stSidebar"] small {
            color: #e2e8f0 !important;
        }

        section[data-testid="stSidebar"] .stImage {
            margin-bottom: 0.1rem !important;
        }

        section[data-testid="stSidebar"] .element-container {
            margin-bottom: 0.25rem !important;
        }

        section[data-testid="stSidebar"] div[data-testid="stSelectbox"] > div > div,
        section[data-testid="stSidebar"] div[data-testid="stTextInput"] input,
        section[data-testid="stSidebar"] div[data-testid="stTextArea"] textarea {
            background: rgba(15, 23, 42, 0.72) !important;
            border: 1px solid rgba(148, 163, 184, 0.22) !important;
            border-radius: 10px !important;
            color: #f8fafc !important;
        }

        section[data-testid="stSidebar"] div[data-testid="stButton"] > button {
            border-radius: 10px !important;
            font-size: 0.9rem !important;
            font-weight: 500 !important;
            border: 1px solid rgba(148, 163, 184, 0.14) !important;
            transition: all 0.15s ease !important;
        }

        section[data-testid="stSidebar"] div[class*="st-key-sidebar_experiment_scroll_region"] {
            background: linear-gradient(180deg, rgba(15, 23, 42, 0.24) 0%, rgba(15, 23, 42, 0.1) 100%);
            border: 1px solid rgba(148, 163, 184, 0.14);
            border-radius: 16px;
            padding: 0.3rem 0.15rem 0.3rem 0.3rem;
            box-shadow: inset 0 1px 0 rgba(255, 255, 255, 0.03);
            scrollbar-gutter: stable both-edges;
        }

        section[data-testid="stSidebar"] div[class*="st-key-sidebar_experiment_scroll_region"] [data-testid="stVerticalBlock"] {
            gap: 0.35rem;
            padding-right: 0.45rem;
            scrollbar-gutter: stable both-edges;
        }

        section[data-testid="stSidebar"] div[class*="st-key-sidebar_experiment_scroll_region"]::-webkit-scrollbar,
        section[data-testid="stSidebar"] div[class*="st-key-sidebar_experiment_scroll_region"] *::-webkit-scrollbar {
            width: 12px;
        }

        section[data-testid="stSidebar"] div[class*="st-key-sidebar_experiment_scroll_region"]::-webkit-scrollbar-track,
        section[data-testid="stSidebar"] div[class*="st-key-sidebar_experiment_scroll_region"] *::-webkit-scrollbar-track {
            background: rgba(15, 23, 42, 0.88);
            border-radius: 999px;
        }

        section[data-testid="stSidebar"] div[class*="st-key-sidebar_experiment_scroll_region"]::-webkit-scrollbar-thumb,
        section[data-testid="stSidebar"] div[class*="st-key-sidebar_experiment_scroll_region"] *::-webkit-scrollbar-thumb {
            background: linear-gradient(180deg, #93c5fd 0%, #60a5fa 100%);
            border-radius: 999px;
            border: 2px solid rgba(15, 23, 42, 0.92);
        }

        section[data-testid="stSidebar"] div[class*="st-key-sidebar_quick_experiment_"] {
            padding-right: 0.45rem !important;
        }

        section[data-testid="stSidebar"] div[class*="st-key-sidebar_quick_experiment_"] button {
            background: rgba(248, 250, 252, 0.98) !important;
            color: #0f172a !important;
            border: 1px solid rgba(203, 213, 225, 0.9) !important;
            box-shadow: inset 0 1px 0 rgba(255, 255, 255, 0.55) !important;
            font-weight: 700 !important;
            letter-spacing: 0.01em !important;
            min-height: 3.15rem !important;
            padding: 0.7rem 0.95rem !important;
        }

        section[data-testid="stSidebar"] div[class*="st-key-sidebar_quick_experiment_"] button div[data-testid="stMarkdownContainer"] {
            width: 100% !important;
        }

        section[data-testid="stSidebar"] div[class*="st-key-sidebar_quick_experiment_"] button div[data-testid="stMarkdownContainer"] p {
            color: #0f172a !important;
            fill: #0f172a !important;
            font-weight: 700 !important;
            line-height: 1.15 !important;
            text-align: center !important;
            white-space: normal !important;
            margin: 0 !important;
        }

        section[data-testid="stSidebar"] div[class*="st-key-sidebar_quick_experiment_"] button:hover {
            background: #ffffff !important;
            border-color: rgba(148, 163, 184, 0.95) !important;
        }

        section[data-testid="stSidebar"] div[class*="st-key-sidebar_quick_experiment_"] button[kind="primary"],
        section[data-testid="stSidebar"] div[class*="st-key-sidebar_quick_experiment_"] button[data-testid="stBaseButton-primary"] {
            background: linear-gradient(135deg, #2563eb 0%, #1d4ed8 100%) !important;
            color: #ffffff !important;
            border-color: rgba(96, 165, 250, 0.45) !important;
            box-shadow: 0 10px 24px rgba(37, 99, 235, 0.22) !important;
        }

        section[data-testid="stSidebar"] div[class*="st-key-sidebar_quick_experiment_"] button[kind="primary"] *,
        section[data-testid="stSidebar"] div[class*="st-key-sidebar_quick_experiment_"] button[data-testid="stBaseButton-primary"] * {
            color: #ffffff !important;
            fill: #ffffff !important;
        }

        section[data-testid="stSidebar"] div[data-testid="stButton"] > button[kind="secondary"],
        section[data-testid="stSidebar"] div[data-testid="stButton"] > button[data-testid="stBaseButton-secondary"] {
            background: rgba(30, 41, 59, 0.72) !important;
            color: #e2e8f0 !important;
            border: 1px solid rgba(148, 163, 184, 0.2) !important;
        }

        section[data-testid="stSidebar"] div[data-testid="stButton"] > button[kind="secondary"] *,
        section[data-testid="stSidebar"] div[data-testid="stButton"] > button[data-testid="stBaseButton-secondary"] * {
            color: #e2e8f0 !important;
            fill: #e2e8f0 !important;
        }

        section[data-testid="stSidebar"] div[data-testid="stButton"] > button[kind="secondary"]:hover,
        section[data-testid="stSidebar"] div[data-testid="stButton"] > button[data-testid="stBaseButton-secondary"]:hover {
            background: rgba(51, 65, 85, 0.9) !important;
            border-color: rgba(148, 163, 184, 0.34) !important;
        }

        section[data-testid="stSidebar"] div[data-testid="stButton"] > button[kind="primary"],
        section[data-testid="stSidebar"] div[data-testid="stButton"] > button[data-testid="stBaseButton-primary"] {
            background: linear-gradient(135deg, #2563eb 0%, #1d4ed8 100%) !important;
            color: #ffffff !important;
            box-shadow: 0 10px 24px rgba(37, 99, 235, 0.22) !important;
        }

        section[data-testid="stSidebar"] div[data-testid="stButton"] > button[kind="primary"] *,
        section[data-testid="stSidebar"] div[data-testid="stButton"] > button[data-testid="stBaseButton-primary"] * {
            color: #ffffff !important;
            fill: #ffffff !important;
        }

        section[data-testid="stSidebar"] div[data-testid="stPopover"] > button {
            border-radius: 10px !important;
            min-height: 2.4rem !important;
        }

        section[data-testid="stSidebar"] div[class*="st-key-sidebar_new_experiment_compact"] button {
            min-width: 36px !important;
            max-width: 36px !important;
            width: 36px !important;
            min-height: 36px !important;
            height: 36px !important;
            padding: 0 !important;
            border-radius: 10px !important;
            font-size: 1.15rem !important;
            font-weight: 700 !important;
        }

        section[data-testid="stSidebar"] div[class*="st-key-sidebar_new_experiment_compact"] button {
            background: linear-gradient(135deg, #2563eb 0%, #1d4ed8 100%) !important;
            color: #ffffff !important;
            border-color: rgba(96, 165, 250, 0.45) !important;
            box-shadow: 0 10px 24px rgba(37, 99, 235, 0.22) !important;
        }

        section[data-testid="stSidebar"] div[class*="st-key-sidebar_new_experiment_compact"] button * {
            color: #ffffff !important;
            fill: #ffffff !important;
        }

        section[data-testid="stSidebar"] div[class*="st-key-sidebar_new_experiment_compact"] button:disabled {
            background: rgba(30, 41, 59, 0.46) !important;
            color: #64748b !important;
            box-shadow: none !important;
        }

        .sidebar-card {
            background: linear-gradient(180deg, rgba(15, 23, 42, 0.84) 0%, rgba(30, 41, 59, 0.72) 100%);
            border: 1px solid rgba(148, 163, 184, 0.18);
            border-radius: 14px;
            padding: 0.72rem 0.8rem;
            margin: 0.25rem 0 0.55rem 0;
        }

        .sidebar-card-title {
            color: #f8fafc;
            font-size: 0.98rem;
            font-weight: 700;
            line-height: 1.25;
        }

        .sidebar-card-meta {
            display: flex;
            flex-wrap: wrap;
            gap: 0.4rem;
            margin-top: 0.55rem;
        }

        .sidebar-pill {
            background: rgba(59, 130, 246, 0.14);
            color: #bfdbfe;
            border: 1px solid rgba(96, 165, 250, 0.2);
            border-radius: 999px;
            padding: 0.18rem 0.55rem;
            font-size: 0.72rem;
            font-weight: 600;
        }

        .sidebar-caption {
            color: #94a3b8;
            font-size: 0.76rem;
            margin-top: 0.35rem;
        }

        .sidebar-toolbar {
            display: flex;
            align-items: center;
            justify-content: space-between;
            margin: 0.2rem 0 0.45rem 0;
        }

        .sidebar-toolbar-title {
            color: #f8fafc;
            font-size: 1rem;
            font-weight: 700;
        }

        .sidebar-section-label {
            color: #cbd5e1;
            font-size: 0.78rem;
            font-weight: 700;
            letter-spacing: 0.04em;
            text-transform: uppercase;
            margin: 0.05rem 0 0.2rem 0;
        }

        .sidebar-scroll-hint {
            color: #dbeafe;
            font-size: 0.76rem;
            font-weight: 600;
            margin-bottom: 0.35rem;
        }
        </style>
        """,
        unsafe_allow_html=True
    )

    st.markdown('<div class="sidebar-toolbar-title">Workspace</div>', unsafe_allow_html=True)

    project_rows = load_sidebar_projects(TEST_USER_ID)
    project_lookup = {row[0]: row for row in project_rows}
    loaded_experiment = st.session_state.get('loaded_experiment')
    current_project_id = st.session_state.get('current_project_id')

    if current_project_id not in project_lookup and loaded_experiment and loaded_experiment.get('project_id') in project_lookup:
        current_project_id = loaded_experiment['project_id']
        st.session_state['current_project_id'] = current_project_id
        st.session_state['current_project_name'] = project_lookup[current_project_id][1]

    st.markdown('<div class="sidebar-section-label">Active Project</div>', unsafe_allow_html=True)
    project_selector_key = "sidebar_project_selector"
    project_options = [None] + [row[0] for row in project_rows]
    selector_value = current_project_id if current_project_id in project_lookup else None
    pending_project_id = st.session_state.pop('_sidebar_pending_project_selector', None)
    if pending_project_id is not None and pending_project_id in project_options:
        st.session_state[project_selector_key] = pending_project_id
    elif project_selector_key not in st.session_state or st.session_state.get(project_selector_key) not in project_options:
        st.session_state[project_selector_key] = selector_value

    project_control_cols = st.columns([0.72, 0.12, 0.16], gap="small")
    with project_control_cols[0]:
        selected_project_id = st.selectbox(
            "Active project",
            options=project_options,
            key=project_selector_key,
            format_func=lambda option: "Select a project" if option is None else (
                f"{project_lookup[option][1]} • {project_lookup[option][6]} experiments • {project_lookup[option][3]}"
            ),
            help="Type to search projects",
            label_visibility="collapsed",
            disabled=not project_rows
        )

    active_control_project_id = (
        selected_project_id if selected_project_id in project_lookup
        else current_project_id if current_project_id in project_lookup
        else None
    )

    with project_control_cols[1]:
        new_experiment_clicked = st.button(
            "+",
            key="sidebar_new_experiment_compact",
            use_container_width=True,
            disabled=active_control_project_id not in project_lookup,
            help="Start a new experiment in the active project"
        )

    with project_control_cols[2]:
        with st.popover("..."):
            with st.form("create_project_form"):
                new_project_name = st.text_input("Project name")
                new_project_description = st.text_area("Description")
                new_project_type = st.selectbox(
                    "Project type",
                    options=["Cathode", "Anode", "Full Cell"],
                    index=2
                )
                create_project_submit = st.form_submit_button("Create Project", use_container_width=True)
                if create_project_submit:
                    if new_project_name and new_project_name.strip():
                        project_id = create_project(TEST_USER_ID, new_project_name.strip(), new_project_description, new_project_type)
                        clear_navigation_caches()
                        st.session_state['_sidebar_pending_project_selector'] = project_id
                        st.session_state['current_project_id'] = project_id
                        st.session_state['current_project_name'] = new_project_name.strip()
                        st.session_state['start_new_experiment'] = False
                        st.rerun()
                    else:
                        st.error("Please enter a project name.")

            if active_control_project_id in project_lookup:
                st.markdown("---")
                active_project = project_lookup[active_control_project_id]
                project_id, project_name, project_desc, project_type, created_date, last_modified, experiment_count = active_project

                if st.button("Rename Project", key=f"sidebar_project_rename_{project_id}", use_container_width=True):
                    st.session_state[f'renaming_project_{project_id}'] = True

                if st.button("Change Project Type", key=f"sidebar_project_type_{project_id}", use_container_width=True):
                    st.session_state[f'changing_project_type_{project_id}'] = True

                if st.button("Delete Project", key=f"sidebar_project_delete_{project_id}", use_container_width=True, type="primary"):
                    st.session_state['confirm_delete_project'] = project_id
                    st.rerun()

    if selected_project_id != current_project_id:
        if selected_project_id is None:
            clear_experiment_editor_state(clear_loaded_experiment=True)
            st.session_state['current_project_id'] = None
            st.session_state['current_project_name'] = None
            st.session_state['start_new_experiment'] = False
        else:
            selected_project = project_lookup[selected_project_id]
            set_active_project(selected_project_id, selected_project[1], start_new_experiment=False)
        current_project_id = st.session_state.get('current_project_id')
        loaded_experiment = st.session_state.get('loaded_experiment')

    if new_experiment_clicked and current_project_id in project_lookup:
        active_project = project_lookup[current_project_id]
        set_active_project(current_project_id, active_project[1], start_new_experiment=True)
        st.rerun()

    if current_project_id in project_lookup:
            active_project = project_lookup[current_project_id]
            project_id, project_name, project_desc, project_type, created_date, last_modified, experiment_count = active_project
            st.markdown(
                f"""
                <div class="sidebar-card">
                    <div class="sidebar-card-title">{html.escape(project_name)}</div>
                    <div class="sidebar-card-meta">
                        <span class="sidebar-pill">{project_type}</span>
                        <span class="sidebar-pill">{experiment_count} experiments</span>
                    </div>
                    <div class="sidebar-caption">Updated {format_nav_date(last_modified)}</div>
                </div>
                """,
                unsafe_allow_html=True
            )

            if st.session_state.get(f'renaming_project_{project_id}', False):
                with st.form(f"rename_project_form_{project_id}"):
                    new_name = st.text_input("Project name", value=project_name)
                    rename_cols = st.columns(2)
                    rename_submit = rename_cols[0].form_submit_button("Save", use_container_width=True)
                    rename_cancel = rename_cols[1].form_submit_button("Cancel", use_container_width=True)
                    if rename_submit:
                        if new_name and new_name.strip() != project_name:
                            rename_project(project_id, new_name.strip())
                            st.session_state['current_project_name'] = new_name.strip()
                            st.session_state[f'renaming_project_{project_id}'] = False
                            clear_navigation_caches()
                            st.rerun()
                        else:
                            st.warning("Enter a different project name.")
                    if rename_cancel:
                        st.session_state[f'renaming_project_{project_id}'] = False
                        st.rerun()

            if st.session_state.get(f'changing_project_type_{project_id}', False):
                with st.form(f"change_project_type_form_{project_id}"):
                    project_type_options = ["Cathode", "Anode", "Full Cell"]
                    new_project_type = st.selectbox(
                        "Project type",
                        options=project_type_options,
                        index=project_type_options.index(project_type)
                    )
                    type_cols = st.columns(2)
                    type_submit = type_cols[0].form_submit_button("Save", use_container_width=True)
                    type_cancel = type_cols[1].form_submit_button("Cancel", use_container_width=True)
                    if type_submit:
                        if new_project_type != project_type:
                            update_project_type(project_id, new_project_type)
                            st.session_state[f'changing_project_type_{project_id}'] = False
                            clear_navigation_caches()
                            st.rerun()
                        else:
                            st.warning("Select a different project type.")
                    if type_cancel:
                        st.session_state[f'changing_project_type_{project_id}'] = False
                        st.rerun()

            st.markdown("#### Experiments")
            experiment_rows = load_sidebar_experiments(project_id)
            current_loaded_id = None
            if loaded_experiment and loaded_experiment.get('project_id') == project_id:
                current_loaded_id = loaded_experiment.get('experiment_id')

            experiment_filter = st.text_input(
                "Filter experiments",
                key=f"sidebar_experiment_filter_{project_id}",
                placeholder="Search by experiment name"
            ).strip().lower()
            sort_key = f"sidebar_experiment_sort_{project_id}"
            valid_sort_options = ["recent", "exp_date", "name_asc", "name_desc"]
            legacy_sort_value = st.session_state.get(sort_key)
            if legacy_sort_value == "name":
                st.session_state[sort_key] = "name_desc"
            elif legacy_sort_value not in valid_sort_options:
                st.session_state[sort_key] = "recent"
            sort_option = st.selectbox(
                "Sort experiments",
                options=valid_sort_options,
                key=sort_key,
                format_func=lambda option: {
                    "recent": "Recently uploaded",
                    "exp_date": "Experiment date",
                    "name_asc": "Name (A-Z)",
                    "name_desc": "Name (Z-A)"
                }[option]
            )

            experiment_items = []
            for experiment_id, experiment_name, created_date, raw_data_json in experiment_rows:
                if experiment_filter and experiment_filter not in experiment_name.lower():
                    continue
                experiment_items.append({
                    "id": experiment_id,
                    "name": experiment_name,
                    "created_date": created_date,
                    "sort_date": get_experiment_sort_date(raw_data_json, created_date)
                })

            if sort_option == "name_asc":
                experiment_items.sort(key=lambda item: item["name"].lower())
            elif sort_option == "name_desc":
                experiment_items.sort(key=lambda item: item["name"].lower(), reverse=True)
            elif sort_option == "exp_date":
                experiment_items.sort(key=lambda item: item["sort_date"] or "", reverse=True)
            else:
                experiment_items.sort(key=lambda item: item["created_date"] or "", reverse=True)

            jump_options = [None] + [item["id"] for item in experiment_items]
            jump_key = f"sidebar_experiment_jump_{project_id}"
            default_jump = current_loaded_id if current_loaded_id in jump_options else None
            pending_experiment_jump = st.session_state.get('_sidebar_pending_experiment_jump')
            if pending_experiment_jump and pending_experiment_jump[0] == project_id and pending_experiment_jump[1] in jump_options:
                st.session_state[jump_key] = pending_experiment_jump[1]
                del st.session_state['_sidebar_pending_experiment_jump']
            elif jump_key not in st.session_state or st.session_state.get(jump_key) not in jump_options:
                st.session_state[jump_key] = default_jump

            selected_experiment_id = st.selectbox(
                "Open experiment",
                options=jump_options,
                key=jump_key,
                format_func=lambda option: "Select an experiment" if option is None else (
                    next(
                        f"{item['name']} • {format_nav_date(item['sort_date'])}"
                        for item in experiment_items if item["id"] == option
                    )
                ),
                help="Type to search experiments"
            )

            if selected_experiment_id and selected_experiment_id != current_loaded_id:
                open_experiment_from_sidebar(selected_experiment_id, project_id, project_name)
                st.rerun()

            if experiment_items:
                st.caption(f"{len(experiment_items)} matching experiment(s)")
                st.markdown('<div class="sidebar-section-label">Browse Results</div>', unsafe_allow_html=True)
                if len(experiment_items) > 6:
                    st.markdown(
                        '<div class="sidebar-scroll-hint">Scroll this list to browse every matching experiment.</div>',
                        unsafe_allow_html=True
                    )
                with st.container(height=420, border=True, key="sidebar_experiment_scroll_region"):
                    for item in experiment_items:
                        button_type = "primary" if item["id"] == current_loaded_id else "secondary"
                        experiment_date = format_nav_date(item['sort_date'])
                        clicked = st.button(
                            f"**{item['name']}** · {experiment_date}",
                            key=f"sidebar_quick_experiment_{item['id']}",
                            use_container_width=True,
                            type=button_type
                        )
                        if clicked and item["id"] != current_loaded_id:
                            open_experiment_from_sidebar(item["id"], project_id, project_name)
                            st.rerun()
            else:
                st.info("No experiments match the current filter.")

            active_experiment = st.session_state.get('loaded_experiment')
            if active_experiment and active_experiment.get('project_id') == project_id:
                active_experiment_id = active_experiment.get('experiment_id')
                active_experiment_name = active_experiment.get('experiment_name', 'Selected experiment')
                st.markdown(
                    f"""
                    <div class="sidebar-card">
                        <div class="sidebar-card-title">{html.escape(active_experiment_name)}</div>
                        <div class="sidebar-caption">Active experiment</div>
                    </div>
                    """,
                    unsafe_allow_html=True
                )

                with st.popover("Active experiment actions"):
                    if st.button("Rename Experiment", key=f"sidebar_exp_rename_{active_experiment_id}", use_container_width=True):
                        st.session_state[f'renaming_experiment_{active_experiment_id}'] = True

                    if st.button("Duplicate Experiment", key=f"sidebar_exp_duplicate_{active_experiment_id}", use_container_width=True):
                        st.session_state['duplicate_experiment'] = (active_experiment_id, active_experiment_name)

                    if st.button("Delete Experiment", key=f"sidebar_exp_delete_{active_experiment_id}", use_container_width=True, type="primary"):
                        st.session_state['confirm_delete_experiment'] = (active_experiment_id, active_experiment_name)
                        st.rerun()

                if st.session_state.get(f'renaming_experiment_{active_experiment_id}', False):
                    with st.form(f"rename_experiment_form_{active_experiment_id}"):
                        new_exp_name = st.text_input("Experiment name", value=active_experiment_name)
                        exp_cols = st.columns(2)
                        exp_submit = exp_cols[0].form_submit_button("Save", use_container_width=True)
                        exp_cancel = exp_cols[1].form_submit_button("Cancel", use_container_width=True)
                        if exp_submit:
                            if new_exp_name and new_exp_name.strip() != active_experiment_name:
                                rename_experiment(active_experiment_id, new_exp_name.strip())
                                st.session_state['loaded_experiment']['experiment_name'] = new_exp_name.strip()
                                st.session_state[f'renaming_experiment_{active_experiment_id}'] = False
                                clear_navigation_caches()
                                st.rerun()
                            else:
                                st.warning("Enter a different experiment name.")
                        if exp_cancel:
                            st.session_state[f'renaming_experiment_{active_experiment_id}'] = False
                            st.rerun()
    st.markdown("---")

    if st.session_state.get('current_project_id'):
        render_preferences_sidebar(st.session_state['current_project_id'])

    st.markdown("---")
    with st.expander("Settings", expanded=False):
        render_tab_settings_section(
            get_available_main_tab_labels(bool(st.session_state.get('current_project_id'))),
            section_label="Main Tabs",
        )

# --- Popover Menu Styling ---
st.markdown(
    """
    <style>
    /* ===== POPOVER MENU STYLING ===== */
    /* Style the popover trigger button (vertical ellipsis) */
    section[data-testid="stSidebar"] div[data-testid="stPopover"] > button {
        min-width: 32px !important;
        max-width: 32px !important;
        height: 32px !important;
        padding: 0 !important;
        background: rgba(71, 85, 105, 0.3) !important;
        border: 1px solid rgba(100, 116, 139, 0.3) !important;
        border-radius: 6px !important;
        color: #94a3b8 !important;
        font-size: 1.2rem !important;
        font-weight: 600 !important;
        display: flex !important;
        align-items: center !important;
        justify-content: center !important;
        transition: all 0.15s ease !important;
    }
    
    section[data-testid="stSidebar"] div[data-testid="stPopover"] > button:hover {
        background: rgba(100, 116, 139, 0.5) !important;
        border-color: rgba(148, 163, 184, 0.5) !important;
        color: #e2e8f0 !important;
        transform: scale(1.05);
    }
    
    /* Popover content panel styling */
    div[data-testid="stPopoverBody"] {
        background: linear-gradient(180deg, #f8fafc 0%, #f1f5f9 100%) !important;
        border: 1px solid rgba(203, 213, 225, 0.8) !important;
        border-radius: 12px !important;
        box-shadow: 0 10px 40px rgba(0, 0, 0, 0.15), 0 4px 12px rgba(0, 0, 0, 0.1) !important;
        padding: 0.75rem !important;
        min-width: 200px !important;
    }
    
    div[data-testid="stPopoverBody"] p,
    div[data-testid="stPopoverBody"] span,
    div[data-testid="stPopoverBody"] strong {
        color: #1e293b !important;
    }
    
    div[data-testid="stPopoverBody"] hr {
        border-color: rgba(148, 163, 184, 0.3) !important;
        margin: 0.5rem 0 !important;
    }
    
    /* Menu item buttons inside popover */
    div[data-testid="stPopoverBody"] div[data-testid="stButton"] > button {
        background: rgba(241, 245, 249, 0.8) !important;
        border: 1px solid rgba(203, 213, 225, 0.6) !important;
        border-radius: 8px !important;
        color: #334155 !important;
        padding: 0.6rem 0.85rem !important;
        margin: 4px 0 !important;
        font-size: 0.9rem !important;
        font-weight: 500 !important;
        text-align: left !important;
        transition: all 0.15s ease !important;
    }
    
    div[data-testid="stPopoverBody"] div[data-testid="stButton"] > button:hover {
        background: rgba(226, 232, 240, 0.9) !important;
        border-color: rgba(148, 163, 184, 0.7) !important;
        transform: translateX(4px);
        color: #1e293b !important;
    }
    
    /* Delete button styling (uses primary type) */
    div[data-testid="stPopoverBody"] div[data-testid="stButton"] > button[kind="primary"] {
        background: rgba(254, 226, 226, 0.6) !important;
        border-color: rgba(252, 165, 165, 0.6) !important;
        color: #dc2626 !important;
    }
    
    div[data-testid="stPopoverBody"] div[data-testid="stButton"] > button[kind="primary"]:hover {
        background: rgba(254, 202, 202, 0.8) !important;
        border-color: rgba(248, 113, 113, 0.7) !important;
        color: #b91c1c !important;
    }
    
    /* ===== GENERAL UI POLISH ===== */
    /* Smooth animations */
    section[data-testid="stSidebar"] * {
        transition: background-color 0.15s ease, border-color 0.15s ease, color 0.15s ease !important;
    }
    
    /* Scrollbar styling for sidebar */
    section[data-testid="stSidebar"]::-webkit-scrollbar,
    section[data-testid="stSidebar"] *::-webkit-scrollbar {
        width: 10px;
        height: 10px;
    }
    
    section[data-testid="stSidebar"]::-webkit-scrollbar-track,
    section[data-testid="stSidebar"] *::-webkit-scrollbar-track {
        background: rgba(15, 23, 42, 0.78);
        border-radius: 999px;
    }
    
    section[data-testid="stSidebar"]::-webkit-scrollbar-thumb,
    section[data-testid="stSidebar"] *::-webkit-scrollbar-thumb {
        background: linear-gradient(180deg, rgba(96, 165, 250, 0.95) 0%, rgba(59, 130, 246, 0.92) 100%);
        border-radius: 999px;
        border: 2px solid rgba(15, 23, 42, 0.88);
    }
    
    section[data-testid="stSidebar"]::-webkit-scrollbar-thumb:hover,
    section[data-testid="stSidebar"] *::-webkit-scrollbar-thumb:hover {
        background: linear-gradient(180deg, rgba(147, 197, 253, 1) 0%, rgba(96, 165, 250, 0.98) 100%);
    }
    </style>
    """,
    unsafe_allow_html=True
)

# Show delete confirmation dialogs when triggered
if st.session_state.get("confirm_delete_project"):
    confirm_delete_project()
else:
    pass

if st.session_state.get("confirm_delete_experiment"):
    confirm_delete_experiment()
else:
    pass

# Handle experiment duplication
if st.session_state.get("duplicate_experiment"):
    experiment_id, experiment_name = st.session_state["duplicate_experiment"]
    try:
        new_experiment_id, new_experiment_name = duplicate_experiment(experiment_id)
        clear_navigation_caches()
        st.success(f"Successfully duplicated '{experiment_name}' as '{new_experiment_name}'! You can now upload new data to this experiment.")
        # Clear the duplication state
        del st.session_state["duplicate_experiment"]
        st.rerun()
    except Exception as e:
        st.error(f"Error duplicating experiment: {str(e)}")
        del st.session_state["duplicate_experiment"]

if 'datasets' not in st.session_state:
    st.session_state['datasets'] = []
datasets = st.session_state.get('datasets', [])
disc_diameter_mm = st.session_state.get('disc_diameter_mm', 15)
experiment_date = st.session_state.get('experiment_date', date.today())
# Ensure experiment_name is always defined
experiment_name = st.session_state.get('sidebar_experiment_name', '') or ''

# Tab selection state
if 'active_main_tab' not in st.session_state:
    st.session_state['active_main_tab'] = 0

# Remove unsupported arguments from st.tabs
# If 'show_cell_inputs_prompt' is set, show a message at the top
if 'show_cell_inputs_prompt' not in st.session_state:
    st.session_state['show_cell_inputs_prompt'] = False

if st.session_state.get('show_cell_inputs_prompt'):
    st.warning('Please click the "Cell Inputs" tab above to start your new experiment.')
    st.session_state['show_cell_inputs_prompt'] = False

# Create tabs - Dashboard is always visible, others depend on project selection
current_project_id = st.session_state.get('current_project_id')
available_main_tab_labels = get_available_main_tab_labels(bool(current_project_id))
visible_main_tab_labels = get_ordered_tab_labels(available_main_tab_labels)

main_tab_map = dict(zip(visible_main_tab_labels, st.tabs(visible_main_tab_labels)))
tab_inputs = main_tab_map["Cell Inputs"]
tab_dashboard = main_tab_map["📊 Dashboard"]
tab_tracking = main_tab_map["⚡ Cycler Tracking"]
tab_cohorts = main_tab_map["🗂️ Cohorts"]
tab_workspace = main_tab_map["🧪 Study Workspace"]
tab_lineage = main_tab_map["🧬 Lineage Explorer"]
tab_batch_builder = main_tab_map["🏗️ Batch Builder"]
tab1 = main_tab_map["Plots"]
tab2 = main_tab_map["Export"]
tab_comparison = main_tab_map.get("Comparison")
tab_master = main_tab_map.get("Master Table")

# --- Dashboard Tab ---
with tab_dashboard:
    st.header("📊 Battery Testing Dashboard")
    st.caption("The Cycler tracking comparison lives in the `⚡ Cycler Tracking` tab.")
    render_dashboard_cohort_snapshot_panel(queue_workspace_callback=queue_study_workspace_snapshot)
    if st.session_state.get('dashboard_cohort_snapshot_id'):
        st.markdown("---")
    
    # Render filter controls in sidebar
    filter_params = render_filter_controls()
    
    # Cache dashboard data for 60 seconds to improve performance
    @st.cache_data(ttl=60, show_spinner=False)
    def load_dashboard_data(user_id: str, filter_params_json: str):
        """Load all dashboard data with caching."""
        import json
        filters = json.loads(filter_params_json)
        
        # Fetch all data
        stats = get_global_statistics(user_id, filters)
        projects = get_project_summaries(user_id, filters)
        top_cells = get_top_performers(
            user_id, 
            metric='retention',
            top_n=5,
            min_cycles=filters.get('min_cycles', 100),
            filter_params=filters
        )
        activity = get_recent_activity(user_id, days=30)
        cells_data = get_cells_with_cycle_data(
            user_id,
            min_cycles=filters.get('min_cycles', 100),
            filter_params=filters
        )
        
        # Generate insights
        dashboard_data = {
            'stats': stats,
            'projects': projects,
            'top_cells': top_cells
        }
        insights = generate_insights(user_id, dashboard_data)
        
        return {
            'stats': stats,
            'projects': projects,
            'top_cells': top_cells,
            'activity': activity,
            'cells_data': cells_data,
            'insights': insights
        }
    
    # Load data with spinner
    with st.spinner("Loading dashboard data..."):
        # Convert filter_params to JSON for caching
        import json
        filter_json = json.dumps(filter_params, default=str)
        
        try:
            data = load_dashboard_data(TEST_USER_ID, filter_json)
        except Exception as e:
            st.error(f"Error loading dashboard data: {str(e)}")
            st.stop()
    
    # Header section with key metrics
    st.subheader("📈 Overview")
    render_dashboard_header(data['stats'])
    
    st.markdown("---")
    
    # Insights section
    st.subheader("💡 Actionable Insights")
    with st.expander("View Recommendations & Alerts", expanded=True):
        render_insights_alerts(data['insights'])
    
    st.markdown("---")
    
    # Two-column layout for project summary and top performers
    col1, col2 = st.columns([1, 1])
    
    with col1:
        st.subheader("📁 Project Summary")
        render_project_summary_grid(data['projects'])
    
    with col2:
        st.subheader("🏆 Top Performers")
        
        # Metric selector for top performers
        metric_options = {
            'retention': 'Capacity Retention',
            'fade_rate': 'Fade Rate',
            'efficiency': 'Coulombic Efficiency'
        }
        selected_metric = st.selectbox(
            "Sort by:",
            options=list(metric_options.keys()),
            format_func=lambda x: metric_options[x],
            key='dashboard_metric_selector'
        )
        
        # Reload top performers with selected metric
        top_cells_by_metric = get_top_performers(
            TEST_USER_ID,
            metric=selected_metric,
            top_n=5,
            min_cycles=filter_params.get('min_cycles', 100),
            filter_params=filter_params
        )
        
        render_top_performers_table(top_cells_by_metric)
    
    st.markdown("---")
    
    # Visualizations section
    st.subheader("📊 Performance Visualizations")
    
    # Tabs for different plots
    viz_tabs = st.tabs(["Retention Curves", "Fade Analysis", "Project Comparison", "Activity"])
    
    with viz_tabs[0]:
        st.markdown("#### Capacity Retention vs. Cycle Number")
        
        col1, col2 = st.columns([3, 1])
        with col2:
            show_avg = st.checkbox("Show Average", value=True, key='dashboard_show_avg')
            max_cells = st.slider(
                "Max cells per project",
                min_value=5,
                max_value=20,
                value=10,
                key='dashboard_max_cells'
            )
        
        with col1:
            if data['cells_data']:
                fig_retention = plot_multi_project_retention(
                    data['cells_data'],
                    group_by='project',
                    show_average=show_avg,
                    max_cells_per_group=max_cells
                )
                st.plotly_chart(fig_retention, use_container_width=True)
            else:
                st.info("No cycling data available. Upload experiments to see retention curves.")
    
    with viz_tabs[1]:
        st.markdown("#### Fade Rate Analysis")
        
        col1, col2 = st.columns([3, 1])
        with col2:
            x_axis_options = {
                'initial_capacity': 'Initial Capacity',
                'temperature': 'Temperature',
                'c_rate': 'C-Rate'
            }
            x_axis = st.selectbox(
                "X-axis:",
                options=list(x_axis_options.keys()),
                format_func=lambda x: x_axis_options[x],
                key='dashboard_fade_x_axis'
            )
            filter_fade_outliers = st.checkbox(
                "Filter outliers",
                value=True,
                key='dashboard_fade_filter_outliers',
                help="Removes physically implausible and statistically extreme points so trendlines remain interpretable."
            )
        
        with col1:
            if data['cells_data']:
                fig_fade = plot_fade_rate_scatter(
                    data['cells_data'],
                    x_axis=x_axis,
                    color_by='project',
                    filter_outliers=filter_fade_outliers
                )
                st.plotly_chart(fig_fade, use_container_width=True)
            else:
                st.info("No cycling data available for fade analysis.")
    
    with viz_tabs[2]:
        st.markdown("#### Project Performance Comparison")
        
        if data['projects']:
            fig_comparison = plot_project_comparison_bar(data['projects'])
            st.plotly_chart(fig_comparison, use_container_width=True)
        else:
            st.info("No project data available.")
    
    with viz_tabs[3]:
        st.markdown("#### Recent Activity Timeline")
        
        if data['activity']:
            fig_activity = plot_activity_timeline(data['activity'])
            st.plotly_chart(fig_activity, use_container_width=True)
        else:
            st.info("No activity in the last 30 days.")

with tab_tracking:
    render_cycler_tracking_tab()

with tab_cohorts:
    render_cohort_explorer(
        open_experiment_callback=open_experiment_from_sidebar,
        queue_comparison_callback=queue_cohort_comparison,
        queue_dashboard_snapshot_callback=queue_dashboard_cohort_snapshot,
        queue_workspace_callback=queue_study_workspace_snapshot,
        current_project_id=st.session_state.get('current_project_id'),
        current_project_name=st.session_state.get('current_project_name'),
    )

with tab_workspace:
    render_study_workspace(
        open_experiment_callback=open_experiment_from_sidebar,
        queue_comparison_callback=queue_cohort_comparison,
    )

with tab_lineage:
    render_lineage_explorer(open_batch_in_cell_inputs_callback=open_batch_in_cell_inputs)

with tab_batch_builder:
    render_batch_builder(
        activate_project_callback=set_active_project,
        current_project_id=st.session_state.get('current_project_id'),
        current_project_name=st.session_state.get('current_project_name'),
    )

# --- Cell Inputs Tab ---
with tab_inputs:
    # If user started a new experiment, clear cell input state
    if st.session_state.get('start_new_experiment'):
        clear_batch_builder_cell_input_state(preserve_request=True)
        project_defaults = get_default_values_for_experiment(st.session_state.get('current_project_id'))
        default_disc_diameter = project_defaults.get('disc_diameter_mm', 15.0)
        # Clear experiment-level session state
        st.session_state['datasets'] = []
        st.session_state['current_experiment_name'] = ''
        st.session_state['current_experiment_date'] = date.today()
        st.session_state['current_disc_diameter_mm'] = default_disc_diameter
        st.session_state['current_group_assignments'] = None
        st.session_state['current_group_names'] = ["Group A", "Group B", "Group C"]
        st.session_state['solids_content'] = 0.0
        st.session_state['pressed_thickness'] = 0.0
        st.session_state['experiment_notes'] = ''
        
        # Clear any remaining cell input session state variables
        keys_to_clear = []
        for key in st.session_state.keys():
            # Clear cell-specific input fields that might have been missed
            if (key.startswith('loading_') or 
                key.startswith('active_') or 
                key.startswith('testnum_') or 
                key.startswith('formation_cycles_') or
                key.startswith('cutoff_lower_') or
                key.startswith('cutoff_upper_') or
                key.startswith('electrolyte_') or
                key.startswith('substrate_') or
                key.startswith('separator_') or  # Added separator_
                key.startswith('formulation_data_') or 
                key.startswith('formulation_saved_') or
                key.startswith('component_dropdown_') or
                key.startswith('component_text_') or
                key.startswith('fraction_') or
                key.startswith('add_row_') or
                key.startswith('delete_row_') or
                key.startswith('multi_file_upload_') or
                key.startswith('assign_all_cells_') or
                key.startswith('use_same_formulation_') or  # Added this
                key == 'datasets' or
                key == 'processed_data_cache' or
                key == 'cache_key'):
                keys_to_clear.append(key)
        
        # Remove the keys
        for key in keys_to_clear:
            del st.session_state[key]
        
        # Reset experiment editor hydration marker and top-level input widgets.
        st.session_state.pop('cell_inputs_loaded_experiment_id', None)
        st.session_state.pop('main_experiment_name', None)
        st.session_state.pop('current_experiment_date', None)
        
        st.session_state['start_new_experiment'] = False
    st.header("Cell Inputs & Experiment Setup")
    st.markdown("---")
    
    # Check if we have a loaded experiment
    loaded_experiment = st.session_state.get('loaded_experiment')
    
    # If a new experiment is being started, always allow editing
    is_new_experiment = not loaded_experiment and st.session_state.get('current_project_id')
    
    if loaded_experiment:
        experiment_data = loaded_experiment['experiment_data']
        cells_data = experiment_data.get('cells', [])
        loaded_experiment_id = loaded_experiment.get('experiment_id')
        initialized_experiment_id = st.session_state.get('cell_inputs_loaded_experiment_id')
        
        # Show different message for experiments with no cells (e.g., duplicates)
        if len(cells_data) == 0:
            st.info(f"Setting up experiment: **{loaded_experiment['experiment_name']}** (ready for data upload)")
        else:
            st.info(f"Editing experiment: **{loaded_experiment['experiment_name']}**")

        tracking_metadata = experiment_data.get('tracking')
        if tracking_metadata:
            tracking_bits = [
                f"Tracked on {tracking_metadata.get('tracking_date')}" if tracking_metadata.get('tracking_date') else None,
                tracking_metadata.get('cycler_channel_text'),
                f"Expected {tracking_metadata.get('tracked_cell_count')} cell(s)" if tracking_metadata.get('tracked_cell_count') else None,
                (
                    f"Manual status: {TRACKING_STATUS_DISPLAY_LABELS.get(tracking_metadata.get('manual_status'), tracking_metadata.get('manual_status'))}"
                    if tracking_metadata.get('manual_status')
                    else None
                ),
            ]
            tracking_text = " | ".join(bit for bit in tracking_bits if bit)
            if tracking_text:
                st.caption(f"Tracking metadata: {tracking_text}")

        ontology_metadata = normalize_ontology_context(experiment_data.get('ontology'))
        if ontology_metadata:
            with st.container(border=True):
                st.markdown("#### Canonical Lineage")
                render_ontology_context_banner(
                    ontology_metadata,
                    key_prefix=f"cell_inputs_{loaded_experiment_id}",
                    legacy_experiment_id=loaded_experiment_id,
                    legacy_experiment_name=loaded_experiment.get('experiment_name'),
                )

        # Only hydrate editor state when loading a different experiment.
        if initialized_experiment_id != loaded_experiment_id:
            parsed_experiment_date = experiment_data.get('experiment_date')
            if isinstance(parsed_experiment_date, str):
                try:
                    parsed_experiment_date = datetime.fromisoformat(parsed_experiment_date).date()
                except:
                    parsed_experiment_date = date.today()
            elif parsed_experiment_date is None:
                parsed_experiment_date = date.today()
            
            st.session_state['current_experiment_name'] = loaded_experiment['experiment_name']
            st.session_state['main_experiment_name'] = loaded_experiment['experiment_name']
            st.session_state['current_experiment_date'] = parsed_experiment_date
            st.session_state['current_disc_diameter_mm'] = experiment_data.get('disc_diameter_mm', 15)
            st.session_state['current_group_assignments'] = experiment_data.get('group_assignments')
            st.session_state['current_group_names'] = experiment_data.get('group_names', ["Group A", "Group B", "Group C"])
            st.session_state['solids_content'] = experiment_data.get('solids_content', 0.0)
            st.session_state['pressed_thickness'] = experiment_data.get('pressed_thickness', 0.0)
            st.session_state['experiment_notes'] = experiment_data.get('experiment_notes', '')
            
            # --- Load cell format data from experiment ---
            st.session_state['current_cell_format'] = experiment_data.get('cell_format', 'Coin')
            if experiment_data.get('cathode_length'):
                st.session_state['current_cathode_length'] = experiment_data.get('cathode_length', 50.0)
            if experiment_data.get('cathode_width'):
                st.session_state['current_cathode_width'] = experiment_data.get('cathode_width', 50.0)
            if experiment_data.get('num_stacked_cells'):
                st.session_state['current_num_stacked_cells'] = experiment_data.get('num_stacked_cells', 1)
            
            # Convert loaded cells data back to datasets format for editing
            loaded_datasets = []
            for cell_data in cells_data:
                has_data = bool(cell_data.get('data_json') or cell_data.get('parquet_path'))
                file_name = cell_data.get('file_name')
                mock_file = None
                if file_name and has_data:
                    mock_file = type('MockFile', (), {
                        'name': file_name,
                        'type': 'text/csv'
                    })()

                loaded_datasets.append({
                    'file': mock_file,
                    'loading': cell_data.get('loading', 20.0),
                    'active': cell_data.get('active_material', 90.0),
                    'testnum': cell_data.get('test_number', cell_data.get('cell_name', '')),
                    'formation_cycles': cell_data.get('formation_cycles', 4),
                    'cutoff_voltage_lower': cell_data.get('cutoff_voltage_lower'),
                    'cutoff_voltage_upper': cell_data.get('cutoff_voltage_upper'),
                    'electrolyte': cell_data.get('electrolyte', '1M LiPF6 1:1:1'),
                    'substrate': cell_data.get('substrate', 'Copper'),
                    'separator': cell_data.get('separator', '25um PP'),
                    'formulation': cell_data.get('formulation', []),
                    'excluded': cell_data.get('excluded', False),
                    'cycler': cell_data.get('cycler'),
                    'channel': cell_data.get('channel'),
                    'cycler_channel': cell_data.get('cycler_channel'),
                    'tracking_placeholder': cell_data.get('tracking_placeholder', False),
                    'uploaded_file_source': False,
                    'has_data': has_data,
                    'file_label': file_name or 'No raw file attached yet'
                })
            
            st.session_state['datasets'] = loaded_datasets
            st.session_state['cell_inputs_loaded_experiment_id'] = loaded_experiment_id
        
        current_experiment_name = st.session_state.get('current_experiment_name', loaded_experiment['experiment_name'])
        current_experiment_date = st.session_state.get('current_experiment_date', date.today())
        current_disc_diameter = st.session_state.get('current_disc_diameter_mm', experiment_data.get('disc_diameter_mm', 15))
        current_group_assignments = st.session_state.get('current_group_assignments')
        current_group_names = st.session_state.get('current_group_names', ["Group A", "Group B", "Group C"])
        datasets = st.session_state.get('datasets', [])
        
    elif is_new_experiment:
        st.info(f"Creating a new experiment in project: **{st.session_state['current_project_name']}**")
        project_defaults = get_default_values_for_experiment(st.session_state.get('current_project_id'))
        st.session_state['cell_inputs_loaded_experiment_id'] = None
        
        if 'current_experiment_name' not in st.session_state:
            st.session_state['current_experiment_name'] = ""
        if 'main_experiment_name' not in st.session_state:
            st.session_state['main_experiment_name'] = ""
        if 'current_experiment_date' not in st.session_state:
            st.session_state['current_experiment_date'] = date.today()
        if 'current_disc_diameter_mm' not in st.session_state:
            st.session_state['current_disc_diameter_mm'] = project_defaults.get('disc_diameter_mm', 15.0)
        if 'current_group_assignments' not in st.session_state:
            st.session_state['current_group_assignments'] = None
        if 'current_group_names' not in st.session_state:
            st.session_state['current_group_names'] = ["Group A", "Group B", "Group C"]
        if 'solids_content' not in st.session_state:
            st.session_state['solids_content'] = 0.0
        if 'pressed_thickness' not in st.session_state:
            st.session_state['pressed_thickness'] = 0.0
        if 'experiment_notes' not in st.session_state:
            st.session_state['experiment_notes'] = ''

        apply_batch_builder_cell_input_request()
        
        current_experiment_name = st.session_state.get('current_experiment_name', "")
        current_experiment_date = st.session_state.get('current_experiment_date', date.today())
        current_disc_diameter = st.session_state.get('current_disc_diameter_mm', project_defaults.get('disc_diameter_mm', 15.0))
        current_group_assignments = st.session_state.get('current_group_assignments')
        current_group_names = st.session_state.get('current_group_names', ["Group A", "Group B", "Group C"])
    else:
        st.info("Create a new experiment or load an existing one from the sidebar")
        st.session_state['cell_inputs_loaded_experiment_id'] = None
        current_experiment_name = st.session_state.get('current_experiment_name', "")
        current_experiment_date = st.session_state.get('current_experiment_date', date.today())
        current_disc_diameter = st.session_state.get('current_disc_diameter_mm', 15)
        current_group_assignments = st.session_state.get('current_group_assignments')
        current_group_names = st.session_state.get('current_group_names', ["Group A", "Group B", "Group C"])
        if 'solids_content' not in st.session_state:
            st.session_state['solids_content'] = 0.0
        if 'pressed_thickness' not in st.session_state:
            st.session_state['pressed_thickness'] = 0.0
        if 'experiment_notes' not in st.session_state:
            st.session_state['experiment_notes'] = ''

    active_batch_builder_template = get_active_batch_builder_template()
    if is_new_experiment and active_batch_builder_template:
        with st.container(border=True):
            st.markdown("#### Canonical Batch Starter")
            st.caption("This experiment will save with canonical ontology context from the Batch Builder.")
            render_ontology_context_banner(
                active_batch_builder_template.get('ontology_context'),
                key_prefix="cell_inputs_batch_builder_template",
                legacy_experiment_name=active_batch_builder_template.get('preferred_experiment_name'),
            )
            batch_bits = [
                active_batch_builder_template.get('project_name'),
                active_batch_builder_template.get('ingestion_hints', {}).get('cell_build_naming'),
            ]
            if any(batch_bits):
                st.caption(" | ".join(bit for bit in batch_bits if bit))
    
    # Experiment metadata inputs
    col1, col2 = st.columns(2)
    with col1:
        experiment_name_input = st.text_input(
            'Experiment Name', 
            value=current_experiment_name if loaded_experiment else experiment_name,
            placeholder='Enter experiment name for file naming',
            key="main_experiment_name"
        )
    
    with col2:
        experiment_date_input = st.date_input(
            "Experiment Date", 
            value=current_experiment_date,
            key="current_experiment_date",
            help="Date associated with this experiment"
        )
    
    # Get current project type to determine input fields
    current_project_id = st.session_state.get('current_project_id')
    project_type = "Full Cell"  # Default
    if current_project_id:
        project_info = get_project_by_id(current_project_id)
        if project_info:
            project_type = project_info[3]  # project_type is the 4th field
    
    # Enhanced Full Cell format selection
    if project_type == "Full Cell":
        st.markdown("#### Cell Configuration")
        
        # Cell format dropdown
        format_options = ["Coin", "Pouch", "Cylindrical", "Prismatic"]
        current_format = st.session_state.get('current_cell_format', 'Coin')
        
        cell_format = st.selectbox(
            'Cell Format',
            options=format_options,
            index=format_options.index(current_format) if current_format in format_options else 0,
            key='cell_format_input',
            help="Select the physical format of the battery cell"
        )
        
        # Store format in session state
        st.session_state['current_cell_format'] = cell_format
        
        if cell_format == "Coin":
            # Show traditional disc diameter input
            disc_diameter_input = st.number_input(
                'Disc Diameter (mm) for Areal Capacity Calculation', 
                min_value=1.0,
                max_value=50.0,
                value=coerce_float_input(current_disc_diameter, 15.0),
                step=1.0,
                help="Diameter of the coin cell disc for areal capacity calculations"
            )
            
        elif cell_format == "Pouch":
            # Show pouch-specific inputs
            st.markdown("##### Pouch Cell Dimensions")
            col1, col2 = st.columns(2)
            
            with col1:
                cathode_length = st.number_input(
                    'Cathode Length (mm)',
                    min_value=1.0,
                    max_value=500.0,
                    value=coerce_float_input(st.session_state.get('current_cathode_length', 50.0), 50.0),
                    step=0.1,
                    key='cathode_length_input',
                    help="Length of the cathode active area"
                )
                
                cathode_width = st.number_input(
                    'Cathode Width (mm)',
                    min_value=1.0,
                    max_value=500.0,
                    value=coerce_float_input(st.session_state.get('current_cathode_width', 50.0), 50.0),
                    step=0.1,
                    key='cathode_width_input',
                    help="Width of the cathode active area"
                )
            
            with col2:
                num_stacked_cells = st.number_input(
                    'Number of Stacked Cells',
                    min_value=1,
                    max_value=100,
                    value=coerce_int_input(st.session_state.get('current_num_stacked_cells', 1), 1),
                    step=1,
                    key='num_stacked_cells_input',
                    help="Number of cells stacked in the pouch configuration"
                )
                
                # Calculate and display total area
                total_area_cm2 = (cathode_length * cathode_width * num_stacked_cells) / 100  # Convert mm² to cm²
                st.metric(
                    label="Total Active Area", 
                    value=f"{total_area_cm2:.2f} cm²",
                    help="Total cathode active area for all stacked cells"
                )
            
            # Store pouch values in session state
            st.session_state['current_cathode_length'] = cathode_length
            st.session_state['current_cathode_width'] = cathode_width
            st.session_state['current_num_stacked_cells'] = num_stacked_cells
            
            # Set disc_diameter_input to None since we're using area calculations instead
            disc_diameter_input = None
            
        elif cell_format in ["Cylindrical", "Prismatic"]:
            # Under construction message
            st.info(f"🚧 **{cell_format} Cell Format - Under Construction**")
            st.markdown(
                f"The {cell_format.lower()} cell format configuration is currently being developed. "
                "Please select 'Coin' or 'Pouch' format for now, or check back in a future update!"
            )
            # Use default disc diameter for now
            disc_diameter_input = current_disc_diameter
            
    else:
        # For non-Full Cell projects (Cathode/Anode), show traditional disc diameter input
        disc_diameter_input = st.number_input(
            'Disc Diameter (mm) for Areal Capacity Calculation', 
            min_value=1.0,
            max_value=50.0,
            value=coerce_float_input(current_disc_diameter, 15.0),
            step=1.0,
            help="Diameter of the electrode disc for areal capacity calculations"
        )
    
    st.markdown("---")
    
    # Cell inputs section
    # Show file upload for new experiments OR loaded experiments with no cells (e.g., duplicates)
    has_cells = loaded_experiment and len(datasets) > 0
    
    if loaded_experiment and has_cells:
        # Pre-compute shared options outside cell loops
        substrate_options = get_substrate_options()
        
        # For loaded experiments, show the cell input fields for editing
        if len(datasets) > 1:
            with st.expander(f'Cell 1: {datasets[0]["testnum"] or "Cell 1"}', expanded=False):
                col1, col2 = st.columns(2)
                with col1:
                    loading_0 = st.number_input(
                        f'Disc loading (mg) for Cell 1', 
                        min_value=0.0, 
                        step=1.0, 
                        value=float(datasets[0]["loading"]),
                        key=f'edit_loading_0'
                    )
                    formation_cycles_0 = st.number_input(
                        f'Formation Cycles for Cell 1', 
                        min_value=0, 
                        step=1, 
                        value=int(datasets[0]["formation_cycles"]),
                        key=f'edit_formation_0'
                    )
                with col2:
                    active_material_0 = st.number_input(
                        f'% Active material for Cell 1', 
                        min_value=0.0, 
                        max_value=100.0, 
                        step=1.0, 
                        value=float(datasets[0]["active"]),
                        key=f'edit_active_0'
                    )
                    test_number_0 = st.text_input(
                        f'Test Number for Cell 1', 
                        value=datasets[0]["testnum"] or "Cell 1",
                        key=f'edit_testnum_0'
                    )
                
                # Electrolyte, Substrate, and Separator selection
                substrate_options = get_substrate_options()
                
                col3, col4, col5 = st.columns(3)
                with col3:
                    from ui_components import render_hybrid_electrolyte_input
                    electrolyte_0 = render_hybrid_electrolyte_input(
                        f'Electrolyte for Cell 1', 
                        default_value=datasets[0]["electrolyte"],
                        key=f'edit_electrolyte_0'
                    )
                with col4:
                    substrate_0 = st.selectbox(
                        f'Substrate for Cell 1', 
                        substrate_options,
                        index=substrate_options.index(datasets[0].get("substrate", "Copper")) if datasets[0].get("substrate") in substrate_options else 0,
                        key=f'edit_substrate_0'
                    )
                with col5:
                    from ui_components import render_hybrid_separator_input
                    separator_0 = render_hybrid_separator_input(
                        f'Separator for Cell 1', 
                        default_value=datasets[0].get("separator", "25um PP"),
                        key=f'edit_separator_0'
                    )
                
                cutoff_lower_default_0 = datasets[0].get("cutoff_voltage_lower")
                cutoff_upper_default_0 = datasets[0].get("cutoff_voltage_upper")
                if cutoff_lower_default_0 is None:
                    cutoff_lower_default_0 = 2.5
                if cutoff_upper_default_0 is None:
                    cutoff_upper_default_0 = 4.2
                
                cutoff_col1, cutoff_col2 = st.columns(2)
                with cutoff_col1:
                    cutoff_lower_0 = st.number_input(
                        "Lower Cutoff Voltage (V) for Cell 1",
                        min_value=0.0,
                        max_value=10.0,
                        step=0.1,
                        value=float(cutoff_lower_default_0),
                        key="edit_cutoff_lower_0",
                    )
                with cutoff_col2:
                    cutoff_upper_0 = st.number_input(
                        "Upper Cutoff Voltage (V) for Cell 1",
                        min_value=0.0,
                        max_value=10.0,
                        step=0.1,
                        value=float(cutoff_upper_default_0),
                        key="edit_cutoff_upper_0",
                    )
                
                # Formulation table
                st.markdown("**Formulation:**")
                from ui_components import render_formulation_table
                # Initialize formulation data if needed
                formulation_key = f'formulation_data_edit_0_loaded'
                if formulation_key not in st.session_state:
                    st.session_state[formulation_key] = datasets[0]["formulation"] if datasets[0]["formulation"] else [{'Component': '', 'Dry Mass Fraction (%)': 0.0}]
                formulation_0 = render_formulation_table(f'edit_0_loaded', project_id, get_project_components)
                
                # Add two buttons: Exclude and Remove
                col_btn1, col_btn2 = st.columns(2)
                
                with col_btn1:
                    exclude_button_disabled = datasets[0].get('excluded', False)
                    if st.button("🚫 Exclude Cell", key=f'exclude_cell_loaded_0', disabled=exclude_button_disabled):
                        datasets[0]['excluded'] = True
                        st.session_state['datasets'] = datasets  # Save to session state
                
                with col_btn2:
                    if st.button("Remove Cell", key=f'remove_cell_loaded_0'):
                        st.session_state[f'confirm_remove_cell_loaded_0'] = True

                # Show excluded status
                if datasets[0].get('excluded', False):
                    st.warning("This cell is excluded from analysis")
                    if st.button("✅ Include Cell", key=f'include_cell_loaded_0'):
                        datasets[0]['excluded'] = False
                        st.session_state['datasets'] = datasets  # Save to session state

                if st.session_state.get(f'confirm_remove_cell_loaded_0', False):
                    st.error("**PERMANENT DELETION** - This will permanently delete the cell data and cannot be undone!")
                    col_confirm1, col_confirm2 = st.columns(2)
                    with col_confirm1:
                        if st.button("Delete Permanently", key=f'confirm_yes_loaded_0', type="primary"):
                            # Actually remove the cell from the datasets
                            datasets.pop(0)
                            st.session_state[f'confirm_remove_cell_loaded_0'] = False
                    with col_confirm2:
                        if st.button("Cancel", key=f'confirm_no_loaded_0'):
                            st.session_state[f'confirm_remove_cell_loaded_0'] = False

                assign_all = st.checkbox('Assign values to all cells', key='assign_all_cells_loaded')
            # Update all datasets with new values
            edited_datasets = []
            for i, dataset in enumerate(datasets):
                # Don't skip excluded cells - still render them but with visual indication
                
                if i == 0:
                    # First cell: preserve original file object and update other fields
                    edited_dataset = {
                        'file': dataset['file'],  # Always preserve original file object
                        'loading': loading_0,
                        'active': active_material_0,
                        'testnum': test_number_0,
                        'formation_cycles': formation_cycles_0,
                        'cutoff_voltage_lower': cutoff_lower_0,
                        'cutoff_voltage_upper': cutoff_upper_0,
                        'electrolyte': electrolyte_0,
                        'substrate': substrate_0,
                        'separator': separator_0,
                        'formulation': formulation_0,
                        'excluded': dataset.get('excluded', False)  # Add this line
                    }
                else:
                    is_excluded = dataset.get('excluded', False)
                    cell_title = f'Cell {i+1}: {dataset["testnum"] or f"Cell {i+1}"}'
                    if is_excluded:
                        cell_title += " ⚠️ EXCLUDED"
                    
                    with st.expander(cell_title, expanded=False):
                        col1, col2 = st.columns(2)
                        if assign_all:
                            loading = loading_0
                            formation_cycles = formation_cycles_0
                            active_material = active_material_0
                            electrolyte = electrolyte_0
                            substrate = substrate_0
                            separator = separator_0
                            cutoff_lower = cutoff_lower_0
                            cutoff_upper = cutoff_upper_0
                            formulation = formulation_0
                            # Test number should remain individual (not assigned to all)
                            test_number = dataset['testnum'] or f'Cell {i+1}'
                        else:
                            with col1:
                                loading = st.number_input(
                                    f'Disc loading (mg) for Cell {i+1}', 
                                    min_value=0.0, 
                                    step=1.0, 
                                    value=float(dataset['loading']),
                                    key=f'edit_loading_{i}'
                                )
                                formation_cycles = st.number_input(
                                    f'Formation Cycles for Cell {i+1}', 
                                    min_value=0, 
                                    step=1, 
                                    value=int(dataset['formation_cycles']),
                                    key=f'edit_formation_{i}'
                                )
                            with col2:
                                active_material = st.number_input(
                                    f'% Active material for Cell {i+1}', 
                                    min_value=0.0, 
                                    max_value=100.0, 
                                    step=1.0, 
                                    value=float(dataset['active']),
                                    key=f'edit_active_{i}'
                                )
                                test_number = st.text_input(
                                    f'Test Number for Cell {i+1}', 
                                    value=dataset['testnum'] or f'Cell {i+1}',
                                    key=f'edit_testnum_{i}'
                                )
                            
                            # Electrolyte, Substrate, and Separator selection
                            substrate_options = get_substrate_options()
                            
                            col3, col4, col5 = st.columns(3)
                            with col3:
                                electrolyte = render_hybrid_electrolyte_input(
                                    f'Electrolyte for Cell {i+1}', 
                                    default_value=dataset['electrolyte'],
                                    key=f'edit_electrolyte_{i}'
                                )
                            with col4:
                                substrate = st.selectbox(
                                    f'Substrate for Cell {i+1}', 
                                    substrate_options,
                                    index=substrate_options.index(dataset.get('substrate', 'Copper')) if dataset.get('substrate') in substrate_options else 0,
                                    key=f'edit_substrate_{i}'
                                )
                            with col5:
                                separator = render_hybrid_separator_input(
                                    f'Separator for Cell {i+1}', 
                                    default_value=dataset.get('separator', '25um PP'),
                                    key=f'edit_separator_{i}'
                                )
                            
                            cutoff_lower_default = dataset.get('cutoff_voltage_lower')
                            cutoff_upper_default = dataset.get('cutoff_voltage_upper')
                            if cutoff_lower_default is None:
                                cutoff_lower_default = cutoff_lower_0
                            if cutoff_upper_default is None:
                                cutoff_upper_default = cutoff_upper_0
                            
                            cutoff_col1, cutoff_col2 = st.columns(2)
                            with cutoff_col1:
                                cutoff_lower = st.number_input(
                                    f'Lower Cutoff Voltage (V) for Cell {i+1}',
                                    min_value=0.0,
                                    max_value=10.0,
                                    step=0.1,
                                    value=float(cutoff_lower_default),
                                    key=f'edit_cutoff_lower_{i}',
                                )
                            with cutoff_col2:
                                cutoff_upper = st.number_input(
                                    f'Upper Cutoff Voltage (V) for Cell {i+1}',
                                    min_value=0.0,
                                    max_value=10.0,
                                    step=0.1,
                                    value=float(cutoff_upper_default),
                                    key=f'edit_cutoff_upper_{i}',
                                )
                            
                            # Formulation table
                            st.markdown("**Formulation:**")
                            from ui_components import render_formulation_table
                            # Initialize formulation data if needed
                            formulation_key = f'formulation_data_edit_{i}_loaded'
                            if formulation_key not in st.session_state:
                                st.session_state[formulation_key] = dataset['formulation'] if dataset['formulation'] else [{'Component': '', 'Dry Mass Fraction (%)': 0.0}]
                            formulation = render_formulation_table(f'edit_{i}_loaded', project_id, get_project_components)
                            
                            # Add two buttons: Exclude and Remove
                            col_btn1, col_btn2 = st.columns(2)
                            
                            with col_btn1:
                                exclude_button_disabled = dataset.get('excluded', False)
                                if st.button("🚫 Exclude Cell", key=f'exclude_cell_loaded_{i}', disabled=exclude_button_disabled):
                                    datasets[i]['excluded'] = True
                                    st.session_state['datasets'] = datasets  # Save to session state
                            
                            with col_btn2:
                                if st.button("Remove Cell", key=f'remove_cell_loaded_{i}'):
                                    st.session_state[f'confirm_remove_cell_loaded_{i}'] = True

                            # Show excluded status
                            if dataset.get('excluded', False):
                                st.warning("This cell is excluded from analysis")
                                if st.button("✅ Include Cell", key=f'include_cell_loaded_{i}'):
                                    datasets[i]['excluded'] = False
                                    st.session_state['datasets'] = datasets  # Save to session state

                            if st.session_state.get(f'confirm_remove_cell_loaded_{i}', False):
                                st.error("**PERMANENT DELETION** - This will permanently delete the cell data and cannot be undone!")
                                col_confirm1, col_confirm2 = st.columns(2)
                                with col_confirm1:
                                    if st.button("Delete Permanently", key=f'confirm_yes_loaded_{i}', type="primary"):
                                        # Actually remove the cell from the datasets
                                        datasets.pop(i)
                                        st.session_state[f'confirm_remove_cell_loaded_{i}'] = False
                                with col_confirm2:
                                    if st.button("Cancel", key=f'confirm_no_loaded_{i}'):
                                        st.session_state[f'confirm_remove_cell_loaded_{i}'] = False
                        
                        # Always preserve original file object, only update other fields
                        edited_dataset = {
                            'file': dataset['file'],  # Always preserve original file object
                            'loading': loading,
                            'active': active_material,
                            'testnum': test_number,
                            'formation_cycles': formation_cycles,
                            'cutoff_voltage_lower': cutoff_lower,
                            'cutoff_voltage_upper': cutoff_upper,
                            'electrolyte': electrolyte,
                            'substrate': substrate,
                            'separator': separator,
                            'formulation': formulation,
                            'excluded': dataset.get('excluded', False),
                            'cycler': dataset.get('cycler'),
                            'channel': dataset.get('channel'),
                            'cycler_channel': dataset.get('cycler_channel'),
                            'tracking_placeholder': dataset.get('tracking_placeholder', False),
                            'uploaded_file_source': dataset.get('uploaded_file_source', False),
                            'has_data': dataset.get('has_data', False),
                            'file_label': dataset.get('file_label')
                        }

                edited_datasets.append(edited_dataset)
            datasets = edited_datasets
            st.session_state['datasets'] = datasets
        else:
            # Only one cell
            for i, dataset in enumerate(datasets):
                # Don't skip excluded cells - still render them but with visual indication
                
                is_excluded = dataset.get('excluded', False)
                cell_title = f'Cell {i+1}: {dataset["testnum"] or f"Cell {i+1}"}'
                if is_excluded:
                    cell_title += " ⚠️ EXCLUDED"
                
                with st.expander(cell_title, expanded=False):
                    col1, col2 = st.columns(2)
                    with col1:
                        loading = st.number_input(
                            f'Disc loading (mg) for Cell {i+1}', 
                            min_value=0.0, 
                            step=1.0, 
                            value=float(dataset['loading']),
                            key=f'edit_loading_{i}'
                        )
                        formation_cycles = st.number_input(
                            f'Formation Cycles for Cell {i+1}', 
                            min_value=0, 
                            step=1, 
                            value=int(dataset['formation_cycles']),
                            key=f'edit_formation_{i}'
                        )
                    with col2:
                        active_material = st.number_input(
                            f'% Active material for Cell {i+1}', 
                            min_value=0.0, 
                            max_value=100.0, 
                            step=1.0, 
                            value=float(dataset['active']),
                            key=f'edit_active_{i}'
                        )
                        test_number = st.text_input(
                            f'Test Number for Cell {i+1}', 
                            value=dataset['testnum'] or f'Cell {i+1}',
                            key=f'edit_testnum_{i}'
                        )
                    
                    # Electrolyte, Substrate, and Separator selection
                    substrate_options = get_substrate_options()
                    
                    col3, col4, col5 = st.columns(3)
                    with col3:
                        electrolyte = render_hybrid_electrolyte_input(
                            f'Electrolyte for Cell {i+1}', 
                            default_value=dataset['electrolyte'],
                            key=f'edit_single_electrolyte_{i}'
                        )
                    with col4:
                        substrate = st.selectbox(
                            f'Substrate for Cell {i+1}', 
                            substrate_options,
                            index=substrate_options.index(dataset.get('substrate', 'Copper')) if dataset.get('substrate') in substrate_options else 0,
                            key=f'edit_single_substrate_{i}'
                        )
                    with col5:
                        separator = render_hybrid_separator_input(
                            f'Separator for Cell {i+1}', 
                            default_value=dataset.get('separator', '25um PP'),
                            key=f'edit_single_separator_{i}'
                        )
                    
                    cutoff_lower_default = dataset.get('cutoff_voltage_lower')
                    cutoff_upper_default = dataset.get('cutoff_voltage_upper')
                    if cutoff_lower_default is None:
                        cutoff_lower_default = 2.5
                    if cutoff_upper_default is None:
                        cutoff_upper_default = 4.2
                    
                    cutoff_col1, cutoff_col2 = st.columns(2)
                    with cutoff_col1:
                        cutoff_lower = st.number_input(
                            f'Lower Cutoff Voltage (V) for Cell {i+1}',
                            min_value=0.0,
                            max_value=10.0,
                            step=0.1,
                            value=float(cutoff_lower_default),
                            key=f'edit_single_cutoff_lower_{i}',
                        )
                    with cutoff_col2:
                        cutoff_upper = st.number_input(
                            f'Upper Cutoff Voltage (V) for Cell {i+1}',
                            min_value=0.0,
                            max_value=10.0,
                            step=0.1,
                            value=float(cutoff_upper_default),
                            key=f'edit_single_cutoff_upper_{i}',
                        )
                    
                    # Formulation table
                    st.markdown("**Formulation:**")
                    from ui_components import render_formulation_table
                    # Initialize formulation data if needed
                    formulation_key = f'formulation_data_edit_single_{i}_loaded'
                    if formulation_key not in st.session_state:
                        st.session_state[formulation_key] = dataset['formulation'] if dataset['formulation'] else [{'Component': '', 'Dry Mass Fraction (%)': 0.0}]
                    formulation = render_formulation_table(f'edit_single_{i}_loaded', project_id, get_project_components)
                    
                    # Add two buttons: Exclude and Remove
                    col_btn1, col_btn2 = st.columns(2)
                    
                    with col_btn1:
                        exclude_button_disabled = dataset.get('excluded', False)
                        if st.button("🚫 Exclude Cell", key=f'exclude_cell_single_{i}', disabled=exclude_button_disabled):
                            datasets[i]['excluded'] = True
                            st.session_state['datasets'] = datasets  # Save to session state
                    
                    with col_btn2:
                        if st.button("Remove Cell", key=f'remove_cell_single_{i}'):
                            st.session_state[f'confirm_remove_cell_single_{i}'] = True

                    # Show excluded status
                    if dataset.get('excluded', False):
                        st.warning("This cell is excluded from analysis")
                        if st.button("✅ Include Cell", key=f'include_cell_single_{i}'):
                            datasets[i]['excluded'] = False
                            st.session_state['datasets'] = datasets  # Save to session state

                    if st.session_state.get(f'confirm_remove_cell_single_{i}', False):
                        st.error("**PERMANENT DELETION** - This will permanently delete the cell data and cannot be undone!")
                        col_confirm1, col_confirm2 = st.columns(2)
                        with col_confirm1:
                            if st.button("Delete Permanently", key=f'confirm_yes_single_{i}', type="primary"):
                                # Actually remove the cell from the datasets  
                                datasets.pop(i)
                                st.session_state[f'confirm_remove_cell_single_{i}'] = False
                        with col_confirm2:
                            if st.button("Cancel", key=f'confirm_no_single_{i}'):
                                st.session_state[f'confirm_remove_cell_single_{i}'] = False

                    # Always preserve original file object, only update other fields
                    edited_dataset = {
                        'file': dataset['file'],  # Always preserve original file object
                        'loading': loading,
                        'active': active_material,
                        'testnum': test_number,
                        'formation_cycles': formation_cycles,
                        'cutoff_voltage_lower': cutoff_lower,
                        'cutoff_voltage_upper': cutoff_upper,
                        'electrolyte': electrolyte,
                        'substrate': substrate,
                        'separator': separator,
                        'formulation': formulation,
                        'excluded': dataset.get('excluded', False),
                        'cycler': dataset.get('cycler'),
                        'channel': dataset.get('channel'),
                        'cycler_channel': dataset.get('cycler_channel'),
                        'tracking_placeholder': dataset.get('tracking_placeholder', False),
                        'uploaded_file_source': dataset.get('uploaded_file_source', False),
                        'has_data': dataset.get('has_data', False),
                        'file_label': dataset.get('file_label')
                    }
                    dataset.update(edited_dataset)
                    
        # --- Add Additional Cells ---
        st.markdown("---")
        st.markdown("#### ➕ Add Additional Cells")
        st.info("Upload new raw data files here to append them to this experiment. They will appear in the list above.")
        
        dynamic_uploader_key = f"append_cells_uploader_{len(datasets)}"
        new_uploaded_files = st.file_uploader(
            "Upload new raw data files", 
            type=['csv', 'xlsx'], 
            accept_multiple_files=True, 
            key=dynamic_uploader_key
        )
        
        if new_uploaded_files:
            if st.button("Append Files to Experiment"):
                from ui_components import int_to_roman
                num_existing = len(datasets)
                placeholder_indexes = [
                    idx for idx, item in enumerate(datasets)
                    if item.get('tracking_placeholder') and not item.get('has_data')
                ]
                
                # Fetch default values from the first cell to auto-populate
                default_loading = datasets[0]['loading'] if datasets else 20.0
                default_active = datasets[0]['active'] if datasets else 90.0
                default_formation = datasets[0]['formation_cycles'] if datasets else 4
                default_cutoff_lower = datasets[0].get('cutoff_voltage_lower') if datasets else None
                default_cutoff_upper = datasets[0].get('cutoff_voltage_upper') if datasets else None
                default_electrolyte = datasets[0]['electrolyte'] if datasets else '1M LiPF6 1:1:1'
                default_substrate = datasets[0].get('substrate', 'Copper') if datasets else 'Copper'
                default_separator = datasets[0].get('separator', '25um PP') if datasets else '25um PP'
                default_formulation = datasets[0].get('formulation', [{'Component': '', 'Dry Mass Fraction (%)': 0.0}]) if datasets else [{'Component': '', 'Dry Mass Fraction (%)': 0.0}]
                assigned_to_placeholders = 0

                for placeholder_index, file in zip(placeholder_indexes, new_uploaded_files):
                    datasets[placeholder_index]['file'] = file
                    datasets[placeholder_index]['uploaded_file_source'] = True
                    datasets[placeholder_index]['file_label'] = file.name
                    assigned_to_placeholders += 1

                remaining_files = list(new_uploaded_files)[assigned_to_placeholders:]

                for file in remaining_files:
                    num_existing += 1
                    test_num_val = f"{experiment_name_input} {int_to_roman(num_existing)}" if experiment_name_input else f'Cell {num_existing}'
                    
                    datasets.append({
                        'file': file,
                        'loading': default_loading,
                        'active': default_active,
                        'testnum': test_num_val,
                        'formation_cycles': default_formation,
                        'cutoff_voltage_lower': default_cutoff_lower,
                        'cutoff_voltage_upper': default_cutoff_upper,
                        'electrolyte': default_electrolyte,
                        'substrate': default_substrate,
                        'separator': default_separator,
                        'formulation': default_formulation,
                        'excluded': False,
                        'cycler': None,
                        'channel': None,
                        'cycler_channel': None,
                        'tracking_placeholder': False,
                        'uploaded_file_source': True,
                        'has_data': False,
                        'file_label': file.name
                    })

                if assigned_to_placeholders:
                    st.success(f"Assigned {assigned_to_placeholders} uploaded file(s) to tracked placeholder cells.")
                st.session_state['datasets'] = datasets
                st.rerun()
    else:
        # New experiment flow - use unified render_cell_inputs
        st.markdown("#### Upload Cell Data Files")
        
        # For duplicated experiments, pre-populate defaults from original experiment
        if loaded_experiment:
            experiment_data = loaded_experiment['experiment_data']
            default_cell_values = experiment_data.get('default_cell_values', {})
            
            if default_cell_values:
                st.info("Using default values from the original experiment. You can modify these as needed.")
                # Pre-populate session state with default values from the original experiment
                if 'loading_0' not in st.session_state:
                    st.session_state['loading_0'] = default_cell_values.get('loading', 20.0)
                if 'active_0' not in st.session_state:
                    st.session_state['active_0'] = default_cell_values.get('active_material', 90.0)
                if 'formation_cycles_0' not in st.session_state:
                    st.session_state['formation_cycles_0'] = default_cell_values.get('formation_cycles', 4)
                if 'cutoff_lower_0' not in st.session_state:
                    st.session_state['cutoff_lower_0'] = default_cell_values.get('cutoff_voltage_lower', 2.5)
                if 'cutoff_upper_0' not in st.session_state:
                    st.session_state['cutoff_upper_0'] = default_cell_values.get('cutoff_voltage_upper', 4.2)
                if 'electrolyte_0' not in st.session_state:
                    st.session_state['electrolyte_0'] = default_cell_values.get('electrolyte', '1M LiPF6 1:1:1')
                if 'substrate_0' not in st.session_state:
                    st.session_state['substrate_0'] = default_cell_values.get('substrate', 'Copper')
                if 'separator_0' not in st.session_state:
                    st.session_state['separator_0'] = default_cell_values.get('separator', '25um PP')
                # Pre-populate formulation data
                if 'formulation_data_0_main_cell_inputs' not in st.session_state:
                    formulation = default_cell_values.get('formulation', [])
                    if formulation:
                        st.session_state['formulation_data_0_main_cell_inputs'] = formulation
        
        current_project_id = st.session_state.get('current_project_id')
        # Pass experiment name for auto-generating cell names with roman numerals
        # Use experiment_name_input directly (defined above) for real-time updates
        datasets, full_cell_data = render_cell_inputs(
            context_key='main_cell_inputs', 
            project_id=current_project_id, 
            get_components_func=get_project_components,
            experiment_name=experiment_name_input if experiment_name_input else '',
            project_type=project_type
        )
        st.session_state['datasets'] = datasets
        if full_cell_data:
            st.session_state['full_cell_data'] = full_cell_data
        # Store original uploaded files separately to prevent loss
        if datasets:
            st.session_state['original_uploaded_files'] = [ds['file'] for ds in datasets if ds.get('file')]
    
    # Group assignment section (if multiple cells)
    enable_grouping = False
    show_averages = False
    group_assignments = current_group_assignments
    group_names = current_group_names
    
    # --- New experiment-level fields ---
    st.markdown('---')
    st.subheader('Experiment Parameters')
    solids_content = st.number_input(
        'Solids Content (%)',
        min_value=0.0, max_value=100.0, step=0.1,
        value=coerce_float_input(st.session_state.get('solids_content', 0.0), 0.0),
        key='solids_content',
        help='Percentage solids in the slurry formulation when the electrode was made.'
    )
    pressed_thickness = st.number_input(
        'Pressed Thickness (um)',
        min_value=0.0, step=0.1,
        value=coerce_float_input(st.session_state.get('pressed_thickness', 0.0), 0.0),
        key='pressed_thickness',
        help='Pressed electrode thickness in microns (um).'
    )
    experiment_notes = st.text_area(
        'Experiment Notes',
        value=st.session_state.get('experiment_notes', ''),
        key='experiment_notes',
        help='Basic notes associated with this experiment.'
    )
    
    if datasets and len([d for d in datasets if d.get('file') or loaded_experiment or is_new_experiment]) > 1:
        st.markdown("---")
        st.markdown("#### 👥 Group Assignment (Optional)")
        enable_grouping = st.checkbox('Assign Cells into Groups?', value=bool(current_group_assignments))
        
        if enable_grouping:
            col1, col2, col3 = st.columns(3)
            with col1:
                group_names[0] = st.text_input('Group A Name', value=group_names[0], key='main_group_name_a')
            with col2:
                group_names[1] = st.text_input('Group B Name', value=group_names[1], key='main_group_name_b')
            with col3:
                group_names[2] = st.text_input('Group C Name', value=group_names[2], key='main_group_name_c')
            
            st.markdown("**Assign each cell to a group:**")
            group_assignments = []
            for i, cell in enumerate(datasets):
                if cell.get('file') or loaded_experiment or is_new_experiment:
                    cell_name = cell['testnum'] or f'Cell {i+1}'
                    default_group = current_group_assignments[i] if (current_group_assignments and i < len(current_group_assignments)) else group_names[0]
                    group = st.radio(
                        f"Assign {cell_name} to group:",
                        [group_names[0], group_names[1], group_names[2], "Exclude"],
                        index=[group_names[0], group_names[1], group_names[2], "Exclude"].index(default_group) if default_group in [group_names[0], group_names[1], group_names[2], "Exclude"] else 0,
                        key=f"main_group_assignment_{i}",
                        horizontal=True
                    )
                    group_assignments.append(group)
            
            show_averages = st.checkbox("Show Group Averages", value=True)
    
    # Update session state with current values
    st.session_state['current_experiment_name'] = experiment_name_input
    
    # Handle different cell formats for Full Cell projects
    if project_type == "Full Cell":
        cell_format = st.session_state.get('current_cell_format', 'Coin')
        if cell_format == "Coin":
            st.session_state['current_disc_diameter_mm'] = disc_diameter_input
        elif cell_format == "Pouch":
            # For pouch cells, we'll store area instead of diameter
            cathode_length = st.session_state.get('current_cathode_length', 50.0)
            cathode_width = st.session_state.get('current_cathode_width', 50.0)
            num_stacked_cells = st.session_state.get('current_num_stacked_cells', 1)
            # Calculate equivalent diameter for backwards compatibility
            total_area_cm2 = (cathode_length * cathode_width * num_stacked_cells) / 100
            equivalent_diameter = 2 * (total_area_cm2 / np.pi) ** 0.5 * 10  # Convert back to mm
            st.session_state['current_disc_diameter_mm'] = equivalent_diameter
        else:
            # For other formats, use current or default value
            st.session_state['current_disc_diameter_mm'] = disc_diameter_input or current_disc_diameter
    else:
        # For non-Full Cell projects, use traditional disc diameter
        st.session_state['current_disc_diameter_mm'] = disc_diameter_input
    
    st.session_state['current_group_assignments'] = group_assignments
    st.session_state['current_group_names'] = group_names
    
    # Render formulation editor modal if needed
    render_formulation_editor_modal()
    
    # Save/Update experiment button
    st.markdown("---")
    if loaded_experiment:
        if st.button("Update Experiment", type="primary", use_container_width=True):
            # Update the loaded experiment with new values
            experiment_id = loaded_experiment['experiment_id']
            project_id = loaded_experiment['project_id']
            
            # Get project type for efficiency calculation
            project_type = "Full Cell"  # Default
            if project_id:
                project_info = get_project_by_id(project_id)
                if project_info:
                    project_type = project_info[3]  # project_type is the 4th field
            
            # Prepare updated cells data with recalculated gravimetric capacities
            updated_cells_data = []
            existing_cells = experiment_data.get('cells', [])
            
            for i, dataset in enumerate(datasets):
                # Check if this is an existing cell or a new one
                if i < len(existing_cells):
                    # Update existing cell
                    original_cell = existing_cells[i]
                    
                    # Read current input values from session state widgets (they have the latest values)
                    widget_loading = st.session_state.get(f'edit_loading_{i}') or st.session_state.get(f'edit_single_loading_{i}')
                    widget_active = st.session_state.get(f'edit_active_{i}') or st.session_state.get(f'edit_single_active_{i}')
                    widget_formation = st.session_state.get(f'edit_formation_{i}') or st.session_state.get(f'edit_single_formation_{i}')
                    widget_testnum = st.session_state.get(f'edit_testnum_{i}') or st.session_state.get(f'edit_single_testnum_{i}')
                    
                    # Use widget values if available, otherwise use dataset values
                    new_loading = widget_loading if widget_loading is not None else dataset.get('loading', 0)
                    new_active = widget_active if widget_active is not None else dataset.get('active', 0)
                    new_formation = widget_formation if widget_formation is not None else dataset.get('formation_cycles', 4)
                    new_testnum = widget_testnum if widget_testnum is not None else dataset.get('testnum', f'Cell {i+1}')
                    
                    # Check if loading or active material has changed
                    original_loading = original_cell.get('loading', 0)
                    original_active = original_cell.get('active_material', 0)
                    
                    uploaded_file = dataset.get('file')
                    uploaded_file_source = dataset.get('uploaded_file_source', False)

                    # Recalculate gravimetric capacities if loading or active material changed
                    updated_data_json = original_cell.get('data_json')
                    updated_file_name = original_cell.get('file_name')
                    processed_cutoff_lower = dataset.get('cutoff_voltage_lower', original_cell.get('cutoff_voltage_lower'))
                    processed_cutoff_upper = dataset.get('cutoff_voltage_upper', original_cell.get('cutoff_voltage_upper'))

                    if uploaded_file_source and uploaded_file:
                        try:
                            temp_dfs = load_and_preprocess_data([dataset], project_type)
                            if temp_dfs and len(temp_dfs) > 0:
                                processed_cell = temp_dfs[0]
                                updated_data_json = processed_cell['df'].to_json()
                                updated_file_name = uploaded_file.name
                                processed_cutoff_lower = processed_cell.get('cutoff_voltage_lower', processed_cutoff_lower)
                                processed_cutoff_upper = processed_cell.get('cutoff_voltage_upper', processed_cutoff_upper)
                                st.info(f"Attached raw file to {new_testnum}")
                            else:
                                st.warning(f"Failed to process uploaded file for {new_testnum}. Keeping metadata only.")
                        except Exception as e:
                            st.error(f"Error processing uploaded file for {new_testnum}: {str(e)}")

                    if (
                        not uploaded_file_source
                        and (new_loading != original_loading or new_active != original_active)
                        and updated_data_json
                    ):
                        try:
                            # Parse the original DataFrame - fix deprecation warning
                            original_df = pd.read_json(StringIO(updated_data_json))
                            
                            # Recalculate gravimetric capacities
                            updated_df = recalculate_gravimetric_capacities(original_df, new_loading, new_active)
                            
                            # Update the data JSON with recalculated values
                            updated_data_json = updated_df.to_json()
                            
                            # Show before/after comparison of first few values
                            if len(updated_df) > 0:
                                try:
                                    first_qdis_old = original_df['Q Dis (mAh/g)'].iloc[0] if 'Q Dis (mAh/g)' in original_df.columns else 'N/A'
                                    first_qdis_new = updated_df['Q Dis (mAh/g)'].iloc[0] if 'Q Dis (mAh/g)' in updated_df.columns else 'N/A'
                                    st.info(f"Recalculated gravimetric capacities for {new_testnum}")
                                    st.info(f"   Changes: Loading: {original_loading}→{new_loading}mg, Active: {original_active}→{new_active}%")
                                    if isinstance(first_qdis_old, (int, float)) and isinstance(first_qdis_new, (int, float)):
                                        st.info(f"   First Cycle Q Dis: {first_qdis_old:.2f}→{first_qdis_new:.2f} mAh/g")
                                    else:
                                        st.info(f"   First Cycle Q Dis: {first_qdis_old}→{first_qdis_new} mAh/g")
                                except Exception:
                                    st.info(f"Recalculated gravimetric capacities for {new_testnum}")
                        except Exception as e:
                            st.warning(f"Could not recalculate capacities for {new_testnum}: {str(e)}")
                    
                    updated_cell = original_cell.copy()
                    
                    # Recalculate porosity if loading changed and we have the required data
                    if (new_loading != original_loading and 
                        pressed_thickness and pressed_thickness > 0 and 
                        dataset.get('formulation') and 
                        disc_diameter_input):
                        try:
                            from porosity_calculations import calculate_porosity_from_experiment_data
                            porosity_data = calculate_porosity_from_experiment_data(
                                disc_mass_mg=new_loading,
                                disc_diameter_mm=disc_diameter_input,
                                pressed_thickness_um=pressed_thickness,
                                formulation=dataset['formulation']
                            )
                            updated_cell['porosity'] = porosity_data['porosity']
                            st.info(f"   Recalculated porosity: {porosity_data['porosity']*100:.1f}%")
                        except Exception as e:
                            st.warning(f"   Could not recalculate porosity for {new_testnum}: {str(e)}")
                    
                    # Read other widget values too
                    widget_electrolyte = st.session_state.get(f'edit_electrolyte_{i}') or st.session_state.get(f'edit_single_electrolyte_{i}')
                    widget_substrate = st.session_state.get(f'edit_substrate_{i}') or st.session_state.get(f'edit_single_substrate_{i}')
                    widget_separator = st.session_state.get(f'edit_separator_{i}') or st.session_state.get(f'edit_single_separator_{i}')
                    new_cutoff_lower = processed_cutoff_lower
                    new_cutoff_upper = processed_cutoff_upper
                    
                    updated_cell.update({
                        'loading': new_loading,
                        'active_material': new_active,
                        'formation_cycles': new_formation,
                        'test_number': new_testnum,
                        'cell_name': new_testnum,
                        'file_name': updated_file_name,
                        'electrolyte': widget_electrolyte if widget_electrolyte is not None else dataset.get('electrolyte', '1M LiPF6 1:1:1'),
                        'substrate': widget_substrate if widget_substrate is not None else dataset.get('substrate', 'Copper'),
                        'separator': widget_separator if widget_separator is not None else dataset.get('separator', '25um PP'),
                        'cutoff_voltage_lower': new_cutoff_lower,
                        'cutoff_voltage_upper': new_cutoff_upper,
                        'formulation': dataset.get('formulation', []),
                        'data_json': updated_data_json,  # Updated with recalculated values
                        'excluded': dataset.get('excluded', False),
                        'cycler': dataset.get('cycler'),
                        'channel': dataset.get('channel'),
                        'cycler_channel': dataset.get('cycler_channel'),
                        'tracking_placeholder': bool(dataset.get('tracking_placeholder', False) and not updated_data_json)
                    })
                    updated_cells_data.append(updated_cell)
                else:
                    # This is a new cell being added to the experiment (e.g., uploading to a duplicate)
                    cell_name = dataset['testnum'] if dataset['testnum'] else f'Cell {i+1}'
                    file_name = dataset['file'].name if dataset.get('file') else f'cell_{i+1}.csv'
                    
                    try:
                        # Process the data to get DataFrame
                        temp_dfs = load_and_preprocess_data([dataset], project_type)
                        if temp_dfs and len(temp_dfs) > 0:
                            processed_cell = temp_dfs[0]
                            df = processed_cell['df']
                            
                            new_cell = {
                                'cell_name': cell_name,
                                'file_name': file_name,
                                'loading': dataset['loading'],
                                'active_material': dataset['active'],
                                'formation_cycles': dataset['formation_cycles'],
                                'test_number': dataset['testnum'],
                                'electrolyte': dataset.get('electrolyte', '1M LiPF6 1:1:1'),
                                'substrate': dataset.get('substrate', 'Copper'),
                                'separator': dataset.get('separator', '25um PP'),
                                'cutoff_voltage_lower': processed_cell.get('cutoff_voltage_lower', dataset.get('cutoff_voltage_lower')),
                                'cutoff_voltage_upper': processed_cell.get('cutoff_voltage_upper', dataset.get('cutoff_voltage_upper')),
                                'formulation': dataset.get('formulation', []),
                                'data_json': df.to_json(),
                                'excluded': dataset.get('excluded', False),
                                'cycler': dataset.get('cycler'),
                                'channel': dataset.get('channel'),
                                'cycler_channel': dataset.get('cycler_channel'),
                                'tracking_placeholder': False
                            }
                            
                            # Calculate porosity if we have the required data
                            if (pressed_thickness and pressed_thickness > 0 and 
                                dataset.get('formulation') and 
                                disc_diameter_input):
                                try:
                                    from porosity_calculations import calculate_porosity_from_experiment_data
                                    porosity_data = calculate_porosity_from_experiment_data(
                                        disc_mass_mg=dataset['loading'],
                                        disc_diameter_mm=disc_diameter_input,
                                        pressed_thickness_um=pressed_thickness,
                                        formulation=dataset['formulation']
                                    )
                                    new_cell['porosity'] = porosity_data['porosity']
                                    st.info(f"   Calculated porosity for {cell_name}: {porosity_data['porosity']*100:.1f}%")
                                except Exception as e:
                                    st.warning(f"   Could not calculate porosity for {cell_name}: {str(e)}")
                            
                            updated_cells_data.append(new_cell)
                            st.info(f"Processed new cell: {cell_name}")
                        else:
                            st.warning(f"Failed to process data for {cell_name}. Skipping this cell.")
                    except Exception as e:
                        st.error(f"Error processing {cell_name}: {str(e)}")
                        continue
            
            try:
                # Prepare cell format data for Full Cell projects
                cell_format_data = {}
                if project_type == "Full Cell":
                    cell_format = st.session_state.get('current_cell_format', 'Coin')
                    cell_format_data['cell_format'] = cell_format
                    if cell_format == "Pouch":
                        cell_format_data['cathode_length'] = st.session_state.get('current_cathode_length', 50.0)
                        cell_format_data['cathode_width'] = st.session_state.get('current_cathode_width', 50.0)
                        cell_format_data['num_stacked_cells'] = st.session_state.get('current_num_stacked_cells', 1)
                
                update_experiment(
                    experiment_id=experiment_id,
                    project_id=project_id,
                    experiment_name=experiment_name_input,
                    experiment_date=experiment_date_input,
                    disc_diameter_mm=st.session_state.get('current_disc_diameter_mm', current_disc_diameter),
                    group_assignments=group_assignments,
                    group_names=group_names,
                    cells_data=updated_cells_data,
                    solids_content=solids_content,
                    pressed_thickness=pressed_thickness,
                    experiment_notes=experiment_notes,
                    cell_format_data=cell_format_data
                )
                clear_navigation_caches()
                
                # Update the loaded experiment in session state
                st.session_state['loaded_experiment']['experiment_name'] = experiment_name_input
                st.session_state['loaded_experiment']['experiment_data'].update({
                    'experiment_date': experiment_date_input.isoformat(),
                    'disc_diameter_mm': disc_diameter_input,
                    'group_assignments': group_assignments,
                    'group_names': group_names,
                    'cells': updated_cells_data,
                    'solids_content': solids_content,
                    'pressed_thickness': pressed_thickness,
                    'experiment_notes': experiment_notes
                })
                
                # Clear any cached processed data to force recalculation
                if 'processed_data_cache' in st.session_state:
                    del st.session_state['processed_data_cache']
                if 'cache_key' in st.session_state:
                    del st.session_state['cache_key']
                
                # Reload the experiment from database to get the updated data
                try:
                    updated_experiment_data = get_experiment_data(experiment_id)
                    if updated_experiment_data:
                        # get_experiment_data returns: id, project_id, cell_name, file_name, loading, active_material, 
                        # formation_cycles, test_number, electrolyte, substrate, separator, data_json, created_date
                        # data_json is at index 11
                        st.session_state['loaded_experiment'] = {
                            'experiment_id': experiment_id,
                            'project_id': project_id,
                            'experiment_name': experiment_name_input,
                            'experiment_data': json.loads(updated_experiment_data[11])  # data_json is at index 11
                        }
                        # Set a flag to indicate that calculations have been updated
                        st.session_state['calculations_updated'] = True
                        st.session_state['update_timestamp'] = datetime.now()
                        st.success("Experiment updated successfully! All calculated values have been refreshed.")
                        st.info("Summary tables, plots, and Master Table will reflect the updated values.")
                    else:
                        st.success("Experiment updated successfully!")
                except Exception as reload_error:
                    st.warning(f"Experiment updated but failed to reload data: {str(reload_error)}")
                    st.success("✅ Experiment updated successfully!")
                
                st.rerun()
            except Exception as e:
                st.error(f"Error updating experiment: {str(e)}")
    
    elif is_new_experiment:
        # Save new experiment (only if we have valid data and a selected project)
        valid_datasets = [ds for ds in datasets if not ds.get('excluded', False) and ds.get('file') and ds.get('loading', 0) > 0 and 0 < ds.get('active', 0) <= 100]
        
        if valid_datasets and st.session_state.get('current_project_id'):
            if st.button("Save New Experiment", type="primary", use_container_width=True):
                current_project_id = st.session_state['current_project_id']
                current_project_name = st.session_state['current_project_name']
                
                # Use experiment name from input or generate one
                exp_name = experiment_name_input if experiment_name_input else f"Experiment_{datetime.now().strftime('%Y%m%d_%H%M%S')}"
                
                # Prepare cells data
                cells_data = []
                for i, ds in enumerate(valid_datasets):
                    cell_name = ds['testnum'] if ds['testnum'] else f'Cell {i+1}'
                    file_name = ds['file'].name if ds['file'] else f'cell_{i+1}.csv'
                    
                    try:
                        # Get project type for efficiency calculation
                        project_type = "Full Cell"  # Default
                        if st.session_state.get('current_project_id'):
                            current_project_id = st.session_state['current_project_id']
                            project_info = get_project_by_id(current_project_id)
                            if project_info:
                                project_type = project_info[3]  # project_type is the 4th field
                        
                        # Process the data to get DataFrame
                        temp_dfs = load_and_preprocess_data([ds], project_type)
                        if temp_dfs and len(temp_dfs) > 0:
                            processed_cell = temp_dfs[0]
                            df = processed_cell['df']
                            
                            cells_data.append({
                                'cell_name': cell_name,
                                'file_name': file_name,
                                'loading': ds['loading'],
                                'active_material': ds['active'],
                                'formation_cycles': ds['formation_cycles'],
                                'test_number': ds['testnum'],
                                'electrolyte': ds.get('electrolyte', '1M LiPF6 1:1:1'),
                                'substrate': ds.get('substrate', 'Copper'),
                                'separator': ds.get('separator', '25um PP'),
                                'cutoff_voltage_lower': processed_cell.get('cutoff_voltage_lower', ds.get('cutoff_voltage_lower')),
                                'cutoff_voltage_upper': processed_cell.get('cutoff_voltage_upper', ds.get('cutoff_voltage_upper')),
                                'formulation': ds.get('formulation', []),
                                'data_json': df.to_json(),
                                'excluded': ds.get('excluded', False)  # Add this line
                            })
                        else:
                            st.warning(f"Failed to process data for {cell_name}. Skipping this cell.")
                    except Exception as e:
                        st.error(f"Error processing {cell_name}: {str(e)}")
                        continue
                
                # Save the experiment
                if cells_data:
                    try:
                        if check_experiment_name_exists(current_project_id, exp_name):
                            experiment_id = get_experiment_by_name(current_project_id, exp_name)
                            # Prepare cell format data for Full Cell projects
                            cell_format_data = {}
                            if project_type == "Full Cell":
                                cell_format = st.session_state.get('current_cell_format', 'Coin')
                                cell_format_data['cell_format'] = cell_format
                                if cell_format == "Pouch":
                                    cell_format_data['cathode_length'] = st.session_state.get('current_cathode_length', 50.0)
                                    cell_format_data['cathode_width'] = st.session_state.get('current_cathode_width', 50.0)
                                    cell_format_data['num_stacked_cells'] = st.session_state.get('current_num_stacked_cells', 1)

                            update_experiment(
                                experiment_id=experiment_id,
                                project_id=current_project_id,
                                experiment_name=exp_name,
                                experiment_date=experiment_date_input,
                                disc_diameter_mm=st.session_state.get('current_disc_diameter_mm', 15),
                                group_assignments=group_assignments,
                                group_names=group_names,
                                cells_data=cells_data,
                                solids_content=solids_content,
                                pressed_thickness=pressed_thickness,
                                experiment_notes=experiment_notes,
                                cell_format_data=cell_format_data,
                                additional_data=get_active_batch_builder_additional_data(),
                            )
                            clear_batch_builder_cell_input_state()
                            clear_navigation_caches()
                            st.success(f"Updated experiment '{exp_name}' in project '{current_project_name}'!")
                        else:
                            # Prepare cell format data for Full Cell projects
                            cell_format_data = {}
                            if project_type == "Full Cell":
                                cell_format = st.session_state.get('current_cell_format', 'Coin')
                                cell_format_data['cell_format'] = cell_format
                                if cell_format == "Pouch":
                                    cell_format_data['cathode_length'] = st.session_state.get('current_cathode_length', 50.0)
                                    cell_format_data['cathode_width'] = st.session_state.get('current_cathode_width', 50.0)
                                    cell_format_data['num_stacked_cells'] = st.session_state.get('current_num_stacked_cells', 1)
                            
                            save_experiment(
                                project_id=current_project_id,
                                experiment_name=exp_name,
                                experiment_date=experiment_date_input,
                                disc_diameter_mm=st.session_state.get('current_disc_diameter_mm', 15),
                                group_assignments=group_assignments,
                                group_names=group_names,
                                cells_data=cells_data,
                                solids_content=solids_content,
                                pressed_thickness=pressed_thickness,
                                experiment_notes=experiment_notes,
                                cell_format_data=cell_format_data,
                                additional_data=get_active_batch_builder_additional_data(),
                            )
                            clear_batch_builder_cell_input_state()
                            clear_navigation_caches()
                            st.success(f"Saved experiment '{exp_name}' with {len(cells_data)} cells in project '{current_project_name}'!")
                        
                        st.rerun()
                    except Exception as e:
                        st.error(f"Failed to save experiment: {str(e)}")
                else:
                    st.error("No valid cell data to save. Please check your files and try again.")
        
        elif valid_datasets and not st.session_state.get('current_project_id'):
            st.warning("Please select a project in the sidebar before saving the experiment.")
        elif not valid_datasets:
            st.info("Upload cell data files and enter valid parameters to save an experiment.")

# --- Data Preprocessing Section ---
# This section is now handled in the Cell Inputs tab
# Data processing now happens in the Cell Inputs tab

# Check if we have loaded experiment data to display
loaded_experiment = st.session_state.get('loaded_experiment')
if loaded_experiment:
    st.markdown("---")
    st.markdown(f"### Loaded Experiment: {loaded_experiment['experiment_name']}")
    
    # Convert saved JSON data back to DataFrames
    loaded_dfs = []
    experiment_data = loaded_experiment['experiment_data']
    cells_data = experiment_data.get('cells', [])
    
    for i, cell_data in enumerate(cells_data):
        cell_name = cell_data.get('cell_name', 'Unknown')
        try:
            # Fix the deprecation warning by using StringIO
            df = pd.read_json(StringIO(cell_data['data_json']))
            
            # Get project type for efficiency recalculation
            project_type = "Full Cell"  # Default
            if st.session_state.get('current_project_id'):
                current_project_id = st.session_state['current_project_id']
                project_info = get_project_by_id(current_project_id)
                if project_info:
                    project_type = project_info[3]  # project_type is the 4th field
            
            # Check if loading or active material have been changed in session state
            # Check both multi-cell and single-cell widget keys
            widget_loading = st.session_state.get(f'edit_loading_{i}') or st.session_state.get(f'edit_single_loading_{i}')
            widget_active = st.session_state.get(f'edit_active_{i}') or st.session_state.get(f'edit_single_active_{i}')
            
            db_loading = cell_data.get('loading', 0)
            db_active = cell_data.get('active_material', 0)
            
            # Use widget values if available, otherwise use database values
            current_loading = widget_loading if widget_loading is not None else db_loading
            current_active = widget_active if widget_active is not None else db_active
            
            # Recalculate gravimetric capacities if loading or active material changed
            if (current_loading != db_loading or current_active != db_active):
                try:
                    df = recalculate_gravimetric_capacities(df, current_loading, current_active)
                except Exception as e:
                    print(f"Error recalculating gravimetric capacities for {cell_name}: {e}")
            
            # Recalculate efficiency based on project type for all projects (was previously only for Anode)
            if 'Q charge (mA.h)' in df.columns and 'Q discharge (mA.h)' in df.columns:
                # Recalculate efficiency for all project types to ensure correctness
                df['Efficiency (-)'] = calculate_efficiency_based_on_project_type(
                    df['Q charge (mA.h)'], 
                    df['Q discharge (mA.h)'], 
                    project_type
                ) / 100  # Convert to decimal for consistency
            
            # Extract electrode data from experiment
            pressed_thickness = experiment_data.get('pressed_thickness', None)
            solids_content = experiment_data.get('solids_content', None)
            porosity = cell_data.get('porosity', None)
            
            # If porosity is not available in cell data, try to calculate it
            if porosity is None or porosity <= 0:
                try:
                    from porosity_calculations import calculate_porosity_from_experiment_data
                    if (cell_data.get('loading') and 
                        experiment_data.get('disc_diameter_mm') and 
                        pressed_thickness and 
                        cell_data.get('formulation')):
                        
                        porosity_data = calculate_porosity_from_experiment_data(
                            disc_mass_mg=cell_data['loading'],
                            disc_diameter_mm=experiment_data['disc_diameter_mm'],
                            pressed_thickness_um=pressed_thickness,
                            formulation=cell_data['formulation']
                        )
                        porosity = porosity_data['porosity']
                except Exception as e:
                    print(f"Error calculating porosity for {cell_name}: {e}")
                    porosity = None
            
            # Get formulation data from cell_data
            formulation = cell_data.get('formulation', [])
            
            loaded_dfs.append({
                'df': df,
                'testnum': cell_data.get('test_number', cell_data.get('cell_name', 'Unknown')),
                'loading': current_loading,  # Use current loading (may be updated from widget)
                'active': current_active,  # Use current active material (may be updated from widget)
                'formation_cycles': cell_data.get('formation_cycles'),
                'project_type': project_type,
                'excluded': cell_data.get('excluded', False),
                # Add electrode data for export functionality
                'pressed_thickness': pressed_thickness,
                'solids_content': solids_content,
                'porosity': porosity,
                # Add formulation data for export functionality
                'formulation': formulation
            })
            
            # Debug info for electrode data loading (can be removed after testing)
            if pressed_thickness or porosity or solids_content:
                print(f"DEBUG: Loaded electrode data for {cell_name}:")
                print(f"  - Pressed thickness: {pressed_thickness}")
                print(f"  - Solids content: {solids_content}")
                print(f"  - Porosity: {porosity}")
        except Exception as e:
            st.error(f"Error loading data for {cell_name}: {str(e)}")
    
    if loaded_dfs:
        # Use loaded data for analysis
        dfs = loaded_dfs
        ready = True
        st.success(f"Loaded {len(loaded_dfs)} cell(s) from saved experiment")
        
        # Update datasets in session state to reflect any widget changes
        # This ensures Summary tables and other components also show updated values
        for i, df_data in enumerate(loaded_dfs):
            if i < len(st.session_state.get('datasets', [])):
                st.session_state['datasets'][i]['loading'] = df_data['loading']
                st.session_state['datasets'][i]['active'] = df_data['active']
        
        # Display experiment metadata
        if experiment_data.get('experiment_date'):
            st.info(f"📅 Experiment Date: {experiment_data['experiment_date']}")
        # LLM Summary Section
        experiment_id = loaded_experiment.get('experiment_id')
        with st.expander("🤖 Generate LLM-Ready Summary", expanded=False):
            st.markdown("""
            Generate a token-efficient summary of this experiment optimized for LLM analysis.
            Includes experiment parameters, cell performance metrics, curve characteristics, 
            and a capacity vs cycle plot image.
            """)
            
            if st.button("Generate Summary", type="primary", use_container_width=True, 
                         key=f'llm_summary_btn_{experiment_id}'):
                with st.spinner("Generating summary and plot..."):
                    try:
                        summary_text, plot_image_base64, stats = generate_experiment_summary(experiment_id)
                        
                        # Display statistics
                        col1, col2, col3, col4 = st.columns(4)
                        with col1:
                            st.metric("Estimated Tokens", f"{stats.get('token_estimate', 0):,}")
                        with col2:
                            st.metric("Characters", f"{stats.get('char_count', 0):,}")
                        with col3:
                            st.metric("Lines", f"{stats.get('line_count', 0):,}")
                        with col4:
                            st.metric("Cells", f"{stats.get('num_cells', 0)}")
                        
                        # Display plot image if available
                        if plot_image_base64:
                            st.markdown("### Capacity vs Cycle Plot")
                            st.markdown("*This plot image can be included in your LLM prompt for visual analysis.*")
                            # Decode and display image
                            import base64
                            from io import BytesIO
                            from PIL import Image
                            img_data = base64.b64decode(plot_image_base64)
                            img = Image.open(BytesIO(img_data))
                            st.image(img, caption=f"{loaded_experiment['experiment_name']} - Capacity vs Cycle", 
                                   use_container_width=True)
                            
                            # Download button for plot
                            st.download_button(
                                label="Download Plot Image",
                                data=img_data,
                                file_name=f"{loaded_experiment['experiment_name']}_capacity_plot.png",
                                mime="image/png",
                                key=f'download_plot_{experiment_id}'
                            )
                        
                        # Display summary text in text area for easy copying
                        st.markdown("### Summary Text (Copy for LLM)")
                        st.text_area(
                            "Experiment Summary",
                            value=summary_text,
                            height=400,
                            key=f'llm_summary_textarea_{experiment_id}',
                            help="Copy this text and paste it into your LLM prompt for analysis. Include the plot image above if using a vision model."
                        )
                        
                        # Copy button using streamlit's clipboard functionality
                        st.code(summary_text, language=None)
                        
                        # Save to session state for later viewing (use different key than widget)
                        st.session_state[f'llm_summary_text_{experiment_id}'] = summary_text
                        st.session_state[f'llm_summary_plot_{experiment_id}'] = plot_image_base64
                        st.session_state[f'llm_summary_stats_{experiment_id}'] = stats
                        
                        st.success("Summary generated! Copy the text above to use with your LLM.")
                        
                    except Exception as e:
                        st.error(f"Error generating summary: {str(e)}")
                        st.exception(e)
            
            # Show cached summary if available
            cached_summary = st.session_state.get(f'llm_summary_text_{experiment_id}')
            if cached_summary:
                if st.button("View Last Summary", key=f'view_summary_{experiment_id}'):
                    cached_stats = st.session_state.get(f'llm_summary_stats_{experiment_id}', {})
                    cached_plot = st.session_state.get(f'llm_summary_plot_{experiment_id}')
                    
                    col1, col2, col3, col4 = st.columns(4)
                    with col1:
                        st.metric("Estimated Tokens", f"{cached_stats.get('token_estimate', 0):,}")
                    with col2:
                        st.metric("Characters", f"{cached_stats.get('char_count', 0):,}")
                    with col3:
                        st.metric("Lines", f"{cached_stats.get('line_count', 0):,}")
                    with col4:
                        st.metric("Cells", f"{cached_stats.get('num_cells', 0)}")
                    
                    if cached_plot:
                        st.markdown("### Capacity vs Cycle Plot")
                        import base64
                        from io import BytesIO
                        from PIL import Image
                        img_data = base64.b64decode(cached_plot)
                        img = Image.open(BytesIO(img_data))
                        st.image(img, caption=f"{loaded_experiment['experiment_name']} - Capacity vs Cycle", 
                               use_container_width=True)
                    
                    st.text_area(
                        "Experiment Summary",
                        value=cached_summary,
                        height=400,
                        key=f'llm_summary_view_textarea_{experiment_id}'
                    )
                    st.code(cached_summary, language=None)
    else:
        st.error("Failed to load experiment data")
        ready = False

# Determine if we have data ready for analysis
if loaded_experiment:
    ready = len(loaded_dfs) > 0
else:
    # For new experiments, check if we have valid uploaded data
    datasets = st.session_state.get('datasets', [])
    # Only include datasets with a real uploaded file for processing
    valid_datasets = []
    for ds in datasets:
        file_obj = ds.get('file')
        if (file_obj and 
            hasattr(file_obj, 'read') and 
            hasattr(file_obj, 'name') and 
            hasattr(file_obj, 'type') and
            ds.get('loading', 0) > 0 and 
            0 < ds.get('active', 0) <= 100):
            # Additional check: ensure it's a real Streamlit UploadedFile, not a mock object
            try:
                # Try to access the file's size property (real uploaded files have this)
                if hasattr(file_obj, 'size') and file_obj.size is not None and file_obj.size > 0:
                    valid_datasets.append(ds)
            except (AttributeError, TypeError):
                # Skip files that don't have proper size attribute or other issues
                continue
    
    # Process uploaded data if we have valid datasets
    if valid_datasets:
        # Create a cache key based on file names and parameters
        cache_key = []
        for ds in valid_datasets:
            if ds.get('file'):
                file_info = f"{ds['file'].name}_{ds['loading']}_{ds['active']}_{ds['formation_cycles']}"
                cache_key.append(file_info)
        cache_key_str = "_".join(cache_key)
        
        # Check if we have cached processed data
        if ('processed_data_cache' in st.session_state and 
            st.session_state.get('cache_key') == cache_key_str):
            dfs = st.session_state['processed_data_cache']
        else:
            # Process data and cache it
            # Final safety check before processing
            safe_datasets = []
            for ds in valid_datasets:
                try:
                    # Test that we can actually read from the file
                    file_obj = ds['file']
                    current_pos = file_obj.tell() if hasattr(file_obj, 'tell') else 0
                    file_obj.seek(0)
                    # Try to read a small sample to verify it's readable
                    sample = file_obj.read(10)
                    file_obj.seek(current_pos)  # Reset position
                    if sample:  # File has content
                        safe_datasets.append(ds)
                except Exception as e:
                    st.warning(f"Skipping invalid file: {ds.get('file', {}).get('name', 'Unknown')} - {str(e)}")
                    continue
            
            if safe_datasets:
                # Get project type for efficiency calculation
                project_type = "Full Cell"  # Default
                if st.session_state.get('current_project_id'):
                    current_project_id = st.session_state['current_project_id']
                    project_info = get_project_by_id(current_project_id)
                    if project_info:
                        project_type = project_info[3]  # project_type is the 4th field
                
                dfs = load_and_preprocess_data(safe_datasets, project_type)
                # After loading and preprocessing, re-attach the latest formation_cycles to each dfs entry
                for i, d in enumerate(dfs):
                    if i < len(safe_datasets):
                        d['formation_cycles'] = safe_datasets[i]['formation_cycles']
                
                # Display file type information
                file_types = [d.get('file_type', 'Unknown') for d in dfs]
                biologic_count = file_types.count('biologic_csv')
                neware_count = file_types.count('neware_xlsx')
                mti_count = file_types.count('mti_xlsx')
                
                # Build info message about processed files
                file_type_parts = []
                if biologic_count > 0:
                    file_type_parts.append(f"{biologic_count} Biologic CSV file(s)")
                if neware_count > 0:
                    file_type_parts.append(f"{neware_count} Neware XLSX file(s)")
                if mti_count > 0:
                    file_type_parts.append(f"{mti_count} MTI XLSX file(s)")
                
                if file_type_parts:
                    st.info(f"Processed {len(dfs)} files: {', '.join(file_type_parts)}")
            else:
                dfs = []
                st.error("No valid files found for processing.")
            
            # Cache the processed data
            st.session_state['processed_data_cache'] = dfs
            st.session_state['cache_key'] = cache_key_str
        
        ready = len(dfs) > 0
    else:
        ready = False

if ready:
    # Use values from Cell Inputs tab
    disc_diameter_mm = st.session_state.get('current_disc_diameter_mm', 15)
    experiment_name = st.session_state.get('current_experiment_name', '')
    group_assignments = st.session_state.get('current_group_assignments')
    group_names = st.session_state.get('current_group_names', ["Group A", "Group B", "Group C"])
    enable_grouping = bool(group_assignments)
    show_averages = enable_grouping
    datasets = st.session_state.get('datasets', [])
    disc_area_cm2 = np.pi * (disc_diameter_mm / 2 / 10) ** 2
    
    # Filter out excluded cells from dfs
    if loaded_experiment:
        dfs = [d for d in loaded_dfs if not d.get('excluded', False)]
    else:
        # For new experiments, we need to filter the processed dfs, not the raw valid_datasets
        # The processed dfs are already cached in st.session_state['processed_data_cache']
        processed_dfs = st.session_state.get('processed_data_cache', [])
        valid_datasets = st.session_state.get('datasets', [])
        
        # Create a mapping of file names to excluded status
        excluded_files = {}
        for ds in valid_datasets:
            if ds.get('file') and hasattr(ds['file'], 'name'):
                excluded_files[ds['file'].name] = ds.get('excluded', False)
        
        # Filter processed dfs based on excluded status
        dfs = []
        for d in processed_dfs:
            # Check if this processed data corresponds to an excluded file
            file_name = d.get('file_name', '')
            if not excluded_files.get(file_name, False):
                dfs.append(d)

    # --- Group Average Curve Calculation for Plotting ---
    group_curves = []
    if enable_grouping and group_assignments is not None:
        group_dfs = [[], [], []]
        for idx, name in enumerate(group_names):
            group_dfs[idx] = [df for df, g in zip(dfs, group_assignments) if g == name]
        def compute_group_avg_curve(group_dfs):
            if not group_dfs:
                return None, None, None, None
            dfs_trimmed = [d['df'] for d in group_dfs]
            x_col = dfs_trimmed[0].columns[0]
            common_cycles = set(dfs_trimmed[0][x_col])
            for df in dfs_trimmed[1:]:
                common_cycles = common_cycles & set(df[x_col])
            common_cycles = sorted(list(common_cycles))
            if not common_cycles:
                return None, None, None, None
            avg_qdis = []
            avg_qchg = []
            avg_eff = []
            for cycle in common_cycles:
                qdis_vals = []
                qchg_vals = []
                eff_vals = []
                for df in dfs_trimmed:
                    row = df[df[x_col] == cycle]
                    if not row.empty:
                        if 'Q Dis (mAh/g)' in row:
                            qdis_vals.append(row['Q Dis (mAh/g)'].values[0])
                        if 'Q Chg (mAh/g)' in row:
                            qchg_vals.append(row['Q Chg (mAh/g)'].values[0])
                        if 'Efficiency (-)' in row and not pd.isnull(row['Efficiency (-)'].values[0]):
                            eff_vals.append(row['Efficiency (-)'].values[0] * 100)
                avg_qdis.append(sum(qdis_vals)/len(qdis_vals) if qdis_vals else None)
                avg_qchg.append(sum(qchg_vals)/len(qchg_vals) if qchg_vals else None)
                avg_eff.append(sum(eff_vals)/len(eff_vals) if eff_vals else None)
            return common_cycles, avg_qdis, avg_qchg, avg_eff
        group_curves = [compute_group_avg_curve(group_dfs[idx]) for idx in range(3)]
    # --- Main Tabs Content ---
    with tab1:
        st.subheader("📈 Cycling Performance Plots")
        
        # Plot Style Toggle - Prominent at top
        col_toggle, col_spacer = st.columns([3, 7])
        with col_toggle:
            plot_style = st.radio(
                "Plot Style",
                options=["📊 Interactive (Plotly)", "📉 Static (Matplotlib)"],
                index=0,
                horizontal=True,
                help="Interactive plots allow zooming, panning, and hovering for details. Static plots are the traditional matplotlib charts.",
                key="plot_style_toggle"
            )
        
        use_interactive = plot_style.startswith("📊")
        
        # Get formation cycles for reference cycle calculation
        formation_cycles = st.session_state.get('current_formation_cycles', 4)
        if ready and datasets:
            # Try to get formation cycles from the first dataset
            if 'formation_cycles' in datasets[0]:
                formation_cycles = datasets[0]['formation_cycles']
        
        # Simplified Plot Controls in expanders
        with st.expander("🎛️ Plot Controls & Customization", expanded=False):
            # Plot Controls
            show_lines, show_efficiency_lines, remove_last_cycle, show_graph_title, show_average_performance, avg_line_toggles, remove_markers, hide_legend, group_plot_toggles, cycle_filter, y_axis_limits, excluded_from_average = render_toggle_section(dfs, enable_grouping=enable_grouping)
            
            # Color customization UI
            custom_colors = render_experiment_color_customization(
                dfs, experiment_name, show_average_performance, 
                enable_grouping, group_names
            )
        
        # Combined plot toggle
        show_combined_plot = st.checkbox(
            "Show Capacity Retention on Secondary Y-Axis",
            value=False,
            key="show_combined_capacity_retention",
            help="Combine Specific Capacity and Capacity Retention into a single graph with dual Y-axes."
        )
        
        # Conditionally show combined plot or separate plots based on toggle
        if show_combined_plot and ready and dfs:
            # Get reference cycle settings
            all_cycles = []
            for d in dfs:
                try:
                    df = d['df']
                    if not df.empty:
                        cycles = df[df.columns[0]].tolist()
                        all_cycles.extend(cycles)
                except:
                    pass
            
            if all_cycles:
                min_cycle = min(all_cycles)
                max_cycle = max(all_cycles)
                
                # Get maximum data length for formation cycles skip limit
                max_data_length = 0
                for d in dfs:
                    try:
                        df = d['df']
                        if not df.empty:
                            max_data_length = max(max_data_length, len(df))
                    except:
                        pass
                
                # Combined Plot Configuration
                config_col1, config_col2 = st.columns([1, 1])
                
                with config_col1:
                    max_skip = max(0, max_data_length - 1) if max_data_length > 0 else 0
                    formation_cycles_skip = st.number_input(
                        "Formation Cycles to Skip",
                        min_value=0,
                        max_value=max_skip,
                        value=0,
                        step=1,
                        key="combined_formation_cycles_skip",
                        help=f"Number of initial cycles to skip when determining the 100% reference capacity."
                    )
                
                with config_col2:
                    retention_threshold = st.slider(
                        "Retention Threshold (%)",
                        min_value=0.0,
                        max_value=100.0,
                        value=80.0,
                        step=5.0,
                        key="combined_retention_threshold",
                        help="Set the threshold line for capacity retention analysis."
                    )
                
                baseline_col1, baseline_col2 = st.columns(2)
                with baseline_col1:
                    show_baseline_100 = st.checkbox('Show 100% baseline', value=True, key='combined_show_baseline_100')
                with baseline_col2:
                    show_baseline_80 = st.checkbox(f'Show {retention_threshold:.0f}% threshold', value=True, key='combined_show_baseline_80')
                
                # Calculate reference cycle based on formation cycles skip
                # The reference cycle index will be formation_cycles_skip (0-based index)
                # But we need the actual cycle number from the data
                reference_cycle_index = formation_cycles_skip
                
                # Get the actual cycle number from the first visible cell
                reference_cycle = None
                for d in dfs:
                    try:
                        df = d['df']
                        if not df.empty and len(df) > reference_cycle_index:
                            cycle_col = df.columns[0]
                            reference_cycle = int(df.iloc[reference_cycle_index][cycle_col])
                            break
                    except:
                        pass
                
                # Fallback to default if not found
                if reference_cycle is None:
                    default_ref_cycle = formation_cycles + 1
                    if default_ref_cycle < min_cycle:
                        default_ref_cycle = min_cycle
                    elif default_ref_cycle > max_cycle:
                        default_ref_cycle = max_cycle
                    reference_cycle = int(default_ref_cycle)
                
                y_axis_preset = st.selectbox(
                    "Retention Y-Axis Range",
                    options=["Auto-scale", "Full Range (0-110%)", "Focused View (70-110%)", "Standard View (50-110%)", "Custom Range"],
                    index=0,
                    key="combined_y_axis_preset"
                )
                
                # Set Y-axis range based on preset
                if y_axis_preset == "Auto-scale":
                    y_axis_min, y_axis_max = None, None  # Will be calculated from capacity range
                elif y_axis_preset == "Full Range (0-110%)":
                    y_axis_min, y_axis_max = 0.0, 110.0
                elif y_axis_preset == "Focused View (70-110%)":
                    y_axis_min, y_axis_max = 70.0, 110.0
                elif y_axis_preset == "Standard View (50-110%)":
                    y_axis_min, y_axis_max = 50.0, 110.0
                else:
                    custom_col1, custom_col2 = st.columns(2)
                    with custom_col1:
                        y_axis_min = st.number_input("Min Y (%)", min_value=0.0, max_value=100.0, value=0.0, step=5.0, key="combined_retention_y_axis_min")
                    with custom_col2:
                        y_axis_max = st.number_input("Max Y (%)", min_value=50.0, max_value=200.0, value=110.0, step=5.0, key="combined_retention_y_axis_max")
                
                # Generate combined plot
                combined_fig = plot_combined_capacity_retention_graph(
                    dfs, show_lines, reference_cycle, formation_cycles, remove_last_cycle,
                    show_graph_title, experiment_name, show_average_performance,
                    avg_line_toggles, remove_markers, hide_legend,
                    retention_threshold=retention_threshold,
                    y_axis_min=y_axis_min,
                    y_axis_max=y_axis_max,
                    show_baseline_line=show_baseline_100,
                    show_threshold_line=show_baseline_80,
                    cycle_filter=cycle_filter,
                    custom_colors=custom_colors,
                    capacity_y_axis_limits=y_axis_limits,
                    formation_cycles_skip=formation_cycles_skip
                )
                st.pyplot(combined_fig)
                
                st.caption(f"Combined view: Specific Capacity (left Y-axis) and Capacity Retention (right Y-axis). Reference: cycle {reference_cycle}.")
            else:
                st.warning("No cycle data available. Please upload data files first.")
        else:
            # Separate plots (when combined toggle is disabled)
            if use_interactive and ready and dfs:
                # Interactive Plotly plots
                st.markdown("### Capacity Plot")
                interactive_cap_fig = plot_interactive_capacity(
                    dfs, show_lines, show_efficiency_lines, remove_last_cycle, 
                    experiment_name, show_average_performance, avg_line_toggles, 
                    group_names, custom_colors, excluded_from_average, cycle_filter, y_axis_limits
                )
                st.plotly_chart(interactive_cap_fig, use_container_width=True)
                
                # Add info about interactive features
                st.info("💡 **Tip**: Hover over data points for cycle, capacity, and retention details. Retention uses the first valid post-formation cycle for each cell and skips clearly anomalous baseline cycles when needed.")
            else:
                # Static matplotlib plot
                fig = plot_capacity_graph(
                    dfs, show_lines, show_efficiency_lines, remove_last_cycle, show_graph_title, experiment_name,
                    show_average_performance, avg_line_toggles, remove_markers, hide_legend,
                    group_a_curve=(group_curves[0][0], group_curves[0][1]) if enable_grouping and group_curves and group_curves[0][0] and group_curves[0][1] and group_plot_toggles.get("Group Q Dis", False) else None,
                    group_b_curve=(group_curves[1][0], group_curves[1][1]) if enable_grouping and group_curves and group_curves[1][0] and group_curves[1][1] and group_plot_toggles.get("Group Q Dis", False) else None,
                    group_c_curve=(group_curves[2][0], group_curves[2][1]) if enable_grouping and group_curves and group_curves[2][0] and group_curves[2][1] and group_plot_toggles.get("Group Q Dis", False) else None,
                    group_a_qchg=(group_curves[0][0], group_curves[0][2]) if enable_grouping and group_curves and group_curves[0][0] and group_curves[0][2] and group_plot_toggles.get("Group Q Chg", False) else None,
                    group_b_qchg=(group_curves[1][0], group_curves[1][2]) if enable_grouping and group_curves and group_curves[1][0] and group_curves[1][2] and group_plot_toggles.get("Group Q Chg", False) else None,
                    group_c_qchg=(group_curves[2][0], group_curves[2][2]) if enable_grouping and group_curves and group_curves[2][0] and group_curves[2][2] and group_plot_toggles.get("Group Q Chg", False) else None,
                    group_a_eff=(group_curves[0][0], group_curves[0][3]) if enable_grouping and group_curves and group_curves[0][0] and group_curves[0][3] and group_plot_toggles.get("Group Efficiency", False) else None,
                    group_b_eff=(group_curves[1][0], group_curves[1][3]) if enable_grouping and group_curves and group_curves[1][0] and group_curves[1][3] and group_plot_toggles.get("Group Efficiency", False) else None,
                    group_c_eff=(group_curves[2][0], group_curves[2][3]) if enable_grouping and group_curves and group_curves[2][0] and group_curves[2][3] and group_plot_toggles.get("Group Efficiency", False) else None,
                    group_names=group_names,
                    cycle_filter=cycle_filter,
                    custom_colors=custom_colors,
                    y_axis_limits=y_axis_limits,
                    excluded_from_average=excluded_from_average
                )
                st.pyplot(fig)
            
            if ready and dfs:
                # Get available cycles from the data to determine valid range for reference cycle
                all_cycles = []
                for d in dfs:
                    try:
                        df = d['df']
                        if not df.empty:
                            cycles = df[df.columns[0]].tolist()
                            all_cycles.extend(cycles)
                    except:
                        pass
                
                if all_cycles:
                    min_cycle = min(all_cycles)
                    max_cycle = max(all_cycles)
                    default_ref_cycle = formation_cycles + 1
                    
                    # Ensure default reference cycle is within valid range
                    if default_ref_cycle < min_cycle:
                        default_ref_cycle = min_cycle
                    elif default_ref_cycle > max_cycle:
                        default_ref_cycle = max_cycle
                    
                    # Reference cycle input controls
                    col1, col2, col3 = st.columns([2, 1, 1])
                    
                    with col1:
                        reference_cycle = st.number_input(
                            "Reference Cycle (100% baseline)",
                            min_value=int(min_cycle),
                            max_value=int(max_cycle),
                            value=int(default_ref_cycle),
                            step=1,
                            key="reference_cycle",
                            help=f"Select which cycle to use as the 100% reference point."
                        )
                    
                    with col2:
                        st.metric("Formation Cycles", formation_cycles)
                    
                    with col3:
                        st.metric("Available Cycles", f"{int(min_cycle)} - {int(max_cycle)}")
                    
                    # Retention Plot Section - Organize controls in expander
                    st.markdown("---")
                    st.markdown("### Capacity Retention Analysis")
                    
                    with st.expander("⚙️ Retention Plot Settings", expanded=False):
                        # Retention plot controls
                        control_col1, control_col2, control_col3 = st.columns([1, 1, 1])
                        
                        with control_col1:
                            retention_threshold = st.slider(
                                "Retention Threshold (%)",
                                min_value=0.0,
                                max_value=100.0,
                                value=80.0,
                                step=5.0,
                                key="retention_threshold"
                            )
                        
                        with control_col2:
                            y_axis_preset = st.selectbox(
                                "Y-Axis Range",
                                options=["Full Range (0-110%)", "Focused View (70-110%)", "Standard View (50-110%)", "Custom Range"],
                                index=0,
                                key="y_axis_preset"
                            )
                        
                        with control_col3:
                            if y_axis_preset == "Custom Range":
                                custom_min = st.number_input("Min Y (%)", min_value=0.0, max_value=100.0, value=st.session_state.get('retention_y_axis_min', 0.0), step=5.0, key="retention_y_axis_min")
                                custom_max = st.number_input("Max Y (%)", min_value=50.0, max_value=200.0, value=st.session_state.get('retention_y_axis_max', 110.0), step=5.0, key="retention_y_axis_max")
                                y_axis_min, y_axis_max = custom_min, custom_max
                            else:
                                if y_axis_preset == "Full Range (0-110%)":
                                    y_axis_min, y_axis_max = 0.0, 110.0
                                elif y_axis_preset == "Focused View (70-110%)":
                                    y_axis_min, y_axis_max = 70.0, 110.0
                                elif y_axis_preset == "Standard View (50-110%)":
                                    y_axis_min, y_axis_max = 50.0, 110.0
                                st.metric("Y-Axis Range", f"{y_axis_min:.0f}% - {y_axis_max:.0f}%")
                        
                        # Retention plot specific options
                        retention_col1, retention_col2 = st.columns(2)
                        with retention_col1:
                            show_baseline_line = st.checkbox(
                                'Show baseline (100%)',
                                value=True,
                                key='retention_baseline'
                            )
                        with retention_col2:
                            show_threshold_line = st.checkbox(
                                f'Show threshold ({retention_threshold:.0f}%)',
                                value=True,
                                key='retention_threshold_line'
                            )
                    
                    # Generate capacity retention plot
                    if use_interactive:
                        # Interactive Plotly retention plot
                        interactive_ret_fig = plot_interactive_retention(
                            dfs, show_lines, reference_cycle, remove_last_cycle, 
                            experiment_name, show_average_performance, custom_colors,
                            retention_threshold, cycle_filter, y_axis_min, y_axis_max
                        )
                        st.plotly_chart(interactive_ret_fig, use_container_width=True)
                    else:
                        # Static matplotlib retention plot
                        retention_fig = plot_capacity_retention_graph(
                            dfs, show_lines, reference_cycle, formation_cycles, remove_last_cycle, 
                            show_graph_title, experiment_name, show_average_performance, 
                            avg_line_toggles, remove_markers, hide_legend,
                            group_a_curve=None,  # Can be extended later for group retention
                            group_b_curve=None,
                            group_c_curve=None,
                            group_names=group_names,
                            retention_threshold=retention_threshold,
                            y_axis_min=y_axis_min,
                            y_axis_max=y_axis_max,
                            show_baseline_line=show_baseline_line,
                            show_threshold_line=show_threshold_line,
                            cycle_filter=cycle_filter,
                            custom_colors=custom_colors
                        )
                        st.pyplot(retention_fig)
                else:
                    st.warning("No cycle data available. Please upload data files first.")
            else:
                st.info("Upload and process data files to see capacity retention analysis.")
        
        # Summary Table Section at the bottom
        if ready and dfs:
            st.markdown("---")
            st.subheader("Summary Statistics")
            
            # Show update notification if calculations were recently updated
            if st.session_state.get('calculations_updated', False):
                update_time = st.session_state.get('update_timestamp')
                if update_time:
                    time_str = update_time.strftime("%H:%M:%S")
                    st.success(f"Values updated at {time_str} - All calculations have been refreshed!")
                    st.session_state['calculations_updated'] = False
            
            # Add toggle for showing average column
            show_average_col = False
            if len(dfs) > 1:
                show_average_col = st.toggle("Show average column", value=True, key="show_average_col_toggle")
            
            # Display summary statistics table
            from ui_components import display_summary_stats
            display_summary_stats(dfs, disc_area_cm2, show_average_col, group_assignments, group_names)
        
        # --- Full Cell Specific Plots ---
        if project_type == 'Full Cell' and ready and dfs:
            st.markdown("---")
            st.header("🔋 Full Cell Performance Analysis")
            st.info("📊 Device-level performance metrics optimized for Full Cell characterization")
            
            # High-Precision Coulombic Efficiency Plot
            st.subheader("⚡ High-Precision Coulombic Efficiency Tracking")
            st.markdown("High-precision CE tracking for cycle life prediction (typically plotted in 99-100% range)")
            
            ce_col1, ce_col2 = st.columns(2)
            with ce_col1:
                ce_y_min = st.number_input(
                    'CE Y-axis Min (%)',
                    min_value=90.0,
                    max_value=100.0,
                    value=99.0,
                    step=0.1,
                    key='ce_y_min',
                    help="Minimum Y-axis value for CE plot"
                )
            with ce_col2:
                ce_y_max = st.number_input(
                    'CE Y-axis Max (%)',
                    min_value=90.0,
                    max_value=100.0,
                    value=100.0,
                    step=0.1,
                    key='ce_y_max',
                    help="Maximum Y-axis value for CE plot"
                )
            
            try:
                from plotting import plot_coulombic_efficiency_precision
                ce_fig = plot_coulombic_efficiency_precision(
                    dfs, show_lines, remove_last_cycle, show_graph_title, experiment_name,
                    remove_markers=remove_markers, hide_legend=hide_legend, 
                    cycle_filter=cycle_filter, custom_colors=custom_colors,
                    y_axis_min=ce_y_min, y_axis_max=ce_y_max
                )
                st.pyplot(ce_fig)
            except Exception as e:
                st.error(f"Error plotting Coulombic Efficiency: {e}")
            
            st.markdown("---")
            
            # Energy Efficiency Plot
            st.subheader("⚡ Energy Efficiency")
            st.markdown("Energy efficiency (E_discharge / E_charge) provides insights into voltage polarization and cycle life")
            
            try:
                from plotting import plot_energy_efficiency
                ee_fig = plot_energy_efficiency(
                    dfs, show_lines, remove_last_cycle, show_graph_title, experiment_name,
                    remove_markers=remove_markers, hide_legend=hide_legend, 
                    cycle_filter=cycle_filter, custom_colors=custom_colors
                )
                st.pyplot(ee_fig)
            except Exception as e:
                st.error(f"Error plotting Energy Efficiency: {e}")
            
            st.markdown("---")
    with tab2:
        st.header("Export & Download")
        
        # Only show export options if data is ready
        if ready:
            stored_experiment_notes = ""
            if dfs:
                if 'experiment_notes' in dfs[0]:
                    stored_experiment_notes = dfs[0]['experiment_notes'] or ""
                else:
                    stored_experiment_notes = st.session_state.get('experiment_notes', "")
            has_experiment_notes = bool(stored_experiment_notes.strip())
            
            include_summary_table = True
            include_electrode_data = False
            include_porosity = False
            include_thickness = False
            include_solids_content = False
            include_formulation = False
            
            current_ref_cycle = st.session_state.get('reference_cycle', 5)
            current_threshold = st.session_state.get('retention_threshold', 80.0)
            
            with st.container(border=True):
                header_cols = st.columns([0.75, 0.25], gap="small")
                with header_cols[0]:
                    st.subheader("PowerPoint")
                    st.caption("Create a clean single-slide summary for the current experiment.")
                with header_cols[1]:
                    st.metric("Visible cells", len(dfs))
                
                control_cols = st.columns([0.58, 0.42], gap="large")
                
                with control_cols[0]:
                    powerpoint_chart = st.radio(
                        "Chart on the slide",
                        options=["Capacity retention", "Capacity comparison"],
                        index=0,
                        horizontal=True,
                        key="export_plot_focus",
                        help="Choose which chart to feature on the PowerPoint slide."
                    )
                    include_retention_plot = powerpoint_chart == "Capacity retention"
                    include_main_plot = not include_retention_plot
                    
                    if has_experiment_notes:
                        include_notes = st.checkbox(
                            "Include experiment notes",
                            value=True,
                            key="export_include_notes",
                            help="Add the notes saved for this experiment."
                        )
                    else:
                        include_notes = False
                        st.caption("No experiment notes are saved for this experiment yet.")
                    
                    if include_retention_plot:
                        st.caption(
                            f"Retention settings stay synced with the Plots tab: reference cycle {current_ref_cycle}, threshold {current_threshold:.0f}%."
                        )
                    else:
                        st.caption("Capacity comparison uses the current plot styling and cell visibility from the Plots tab.")
                
                with control_cols[1]:
                    st.markdown("**Included on the slide**")
                    slide_contents = [
                        "Summary metrics table",
                        "Experiment metadata",
                        "Selected chart",
                    ]
                    if include_notes:
                        slide_contents.append("Experiment notes")
                    
                    for item in slide_contents:
                        st.markdown(f"- {item}")
                    
                    summary_cols = st.columns(2)
                    with summary_cols[0]:
                        st.metric("Chart", "Retention" if include_retention_plot else "Capacity")
                    with summary_cols[1]:
                        st.metric("Notes", "Included" if include_notes else "Off")
                
                from export import export_powerpoint
                
                try:
                    with st.spinner("Preparing PowerPoint..."):
                        pptx_bytes, pptx_file_name = export_powerpoint(
                            dfs=dfs,
                            show_averages=show_average_performance,
                            experiment_name=experiment_name,
                            show_lines=show_lines,
                            show_efficiency_lines=show_efficiency_lines,
                            remove_last_cycle=remove_last_cycle,
                            include_summary_table=include_summary_table,
                            include_main_plot=include_main_plot,
                            include_retention_plot=include_retention_plot,
                            include_notes=include_notes,
                            include_electrode_data=include_electrode_data,
                            include_porosity=include_porosity,
                            include_thickness=include_thickness,
                            include_solids_content=include_solids_content,
                            include_formulation=include_formulation,
                            experiment_notes=stored_experiment_notes if include_notes else "",
                            retention_threshold=current_threshold,
                            reference_cycle=current_ref_cycle,
                            formation_cycles=dfs[0].get('formation_cycles', 4) if dfs else st.session_state.get('current_formation_cycles', 4),
                            retention_show_lines=show_lines,
                            retention_remove_markers=remove_markers,
                            retention_hide_legend=hide_legend,
                            retention_show_title=show_graph_title,
                            show_baseline_line=st.session_state.get('retention_show_baseline', True),
                            show_threshold_line=st.session_state.get('retention_show_threshold', True),
                            y_axis_min=st.session_state.get('y_axis_min', 0.0),
                            y_axis_max=st.session_state.get('y_axis_max', 110.0),
                            show_graph_title=st.session_state.get('show_graph_title', True),
                            show_average_performance=show_average_performance,
                            avg_line_toggles=st.session_state.get('avg_line_toggles', {}),
                            remove_markers=st.session_state.get('remove_markers', False),
                            hide_legend=st.session_state.get('hide_legend', False)
                        )
                    
                    st.download_button(
                        "Download PowerPoint",
                        data=pptx_bytes,
                        file_name=pptx_file_name,
                        mime='application/vnd.openxmlformats-officedocument.presentationml.presentation',
                        key='download_enhanced_pptx',
                        use_container_width=True
                    )
                    st.caption(f"Ready: {pptx_file_name}")
                except Exception as e:
                    st.error(f"Error generating PowerPoint: {str(e)}")
                    st.error("Please check your data and settings, then try again.")
            
            # Project-Level Export Section
            if current_project_id:
                with st.container(border=True):
                    st.subheader("Project PowerPoint")
                    st.caption("Build one deck with a slide for every experiment in the current project using the same chart choice above.")
                    
                    if st.button("Build Project PowerPoint", type="secondary", use_container_width=True):
                        try:
                            from database import get_all_project_experiments_data, get_project_by_id
                            from export import export_powerpoint
                            from io import BytesIO
                            from pptx import Presentation
                            import json
                            from io import StringIO
                            from data_processing import calculate_efficiency_based_on_project_type
                            
                            # Get all experiments for the project, sorted by creation date
                            all_experiments_data = get_all_project_experiments_data(current_project_id)
                            
                            if not all_experiments_data:
                                st.error("No experiments found in this project.")
                            else:
                                # Sort by creation date (chronologically)
                                # Handle None values by using a far-future date for sorting
                                from datetime import datetime
                                def get_sort_key(exp_data):
                                    created_date = exp_data[13] if len(exp_data) > 13 else None  # created_date is index 13
                                    if created_date is None:
                                        return datetime.max  # Put None dates at the end
                                    # If it's already a datetime, use it directly
                                    if isinstance(created_date, datetime):
                                        return created_date
                                    # If it's a string, try to parse it
                                    if isinstance(created_date, str):
                                        try:
                                            return datetime.fromisoformat(created_date.replace('Z', '+00:00'))
                                        except:
                                            return datetime.max
                                    return datetime.max
                                
                                all_experiments_data.sort(key=get_sort_key)
                                
                                st.info(f"Found {len(all_experiments_data)} experiment(s). Generating slides...")
                                
                                # Create a new presentation for the project
                                project_prs = Presentation()
                                project_prs.slide_layouts[6]  # Blank layout
                                
                                # Get project info
                                project_info = get_project_by_id(current_project_id)
                                project_name = project_info[1] if project_info else "Project"
                                project_type = project_info[3] if project_info and len(project_info) > 3 else "Full Cell"
                                
                                # Process each experiment
                                experiments_processed = 0
                                for exp_data in all_experiments_data:
                                    exp_id, exp_name, file_name, loading, active_material, formation_cycles, test_number, electrolyte, substrate, separator, formulation_json, data_json, created_date, porosity, experiment_notes, cutoff_voltage_lower, cutoff_voltage_upper = exp_data
                                    
                                    try:
                                        # Parse experiment data
                                        parsed_data = json.loads(data_json)
                                        
                                        # Load cells data
                                        cells_data = parsed_data.get('cells', [])
                                        if not cells_data:
                                            # Single cell experiment
                                            cells_data = [{
                                                'data_json': data_json,
                                                'cell_name': exp_name,
                                                'test_number': test_number,
                                                'loading': loading,
                                                'active_material': active_material,
                                                'formation_cycles': formation_cycles,
                                                'formulation': json.loads(formulation_json) if formulation_json else []
                                            }]
                                        
                                        # Build dfs structure for this experiment
                                        exp_dfs = []
                                        for cell_data in cells_data:
                                            if cell_data.get('excluded', False):
                                                continue
                                            
                                            cell_data_json = cell_data.get('data_json', '')
                                            if not cell_data_json:
                                                continue
                                            
                                            df = pd.read_json(StringIO(cell_data_json))
                                            
                                            # Recalculate efficiency based on project type
                                            if 'Q charge (mA.h)' in df.columns and 'Q discharge (mA.h)' in df.columns:
                                                df['Efficiency (-)'] = calculate_efficiency_based_on_project_type(
                                                    df['Q charge (mA.h)'], 
                                                    df['Q discharge (mA.h)'], 
                                                    project_type
                                                ) / 100
                                            
                                            # Get formulation
                                            formulation = cell_data.get('formulation', [])
                                            if not formulation and formulation_json:
                                                try:
                                                    formulation = json.loads(formulation_json)
                                                except:
                                                    formulation = []
                                            
                                            exp_dfs.append({
                                                'df': df,
                                                'testnum': cell_data.get('test_number', cell_data.get('cell_name', 'Unknown')),
                                                'loading': cell_data.get('loading', loading),
                                                'active': cell_data.get('active_material', active_material),
                                                'formation_cycles': cell_data.get('formation_cycles', formation_cycles),
                                                'project_type': project_type,
                                                'excluded': False,
                                                'pressed_thickness': parsed_data.get('pressed_thickness'),
                                                'solids_content': parsed_data.get('solids_content'),
                                                'porosity': cell_data.get('porosity', porosity),
                                                'formulation': formulation
                                            })
                                        
                                        if not exp_dfs:
                                            st.warning(f"Skipping {exp_name}: No valid cell data found.")
                                            continue
                                        
                                        # Add slides for this experiment to the project presentation
                                        export_powerpoint(
                                            dfs=exp_dfs,
                                            show_averages=True,
                                            experiment_name=exp_name,
                                            show_lines=show_lines,
                                            show_efficiency_lines=show_efficiency_lines,
                                            remove_last_cycle=remove_last_cycle,
                                            include_summary_table=include_summary_table,
                                            include_main_plot=include_main_plot,
                                            include_retention_plot=include_retention_plot,
                                            include_notes=include_notes,
                                            include_electrode_data=include_electrode_data,
                                            include_porosity=include_porosity,
                                            include_thickness=include_thickness,
                                            include_solids_content=include_solids_content,
                                            include_formulation=include_formulation,
                                            experiment_notes=(experiment_notes or "") if include_notes else "",
                                            retention_threshold=current_threshold,
                                            reference_cycle=current_ref_cycle,
                                            formation_cycles=formation_cycles or 4,
                                            retention_show_lines=show_lines,
                                            retention_remove_markers=remove_markers,
                                            retention_hide_legend=hide_legend,
                                            retention_show_title=show_graph_title,
                                            show_baseline_line=st.session_state.get('retention_show_baseline', True),
                                            show_threshold_line=st.session_state.get('retention_show_threshold', True),
                                            y_axis_min=st.session_state.get('y_axis_min', 0.0),
                                            y_axis_max=st.session_state.get('y_axis_max', 110.0),
                                            show_graph_title=st.session_state.get('show_graph_title', True),
                                            show_average_performance=show_average_performance,
                                            avg_line_toggles=st.session_state.get('avg_line_toggles', {}),
                                            remove_markers=st.session_state.get('remove_markers', False),
                                            hide_legend=st.session_state.get('hide_legend', False),
                                            existing_prs=project_prs  # Append to project presentation
                                        )
                                        
                                        experiments_processed += 1
                                        
                                    except Exception as e:
                                        st.warning(f"Error processing experiment {exp_name}: {str(e)}")
                                        import logging
                                        logging.error(f"Error processing experiment {exp_name}: {e}")
                                        continue
                                
                                if experiments_processed > 0:
                                    # Save project presentation
                                    project_bio = BytesIO()
                                    project_prs.save(project_bio)
                                    project_bio.seek(0)
                                    
                                    st.success(f"Project export completed! Processed {experiments_processed} experiment(s).")
                                    st.download_button(
                                        "Download Project PowerPoint",
                                        data=project_bio,
                                        file_name=f"{project_name} Project Export.pptx",
                                        mime='application/vnd.openxmlformats-officedocument.presentationml.presentation',
                                        key='download_project_pptx',
                                        use_container_width=True
                                    )
                                else:
                                    st.error("No experiments could be processed for export.")
                        
                        except Exception as e:
                            st.error(f"Error exporting project: {str(e)}")
                            import traceback
                            st.error(traceback.format_exc())
            
            with st.container(border=True):
                st.subheader("Excel")
                st.caption("Download the processed data and summary sheets for the current experiment.")
                
                from export import export_excel
                
                try:
                    with st.spinner("Preparing Excel workbook..."):
                        excel_bytes, excel_file_name = export_excel(dfs, show_average_performance, experiment_name)
                    
                    st.download_button(
                        "Download Excel",
                        data=excel_bytes,
                        file_name=excel_file_name,
                        mime='application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',
                        key='download_enhanced_excel',
                        use_container_width=True
                    )
                    st.caption(f"Ready: {excel_file_name}")
                except Exception as e:
                    st.error(f"Error generating Excel file: {str(e)}")
                    st.error("Please check your data and try again.")
        else:
            with st.container(border=True):
                st.subheader("Exports unlock after data is processed")
                st.caption(
                    "Load an experiment from the sidebar, or create one in Cell Inputs and upload data to enable downloads."
                )
                st.markdown("- PowerPoint summary slides")
                st.markdown("- Excel workbooks with charts")
                st.markdown("- Project-wide PowerPoint export")

if not ready:
    with tab1:
        st.subheader("📈 Cycling Performance Plots")
        with st.container(border=True):
            st.markdown(
                """
                <div style="padding: 2.5rem 1rem; text-align: center;">
                    <div style="font-size: 2rem; line-height: 1; margin-bottom: 0.75rem;">📈</div>
                    <div style="font-size: 1.2rem; font-weight: 600; margin-bottom: 0.4rem;">Open an experiment to view plots</div>
                    <div style="color: #5f6b7a; max-width: 32rem; margin: 0 auto;">
                        Load an existing experiment from the sidebar, or set up a new experiment in Cell Inputs and upload data to generate plots.
                    </div>
                </div>
                """,
                unsafe_allow_html=True,
            )

# --- Comparison Tab ---
if tab_comparison and current_project_id:
    with tab_comparison:
        current_project_name = st.session_state.get('current_project_name', 'Selected Project')
        st.caption(f"Project: {current_project_name}")
        
        # Get all experiments data for this project
        all_experiments_data = get_all_project_experiments_data(current_project_id)
        
        if not all_experiments_data:
            st.info("No experiments found in this project. Create experiments to see comparison data.")
        else:
            # Extract experiment names for selection
            experiment_options = []
            experiment_dict = {}
            
            for exp_data in all_experiments_data:
                exp_id, exp_name, file_name, loading, active_material, formation_cycles, test_number, electrolyte, substrate, separator, formulation_json, data_json, created_date, porosity, experiment_notes, cutoff_voltage_lower, cutoff_voltage_upper = exp_data
                experiment_options.append(exp_name)
                experiment_dict[exp_name] = exp_data

            if 'comparison_selected_experiments' not in st.session_state:
                st.session_state['comparison_selected_experiments'] = []

            # Build a lookup that maps stripped names to actual DB names
            # so prefilled selections (which strip whitespace) still match
            stripped_to_actual = {name.strip(): name for name in experiment_options}

            sanitized_selection = []
            for experiment_name in st.session_state.get('comparison_selected_experiments', []):
                if experiment_name in experiment_options:
                    sanitized_selection.append(experiment_name)
                elif experiment_name.strip() in stripped_to_actual:
                    sanitized_selection.append(stripped_to_actual[experiment_name.strip()])

            if sanitized_selection != st.session_state.get('comparison_selected_experiments', []):
                st.session_state['comparison_selected_experiments'] = sanitized_selection

            comparison_prefill_notice = st.session_state.pop('comparison_prefill_notice', None)
            if comparison_prefill_notice:
                st.info(comparison_prefill_notice)
            
            # Experiment Selection
            selected_experiments = st.multiselect(
                "Select experiments to compare",
                options=experiment_options,
                key='comparison_selected_experiments',
                help="Select two or more experiments to compare"
            )
            
            if len(selected_experiments) < 2:
                st.warning("Please select at least 2 experiments to enable comparison.")
            else:
                st.caption(f"Comparing {len(selected_experiments)} experiments")
                
                # Process selected experiments data
                comparison_data = []
                individual_cells_comparison = []
                
                for exp_name in selected_experiments:
                    exp_data = experiment_dict[exp_name]
                    exp_id, exp_name, file_name, loading, active_material, formation_cycles, test_number, electrolyte, substrate, separator, formulation_json, data_json, created_date, porosity, experiment_notes, cutoff_voltage_lower, cutoff_voltage_upper = exp_data
                    
                    try:
                        parsed_data = json.loads(data_json)
                        
                        # Check if this is a multi-cell experiment or single cell
                        if 'cells' in parsed_data:
                            # Multi-cell experiment
                            cells_data = parsed_data['cells']
                            disc_diameter = parsed_data.get('disc_diameter_mm', 15)
                            disc_area_cm2 = np.pi * (disc_diameter / 2 / 10) ** 2
                            
                            experiment_cells = []
                            for cell_data in cells_data:
                                # Skip excluded cells
                                if cell_data.get('excluded', False):
                                    continue
                                    
                                try:
                                    df = pd.read_json(StringIO(cell_data['data_json']))
                                    
                                    # Get project type for efficiency calculation
                                    project_type = "Full Cell"  # Default
                                    if st.session_state.get('current_project_id'):
                                        current_project_id = st.session_state['current_project_id']
                                        project_info = get_project_by_id(current_project_id)
                                        if project_info:
                                            project_type = project_info[3]  # project_type is the 4th field
                                    
                                    cell_summary = calculate_cell_summary(df, cell_data, disc_area_cm2, project_type)
                                    cell_summary['experiment_name'] = exp_name
                                    cell_summary['experiment_date'] = parsed_data.get('experiment_date', created_date)
                                    # Add formulation data to cell summary
                                    if 'formulation' in cell_data:
                                        cell_summary['formulation_json'] = json.dumps(cell_data['formulation'])
                                    experiment_cells.append(cell_summary)
                                    individual_cells_comparison.append(cell_summary)
                                except Exception as e:
                                    continue
                            
                            # Calculate experiment average
                            if experiment_cells:
                                exp_summary = calculate_experiment_average(experiment_cells, exp_name, parsed_data.get('experiment_date', created_date))
                                # Add formulation data to experiment summary (use first cell's formulation as representative)
                                if experiment_cells and 'formulation_json' in experiment_cells[0]:
                                    exp_summary['formulation_json'] = experiment_cells[0]['formulation_json']
                                # Add porosity data to experiment summary (use average from cells)
                                porosity_values = [cell.get('porosity') for cell in experiment_cells if cell.get('porosity') is not None]
                                if porosity_values:
                                    exp_summary['porosity'] = sum(porosity_values) / len(porosity_values)
                                # Add pressed thickness data to experiment summary
                                exp_summary['pressed_thickness'] = parsed_data.get('pressed_thickness')
                                # Add disc diameter data to experiment summary
                                exp_summary['disc_diameter_mm'] = disc_diameter
                                # Add experiment notes to experiment summary
                                exp_summary['experiment_notes'] = experiment_notes
                                comparison_data.append(exp_summary)
                        else:
                            # Legacy single cell experiment
                            df = pd.read_json(StringIO(data_json))
                            
                            # Get project type for efficiency calculation
                            project_type = "Full Cell"  # Default
                            if st.session_state.get('current_project_id'):
                                current_project_id = st.session_state['current_project_id']
                                project_info = get_project_by_id(current_project_id)
                                if project_info:
                                    project_type = project_info[3]  # project_type is the 4th field
                            
                            cell_summary = calculate_cell_summary(df, {
                                'cell_name': test_number or exp_name,
                                'loading': loading,
                                'active_material': active_material,
                                'formation_cycles': formation_cycles,
                                'test_number': test_number
                            }, np.pi * (15 / 2 / 10) ** 2, project_type)  # Default disc size
                            cell_summary['experiment_name'] = exp_name
                            cell_summary['experiment_date'] = created_date
                            # Add formulation data to cell summary
                            if formulation_json:
                                cell_summary['formulation_json'] = formulation_json
                            individual_cells_comparison.append(cell_summary)
                            
                            # Also add as experiment summary (since it's a single cell)
                            exp_summary = cell_summary.copy()
                            exp_summary['cell_name'] = f"{exp_name} (Single Cell)"
                            comparison_data.append(exp_summary)
                            
                    except Exception as e:
                        st.error(f"Error processing experiment {exp_name}: {str(e)}")
                        continue
                
                # Generate comparison visualizations and tables
                if comparison_data:
                    # Create two columns for better layout
                    col1, col2 = st.columns([2, 1])
                    
                    with col1:
                        st.subheader("Comparison Visualization")
                        
                        # Plot selection
                        plot_types = st.multiselect(
                            "Metrics",
                            options=["Reversible Capacity", "Coulombic Efficiency", "First Discharge Capacity", 
                                     "First Cycle Efficiency", "Cycle Life (80%)", "Areal Capacity"],
                            default=["Reversible Capacity", "First Cycle Efficiency"]
                        )
                        
                        if plot_types:
                            # Map plot types to data keys
                            plot_mapping = {
                                "Reversible Capacity": ("reversible_capacity", "mAh/g"),
                                "Coulombic Efficiency": ("coulombic_efficiency", "%"),
                                "First Discharge Capacity": ("first_discharge", "mAh/g"),
                                "First Cycle Efficiency": ("first_efficiency", "%"),
                                "Cycle Life (80%)": ("cycle_life_80", "cycles"),
                                "Areal Capacity": ("areal_capacity", "mAh/cm²")
                            }
                            
                            # Render color/name logic is already collected below for the capacity plots,
                            # but we can grab st.session_state if it exists for consistency, or pass empty.
                            custom_names = st.session_state.get('comp_custom_names', {})
                            
                            # Create interactive comparison plot
                            try:
                                fig = plot_interactive_comparison_metrics(
                                    comparison_data=comparison_data,
                                    selected_metrics=plot_types,
                                    plot_mapping=plot_mapping,
                                    custom_names=custom_names
                                )
                                st.plotly_chart(fig, use_container_width=True)
                            except Exception as e:
                                st.error(f"Error generating comparison plot: {str(e)}")
                        else:
                            st.warning("Please select at least one metric to compare.")
                    
                    with col2:
                        st.subheader("Quick Stats")
                        
                        st.metric("Experiments", len(selected_experiments))
                        st.metric("Total Cells", len(individual_cells_comparison))
                        
                        # Show best performer for first selected metric if applicable
                        if plot_types:
                            primary_metric = plot_types[0]
                            data_key, unit = plot_mapping[primary_metric]
                            
                            exp_names = []
                            values = []
                            for exp in comparison_data:
                                val = exp.get(data_key)
                                if val is not None:
                                    exp_names.append(exp['experiment_name'])
                                    values.append(val)
                                    
                            if values and exp_names:
                                best_idx = np.argmax(values)
                                best_exp = exp_names[best_idx]
                                best_value = values[best_idx]
                                st.metric(f"Best {primary_metric}", f"{best_value:.2f} {unit}")
                    
                    # Capacity comparison plot section
                    st.subheader("Capacity Data Comparison")
                    
                    # Prepare experiment data for plotting
                    experiments_plot_data = []
                    for exp_name in selected_experiments:
                        exp_data = experiment_dict[exp_name]
                        exp_id, exp_name, file_name, loading, active_material, formation_cycles, test_number, electrolyte, substrate, separator, formulation_json, data_json, created_date, porosity, experiment_notes, cutoff_voltage_lower, cutoff_voltage_upper = exp_data
                        
                        try:
                            parsed_data = json.loads(data_json)
                            dfs = []
                            
                            # Check if this is a multi-cell experiment or single cell
                            if 'cells' in parsed_data:
                                # Multi-cell experiment - cells_data is a list
                                cells_data = parsed_data['cells']
                                
                                # Process each cell in the experiment
                                for cell_data in cells_data:
                                    if cell_data.get('excluded', False):
                                        continue  # Skip excluded cells
                                    
                                    if 'data_json' in cell_data:
                                        df = pd.read_json(StringIO(cell_data['data_json']))
                                        test_num = cell_data.get('test_number', cell_data.get('testnum', f'Cell {len(dfs)+1}'))
                                        dfs.append({
                                            'df': df,
                                            'testnum': test_num,
                                            'loading': cell_data.get('loading', loading),
                                            'active_material': cell_data.get('active_material', active_material),
                                            'formation_cycles': cell_data.get('formation_cycles', formation_cycles)
                                        })
                            else:
                                # Single cell experiment - data_json is at the top level
                                df = pd.read_json(StringIO(data_json))
                                test_num = test_number or f'Cell 1'
                                dfs.append({
                                    'df': df,
                                    'testnum': test_num,
                                    'loading': loading,
                                    'active_material': active_material,
                                    'formation_cycles': formation_cycles
                                })
                            
                            if dfs:  # Only add if we have valid data
                                experiments_plot_data.append({
                                    'experiment_name': exp_name,
                                    'dfs': dfs
                                })
                        except Exception as e:
                            st.warning(f"Could not load plotting data for {exp_name}: {str(e)}")
                            # Add debug info
                            st.info(f"Debug info for {exp_name}: data_json type = {type(data_json)}, length = {len(str(data_json)) if data_json else 'None'}")
                            continue
                    
                    if len(experiments_plot_data) >= 1:
                        # Add toggle for interactive vs static
                        col_toggle, col_spacer = st.columns([3, 7])
                        with col_toggle:
                            plot_style_comp = st.radio(
                                "Plot Style",
                                options=["📊 Interactive (Plotly)", "📉 Static (Matplotlib)"],
                                index=0,
                                horizontal=True,
                                help="Interactive plots allow zooming, panning, and hovering for details.",
                                key="plot_style_comps"
                            )
                        use_interactive_comp = plot_style_comp.startswith("📊")

                        # Render plot options
                        show_lines, show_efficiency_lines, remove_last_cycle, show_graph_title, show_average_performance, avg_line_toggles, remove_markers, hide_legend, cycle_filter, y_axis_limits, custom_title, excluded_from_average = render_comparison_plot_options(experiments_plot_data)
                        
                        # Render color customization UI
                        custom_colors = render_comparison_color_customization(
                            experiments_plot_data, 
                            show_average_performance
                        )
                        
                        # Render name customization UI
                        custom_names = render_comparison_name_customization(
                            experiments_plot_data, 
                            show_average_performance
                        )
                        
                        # Generate the comparison plot
                        try:
                            if use_interactive_comp:
                                comparison_fig = plot_interactive_comparison_capacity(
                                    experiments_plot_data,
                                    show_lines,
                                    show_efficiency_lines,
                                    remove_last_cycle,
                                    show_graph_title,
                                    show_average_performance,
                                    avg_line_toggles,
                                    hide_legend,
                                    cycle_filter,
                                    custom_colors,
                                    y_axis_limits,
                                    custom_names,
                                    custom_title,
                                    excluded_from_average
                                )
                                # Display the interactive plot
                                st.plotly_chart(comparison_fig, use_container_width=True)
                                st.info("💡 **Tip**: Hover over data points for cycle, capacity, and retention details. Retention uses the first valid post-formation cycle for each cell and skips clearly anomalous baseline cycles when needed.")
                            else:
                                comparison_fig = plot_comparison_capacity_graph(
                                    experiments_plot_data,
                                    show_lines,
                                    show_efficiency_lines,
                                    remove_last_cycle,
                                    show_graph_title,
                                    show_average_performance,
                                    avg_line_toggles,
                                    remove_markers,
                                    hide_legend,
                                    cycle_filter,
                                    custom_colors,
                                    y_axis_limits,
                                    custom_names,
                                    custom_title,
                                    excluded_from_average
                                )
                                
                                # Display the plot
                                st.pyplot(comparison_fig)
                                
                                # Export option for the comparison plot
                                buf = io.BytesIO()
                                comparison_fig.savefig(buf, format='png', dpi=300, bbox_inches='tight')
                                buf.seek(0)
                                st.download_button(
                                    label="Download Plot",
                                    data=buf,
                                    file_name="capacity_comparison_plot.png",
                                    mime="image/png"
                                )
                            
                        except Exception as e:
                            st.error(f"Error generating comparison plot: {str(e)}")
                    else:
                        st.warning("No valid experiment data available for capacity plotting.")
                    
                    # Summary comparison table
                    st.subheader("Comparison Summary Table")
                    
                    # Table filter options
                    with st.expander("Table Options", expanded=False):
                        show_columns = st.multiselect(
                            "Select metrics to display",
                            ["Experiment", "Reversible Capacity (mAh/g)", "Coulombic Efficiency (%)", 
                             "First Discharge (mAh/g)", "First Efficiency (%)", 
                             "Cycle Life (80%)", "Areal Capacity (mAh/cm²)", "Active Material (%)", "Date"],
                            default=["Experiment", "Reversible Capacity (mAh/g)", "Coulombic Efficiency (%)", 
                                   "First Discharge (mAh/g)", "Cycle Life (80%)"]
                        )
                    
                    if show_columns:
                        # Create comparison DataFrame
                        comparison_df_data = []
                        for exp in comparison_data:
                            row = {
                                'Experiment': exp['experiment_name'],
                                'Reversible Capacity (mAh/g)': exp.get('reversible_capacity', 'N/A'),
                                'Coulombic Efficiency (%)': exp.get('coulombic_efficiency', 'N/A'),
                                'First Discharge (mAh/g)': exp.get('first_discharge', 'N/A'),
                                'First Efficiency (%)': exp.get('first_efficiency', 'N/A'),
                                'Cycle Life (80%)': exp.get('cycle_life_80', 'N/A'),
                                'Areal Capacity (mAh/cm²)': exp.get('areal_capacity', 'N/A'),
                                'Active Material (%)': f"{exp.get('active_material', 'N/A'):.1f}" if exp.get('active_material') is not None and exp.get('active_material') != 'N/A' else 'N/A',
                                'Date': exp.get('experiment_date', 'N/A')
                            }
                            comparison_df_data.append(row)
                        
                        comparison_df = pd.DataFrame(comparison_df_data)
                        
                        # Filter to selected columns
                        available_columns = [col for col in show_columns if col in comparison_df.columns]
                        if available_columns:
                            filtered_df = comparison_df[available_columns]
                            st.dataframe(filtered_df, use_container_width=True)
                            
                            # Export option for table
                            csv_data = filtered_df.to_csv(index=False)
                            st.download_button(
                                label="Download Table (CSV)",
                                data=csv_data,
                                file_name="experiment_comparison.csv",
                                mime="text/csv"
                            )
                        else:
                            st.warning("No valid columns selected for display.")
                    else:
                        st.warning("Please select at least one column to display.")
                    
                    # Individual cells comparison (optional detailed view)
                    if individual_cells_comparison:
                        with st.expander("Individual Cells Detailed Comparison", expanded=False):
                            
                            # Create individual cells DataFrame
                            individual_df_data = []
                            for cell in individual_cells_comparison:
                                row = {
                                    'Experiment': cell.get('experiment_name', 'Unknown'),
                                    'Cell Name': cell['cell_name'],
                                    'Reversible Capacity (mAh/g)': cell.get('reversible_capacity', 'N/A'),
                                    'Coulombic Efficiency (%)': cell.get('coulombic_efficiency', 'N/A'),
                                    'First Discharge (mAh/g)': cell.get('first_discharge', 'N/A'),
                                    'First Efficiency (%)': cell.get('first_efficiency', 'N/A'),
                                    'Cycle Life (80%)': cell.get('cycle_life_80', 'N/A'),
                                    'Areal Capacity (mAh/cm²)': cell.get('areal_capacity', 'N/A'),
                                    'Loading (mg)': cell.get('loading', 'N/A')
                                }
                                individual_df_data.append(row)
                            
                            individual_df = pd.DataFrame(individual_df_data)
                            st.dataframe(individual_df, use_container_width=True)
                            
                            # Export option for individual cells
                            individual_csv = individual_df.to_csv(index=False)
                            st.download_button(
                                label="Download Individual Cells (CSV)",
                                data=individual_csv,
                                file_name="individual_cells_comparison.csv",
                                mime="text/csv"
                            )
                    else:
                        st.error("No valid data found for selected experiments.")
                else:
                    st.warning("Could not load comparison data for the selected experiments. Try re-selecting them.")

            # Formulation-Based Comparison Section
            st.subheader("Formulation-Based Comparison")
            
            # Get formulation summary for the project
            formulation_summary = get_formulation_summary(current_project_id)
            
            if not formulation_summary:
                st.info("No formulation data found in this project. Add formulations to your experiments to enable formulation-based comparisons.")
            else:
                # Component selection
                component_names = sorted(list(formulation_summary.keys()))
                
                col1, col2 = st.columns([2, 1])
                with col1:
                    selected_component = st.selectbox(
                        "Select formulation component",
                        options=component_names
                    )
                
                with col2:
                    if selected_component:
                        stats = formulation_summary[selected_component]
                        st.metric(f"{selected_component} Range", f"{stats['min']:.1f} - {stats['max']:.1f}%")
                
                if selected_component:
                    # Get experiments with this component
                    matching_experiments = get_experiments_by_formulation_component(
                        current_project_id, selected_component
                    )
                    
                    if not matching_experiments:
                        st.warning(f"No experiments found with {selected_component} in their formulation.")
                    else:
                        # Filter options
                        with st.expander("Filter Options", expanded=False):
                            filter_col1, filter_col2 = st.columns(2)
                            with filter_col1:
                                min_pct = st.number_input(f"Minimum {selected_component} %", min_value=0.0, max_value=100.0, value=0.0, step=0.5)
                            with filter_col2:
                                max_pct = st.number_input(f"Maximum {selected_component} %", min_value=0.0, max_value=100.0, value=100.0, step=0.5)
                        
                        # Apply filters
                        if min_pct > 0 or max_pct < 100:
                            filtered_experiments = get_experiments_by_formulation_component(
                                current_project_id, selected_component,
                                min_percentage=min_pct if min_pct > 0 else None,
                                max_percentage=max_pct if max_pct < 100 else None
                            )
                        else:
                            filtered_experiments = matching_experiments
                        
                        if filtered_experiments:
                            st.caption(f"Found {len(filtered_experiments)} experiments")
                            
                            # Create comparison DataFrame
                            comparison_df = create_formulation_comparison_dataframe(
                                filtered_experiments, selected_component
                            )
                            
                            if not comparison_df.empty:
                                # Initialize excluded metric values in session state
                                # Store as set of tuples: (experiment_name, metric_name)
                                exclusion_key = f'formulation_excluded_{selected_component}_{current_project_id}'
                                if exclusion_key not in st.session_state:
                                    st.session_state[exclusion_key] = set()
                                
                                excluded_values = st.session_state[exclusion_key]
                                
                                # Initialize overwritten values in session state
                                # Store as dict: {(experiment_name, metric_name): overwritten_value}
                                overwrite_key = f'formulation_overwritten_{selected_component}_{current_project_id}'
                                if overwrite_key not in st.session_state:
                                    st.session_state[overwrite_key] = {}
                                
                                overwritten_values = st.session_state[overwrite_key]
                                
                                # Visualization: Performance vs Component Percentage
                                st.subheader(f"Performance vs {selected_component} Percentage")
                                
                                # Metric selection for Y-axis
                                metric_options = {
                                'Reversible Capacity (mAh/g)': 'Reversible Capacity (mAh/g)',
                                'First Discharge (mAh/g)': 'First Discharge (mAh/g)',
                                'First Efficiency (%)': 'First Efficiency (%)',
                                'Cycle Life': 'Cycle Life',
                                'Porosity (%)': 'Porosity (%)'
                                }
                                
                                selected_metric = st.selectbox("Performance metric", options=list(metric_options.keys()))
                                
                                metric_col = metric_options[selected_metric]
                                
                                # Create a copy for plotting with overwritten values applied
                                plot_df = comparison_df.copy()
                                
                                # Apply overwritten values
                                for idx, row in plot_df.iterrows():
                                    exp_name = row['Experiment']
                                    overwrite_key_tuple = (exp_name, metric_col)
                                    if overwrite_key_tuple in overwritten_values:
                                        plot_df.at[idx, metric_col] = overwritten_values[overwrite_key_tuple]
                                
                                # Filter out excluded metric values and rows with missing data
                                def is_excluded(exp_name, metric):
                                    return (exp_name, metric) in excluded_values
                                
                                plot_df = plot_df[
                                    (~plot_df.apply(lambda row: is_excluded(row['Experiment'], metric_col), axis=1)) &
                                    plot_df[metric_col].notna() & 
                                    (plot_df[metric_col] != 'N/A')
                                ].copy()
                                
                                if not plot_df.empty:
                                    # Convert metric column to numeric if needed
                                    plot_df[metric_col] = pd.to_numeric(plot_df[metric_col], errors='coerce')
                                    plot_df = plot_df.dropna(subset=[metric_col, 'Component %'])
                                    
                                    if not plot_df.empty:
                                        # Create scatter plot
                                        fig, ax = plt.subplots(figsize=(10, 6))
                                        
                                        scatter = ax.scatter(
                                            plot_df['Component %'],
                                            plot_df[metric_col],
                                            s=100,
                                            alpha=0.6,
                                            c=plot_df['Component %'],
                                            cmap='viridis',
                                            edgecolors='black',
                                            linewidths=1
                                        )
                                        
                                        # Add trend line
                                        z = np.polyfit(plot_df['Component %'], plot_df[metric_col], 1)
                                        p = np.poly1d(z)
                                        ax.plot(
                                            plot_df['Component %'],
                                            p(plot_df['Component %']),
                                            "r--",
                                            alpha=0.5,
                                            label=f'Trend: y = {z[0]:.2f}x + {z[1]:.2f}'
                                        )
                                        
                                        ax.set_xlabel(f'{selected_component} (%)', fontsize=12, fontweight='bold')
                                        ax.set_ylabel(selected_metric, fontsize=12, fontweight='bold')
                                        ax.set_title(f'{selected_metric} vs {selected_component} Percentage', fontsize=14, fontweight='bold')
                                        ax.grid(True, alpha=0.3)
                                        ax.legend()
                                        
                                        # Add colorbar
                                        cbar = plt.colorbar(scatter, ax=ax)
                                        cbar.set_label(f'{selected_component} (%)', fontsize=10)
                                        
                                        plt.tight_layout()
                                        st.pyplot(fig)
                                        
                                        # Export plot
                                        buf = io.BytesIO()
                                        fig.savefig(buf, format='png', dpi=300, bbox_inches='tight')
                                        buf.seek(0)
                                        st.download_button(
                                            label=f"Download Plot",
                                            data=buf,
                                            file_name=f"{selected_metric.replace(' ', '_')}_vs_{selected_component.replace(' ', '_')}.png",
                                            mime="image/png"
                                        )
                                    else:
                                        st.warning(f"No valid data points for {selected_metric} comparison.")
                                else:
                                    st.warning(f"No data available for {selected_metric} in the selected experiments.")
                                
                                # Simple comparison table - clean and user-friendly
                                st.subheader("Detailed Comparison Table")

                                # Table options
                                with st.expander("🔧 Table Options", expanded=False):
                                    show_columns = st.multiselect(
                                        "Select metrics to display:",
                                        ['Experiment', 'Component %', 'Reversible Capacity (mAh/g)',
                                         'First Discharge (mAh/g)', 'First Efficiency (%)',
                                         'Cycle Life', 'Loading (mg)', 'Active Material (%)',
                                         'Porosity (%)', 'Date'],
                                        default=['Experiment', 'Component %', 'Reversible Capacity (mAh/g)',
                                                'First Discharge (mAh/g)', 'Cycle Life'],
                                        help="Choose which columns to display in the comparison table"
                                    )

                                if show_columns:
                                    # Filter to selected columns and sort by component percentage
                                    available_columns = [col for col in show_columns if col in comparison_df.columns]
                                    if available_columns:
                                        filtered_df = comparison_df[available_columns].sort_values('Component %')
                                        st.dataframe(filtered_df, use_container_width=True)

                                        # Export option for table
                                        csv_data = filtered_df.to_csv(index=False)
                                        st.download_button(
                                            label="Download Table (CSV)",
                                            data=csv_data,
                                            file_name=f"{selected_component.replace(' ', '_')}_comparison.csv",
                                            mime="text/csv"
                                        )
                                    else:
                                        st.warning("No valid columns selected for display.")
                                else:
                                    st.warning("Please select at least one column to display.")

                                # Optional individual cells comparison
                                with st.expander("🔬 Individual Cells Detailed Comparison", expanded=False):
                                    st.markdown("**All individual cells from selected experiments:**")

                                    # Create individual cells DataFrame
                                    individual_df_data = []
                                    for exp in filtered_experiments:
                                        exp_id, project_id, cell_name, file_name, loading, active_material, \
                                        formation_cycles, test_number, electrolyte, substrate, separator, \
                                        formulation_json, data_json, solids_content, pressed_thickness, \
                                        experiment_notes, created_date, porosity = exp
                                        
                                        exp_name = cell_name
                                        component_pct = extract_formulation_component_from_experiment(exp, selected_component)

                                        # Extract cell data from data_json
                                        reversible_capacity = None
                                        first_discharge = None
                                        first_efficiency = None
                                        cycle_life = None
                                        cell_loading = loading
                                        cell_active_material = active_material
                                        
                                        if data_json:
                                            try:
                                                parsed_data = json.loads(data_json)
                                                formation_cycles = formation_cycles or 4
                                                
                                                if 'cells' in parsed_data:
                                                    # Multi-cell experiment - process each cell
                                                    for cell in parsed_data['cells']:
                                                        if cell.get('excluded', False):
                                                            continue
                                                        
                                                        # Reset metrics for each cell
                                                        cell_reversible_capacity = None
                                                        cell_first_discharge = None
                                                        cell_first_efficiency = None
                                                        cell_cycle_life = None
                                                        cell_loading = loading  # Start with experiment-level loading
                                                        cell_active_material = active_material  # Start with experiment-level active_material
                                                        
                                                        # Extract loading and active_material from cell data
                                                        if cell.get('loading') is not None:
                                                            cell_loading = cell.get('loading')
                                                        if cell.get('active_material') is not None:
                                                            cell_active_material = cell.get('active_material')
                                                        
                                                        if 'data_json' in cell:
                                                            try:
                                                                df = pd.read_json(StringIO(cell['data_json']))
                                                                
                                                                # Get first discharge capacity (max of first 3 cycles)
                                                                if 'Q Dis (mAh/g)' in df.columns:
                                                                    first_three = df['Q Dis (mAh/g)'].head(3).tolist()
                                                                    if first_three:
                                                                        cell_first_discharge = max(first_three)
                                                                    
                                                                    # Get first post-formation cycle (reversible capacity)
                                                                    if len(df) > formation_cycles:
                                                                        cell_reversible_capacity = df['Q Dis (mAh/g)'].iloc[formation_cycles]
                                                                
                                                                # Get first cycle efficiency
                                                                if 'Efficiency (-)' in df.columns and len(df) > 0:
                                                                    first_eff = df['Efficiency (-)'].iloc[0]
                                                                    if first_eff is not None:
                                                                        try:
                                                                            cell_first_efficiency = float(first_eff) * 100
                                                                        except (ValueError, TypeError):
                                                                            pass
                                                                
                                                                # Calculate cycle life (80% threshold)
                                                                if 'Q Dis (mAh/g)' in df.columns and len(df) > formation_cycles:
                                                                    post_formation = df.iloc[formation_cycles:]
                                                                    if not post_formation.empty:
                                                                        initial_capacity = post_formation['Q Dis (mAh/g)'].iloc[0]
                                                                        if initial_capacity > 0:
                                                                            threshold = 0.8 * initial_capacity
                                                                            below_threshold = post_formation[post_formation['Q Dis (mAh/g)'] < threshold]
                                                                            if not below_threshold.empty:
                                                                                cell_cycle_life = int(post_formation.index[below_threshold.index[0]])
                                                            except Exception:
                                                                pass
                                                        
                                                        # Add row for this cell
                                                        row = {
                                                            'Experiment': exp_name,
                                                            'Component %': f"{component_pct:.1f}%" if component_pct else 'N/A',
                                                            'Reversible Capacity (mAh/g)': cell_reversible_capacity if cell_reversible_capacity is not None else 'N/A',
                                                            'First Discharge (mAh/g)': cell_first_discharge if cell_first_discharge is not None else 'N/A',
                                                            'First Efficiency (%)': f"{cell_first_efficiency:.2f}%" if cell_first_efficiency is not None else 'N/A',
                                                            'Cycle Life': cell_cycle_life if cell_cycle_life is not None else 'N/A',
                                                            'Loading (mg)': cell_loading if cell_loading is not None else 'N/A',
                                                            'Active Material (%)': cell_active_material if cell_active_material is not None else 'N/A',
                                                            'Porosity (%)': f"{porosity * 100:.2f}%" if porosity is not None else 'N/A'
                                                        }
                                                        individual_df_data.append(row)
                                                else:
                                                    # Legacy single cell experiment
                                                    try:
                                                        df = pd.read_json(StringIO(data_json))
                                                        formation_cycles = formation_cycles or 4
                                                        
                                                        if 'Q Dis (mAh/g)' in df.columns:
                                                            # First discharge (max of first 3)
                                                            first_three = df['Q Dis (mAh/g)'].head(3).tolist()
                                                            if first_three:
                                                                first_discharge = max(first_three)
                                                            
                                                            # Reversible capacity
                                                            if len(df) > formation_cycles:
                                                                reversible_capacity = df['Q Dis (mAh/g)'].iloc[formation_cycles]
                                                        
                                                        # First cycle efficiency
                                                        if 'Efficiency (-)' in df.columns and len(df) > 0:
                                                            first_eff = df['Efficiency (-)'].iloc[0]
                                                            if first_eff is not None:
                                                                try:
                                                                    first_efficiency = float(first_eff) * 100
                                                                except (ValueError, TypeError):
                                                                    pass
                                                        
                                                        # Cycle life
                                                        if 'Q Dis (mAh/g)' in df.columns and len(df) > formation_cycles:
                                                            post_formation = df.iloc[formation_cycles:]
                                                            if not post_formation.empty:
                                                                initial_capacity = post_formation['Q Dis (mAh/g)'].iloc[0]
                                                                if initial_capacity > 0:
                                                                    threshold = 0.8 * initial_capacity
                                                                    below_threshold = post_formation[post_formation['Q Dis (mAh/g)'] < threshold]
                                                                    if not below_threshold.empty:
                                                                        cycle_life = int(post_formation.index[below_threshold.index[0]])
                                                    except Exception:
                                                        pass
                                                    
                                                    # Add row for single cell experiment
                                                    row = {
                                                        'Experiment': exp_name,
                                                        'Component %': f"{component_pct:.1f}%" if component_pct else 'N/A',
                                                        'Reversible Capacity (mAh/g)': reversible_capacity if reversible_capacity is not None else 'N/A',
                                                        'First Discharge (mAh/g)': first_discharge if first_discharge is not None else 'N/A',
                                                        'First Efficiency (%)': f"{first_efficiency:.2f}%" if first_efficiency is not None else 'N/A',
                                                        'Cycle Life': cycle_life if cycle_life is not None else 'N/A',
                                                        'Loading (mg)': cell_loading if cell_loading is not None else 'N/A',
                                                        'Active Material (%)': cell_active_material if cell_active_material is not None else 'N/A',
                                                        'Porosity (%)': f"{porosity * 100:.2f}%" if porosity is not None else 'N/A'
                                                    }
                                                    individual_df_data.append(row)
                                            except Exception:
                                                # If data_json parsing fails, still add a row with available data
                                                row = {
                                                    'Experiment': exp_name,
                                                    'Component %': f"{component_pct:.1f}%" if component_pct else 'N/A',
                                                    'Reversible Capacity (mAh/g)': 'N/A',
                                                    'First Discharge (mAh/g)': 'N/A',
                                                    'First Efficiency (%)': 'N/A',
                                                    'Cycle Life': 'N/A',
                                                    'Loading (mg)': cell_loading if cell_loading is not None else 'N/A',
                                                    'Active Material (%)': cell_active_material if cell_active_material is not None else 'N/A',
                                                    'Porosity (%)': f"{porosity * 100:.2f}%" if porosity is not None else 'N/A'
                                                }
                                                individual_df_data.append(row)
                                        else:
                                            # No data_json available, add row with available data
                                            row = {
                                                'Experiment': exp_name,
                                                'Component %': f"{component_pct:.1f}%" if component_pct else 'N/A',
                                                'Reversible Capacity (mAh/g)': 'N/A',
                                                'First Discharge (mAh/g)': 'N/A',
                                                'First Efficiency (%)': 'N/A',
                                                'Cycle Life': 'N/A',
                                                'Loading (mg)': cell_loading if cell_loading is not None else 'N/A',
                                                'Active Material (%)': cell_active_material if cell_active_material is not None else 'N/A',
                                                'Porosity (%)': f"{porosity * 100:.2f}%" if porosity is not None else 'N/A'
                                        }
                                        individual_df_data.append(row)

                                    if individual_df_data:
                                        individual_df = pd.DataFrame(individual_df_data)
                                        st.dataframe(individual_df, use_container_width=True)

                                        # Export option for individual cells
                                        individual_csv = individual_df.to_csv(index=False)
                                        st.download_button(
                                            label="Download Individual Cells (CSV)",
                                            data=individual_csv,
                                            file_name=f"{selected_component.replace(' ', '_')}_individual_cells.csv",
                                            mime="text/csv"
                                        )
                                    else:
                                        st.info("No individual cell data available for the selected experiments.")

                                # Grouped analysis
                                st.subheader("Grouped Analysis")
                                st.caption(f"Experiments grouped by {selected_component} percentage ranges")

                                grouped = group_experiments_by_formulation_range(
                                    filtered_experiments, selected_component, range_size=5.0
                                )

                                if grouped:
                                    for range_label in sorted(grouped.keys(), key=lambda x: float(x.split('-')[0])):
                                        experiments_in_range = grouped[range_label]
                                        with st.expander(f"{range_label} ({len(experiments_in_range)} experiments)"):
                                            for exp in experiments_in_range:
                                                exp_name = exp[2]  # cell_name
                                                component_pct = extract_formulation_component_from_experiment(exp, selected_component)
                                                st.write(f"• **{exp_name}**: {component_pct:.1f}% {selected_component}")
                            else:
                                st.warning("Could not create comparison table. Check that experiments have valid formulation and performance data.")
                        else:
                            st.warning(f"No experiments found with {selected_component} in the specified range ({min_pct:.1f}% - {max_pct:.1f}%).")

# --- Master Table Tab ---
if tab_master and current_project_id:
    with tab_master:
        st.header("Master Table")
        current_project_name = st.session_state.get('current_project_name', 'Selected Project')
        st.markdown(f"**Project:** {current_project_name}")
        
        # Show update notification if calculations were recently updated
        if st.session_state.get('calculations_updated', False):
            update_time = st.session_state.get('update_timestamp')
            if update_time:
                time_str = update_time.strftime("%H:%M:%S")
                st.info(f"Data refreshed at {time_str} - Master table shows updated calculations!")
        
        st.markdown("---")
        
        # Get all experiments data for this project
        all_experiments_data = get_all_project_experiments_data(current_project_id)
        
        if not all_experiments_data:
            st.info("No experiments found in this project. Create experiments to see master table data.")
        else:
            # Process experiment data
            experiment_summaries = []
            individual_cells = []
            
            for exp_data in all_experiments_data:
                exp_id, exp_name, file_name, loading, active_material, formation_cycles, test_number, electrolyte, substrate, separator, formulation_json, data_json, created_date, porosity, experiment_notes, cutoff_voltage_lower, cutoff_voltage_upper = exp_data
                
                # If substrate or separator are None from database, we'll extract them from JSON data
                extracted_substrate = substrate
                extracted_separator = separator
                
                try:
                    parsed_data = json.loads(data_json)
                    
                    # Check if this is a multi-cell experiment or single cell
                    if 'cells' in parsed_data:
                        # Multi-cell experiment
                        cells_data = parsed_data['cells']
                        disc_diameter = parsed_data.get('disc_diameter_mm', 15)
                        disc_area_cm2 = np.pi * (disc_diameter / 2 / 10) ** 2
                        
                        experiment_cells = []
                        for cell_data in cells_data:
                            if cell_data.get('excluded', False):
                                continue
                            try:
                                df = pd.read_json(StringIO(cell_data['data_json']))
                                
                                # Get project type for efficiency calculation
                                project_type = "Full Cell"  # Default
                                if st.session_state.get('current_project_id'):
                                    current_project_id = st.session_state['current_project_id']
                                    project_info = get_project_by_id(current_project_id)
                                    if project_info:
                                        project_type = project_info[3]  # project_type is the 4th field
                                
                                cell_summary = calculate_cell_summary(df, cell_data, disc_area_cm2, project_type)
                                cell_summary['experiment_name'] = exp_name
                                cell_summary['experiment_date'] = parsed_data.get('experiment_date', created_date)
                                # Add pressed thickness data from experiment
                                cell_summary['pressed_thickness'] = parsed_data.get('pressed_thickness')
                                # Add disc diameter data from experiment
                                cell_summary['disc_diameter_mm'] = disc_diameter
                                # Add electrolyte, substrate, and separator data to cell summary
                                cell_summary['electrolyte'] = cell_data.get('electrolyte', 'N/A')
                                cell_summary['substrate'] = cell_data.get('substrate', 'N/A')
                                cell_summary['separator'] = cell_data.get('separator', 'N/A')
                                # Add cutoff voltages to cell summary (fall back to experiment-level from DB if not in cell)
                                cell_summary['cutoff_voltage_lower'] = cell_data.get('cutoff_voltage_lower') if cell_data.get('cutoff_voltage_lower') is not None else cutoff_voltage_lower
                                cell_summary['cutoff_voltage_upper'] = cell_data.get('cutoff_voltage_upper') if cell_data.get('cutoff_voltage_upper') is not None else cutoff_voltage_upper
                                # Add formulation data to cell summary
                                if 'formulation' in cell_data:
                                    cell_summary['formulation_json'] = json.dumps(cell_data['formulation'])
                                # Add porosity data from cell_data if available, or recalculate if missing
                                if 'porosity' in cell_data and cell_data['porosity'] is not None and cell_data['porosity'] > 0:
                                    cell_summary['porosity'] = cell_data['porosity']
                                else:
                                    # Recalculate porosity if missing or invalid
                                    try:
                                        from porosity_calculations import calculate_porosity_from_experiment_data
                                        if (cell_data.get('loading') and 
                                            disc_diameter and 
                                            parsed_data.get('pressed_thickness') and 
                                            cell_data.get('formulation')):
                                            
                                            porosity_data = calculate_porosity_from_experiment_data(
                                                disc_mass_mg=cell_data['loading'],
                                                disc_diameter_mm=disc_diameter,
                                                pressed_thickness_um=parsed_data['pressed_thickness'],
                                                formulation=cell_data['formulation']
                                            )
                                            cell_summary['porosity'] = porosity_data['porosity']
                                        else:
                                            cell_summary['porosity'] = None
                                    except Exception:
                                        cell_summary['porosity'] = None
                                experiment_cells.append(cell_summary)
                                individual_cells.append(cell_summary)
                            except Exception as e:
                                continue
                        
                        # Calculate experiment average
                        if experiment_cells:
                            exp_summary = calculate_experiment_average(experiment_cells, exp_name, parsed_data.get('experiment_date', created_date))
                            # Add formulation data to experiment summary (use first cell's formulation as representative)
                            if experiment_cells and 'formulation_json' in experiment_cells[0]:
                                exp_summary['formulation_json'] = experiment_cells[0]['formulation_json']
                            # Add porosity data to experiment summary (use average from cells)
                            porosity_values = [cell.get('porosity') for cell in experiment_cells if cell.get('porosity') is not None]
                            if porosity_values:
                                exp_summary['porosity'] = sum(porosity_values) / len(porosity_values)
                            # Add pressed thickness data to experiment summary
                            exp_summary['pressed_thickness'] = parsed_data.get('pressed_thickness')
                            # Add disc diameter data to experiment summary
                            exp_summary['disc_diameter_mm'] = disc_diameter
                            # Add electrolyte, substrate, and separator data to experiment summary (use first cell's values as representative)
                            if experiment_cells:
                                exp_summary['electrolyte'] = experiment_cells[0].get('electrolyte', 'N/A')
                                exp_summary['substrate'] = experiment_cells[0].get('substrate', 'N/A')
                                exp_summary['separator'] = experiment_cells[0].get('separator', 'N/A')
                                # Add cutoff voltages to experiment summary (use first cell or fall back to experiment-level from DB)
                                exp_summary['cutoff_voltage_lower'] = experiment_cells[0].get('cutoff_voltage_lower') if experiment_cells[0].get('cutoff_voltage_lower') is not None else cutoff_voltage_lower
                                exp_summary['cutoff_voltage_upper'] = experiment_cells[0].get('cutoff_voltage_upper') if experiment_cells[0].get('cutoff_voltage_upper') is not None else cutoff_voltage_upper
                            # Add experiment notes to experiment summary
                            exp_summary['experiment_notes'] = experiment_notes
                            experiment_summaries.append(exp_summary)
                    
                    else:
                        # Legacy single cell experiment
                        df = pd.read_json(StringIO(data_json))
                        
                        # Get project type for efficiency calculation
                        project_type = "Full Cell"  # Default
                        if st.session_state.get('current_project_id'):
                            current_project_id = st.session_state['current_project_id']
                            project_info = get_project_by_id(current_project_id)
                            if project_info:
                                project_type = project_info[3]  # project_type is the 4th field
                        
                        cell_summary = calculate_cell_summary(df, {
                            'cell_name': test_number or exp_name,
                            'loading': loading,
                            'active_material': active_material,
                            'formation_cycles': formation_cycles,
                            'test_number': test_number
                        }, np.pi * (15 / 2 / 10) ** 2, project_type)  # Default disc size
                        cell_summary['experiment_name'] = exp_name
                        cell_summary['experiment_date'] = created_date
                        # Add electrolyte, substrate, and separator data to cell summary
                        cell_summary['electrolyte'] = electrolyte if electrolyte else 'N/A'
                        cell_summary['substrate'] = extracted_substrate if extracted_substrate else 'N/A'
                        cell_summary['separator'] = extracted_separator if extracted_separator else 'N/A'
                        # Add cutoff voltages to cell summary
                        cell_summary['cutoff_voltage_lower'] = cutoff_voltage_lower
                        cell_summary['cutoff_voltage_upper'] = cutoff_voltage_upper
                        # Add formulation data to cell summary
                        if formulation_json:
                            cell_summary['formulation_json'] = formulation_json
                        # Add porosity data from database if available, or recalculate if missing
                        if porosity is not None and porosity > 0:
                            cell_summary['porosity'] = porosity
                        else:
                            # Recalculate porosity for legacy experiments if missing or invalid
                            try:
                                from porosity_calculations import calculate_porosity_from_experiment_data
                                if (loading and 
                                    disc_diameter and 
                                    formulation_json):
                                    
                                    # Parse formulation data
                                    formulation_data = json.loads(formulation_json)
                                    
                                    # Get pressed thickness from database
                                    conn = get_db_connection()
                                    cursor = conn.cursor()
                                    cursor.execute('SELECT pressed_thickness FROM cell_experiments WHERE id = ?', (exp_id,))
                                    result = cursor.fetchone()
                                    conn.close()
                                    pressed_thickness = result[0] if result and result[0] is not None else None
                                    
                                    if pressed_thickness:
                                        porosity_data = calculate_porosity_from_experiment_data(
                                            disc_mass_mg=loading,
                                            disc_diameter_mm=disc_diameter,
                                            pressed_thickness_um=pressed_thickness,
                                            formulation=formulation_data
                                        )
                                        cell_summary['porosity'] = porosity_data['porosity']
                                    else:
                                        cell_summary['porosity'] = None
                                else:
                                    cell_summary['porosity'] = None
                            except Exception:
                                cell_summary['porosity'] = None
                        # Add pressed thickness data from database if available
                        conn = get_db_connection()
                        cursor = conn.cursor()
                        cursor.execute('SELECT pressed_thickness FROM cell_experiments WHERE id = ?', (exp_id,))
                        result = cursor.fetchone()
                        conn.close()
                        if result and result[0] is not None:
                            cell_summary['pressed_thickness'] = result[0]
                        # Add disc diameter data (default to 15mm for legacy experiments)
                        cell_summary['disc_diameter_mm'] = 15
                        individual_cells.append(cell_summary)
                        
                        # Also add as experiment summary (since it's a single cell)
                        exp_summary = cell_summary.copy()
                        exp_summary['cell_name'] = f"{exp_name} (Single Cell)"
                        # Add electrolyte, substrate, and separator data to experiment summary
                        exp_summary['electrolyte'] = electrolyte if electrolyte else 'N/A'
                        exp_summary['substrate'] = extracted_substrate if extracted_substrate else 'N/A'
                        exp_summary['separator'] = extracted_separator if extracted_separator else 'N/A'
                        # Add cutoff voltages to experiment summary
                        exp_summary['cutoff_voltage_lower'] = cutoff_voltage_lower
                        exp_summary['cutoff_voltage_upper'] = cutoff_voltage_upper
                        # Add porosity data from database if available, or recalculate if missing
                        if porosity is not None and porosity > 0:
                            exp_summary['porosity'] = porosity
                        else:
                            # Recalculate porosity for legacy experiments if missing or invalid
                            try:
                                from porosity_calculations import calculate_porosity_from_experiment_data
                                if (loading and 
                                    disc_diameter and 
                                    formulation_json):
                                    
                                    # Parse formulation data
                                    formulation_data = json.loads(formulation_json)
                                    
                                    # Get pressed thickness from database
                                    conn = get_db_connection()
                                    cursor = conn.cursor()
                                    cursor.execute('SELECT pressed_thickness FROM cell_experiments WHERE id = ?', (exp_id,))
                                    result = cursor.fetchone()
                                    conn.close()
                                    pressed_thickness = result[0] if result and result[0] is not None else None
                                    
                                    if pressed_thickness:
                                        porosity_data = calculate_porosity_from_experiment_data(
                                            disc_mass_mg=loading,
                                            disc_diameter_mm=disc_diameter,
                                            pressed_thickness_um=pressed_thickness,
                                            formulation=formulation_data
                                        )
                                        exp_summary['porosity'] = porosity_data['porosity']
                                    else:
                                        exp_summary['porosity'] = None
                                else:
                                    exp_summary['porosity'] = None
                            except Exception:
                                exp_summary['porosity'] = None
                        # Add pressed thickness data from database if available
                        # For legacy experiments, we need to get this from the database
                        conn = get_db_connection()
                        cursor = conn.cursor()
                        cursor.execute('SELECT pressed_thickness FROM cell_experiments WHERE id = ?', (exp_id,))
                        result = cursor.fetchone()
                        conn.close()
                        if result and result[0] is not None:
                            exp_summary['pressed_thickness'] = result[0]
                        # Add disc diameter data (default to 15mm for legacy experiments)
                        exp_summary['disc_diameter_mm'] = 15
                        # Add experiment notes to experiment summary
                        exp_summary['experiment_notes'] = experiment_notes
                        experiment_summaries.append(exp_summary)
                        
                except Exception as e:
                    st.error(f"Error processing experiment {exp_name}: {str(e)}")
                    continue
            
            # ===========================
            # Automated Anomaly Detection & Flagging
            # ===========================
            from cell_flags import analyze_cell_for_flags, get_experiment_context
            
            all_flags = {}  # Dictionary mapping cell_name to list of flags
            
            if individual_cells:
                # Build experiment context for statistical comparison
                experiment_context = get_experiment_context(individual_cells)
                
                # Analyze each cell for anomalies
                for exp_data in all_experiments_data:
                    exp_id, exp_name, file_name, loading, active_material, formation_cycles, test_number, electrolyte, substrate, separator, formulation_json, data_json, created_date, porosity, experiment_notes, cutoff_voltage_lower, cutoff_voltage_upper = exp_data
                    
                    try:
                        parsed_data = json.loads(data_json)
                        
                        # Check if this is a multi-cell experiment or single cell
                        if 'cells' in parsed_data:
                            # Multi-cell experiment
                            cells_data = parsed_data['cells']
                            disc_diameter = parsed_data.get('disc_diameter_mm', 15)
                            disc_area_cm2 = np.pi * (disc_diameter / 2 / 10) ** 2
                            
                            for cell_data in cells_data:
                                if cell_data.get('excluded', False):
                                    continue
                                try:
                                    df = pd.read_json(StringIO(cell_data['data_json']))
                                    cell_name = cell_data.get('test_number') or cell_data.get('cell_name', 'Unknown')
                                    
                                    # Find corresponding cell_summary from individual_cells
                                    cell_summary = next((c for c in individual_cells if c['cell_name'] == cell_name), None)
                                    
                                    if cell_summary:
                                        # Analyze cell for flags
                                        flags = analyze_cell_for_flags(df, cell_summary, experiment_context)
                                        all_flags[cell_name] = flags
                                except Exception:
                                    continue
                        
                        else:
                            # Legacy single cell experiment
                            df = pd.read_json(StringIO(data_json))
                            cell_name = test_number or exp_name
                            
                            # Find corresponding cell_summary
                            cell_summary = next((c for c in individual_cells if c['cell_name'] == cell_name), None)
                            
                            if cell_summary:
                                # Analyze cell for flags
                                flags = analyze_cell_for_flags(df, cell_summary, experiment_context)
                                all_flags[cell_name] = flags
                    
                    except Exception:
                        continue
            
            # Import flag display functions
            from display_components import display_cell_flags_summary, display_detailed_flags_section
            
            # Section 1: Average Cell Data per Experiment
            st.markdown("### 📊 Section 1: Experiment Summary")
            if experiment_summaries:
                exp_count = len(experiment_summaries)
                st.markdown(f"**{exp_count} experiment(s)** in this project • Showing averaged data per experiment")
                with st.expander("🎯 Customize & View Table", expanded=True):
                    display_experiment_summaries_table(experiment_summaries, all_flags)
            else:
                st.info("💡 No experiment summary data available. Create experiments to see data here.")
            
            st.markdown("---")
            
            # Section 2: All Individual Cells Data
            st.markdown("### 🧪 Section 2: Individual Cell Details")
            if individual_cells:
                cell_count = len(individual_cells)
                st.markdown(f"**{cell_count} cell(s)** tracked • Detailed data for each individual cell")
                with st.expander("🎯 Customize & View Table", expanded=False):
                    display_individual_cells_table(individual_cells, all_flags)
            else:
                st.info("💡 No individual cell data available. Upload cell data to experiments to see details here.")
            
            st.markdown("---")
            
            # Section 3: Best Performing Cells Analysis
            with st.expander("### 🏅 Section 3: Best Performing Cells Analysis", expanded=True):
                display_best_performers_analysis(individual_cells)
            
            st.markdown("---")
            
            # Automated Anomaly Detection & Flagging Section (at bottom for easy reference)
            if all_flags:
                display_cell_flags_summary(all_flags)
                st.markdown("---")
                display_detailed_flags_section(all_flags)

# --- Data Preprocessing Section ---

# Add this function after the imports at the top of the file
def recalculate_gravimetric_capacities(df, new_loading, new_active_material):
    """
    Recalculate gravimetric capacities when loading or active material values change.
    Returns a new DataFrame with updated Q Chg (mAh/g) and Q Dis (mAh/g) values.
    """
    try:
        # Create a copy of the DataFrame to avoid modifying the original
        updated_df = df.copy()
        
        # Calculate new active mass
        active_mass = (new_loading / 1000) * (new_active_material / 100)
        if active_mass <= 0:
            raise ValueError("Active mass must be greater than 0. Check loading and active material values.")
        
        # Recalculate gravimetric capacities
        if 'Q charge (mA.h)' in updated_df.columns:
            updated_df['Q Chg (mAh/g)'] = updated_df['Q charge (mA.h)'] / active_mass
        if 'Q discharge (mA.h)' in updated_df.columns:
            updated_df['Q Dis (mAh/g)'] = updated_df['Q discharge (mA.h)'] / active_mass
        
        return updated_df
    except Exception as e:
        st.error(f"Error recalculating gravimetric capacities: {str(e)}")
        return df  # Return original DataFrame if calculation fails
