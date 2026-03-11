 # export.py
from typing import List, Dict, Any, Tuple, Optional
import io
import json
import logging
import tempfile
import os
import traceback
from openpyxl import load_workbook
from openpyxl.chart import LineChart, Reference
import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import matplotlib.pyplot as plt
from matplotlib.figure import Figure
import numpy as np

# Set up comprehensive logging
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(name)s - %(levelname)s - %(message)s'
)
logger = logging.getLogger(__name__)

def safe_cycle_life_calculation(df_cell: pd.DataFrame, formation_cycles: int = 4) -> Optional[int]:
    """
    Calculate cycle life using the same logic as the main application.
    Returns the cycle number where capacity drops below 80% of initial capacity.
    """
    try:
        logger.debug(f"Calculating cycle life for cell with {len(df_cell)} cycles, formation_cycles={formation_cycles}")
        
        # Get discharge capacity series, excluding formation cycles
        if len(df_cell) <= formation_cycles:
            logger.warning(f"Not enough cycles for cycle life calculation: {len(df_cell)} <= {formation_cycles}")
            return None
            
        # Get the first cycle after formation as reference
        post_formation_data = df_cell.iloc[formation_cycles:]
        if post_formation_data.empty:
            logger.warning("No post-formation data available")
            return None
            
        # Get discharge capacity data
        qdis_col = 'Q Dis (mAh/g)'
        if qdis_col not in df_cell.columns:
            logger.warning(f"Discharge capacity column '{qdis_col}' not found in data")
            return None
            
        qdis_data = pd.to_numeric(post_formation_data[qdis_col], errors='coerce')
        qdis_data = qdis_data.dropna()
        
        if qdis_data.empty:
            logger.warning("No valid discharge capacity data found")
            return None
            
        # Use the first post-formation cycle as reference
        initial_capacity = qdis_data.iloc[0]
        if initial_capacity <= 0:
            logger.warning(f"Invalid initial capacity: {initial_capacity}")
            return None
            
        threshold = 0.8 * initial_capacity
        logger.debug(f"Initial capacity: {initial_capacity}, threshold: {threshold}")
        
        # Find first cycle below threshold
        below_threshold = qdis_data < threshold
        if below_threshold.any():
            first_below_idx = below_threshold.idxmax()
            # Get the actual cycle number from the original dataframe
            cycle_number = df_cell.iloc[first_below_idx][df_cell.columns[0]]
            logger.info(f"Cycle life calculated: {int(cycle_number)} cycles")
            return int(cycle_number)
        else:
            # If no cycle below threshold, return the last cycle
            last_cycle = df_cell.iloc[-1][df_cell.columns[0]]
            logger.info(f"No cycle below threshold, using last cycle: {int(last_cycle)}")
            return int(last_cycle)
            
    except Exception as e:
        logger.error(f"Error calculating cycle life: {e}")
        logger.error(f"Traceback: {traceback.format_exc()}")
        return None

def get_cell_metrics(df_cell: pd.DataFrame, formation_cycles: int = 4) -> Dict[str, Any]:
    """
    Calculate all cell metrics using the same logic as the main application.
    """
    metrics = {}
    
    try:
        logger.debug(f"Calculating metrics for cell with {len(df_cell)} cycles")
        
        # 1st Cycle Discharge Capacity (max of first 3 cycles)
        first_three_qdis = df_cell['Q Dis (mAh/g)'].head(3).tolist()
        max_qdis = max(first_three_qdis) if first_three_qdis else None
        metrics['max_qdis'] = max_qdis
        metrics['qdis_str'] = f"{max_qdis:.1f}" if isinstance(max_qdis, (int, float)) else "N/A"
        
        # First Cycle Efficiency
        if 'Efficiency (-)' in df_cell.columns and not df_cell['Efficiency (-)'].empty:
            first_cycle_eff = df_cell['Efficiency (-)'].iloc[0]
            try:
                eff_pct = float(first_cycle_eff) * 100
                metrics['eff_pct'] = eff_pct
                metrics['eff_str'] = f"{eff_pct:.1f}%"
            except (ValueError, TypeError):
                metrics['eff_pct'] = None
                metrics['eff_str'] = "N/A"
        else:
            metrics['eff_pct'] = None
            metrics['eff_str'] = "N/A"
        
        # Cycle Life (80%)
        cycle_life = safe_cycle_life_calculation(df_cell, formation_cycles)
        metrics['cycle_life'] = cycle_life
        metrics['cycle_life_str'] = f"{cycle_life}" if cycle_life is not None else "N/A"
        
        # Reversible Capacity (first cycle after formation)
        if len(df_cell) > formation_cycles:
            reversible_capacity = df_cell['Q Dis (mAh/g)'].iloc[formation_cycles]
            metrics['reversible_capacity'] = reversible_capacity
            metrics['reversible_str'] = f"{reversible_capacity:.1f}" if isinstance(reversible_capacity, (int, float)) else "N/A"
        else:
            metrics['reversible_capacity'] = None
            metrics['reversible_str'] = "N/A"
        
        # Coulombic Efficiency (average of post-formation cycles)
        eff_col = 'Efficiency (-)'
        qdis_col = 'Q Dis (mAh/g)'
        qch_col = 'Q Ch (mAh/g)'
        n_cycles = len(df_cell)
        
        # Coulombic efficiency: prefer Q Dis / Q Ch per cycle; fall back to provided efficiency column
        post_ce_values = []
        try:
            start_idx = formation_cycles if n_cycles > formation_cycles else 0
            
            if qdis_col in df_cell.columns and qch_col in df_cell.columns:
                qdis_series = pd.to_numeric(df_cell[qdis_col], errors='coerce')
                qch_series = pd.to_numeric(df_cell[qch_col], errors='coerce')
                
                for i in range(start_idx, n_cycles):
                    qdis = qdis_series.iloc[i]
                    qch = qch_series.iloc[i]
                    if pd.notna(qdis) and pd.notna(qch) and qch > 0:
                        post_ce_values.append((qdis / qch) * 100)
            
            elif eff_col in df_cell.columns:
                eff_series = pd.to_numeric(df_cell[eff_col], errors='coerce')
                post_ce_values = [val * 100 for val in eff_series.iloc[start_idx:n_cycles] if pd.notna(val) and val > 0]
            
            if post_ce_values:
                avg_eff = float(np.mean(post_ce_values))
                metrics['coulombic_eff'] = avg_eff
                metrics['coulombic_str'] = f"{avg_eff:.1f}%"
            else:
                metrics['coulombic_eff'] = None
                metrics['coulombic_str'] = "N/A"
        except Exception as e:
            logger.warning(f"Error calculating coulombic efficiency: {e}")
            metrics['coulombic_eff'] = None
            metrics['coulombic_str'] = "N/A"
        
        logger.debug(f"Metrics calculated successfully: {metrics}")
        return metrics
        
    except Exception as e:
        logger.error(f"Error calculating cell metrics: {e}")
        logger.error(f"Traceback: {traceback.format_exc()}")
        # Return safe defaults
        return {
            'max_qdis': None, 'qdis_str': "N/A",
            'eff_pct': None, 'eff_str': "N/A",
            'cycle_life': None, 'cycle_life_str': "N/A",
            'reversible_capacity': None, 'reversible_str': "N/A",
            'coulombic_eff': None, 'coulombic_str': "N/A"
        }

def create_main_plot_from_session_state(
    dfs: List[Dict[str, Any]],
    show_lines: Dict[str, bool],
    show_efficiency_lines: Dict[str, bool],
    remove_last_cycle: bool,
    show_graph_title: bool,
    experiment_name: str,
    show_average_performance: bool,
    avg_line_toggles: Dict[str, bool],
    remove_markers: bool,
    hide_legend: bool
) -> Figure:
    """
    Create the main capacity plot using the exact same logic as the Plots tab.
    """
    try:
        logger.info("Creating main capacity plot from session state...")
        logger.debug(f"Plot parameters: show_lines={show_lines}, remove_last_cycle={remove_last_cycle}, show_graph_title={show_graph_title}")
        
        # Import the plotting function from plotting.py
        from plotting import plot_capacity_graph
        
        # Create the plot using the same function as the Plots tab
        fig = plot_capacity_graph(
            dfs, show_lines, show_efficiency_lines, remove_last_cycle,
            show_graph_title, experiment_name, show_average_performance,
            avg_line_toggles, remove_markers, hide_legend
        )
        
        logger.info("Main capacity plot created successfully")
        return fig
        
    except Exception as e:
        logger.error(f"Error creating main capacity plot: {e}")
        logger.error(f"Traceback: {traceback.format_exc()}")
        
        # Fallback: create a simple error plot
        fig, ax = plt.subplots(figsize=(8, 5))
        ax.text(0.5, 0.5, 'Main plot generation failed', ha='center', va='center', transform=ax.transAxes)
        ax.set_xlabel('Cycle')
        ax.set_ylabel('Capacity (mAh/g)')
        ax.set_title('Capacity Plot (Error)')
        return fig

def create_retention_plot_from_session_state(
    dfs: List[Dict[str, Any]],
    show_lines: Dict[str, bool],
    reference_cycle: int,
    formation_cycles: int,
    remove_last_cycle: bool,
    show_title: bool,
    experiment_name: str,
    retention_threshold: float,
    y_axis_min: float,
    y_axis_max: float,
    show_baseline_line: bool,
    show_threshold_line: bool,
    remove_markers: bool,
    hide_legend: bool
) -> Figure:
    """
    Create the capacity retention plot using the exact same logic as the Plots tab.
    """
    try:
        logger.info("Creating capacity retention plot from session state...")
        logger.debug(f"Retention plot parameters: reference_cycle={reference_cycle}, retention_threshold={retention_threshold}, y_axis_min={y_axis_min}, y_axis_max={y_axis_max}")
        
        # Validate input data
        if not dfs:
            logger.error("No dataframes provided for retention plot")
            raise ValueError("No dataframes provided for retention plot")
        
        # Check if reference cycle exists in data
        for i, d in enumerate(dfs):
            df = d['df']
            if reference_cycle not in df[df.columns[0]].values:
                logger.warning(f"Reference cycle {reference_cycle} not found in cell {i+1}")
        
        # Import the plotting function from plotting.py
        from plotting import plot_capacity_retention_graph
        
        # Create the plot using the same function as the Plots tab
        fig = plot_capacity_retention_graph(
            dfs, show_lines, reference_cycle, formation_cycles, remove_last_cycle,
            show_title, experiment_name, False, {}, remove_markers, hide_legend,
            None, None, None, None, retention_threshold, y_axis_min, y_axis_max,
            show_baseline_line, show_threshold_line
        )
        
        logger.info("Capacity retention plot created successfully")
        return fig
        
    except Exception as e:
        logger.error(f"Error creating capacity retention plot: {e}")
        logger.error(f"Traceback: {traceback.format_exc()}")
        
        # Fallback: create a simple error plot
        fig, ax = plt.subplots(figsize=(8, 5))
        ax.text(0.5, 0.5, 'Retention plot generation failed', ha='center', va='center', transform=ax.transAxes)
        ax.set_xlabel('Cycle')
        ax.set_ylabel('Capacity Retention (%)')
        ax.set_title('Capacity Retention Plot (Error)')
        return fig

def save_figure_to_temp_file(fig: Figure, dpi: int = 300) -> str:
    """
    Save a matplotlib figure to a temporary file and return the path.
    Handles cleanup and error logging.
    """
    temp_file = None
    try:
        logger.debug("Saving figure to temporary file...")
        
        # Create temporary file
        temp_file = tempfile.NamedTemporaryFile(delete=False, suffix='.png')
        temp_path = temp_file.name
        temp_file.close()
        
        # Save figure with high quality settings
        fig.savefig(temp_path, format='png', bbox_inches='tight', dpi=dpi, 
                   facecolor='white', edgecolor='none', pad_inches=0.1)
        
        # Explicitly close all matplotlib resources to prevent memory issues
        plt.close(fig)
        plt.close('all')  # Close any remaining figures
        
        # Verify file was created and has content
        if os.path.exists(temp_path) and os.path.getsize(temp_path) > 0:
            logger.info(f"Figure saved successfully to: {temp_path} ({os.path.getsize(temp_path)} bytes)")
            return temp_path
        else:
            raise ValueError("Figure file was not created or is empty")
        
    except Exception as e:
        logger.error(f"Error saving figure to temp file: {e}")
        logger.error(f"Traceback: {traceback.format_exc()}")
        if temp_file:
            temp_file.close()
        raise

def add_picture_to_slide_safely(slide, image_path: str, left: float, top: float, width: float, height: float) -> bool:
    """
    Safely add a picture to a slide with comprehensive error handling.
    """
    try:
        logger.debug(f"Adding image to slide: {image_path}")
        
        # Validate file exists and has content
        if not os.path.exists(image_path):
            logger.error(f"Image file does not exist: {image_path}")
            return False
        
        file_size = os.path.getsize(image_path)
        if file_size == 0:
            logger.error(f"Image file is empty: {image_path}")
            return False
        
        logger.debug(f"Image file size: {file_size} bytes")
        
        # Add picture to slide
        slide.shapes.add_picture(image_path, Inches(left), Inches(top), width=Inches(width), height=Inches(height))
        logger.info(f"Successfully added image to slide: {image_path}")
        return True
        
    except Exception as e:
        logger.error(f"Error adding image to slide: {e}")
        logger.error(f"Traceback: {traceback.format_exc()}")
        return False

def cleanup_temp_files(temp_files: List[str]):
    """
    Clean up temporary files with error handling.
    """
    for temp_file in temp_files:
        try:
            if os.path.exists(temp_file):
                os.unlink(temp_file)
                logger.debug(f"Cleaned up temp file: {temp_file}")
        except Exception as e:
            logger.error(f"Error cleaning up temp file {temp_file}: {e}")

def add_formatted_table_to_slide(slide, table_data, left, top, width, height, header_color=None, avg_row_color=None):
    """
    Add a formatted table to a slide with proper styling.
    """
    try:
        rows = len(table_data)
        cols = len(table_data[0]) if table_data else 0
        if rows == 0 or cols == 0:
            return None
        
        # Calculate dynamic height based on number of rows
        row_height = max(0.3, min(0.5, height / rows))
        table_height = max(height, row_height * rows)
        
        table_shape = slide.shapes.add_table(rows, cols, Inches(left), Inches(top), Inches(width), Inches(table_height))
        table = table_shape.table
        
        # Set column widths
        col_width = width / cols
        for col_idx in range(cols):
            table.columns[col_idx].width = Inches(col_width)
        
        # Populate table with data
        for i, row_data in enumerate(table_data):
            for j, cell_data in enumerate(row_data):
                if j < cols:  # Safety check
                    cell = table.cell(i, j)
                    cell.text = str(cell_data)
                    
                    # Set vertical alignment
                    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
                    
                    # Format paragraphs
                    for para in cell.text_frame.paragraphs:
                        para.alignment = PP_ALIGN.CENTER
                        for run in para.runs:
                            run.font.size = Pt(11)
                            run.font.name = 'Calibri'
                    
                    # Header row styling
                    if i == 0:
                        cell.fill.solid()
                        header_color_rgb = header_color if header_color else RGBColor(68, 114, 196)
                        cell.fill.fore_color.rgb = header_color_rgb
                        for para in cell.text_frame.paragraphs:
                            for run in para.runs:
                                run.font.color.rgb = RGBColor(255, 255, 255)
                                run.font.bold = True
                                run.font.size = Pt(12)
                    
                    # Average row styling (if exists and avg_row_color provided)
                    # Check if this is the last row and contains "Average" or "AVERAGE" in the first cell
                    elif avg_row_color and i == len(table_data) - 1:
                        # Check first cell of the row for average indicators
                        first_cell_text = str(table_data[i][0]).upper() if table_data[i] else ""
                        if "AVERAGE" in first_cell_text or "PERFORMANCE" in first_cell_text:
                            cell.fill.solid()
                            cell.fill.fore_color.rgb = avg_row_color
                            for para in cell.text_frame.paragraphs:
                                for run in para.runs:
                                    run.font.bold = True
                                    run.font.size = Pt(11)
        
        return table
    except Exception as e:
        logger.error(f"Error adding formatted table: {e}")
        return None

def add_formulation_table_to_slide(slide, formulation_data, left, top, width, height):
    """
    Add a formulation component table to a slide.
    Handles multiple key variations for legacy data compatibility.
    """
    try:
        if not formulation_data or len(formulation_data) == 0:
            logger.warning("No formulation data provided")
            return None
        
        # Prepare table data
        table_data = [["Component", "Dry Mass Fraction (%)"]]
        
        for item in formulation_data:
            if not isinstance(item, dict):
                logger.warning(f"Skipping non-dict item in formulation data: {item}")
                continue
            
            # Try multiple key variations for component (handle legacy data)
            component = (item.get('Component') or 
                        item.get('component') or 
                        item.get('Component Name') or
                        item.get('component_name') or
                        '')
            
            # Try multiple key variations for dry mass fraction (handle legacy data)
            dry_mass_fraction = (item.get('Dry Mass Fraction (%)') or
                                item.get('dry_mass_fraction') or
                                item.get('Value') or
                                item.get('value') or
                                item.get('Dry Mass Fraction') or
                                item.get('dry_mass_fraction_pct') or
                                None)
            
            # Handle None, empty string, or invalid values
            if dry_mass_fraction is None or dry_mass_fraction == '':
                dry_mass_fraction_str = 'N/A'
            elif isinstance(dry_mass_fraction, (int, float)):
                # Ensure it's a valid number
                if dry_mass_fraction < 0 or dry_mass_fraction > 100:
                    logger.warning(f"Invalid dry mass fraction value: {dry_mass_fraction}, using as-is")
                dry_mass_fraction_str = f"{dry_mass_fraction:.1f}"
            elif isinstance(dry_mass_fraction, str):
                # Try to parse string values (e.g., "50.0", "50.0%", "None")
                dry_mass_fraction_clean = dry_mass_fraction.strip().rstrip('%')
                if dry_mass_fraction_clean.lower() in ('none', 'null', '', 'n/a'):
                    dry_mass_fraction_str = 'N/A'
                else:
                    try:
                        # Try to convert to float and format
                        fraction_float = float(dry_mass_fraction_clean)
                        dry_mass_fraction_str = f"{fraction_float:.1f}"
                    except (ValueError, TypeError):
                        # If conversion fails, use the string as-is
                        logger.warning(f"Could not convert dry mass fraction '{dry_mass_fraction}' to number, using as string")
                        dry_mass_fraction_str = dry_mass_fraction_clean
            else:
                # Fallback for any other type
                dry_mass_fraction_str = str(dry_mass_fraction) if dry_mass_fraction else 'N/A'
            
            # Only add row if component is not empty
            if component:
                table_data.append([component, dry_mass_fraction_str])
                logger.debug(f"Added formulation row: {component} = {dry_mass_fraction_str}")
            else:
                logger.warning(f"Skipping row with empty component: {item}")
        
        if len(table_data) <= 1:
            logger.warning("No valid formulation components found after processing")
            return None
        
        # Calculate table height
        rows = len(table_data)
        row_height = 0.4
        table_height = max(height, row_height * rows)
        
        # Add table
        table = add_formatted_table_to_slide(
            slide, table_data, left, top, width, table_height,
            header_color=RGBColor(68, 114, 196)
        )
        
        logger.info(f"Formulation table added with {rows} rows (including header)")
        return table
        
    except Exception as e:
        logger.error(f"Error adding formulation table: {e}")
        logger.error(f"Traceback: {traceback.format_exc()}")
        return None

def validate_powerpoint_structure(prs: Presentation) -> bool:
    """
    Validate that the PowerPoint presentation has the expected structure.
    """
    try:
        if not prs.slides:
            logger.error("PowerPoint has no slides")
            return False
        
        slide = prs.slides[0]
        if not slide.shapes:
            logger.error("PowerPoint slide has no shapes")
            return False
        
        logger.info(f"PowerPoint validation passed: {len(prs.slides)} slides, {len(slide.shapes)} shapes")
        return True
        
    except Exception as e:
        logger.error(f"Error validating PowerPoint structure: {e}")
        return False

def export_powerpoint(
    dfs: List[Dict[str, Any]], 
    show_averages: bool, 
    experiment_name: str, 
    show_lines: Dict[str, bool], 
    show_efficiency_lines: Dict[str, bool], 
    remove_last_cycle: bool,
    # Advanced slide content control
    include_summary_table: bool = True,
    include_main_plot: bool = True,
    include_retention_plot: bool = True,
    include_notes: bool = False,
    include_electrode_data: bool = False,
    include_porosity: bool = False,
    include_thickness: bool = False,
    include_solids_content: bool = False,
    include_formulation: bool = False,
    experiment_notes: str = "",
    # Parameters for retention plot
    retention_threshold: float = 80.0,
    reference_cycle: int = 5,
    formation_cycles: int = 4,
    retention_show_lines: Dict[str, bool] = None,
    retention_remove_markers: bool = False,
    retention_hide_legend: bool = False,
    retention_show_title: bool = True,
    show_baseline_line: bool = True,
    show_threshold_line: bool = True,
    y_axis_min: float = 0.0,
    y_axis_max: float = 110.0,
    # Plot customization parameters (from session state)
    show_graph_title: bool = True,
    show_average_performance: bool = False,
    avg_line_toggles: Dict[str, bool] = None,
    remove_markers: bool = False,
    hide_legend: bool = False,
    # Optional: existing presentation to append slides to
    existing_prs: Optional[Presentation] = None
) -> Tuple[io.BytesIO, str]:
    """
    Enhanced PowerPoint export with a highly dense, professional layout matching the Example Slide.
    All data is neatly aggregated onto a single unified status slide.
    """
    temp_files = []
    
    try:
        logger.info("=== Starting PowerPoint Export (Single Slide Dense Layout) ===")
        if not dfs:
            raise ValueError("No dataframes provided for export")
            
        if retention_show_lines is None:
            retention_show_lines = show_lines
        if avg_line_toggles is None:
            avg_line_toggles = {}
            
        if existing_prs is not None:
            prs = existing_prs
        else:
            prs = Presentation()
            
        # Determine slide headers
        if experiment_name:
            project_name = dfs[0].get('project_name', 'Unknown Project') if dfs else 'Unknown Project'
            title_text = f"Project: {project_name} – EXP ID: {experiment_name}"
        else:
            testnum = dfs[0].get('testnum', 'Cell 1') if dfs else 'Cell 1'
            title_text = f"Project: Unknown – EXP ID: {testnum}"
            
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
        
        # 1. ADD LOGO
        logo_path = "logo.png"
        if os.path.exists(logo_path):
            add_picture_to_slide_safely(slide, logo_path, left=0.2, top=0.2, width=0.8, height=0.6)
            
        # 2. ADD MAIN TITLE
        title_box = slide.shapes.add_textbox(Inches(1.2), Inches(0.2), Inches(8.5), Inches(0.5))
        title_tf = title_box.text_frame
        title_tf.clear()
        p = title_tf.add_paragraph()
        p.text = title_text
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.name = 'Times New Roman'
        p.alignment = PP_ALIGN.LEFT
        
        # Extract variables for Subtitle
        active_material = dfs[0].get('active_material', 'Unknown') if dfs else 'Unknown'
        subtitle_text = f"CHEMISTRY: {active_material} | KEY VARIABLES: Custom"
        
        subtitle_box = slide.shapes.add_textbox(Inches(1.2), Inches(0.6), Inches(8.5), Inches(0.4))
        subtitle_tf = subtitle_box.text_frame
        subtitle_tf.clear()
        p2 = subtitle_tf.add_paragraph()
        p2.text = subtitle_text
        p2.font.size = Pt(18)
        p2.font.bold = True
        p2.font.name = 'Times New Roman'
        
        # 3. ADD VERTICAL DIVIDER LINE
        try:
            from pptx.enum.shapes import MSO_CONNECTOR
            line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(5.0), Inches(1.2), Inches(5.0), Inches(7.0))
            line.line.color.rgb = RGBColor(0, 0, 0)
            line.line.width = Pt(1.5)
        except ImportError:
            pass
        
        # 4. LEFT COLUMN: SPECIFICATIONS & PROTOCOL (x=0.2 to 4.8)
        left_col_x = 0.2
        
        header_l_box = slide.shapes.add_textbox(Inches(left_col_x), Inches(1.2), Inches(4.5), Inches(0.4))
        hl_tf = header_l_box.text_frame
        hl_tf.clear()
        hl_p = hl_tf.add_paragraph()
        hl_p.text = "SPECIFICATIONS & PROTOCOL:"
        hl_p.font.size = Pt(16)
        hl_p.font.bold = True
        hl_p.font.name = 'Times New Roman'
        
        spec_box = slide.shapes.add_textbox(Inches(left_col_x), Inches(1.6), Inches(4.6), Inches(5.8))
        spec_tf = spec_box.text_frame
        spec_tf.clear()
        spec_tf.word_wrap = True
        
        def add_section_header(text):
            p = spec_tf.add_paragraph()
            p.text = text
            p.font.size = Pt(14)
            p.font.bold = True
            p.font.name = 'Times New Roman'
            
        def add_bullet(label, value):
            p = spec_tf.add_paragraph()
            p.font.size = Pt(12)
            p.font.name = 'Times New Roman'
            p.level = 1
            
            run = p.add_run()
            run.text = f"{label}: "
            run.font.bold = True
            
            run2 = p.add_run()
            run2.text = str(value)
            run2.font.bold = False
            
        # Section 1: Cell & Electrode Metadata
        add_section_header("1. Cell & Electrode Metadata")
        
        formulation_data = dfs[0].get('formulation', []) if dfs else []
        formulation_str = "Unknown"
        if formulation_data and isinstance(formulation_data, list):
            comps = []
            for item in formulation_data:
                if isinstance(item, dict):
                    comp = item.get('Component', item.get('component', ''))
                    val = item.get('Dry Mass Fraction (%)', item.get('Value', ''))
                    if comp and val:
                        comps.append(f"{comp} ({val}%)")
            if comps:
                formulation_str = " | ".join(comps)
                
        add_bullet("Cathode", formulation_str)
        
        loading = dfs[0].get('loading', 'N/A') if dfs else 'N/A'
        thickness = dfs[0].get('pressed_thickness', 'N/A') if dfs else 'N/A'
        add_bullet("Loading", f"{loading} mg/cm²; {thickness} µm thick")
        
        substrate = dfs[0].get('substrate', 'Unknown') if dfs else 'Unknown'
        add_bullet("Current Collector", substrate)
        
        electrolyte = dfs[0].get('electrolyte', 'Unknown') if dfs else 'Unknown'
        add_bullet("Electrolyte", electrolyte)
        
        add_bullet("E/C Ratio", "Not recorded in DB")
        
        separator = dfs[0].get('separator', 'Unknown') if dfs else 'Unknown'
        add_bullet("Separator", separator)
        
        # Add spacing
        spec_tf.add_paragraph()
        
        # Section 2: Testing Protocol
        add_section_header("2. Testing Protocol")
        
        v_low = dfs[0].get('cutoff_voltage_lower', '3.0') if dfs else '3.0'
        v_high = dfs[0].get('cutoff_voltage_upper', '4.2') if dfs else '4.2'
        add_bullet("Voltage Window", f"{v_low}-{v_high}V vs. Li/Li+")
        
        form_cycles = dfs[0].get('formation_cycles', 4) if dfs else 4
        add_bullet("Formation", f"{form_cycles} cycles")
        
        add_bullet("Cycling Rate", "Not recorded in DB")
        
        # Add spacing
        spec_tf.add_paragraph()
        
        # Section 3: Conclusion
        add_section_header("3. Conclusion")
        if experiment_notes:
            for line in experiment_notes.split('\n'):
                if line.strip():
                    p = spec_tf.add_paragraph()
                    p.text = line.strip()
                    p.font.size = Pt(12)
                    p.font.name = 'Times New Roman'
                    p.level = 1
        else:
            p = spec_tf.add_paragraph()
            p.text = "No notes provided."
            p.font.size = Pt(12)
            p.font.name = 'Times New Roman'
            p.level = 1
            
        # 5. RIGHT COLUMN: KEY RESULTS (x=5.2 to x=9.8)
        right_col_x = 5.2
        
        header_r_box = slide.shapes.add_textbox(Inches(right_col_x), Inches(1.2), Inches(4.5), Inches(0.4))
        hr_tf = header_r_box.text_frame
        hr_tf.clear()
        hr_p = hr_tf.add_paragraph()
        hr_p.text = "KEY RESULTS"
        hr_p.font.size = Pt(16)
        hr_p.font.bold = True
        hr_p.font.name = 'Times New Roman'
        
        # Right Col: Plots
        plot_y = 1.7
        if include_retention_plot or include_main_plot:
            try:
                # Prioritize retention plot for the single slide format as seen in example
                if include_retention_plot:
                    fig = create_retention_plot_from_session_state(
                        dfs, retention_show_lines, reference_cycle, formation_cycles, remove_last_cycle,
                        retention_show_title, experiment_name, retention_threshold, y_axis_min, y_axis_max,
                        show_baseline_line, show_threshold_line, retention_remove_markers, retention_hide_legend
                    )
                else:
                    fig = create_main_plot_from_session_state(
                        dfs, show_lines, show_efficiency_lines, remove_last_cycle,
                        show_graph_title, experiment_name, show_average_performance,
                        avg_line_toggles, remove_markers, hide_legend
                    )
                    
                if fig is not None:
                    img_path = save_figure_to_temp_file(fig)
                    temp_files.append(img_path)
                    add_picture_to_slide_safely(slide, img_path, right_col_x, plot_y, 4.6, 3.2)
            except Exception as e:
                logger.error(f"Error generating plot for slide: {e}")
                
        # Right Col: Table
        table_y = 5.0
        table_data = [
            ["Metric", "Benchmark", "Value", "Score"]
        ]
        
        df_cell = dfs[0]['df'] if dfs else pd.DataFrame()
        metrics = get_cell_metrics(df_cell, formation_cycles) if not df_cell.empty else {}
        
        areal_cap = "N/A"
        loading_val = dfs[0].get('loading', None) if dfs else None
        if loading_val and isinstance(loading_val, (int, float)) and metrics.get('max_qdis'):
            try:
                areal_cap = f"{(loading_val * float(metrics['max_qdis']) / 1000):.2f}"
            except Exception:
                pass
            
        table_data.extend([
            ["1st Discharge Cap. (mAh/g)", ">200", metrics.get('qdis_str', 'N/A'), ""],
            ["Reversible Capacity (mAh/g)", ">180", metrics.get('reversible_str', 'N/A'), ""],
            ["Areal Cap. (mAh/cm2)", "N/A", areal_cap, ""],
            ["FCE (%)", ">95", metrics.get('eff_str', 'N/A'), ""],
            ["Capacity Retention(%)", ">80", f"{metrics.get('cycle_life_str', 'N/A')} cycles", ""],
            ["Avg. Post-Form. CE (%)", ">99.9", metrics.get('coulombic_str', 'N/A'), ""]
        ])
        
        # Build Table
        table = add_formatted_table_to_slide(
            slide, table_data, 
            left=right_col_x, top=table_y, width=4.6, height=2.2,
            header_color=RGBColor(38, 77, 107)  # Dark teal from example
        )
        
        # Style table to match example
        if table:
            for row_idx, row in enumerate(table.rows):
                for col_idx, cell in enumerate(row.cells):
                    for paragraph in cell.text_frame.paragraphs:
                        paragraph.font.name = 'Times New Roman'
                        paragraph.font.size = Pt(10) if row_idx > 0 else Pt(11)
                        if row_idx > 0:
                            # Light gray alternating rows
                            if row_idx % 2 == 1:
                                cell.fill.solid()
                                cell.fill.fore_color.rgb = RGBColor(240, 240, 240)
                                
        if existing_prs is None:
            if not validate_powerpoint_structure(prs):
                raise ValueError("PowerPoint structure validation failed")
            
            pptx_bytes = io.BytesIO()
            prs.save(pptx_bytes)
            pptx_bytes.seek(0)
            
            testnum = dfs[0].get('testnum', 'Summary') if dfs else 'Summary'
            if experiment_name:
                pptx_file_name = f'{experiment_name} PPTX Export.pptx'
            else:
                pptx_file_name = f'{testnum} Export.pptx'
            
            return pptx_bytes, pptx_file_name
        else:
            return None, ""
            
    except Exception as e:
        logger.error(f"Error: {e}")
        logger.error(f"Traceback: {traceback.format_exc()}")
        raise
    finally:
        cleanup_temp_files(temp_files)

def export_excel(dfs: List[Dict[str, Any]], show_averages: bool, experiment_name: str) -> Tuple[io.BytesIO, str]:
    """
    Export data to Excel format with charts.
    """
    try:
        logger.info("Starting Excel export...")
        
        # Create Excel workbook
        output = io.BytesIO()
        
        with pd.ExcelWriter(output, engine='openpyxl') as writer:
            # Summary sheet
            summary_data = []
            headers = ["Cell", "1st Cycle Discharge Capacity (mAh/g)", "First Cycle Efficiency (%)", "Cycle Life (80%)", "Reversible Capacity (mAh/g)", "Coulombic Efficiency (%)"]
            summary_data.append(headers)
            
            for i, d in enumerate(dfs):
                df_cell = d['df']
                cell_name = d.get('testnum', f'Cell {i+1}') or f'Cell {i+1}'
                
                # Calculate metrics using the same logic as the app
                metrics = get_cell_metrics(df_cell, 4)  # Default formation cycles
                
                summary_data.append([
                    cell_name, 
                    metrics['qdis_str'], 
                    metrics['eff_str'], 
                    metrics['cycle_life_str'],
                    metrics['reversible_str'],
                    metrics['coulombic_str']
                ])
            
            # Always add Average Performance row when there are multiple cells
            if len(dfs) > 1:
                # Calculate averages
                avg_metrics = {
                    'max_qdis': [], 'eff_pct': [], 'cycle_life': [],
                    'reversible_capacity': [], 'coulombic_eff': []
                }
                
                for d in dfs:
                    df_cell = d['df']
                    metrics = get_cell_metrics(df_cell, 4)
                    
                    if metrics['max_qdis'] is not None:
                        avg_metrics['max_qdis'].append(metrics['max_qdis'])
                    if metrics['eff_pct'] is not None:
                        avg_metrics['eff_pct'].append(metrics['eff_pct'])
                    if metrics['cycle_life'] is not None:
                        avg_metrics['cycle_life'].append(metrics['cycle_life'])
                    if metrics['reversible_capacity'] is not None:
                        avg_metrics['reversible_capacity'].append(metrics['reversible_capacity'])
                    if metrics['coulombic_eff'] is not None:
                        avg_metrics['coulombic_eff'].append(metrics['coulombic_eff'])
                
                # Calculate final averages
                avg_row = ["Average Performance"]
                avg_row.append(f"{np.mean(avg_metrics['max_qdis']):.1f}" if avg_metrics['max_qdis'] else "N/A")
                avg_row.append(f"{np.mean(avg_metrics['eff_pct']):.1f}%" if avg_metrics['eff_pct'] else "N/A")
                avg_row.append(f"{np.mean(avg_metrics['cycle_life']):.0f}" if avg_metrics['cycle_life'] else "N/A")
                avg_row.append(f"{np.mean(avg_metrics['reversible_capacity']):.1f}" if avg_metrics['reversible_capacity'] else "N/A")
                avg_row.append(f"{np.mean(avg_metrics['coulombic_eff']):.1f}%" if avg_metrics['coulombic_eff'] else "N/A")
                
                summary_data.append(avg_row)
            
            # Create summary DataFrame
            summary_df = pd.DataFrame(summary_data[1:], columns=summary_data[0])
            summary_df.to_excel(writer, sheet_name='Summary', index=False)
            
            # Data sheets for each cell
            for i, d in enumerate(dfs):
                df = d['df']
                cell_name = d.get('testnum', f'Cell {i+1}') or f'Cell {i+1}'
                sheet_name = f'Cell_{i+1}' if len(cell_name) > 31 else cell_name
                df.to_excel(writer, sheet_name=sheet_name, index=False)
        
        output.seek(0)
        
        # Generate filename
        if experiment_name:
            if len(dfs) == 1:
                testnum = dfs[0].get('testnum', 'Cell 1') or 'Cell 1'
                excel_file_name = f'{experiment_name} {testnum} Data.xlsx'
            else:
                excel_file_name = f'{experiment_name} Cell Comparison Data.xlsx'
        else:
            if len(dfs) == 1:
                testnum = dfs[0].get('testnum', 'Cell 1') or 'Cell 1'
                excel_file_name = f'{testnum} Data.xlsx'
            else:
                excel_file_name = 'Cell Comparison Data.xlsx'
        
        logger.info(f"Excel export completed successfully: {excel_file_name}")
        return output, excel_file_name
        
    except Exception as e:
        logger.error(f"Error in Excel export: {e}")
        logger.error(f"Traceback: {traceback.format_exc()}")
        raise