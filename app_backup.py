import streamlit as st
import pandas as pd
import io
from openpyxl import load_workbook
import matplotlib.pyplot as plt
import numpy as np
from datetime import datetime, date
import re
import sqlite3
import json
import os
from pathlib import Path

# Database setup
from contextlib import contextmanager

@contextmanager
def get_db_connection():
    """Context manager for database connections to ensure proper cleanup."""
    conn = sqlite3.connect('cellscope.db')
    try:
        yield conn
    except Exception as e:
        conn.rollback()
        raise
    finally:
        conn.close()

def init_database():
    """Initialize the database and create tables if they don't exist."""
    with get_db_connection() as conn:
        cursor = conn.cursor()
        
        # Create projects table
        cursor.execute('''
            CREATE TABLE IF NOT EXISTS projects (
                id INTEGER PRIMARY KEY AUTOINCREMENT,
                user_id TEXT NOT NULL,
                name TEXT NOT NULL,
                description TEXT,
                created_date TIMESTAMP DEFAULT CURRENT_TIMESTAMP,
                last_modified TIMESTAMP DEFAULT CURRENT_TIMESTAMP
            )
        ''')
        
        # Create cell experiments table
        cursor.execute('''
            CREATE TABLE IF NOT EXISTS cell_experiments (
                id INTEGER PRIMARY KEY AUTOINCREMENT,
                project_id INTEGER NOT NULL,
                cell_name TEXT NOT NULL,
                file_name TEXT,
                loading REAL,
                active_material REAL,
                formation_cycles INTEGER,
                test_number TEXT,
                electrolyte TEXT,
                formulation_json TEXT,
                data_json TEXT,
                created_date TIMESTAMP DEFAULT CURRENT_TIMESTAMP,
                FOREIGN KEY (project_id) REFERENCES projects (id)
            )
        ''')
        
        # Create cell recipes table (for future use)
        cursor.execute('''
            CREATE TABLE IF NOT EXISTS cell_recipes (
                id INTEGER PRIMARY KEY AUTOINCREMENT,
                project_id INTEGER NOT NULL,
                recipe_name TEXT NOT NULL,
                recipe_data TEXT,
                created_date TIMESTAMP DEFAULT CURRENT_TIMESTAMP,
                FOREIGN KEY (project_id) REFERENCES projects (id)
            )
        ''')
        
        conn.commit()

# Initialize database
init_database()

def migrate_database():
    """Migrate database to add missing columns if they don't exist."""
    with get_db_connection() as conn:
        cursor = conn.cursor()
        
        # Check if electrolyte column exists
        cursor.execute("PRAGMA table_info(cell_experiments)")
        columns = [column[1] for column in cursor.fetchall()]
        
        # Add electrolyte column if it doesn't exist
        if 'electrolyte' not in columns:
            try:
                cursor.execute('ALTER TABLE cell_experiments ADD COLUMN electrolyte TEXT')
                print("Added electrolyte column to cell_experiments table")
            except sqlite3.OperationalError as e:
                print(f"Error adding electrolyte column: {e}")
        
        # Add formulation_json column if it doesn't exist
        if 'formulation_json' not in columns:
            try:
                cursor.execute('ALTER TABLE cell_experiments ADD COLUMN formulation_json TEXT')
                print("Added formulation_json column to cell_experiments table")
            except sqlite3.OperationalError as e:
                print(f"Error adding formulation_json column: {e}")
        
        conn.commit()

# Run database migration
migrate_database()

def get_project_components(project_id):
    """Get all unique components used in formulations within a project."""
    with get_db_connection() as conn:
        cursor = conn.cursor()
        cursor.execute('''
            SELECT formulation_json, data_json FROM cell_experiments 
            WHERE project_id = ? AND (formulation_json IS NOT NULL OR data_json IS NOT NULL)
        ''', (project_id,))
        results = cursor.fetchall()
        
        components = set()
        
        for formulation_json, data_json in results:
            # Check formulation_json field first
            if formulation_json:
                try:
                    formulation_data = json.loads(formulation_json)
                    if isinstance(formulation_data, list):
                        for item in formulation_data:
                            if isinstance(item, dict) and item.get('Component'):
                                components.add(item['Component'].strip())
                except (json.JSONDecodeError, AttributeError):
                    pass
            
            # Check formulation data within data_json (for multi-cell experiments)
            if data_json:
                try:
                    data = json.loads(data_json)
                    if 'cells' in data:
                        for cell in data['cells']:
                            formulation = cell.get('formulation', [])
                            if isinstance(formulation, list):
                                for item in formulation:
                                    if isinstance(item, dict) and item.get('Component'):
                                        components.add(item['Component'].strip())
                except (json.JSONDecodeError, AttributeError):
                    pass
        
        # Remove empty strings and return sorted list
        return sorted([comp for comp in components if comp])

# Test user (for now)
TEST_USER_ID = "admin"

def get_user_projects(user_id):
    """Get all projects for a user."""
    with get_db_connection() as conn:
        cursor = conn.cursor()
        cursor.execute('SELECT id, name, description, created_date, last_modified FROM projects WHERE user_id = ? ORDER BY last_modified DESC', (user_id,))
        return cursor.fetchall()

def create_project(user_id, name, description=""):
    """Create a new project."""
    with get_db_connection() as conn:
        cursor = conn.cursor()
        cursor.execute('INSERT INTO projects (user_id, name, description) VALUES (?, ?, ?)', (user_id, name, description))
        project_id = cursor.lastrowid
        conn.commit()
        return project_id

def save_cell_experiment(project_id, cell_name, file_name, loading, active_material, formation_cycles, test_number, df):
    """Save cell experiment data to database."""
    with get_db_connection() as conn:
        cursor = conn.cursor()
        
        # Convert DataFrame to JSON for storage
        data_json = df.to_json()
        
        cursor.execute('''
            INSERT INTO cell_experiments 
            (project_id, cell_name, file_name, loading, active_material, formation_cycles, test_number, data_json)
            VALUES (?, ?, ?, ?, ?, ?, ?, ?)
        ''', (project_id, cell_name, file_name, loading, active_material, formation_cycles, test_number, data_json))
        
        # Update project last_modified
        cursor.execute('UPDATE projects SET last_modified = CURRENT_TIMESTAMP WHERE id = ?', (project_id,))
        
        conn.commit()

def update_cell_experiment(experiment_id, cell_name, file_name, loading, active_material, formation_cycles, test_number, df, project_id):
    """Update an existing cell experiment in the database."""
    with get_db_connection() as conn:
        cursor = conn.cursor()
        
        # Convert DataFrame to JSON for storage
        data_json = df.to_json()
        
        cursor.execute('''
            UPDATE cell_experiments 
            SET cell_name = ?, file_name = ?, loading = ?, active_material = ?, 
                formation_cycles = ?, test_number = ?, data_json = ?
            WHERE id = ?
        ''', (cell_name, file_name, loading, active_material, formation_cycles, test_number, data_json, experiment_id))
        
        # Update project last_modified
        cursor.execute('UPDATE projects SET last_modified = CURRENT_TIMESTAMP WHERE id = ?', (project_id,))
        
        conn.commit()

def get_experiment_by_name_and_file(project_id, cell_name, file_name):
    """Get experiment ID by cell name and file name."""
    with get_db_connection() as conn:
        cursor = conn.cursor()
        cursor.execute('''
            SELECT id FROM cell_experiments 
            WHERE project_id = ? AND cell_name = ? AND file_name = ?
        ''', (project_id, cell_name, file_name))
        result = cursor.fetchone()
        return result[0] if result else None

def get_project_experiments(project_id):
    """Get all experiments for a project."""
    with get_db_connection() as conn:
        cursor = conn.cursor()
        cursor.execute('''
            SELECT id, cell_name, file_name, data_json, created_date
            FROM cell_experiments 
            WHERE project_id = ? 
            ORDER BY created_date
        ''', (project_id,))
        return cursor.fetchall()

def check_experiment_exists(project_id, cell_name, file_name):
    """Check if an experiment already exists in the project (for overwrite logic)."""
    with get_db_connection() as conn:
        cursor = conn.cursor()
        cursor.execute('''
            SELECT id FROM cell_experiments 
            WHERE project_id = ? AND cell_name = ? AND file_name = ?
        ''', (project_id, cell_name, file_name))
        result = cursor.fetchone()
        return result is not None

def get_experiment_data(experiment_id):
    """Get detailed data for a specific experiment."""
    with get_db_connection() as conn:
        cursor = conn.cursor()
        cursor.execute('''
            SELECT id, project_id, cell_name, file_name, loading, active_material, formation_cycles, test_number, data_json, created_date
            FROM cell_experiments 
            WHERE id = ?
        ''', (experiment_id,))
        return cursor.fetchone()

def delete_cell_experiment(experiment_id):
    """Delete a cell experiment from the database."""
    with get_db_connection() as conn:
        cursor = conn.cursor()
        
        # Get project_id before deleting
        cursor.execute('SELECT project_id FROM cell_experiments WHERE id = ?', (experiment_id,))
        result = cursor.fetchone()
        project_id = result[0] if result else None
        
        # Delete the experiment
        cursor.execute('DELETE FROM cell_experiments WHERE id = ?', (experiment_id,))
        
        # Update project last_modified if project exists
        if project_id:
            cursor.execute('UPDATE projects SET last_modified = CURRENT_TIMESTAMP WHERE id = ?', (project_id,))
        
        conn.commit()

def delete_project(project_id):
    """Delete a project and all its experiments from the database."""
    with get_db_connection() as conn:
        cursor = conn.cursor()
        
        # Delete all experiments in the project first
        cursor.execute('DELETE FROM cell_experiments WHERE project_id = ?', (project_id,))
        
        # Delete the project
        cursor.execute('DELETE FROM projects WHERE id = ?', (project_id,))
        
        conn.commit()

def rename_project(project_id, new_name):
    """Rename a project in the database."""
    with get_db_connection() as conn:
        cursor = conn.cursor()
        cursor.execute('UPDATE projects SET name = ?, last_modified = CURRENT_TIMESTAMP WHERE id = ?', (new_name, project_id))
        conn.commit()

def rename_experiment(experiment_id, new_name):
    """Rename an experiment in the database."""
    with get_db_connection() as conn:
        cursor = conn.cursor()
        cursor.execute('UPDATE cell_experiments SET cell_name = ?, file_name = ? WHERE id = ?', (new_name, f"{new_name}.json", experiment_id))
        conn.commit()

# New experiment-based functions
def save_experiment(project_id, experiment_name, experiment_date, disc_diameter_mm, group_assignments, group_names, cells_data):
    """Save a complete experiment with all its cells."""
    with get_db_connection() as conn:
        cursor = conn.cursor()
        
        # Convert group data to JSON
        group_assignments_json = json.dumps(group_assignments) if group_assignments else None
        group_names_json = json.dumps(group_names) if group_names else None
        
        # Collect all unique components from all cells' formulations
        unique_components = set()
        for cell in cells_data:
            for item in cell.get('formulation', []):
                comp = item.get('Component')
                if comp:
                    unique_components.add(comp.strip())
        formulation_json = json.dumps([{'Component': c} for c in sorted(unique_components)]) if unique_components else None
        
        # Create the experiment record
        cursor.execute('''
            INSERT INTO cell_experiments 
            (project_id, cell_name, file_name, loading, active_material, formation_cycles, test_number, electrolyte, formulation_json, data_json)
            VALUES (?, ?, ?, ?, ?, ?, ?, ?, ?, ?)
        ''', (project_id, experiment_name, f"{experiment_name}.json", None, None, None, None, None, formulation_json, json.dumps({
            'experiment_date': experiment_date.isoformat() if experiment_date else None,
            'disc_diameter_mm': disc_diameter_mm,
            'group_assignments': group_assignments,
            'group_names': group_names,
            'cells': cells_data
        })))
        
        # Update project last_modified
        cursor.execute('UPDATE projects SET last_modified = CURRENT_TIMESTAMP WHERE id = ?', (project_id,))
        
        experiment_id = cursor.lastrowid
        conn.commit()
        return experiment_id

def update_experiment(experiment_id, project_id, experiment_name, experiment_date, disc_diameter_mm, group_assignments, group_names, cells_data):
    """Update an existing experiment."""
    with get_db_connection() as conn:
        cursor = conn.cursor()
        
        # Convert group data to JSON
        group_assignments_json = json.dumps(group_assignments) if group_assignments else None
        group_names_json = json.dumps(group_names) if group_names else None
        
        # Update the experiment record
        cursor.execute('''
            UPDATE cell_experiments 
            SET cell_name = ?, file_name = ?, data_json = ?
            WHERE id = ?
        ''', (experiment_name, f"{experiment_name}.json", json.dumps({
            'experiment_date': experiment_date.isoformat() if experiment_date else None,
            'disc_diameter_mm': disc_diameter_mm,
            'group_assignments': group_assignments,
            'group_names': group_names,
            'cells': cells_data
        }), experiment_id))
        
        # Update project last_modified
        cursor.execute('UPDATE projects SET last_modified = CURRENT_TIMESTAMP WHERE id = ?', (project_id,))
        
        conn.commit()

def check_experiment_name_exists(project_id, experiment_name):
    """Check if an experiment with this name already exists in the project."""
    with get_db_connection() as conn:
        cursor = conn.cursor()
        cursor.execute('''
            SELECT id FROM cell_experiments 
            WHERE project_id = ? AND cell_name = ?
        ''', (project_id, experiment_name))
        result = cursor.fetchone()
        return result is not None

def get_experiment_by_name(project_id, experiment_name):
    """Get experiment ID by experiment name."""
    with get_db_connection() as conn:
        cursor = conn.cursor()
        cursor.execute('''
            SELECT id FROM cell_experiments 
            WHERE project_id = ? AND cell_name = ?
        ''', (project_id, experiment_name))
        result = cursor.fetchone()
        return result[0] if result else None

def get_all_project_experiments_data(project_id):
    """Get all experiments data for a project for Master Table analysis."""
    with get_db_connection() as conn:
        cursor = conn.cursor()
        cursor.execute('''
            SELECT id, cell_name, file_name, loading, active_material, formation_cycles, test_number, electrolyte, formulation_json, data_json, created_date
            FROM cell_experiments 
            WHERE project_id = ? 
            ORDER BY created_date
        ''', (project_id,))
        return cursor.fetchall()

def extract_date_from_filename(filename):
    match = re.search(r'(\d{6})', filename)
    if match:
        yymmdd = match.group(1)
        try:
            parsed_date = datetime.strptime(yymmdd, "%y%m%d").date()
            return parsed_date
        except ValueError:
            return None
    return None

def calculate_cell_summary(df, cell_data, disc_area_cm2):
    """Calculate summary statistics for a single cell."""
    try:
        # Basic cell info
        cell_name = cell_data.get('test_number') or cell_data.get('cell_name', 'Unknown')
        loading = cell_data.get('loading', 0)
        active_material = cell_data.get('active_material', 0)
        formation_cycles = cell_data.get('formation_cycles', 4)
        
        # 1st Cycle Discharge Capacity (mAh/g)
        first_three_qdis = df['Q Dis (mAh/g)'].head(3).tolist()
        max_qdis = max(first_three_qdis) if first_three_qdis else None
        
        # First Cycle Efficiency (%)
        eff_pct = None
        if 'Efficiency (-)' in df.columns and not df['Efficiency (-)'].empty:
            first_cycle_eff = df['Efficiency (-)'].iloc[0]
            try:
                eff_pct = float(first_cycle_eff) * 100
            except (ValueError, TypeError):
                eff_pct = None
        
        # Cycle Life (80%)
        cycle_life_80 = None
        try:
            qdis_series = df['Q Dis (mAh/g)'].dropna()
            cycle_index_series = df[df.columns[0]].iloc[qdis_series.index]
            cycle_life_80 = calculate_cycle_life_80(qdis_series, cycle_index_series)
        except:
            pass
        
        # Initial Areal Capacity (mAh/cm²)
        areal_capacity = None
        try:
            from ui_components import get_initial_areal_capacity
            areal_capacity, _, _, _ = get_initial_areal_capacity(df, disc_area_cm2)
        except:
            pass
        
        # Reversible Capacity (mAh/g)
        reversible_capacity = None
        if len(df) > formation_cycles:
            reversible_capacity = df['Q Dis (mAh/g)'].iloc[formation_cycles]
        
        # Coulombic Efficiency (post-formation, %)
        ceff_avg = None
        try:
            eff_col = 'Efficiency (-)'
            qdis_col = 'Q Dis (mAh/g)'
            n_cycles = len(df)
            ceff_values = []
            if eff_col in df.columns and qdis_col in df.columns and n_cycles > formation_cycles+1:
                prev_qdis = df[qdis_col].iloc[formation_cycles]
                prev_eff = df[eff_col].iloc[formation_cycles]
                for i in range(formation_cycles+1, n_cycles):
                    curr_qdis = df[qdis_col].iloc[i]
                    curr_eff = df[eff_col].iloc[i]
                    try:
                        pq = float(prev_qdis)
                        cq = float(curr_qdis)
                        pe = float(prev_eff)
                        ce = float(curr_eff)
                        if pq > 0 and (cq < 0.95 * pq or ce < 0.95 * pe):
                            break
                        ceff_values.append(ce)
                        prev_qdis = cq
                        prev_eff = ce
                    except (ValueError, TypeError):
                        continue
            if ceff_values:
                ceff_avg = sum(ceff_values) / len(ceff_values) * 100
        except:
            pass
        
        return {
            'cell_name': cell_name,
            'loading': loading,
            'active_material': active_material,
            'formation_cycles': formation_cycles,
            'first_discharge': max_qdis,
            'first_efficiency': eff_pct,
            'cycle_life_80': cycle_life_80,
            'areal_capacity': areal_capacity,
            'reversible_capacity': reversible_capacity,
            'coulombic_efficiency': ceff_avg
        }
    except Exception as e:
        # Return basic info if calculation fails
        return {
            'cell_name': cell_data.get('test_number') or cell_data.get('cell_name', 'Unknown'),
            'loading': cell_data.get('loading', 0),
            'active_material': cell_data.get('active_material', 0),
            'formation_cycles': cell_data.get('formation_cycles', 4),
            'first_discharge': None,
            'first_efficiency': None,
            'cycle_life_80': None,
            'areal_capacity': None,
            'reversible_capacity': None,
            'coulombic_efficiency': None
        }

def calculate_experiment_average(experiment_cells, exp_name, exp_date):
    """Calculate average values for an experiment."""
    if not experiment_cells:
        return None
    
    # Calculate averages for numeric fields
    numeric_fields = ['first_discharge', 'first_efficiency', 'cycle_life_80', 'areal_capacity', 'reversible_capacity', 'coulombic_efficiency']
    averages = {}
    
    for field in numeric_fields:
        values = [cell[field] for cell in experiment_cells if cell[field] is not None]
        averages[field] = sum(values) / len(values) if values else None
    
    return {
        'cell_name': f"{exp_name} (Avg)",
        'experiment_name': exp_name,
        'experiment_date': exp_date,
        'cell_count': len(experiment_cells),
        'loading': sum(cell['loading'] for cell in experiment_cells) / len(experiment_cells),
        'active_material': sum(cell['active_material'] for cell in experiment_cells) / len(experiment_cells),
        'formation_cycles': int(sum(cell['formation_cycles'] for cell in experiment_cells) / len(experiment_cells)),
        **averages
    }

def display_experiment_summaries_table(experiment_summaries):
    """Display the experiment summaries table with custom column order."""
    if not experiment_summaries:
        return
    # Custom column order for Section 1
    col_order = [
        'Experiment',
        'Reversible Capacity (mAh/g)',
        'Coulombic Efficiency (%)',
        'Areal Capacity (mAh/cm²)',
        '1st Discharge (mAh/g)',
        'First Efficiency (%)',
        'Cycle Life (80%)',
        'Date'
    ]
    df_data = []
    for exp in experiment_summaries:
        row = {
            'Experiment': exp['experiment_name'],
            'Date': exp.get('experiment_date', 'N/A'),
            '1st Discharge (mAh/g)': f"{exp['first_discharge']:.1f}" if exp['first_discharge'] else "N/A",
            'First Efficiency (%)': f"{exp['first_efficiency']:.1f}" if exp['first_efficiency'] else "N/A",
            'Cycle Life (80%)': f"{exp['cycle_life_80']:.0f}" if exp['cycle_life_80'] else "N/A",
            'Areal Capacity (mAh/cm²)': f"{exp['areal_capacity']:.3f}" if exp['areal_capacity'] else "N/A",
            'Reversible Capacity (mAh/g)': f"{exp['reversible_capacity']:.1f}" if exp['reversible_capacity'] else "N/A",
            'Coulombic Efficiency (%)': f"{exp['coulombic_efficiency']:.2f}" if exp['coulombic_efficiency'] else "N/A"
        }
        df_data.append(row)
    df = pd.DataFrame(df_data)
    df = df[[col for col in col_order if col in df.columns]]
    st.dataframe(df, use_container_width=True)

def display_individual_cells_table(individual_cells):
    """Display the individual cells table with custom column order."""
    if not individual_cells:
        return
    # Custom column order for Section 2
    col_order = [
        'Cell Name',
        'Reversible Capacity (mAh/g)',
        'Coulombic Efficiency (%)',
        'Areal Capacity (mAh/cm²)',
        '1st Discharge (mAh/g)',
        'First Efficiency (%)',
        'Cycle Life (80%)',
        'Active Material (%)',
        'Loading (mg)',
        'Date',
        'Experiment'
    ]
    df_data = []
    for cell in individual_cells:
        row = {
            'Experiment': cell.get('experiment_name', 'N/A'),
            'Cell Name': cell['cell_name'],
            'Date': cell.get('experiment_date', 'N/A'),
            'Loading (mg)': f"{cell['loading']:.1f}" if cell['loading'] else "N/A",
            'Active Material (%)': f"{cell['active_material']:.1f}" if cell['active_material'] else "N/A",
            '1st Discharge (mAh/g)': f"{cell['first_discharge']:.1f}" if cell['first_discharge'] else "N/A",
            'First Efficiency (%)': f"{cell['first_efficiency']:.1f}" if cell['first_efficiency'] else "N/A",
            'Cycle Life (80%)': f"{cell['cycle_life_80']:.0f}" if cell['cycle_life_80'] else "N/A",
            'Areal Capacity (mAh/cm²)': f"{cell['areal_capacity']:.3f}" if cell['areal_capacity'] else "N/A",
            'Reversible Capacity (mAh/g)': f"{cell['reversible_capacity']:.1f}" if cell['reversible_capacity'] else "N/A",
            'Coulombic Efficiency (%)': f"{cell['coulombic_efficiency']:.2f}" if cell['coulombic_efficiency'] else "N/A"
        }
        df_data.append(row)
    df = pd.DataFrame(df_data)
    df = df[[col for col in col_order if col in df.columns]]
    st.dataframe(df, use_container_width=True)

def display_best_performers_analysis(individual_cells):
    """Display best performing cells analysis."""
    if not individual_cells:
        return
    
    # Filter cells with valid data
    valid_cells = [cell for cell in individual_cells if any([
        cell['first_discharge'], cell['first_efficiency'], cell['cycle_life_80'], 
        cell['areal_capacity'], cell['reversible_capacity'], cell['coulombic_efficiency']
    ])]
    
    if not valid_cells:
        st.info("No valid data for performance analysis.")
        return
    
    # Find best performers for each metric
    metrics = {
        'Highest 1st Discharge Capacity': ('first_discharge', lambda x: x, 'mAh/g'),
        'Highest First Cycle Efficiency': ('first_efficiency', lambda x: x, '%'),
        'Longest Cycle Life': ('cycle_life_80', lambda x: x, 'cycles'),
        'Highest Areal Capacity': ('areal_capacity', lambda x: x, 'mAh/cm²'),
        'Highest Reversible Capacity': ('reversible_capacity', lambda x: x, 'mAh/g'),
        'Highest Coulombic Efficiency': ('coulombic_efficiency', lambda x: x, '%')
    }
    
    st.markdown("#### 🥇 Best Individual Performers by Metric")
    
    cols = st.columns(2)
    for i, (metric_name, (field, transform, unit)) in enumerate(metrics.items()):
        with cols[i % 2]:
            valid_for_metric = [cell for cell in valid_cells if cell[field] is not None]
            if valid_for_metric:
                best_cell = max(valid_for_metric, key=lambda x: transform(x[field]))
                value = transform(best_cell[field])
                st.metric(
                    label=metric_name,
                    value=f"{value:.2f} {unit}" if isinstance(value, float) else f"{value} {unit}",
                    help=f"Cell: {best_cell['cell_name']} from {best_cell.get('experiment_name', 'Unknown')}"
                )
    
    st.markdown("#### 🏆 Overall Best Performer")
    
    # Calculate normalized scores for overall performance
    # Normalize each metric to 0-1 scale and sum them
    performance_scores = []
    
    for cell in valid_cells:
        score = 0
        valid_metrics = 0
        
        # Normalize each metric (higher is better for all)
        for field, transform, unit in metrics.values():
            if cell[field] is not None:
                all_values = [c[field] for c in valid_cells if c[field] is not None]
                if all_values:
                    min_val = min(all_values)
                    max_val = max(all_values)
                    if max_val > min_val:
                        normalized = (transform(cell[field]) - min_val) / (max_val - min_val)
                        score += normalized
                        valid_metrics += 1
        
        if valid_metrics > 0:
            avg_score = score / valid_metrics
            performance_scores.append((cell, avg_score, valid_metrics))
    
    if performance_scores:
        # Sort by average score descending
        performance_scores.sort(key=lambda x: x[1], reverse=True)
        
        best_overall = performance_scores[0]
        best_cell, best_score, metrics_count = best_overall
        
        col1, col2 = st.columns([2, 1])
        with col1:
            st.success(f"🏆 **{best_cell['cell_name']}** from experiment **{best_cell.get('experiment_name', 'Unknown')}**")
            st.markdown(f"**Overall Performance Score:** {best_score:.2f}/1.00 (based on {metrics_count} metrics)")
        
        with col2:
            # Show top 3 performers
            st.markdown("**Top 3 Overall:**")
            for i, (cell, score, _) in enumerate(performance_scores[:3]):
                medal = ["🥇", "🥈", "🥉"][i]
                st.write(f"{medal} {cell['cell_name']} ({score:.2f})")
    
    else:
        st.info("Insufficient data for overall performance ranking.")

# --- Helper function for cycle life calculation ---
def calculate_cycle_life_80(qdis_series, cycle_index_series):
    # Use max of cycles 3 and 4 as initial, or last available if <4 cycles
    if len(qdis_series) >= 4:
        initial_qdis = max(qdis_series.iloc[2], qdis_series.iloc[3])
    elif len(qdis_series) > 0:
        initial_qdis = qdis_series.iloc[-1]
    else:
        return None
    threshold = 0.8 * initial_qdis
    below_threshold = qdis_series <= threshold
    if below_threshold.any():
        first_below_idx = below_threshold.idxmin()
        return int(cycle_index_series.iloc[first_below_idx])
    else:
        return int(cycle_index_series.iloc[-1])

def get_qdis_series(df_cell):
    qdis_raw = df_cell['Q Dis (mAh/g)']
    if pd.api.types.is_scalar(qdis_raw):
        return pd.Series([qdis_raw]).dropna()
    else:
        return pd.Series(qdis_raw).dropna()
from data_processing import load_and_preprocess_data
from ui_components import render_toggle_section, display_summary_stats, display_averages, render_cell_inputs, get_initial_areal_capacity
from plotting import plot_capacity_graph

# =============================
# Battery Data Gravimetric Capacity Calculator App
# =============================

# --- Top Bar ---
with st.container():
    col1, col2, col3 = st.columns([2, 1, 1])
    
    with col1:
        st.title("🔋 CellScope")
    
    with col2:
        # Show current experiment status
        loaded_experiment = st.session_state.get('loaded_experiment')
        if loaded_experiment:
            st.info(f"📊 {loaded_experiment['experiment_name']}")
    
    with col3:
        # Save button for loaded experiments
        if loaded_experiment:
            if st.button("💾 Save", key="save_changes_btn"):
                # Get current experiment data
                experiment_data = loaded_experiment['experiment_data']
                experiment_id = loaded_experiment['experiment_id']
                project_id = loaded_experiment['project_id']
                
                # Get current values from session state or use loaded values
                current_experiment_date = st.session_state.get('current_experiment_date', experiment_data.get('experiment_date'))
                current_disc_diameter = st.session_state.get('current_disc_diameter_mm', experiment_data.get('disc_diameter_mm'))
                current_group_assignments = st.session_state.get('current_group_assignments', experiment_data.get('group_assignments'))
                current_group_names = st.session_state.get('current_group_names', experiment_data.get('group_names'))
                
                # Convert date string to date object if needed
                if isinstance(current_experiment_date, str):
                    try:
                        current_experiment_date = datetime.fromisoformat(current_experiment_date).date()
                    except:
                        current_experiment_date = date.today()
                
                # Update the experiment with current data
                update_experiment(
                    experiment_id=experiment_id,
                    project_id=project_id,
                    experiment_name=loaded_experiment['experiment_name'],
                    experiment_date=current_experiment_date,
                    disc_diameter_mm=current_disc_diameter,
                    group_assignments=current_group_assignments,
                    group_names=current_group_names,
                    cells_data=experiment_data['cells']  # Keep original cell data
                )
                st.success("✅ Changes saved!")
                st.rerun()

st.markdown("---")

# --- Sidebar ---
with st.sidebar:
    # Restore logo at the top
    try:
        st.image("logo.png", width=150)
    except:
        st.image("https://placehold.co/150x80?text=Logo", width=150)
    st.markdown("---")
    st.markdown("#### 📁 Projects")
    user_projects = get_user_projects(TEST_USER_ID)
    if user_projects:
        for p in user_projects:
            project_id, project_name, project_desc, created_date, last_modified = p
            project_expanded = st.session_state.get(f'project_expanded_{project_id}', False)
            
            # Project row with dropdown arrow and three dots
            project_cols = st.columns([0.1, 0.7, 0.2])
            with project_cols[0]:
                # Dropdown arrow
                arrow = "▼" if project_expanded else "▶"
                if st.button(arrow, key=f'project_toggle_{project_id}', help="Expand/Collapse"):
                    st.session_state[f'project_expanded_{project_id}'] = not project_expanded
                    st.rerun()
            
            with project_cols[1]:
                # Project name button
                if st.button(project_name, key=f'project_select_{project_id}', use_container_width=True):
                    st.session_state['current_project_id'] = project_id
                    st.session_state['current_project_name'] = project_name
                    st.session_state[f'project_expanded_{project_id}'] = True
                    st.rerun()
            
            with project_cols[2]:
                menu_open = st.session_state.get(f'project_menu_open_{project_id}', False)
                if st.button('⋯', key=f'project_menu_btn_{project_id}', help="Project options"):
                    # Close all other menus
                    for p2 in user_projects:
                        st.session_state[f'project_menu_open_{p2[0]}'] = False
                    st.session_state[f'project_menu_open_{project_id}'] = not menu_open
                    st.rerun()
                # Simple vertical dropdown menu
                if menu_open:
                    with st.container():
                        if st.button('New Experiment', key=f'project_new_exp_{project_id}_menu', use_container_width=True):
                            st.session_state['current_project_id'] = project_id
                            st.session_state['current_project_name'] = project_name
                            st.session_state['start_new_experiment'] = True
                            if 'loaded_experiment' in st.session_state:
                                del st.session_state['loaded_experiment']
                            st.session_state[f'project_menu_open_{project_id}'] = False
                            st.session_state['show_cell_inputs_prompt'] = True
                            st.rerun()
                        if st.button('Rename', key=f'project_rename_{project_id}_menu', use_container_width=True):
                            st.session_state[f'renaming_project_{project_id}'] = True
                            st.session_state[f'project_menu_open_{project_id}'] = False
                            st.rerun()
                        if st.button('Delete', key=f'project_delete_{project_id}_menu', use_container_width=True):
                            st.session_state['confirm_delete_project'] = project_id
                            st.session_state[f'project_menu_open_{project_id}'] = False
                            st.rerun()
            
            # Inline rename for project
            if st.session_state.get(f'renaming_project_{project_id}', False):
                rename_cols = st.columns([0.8, 0.2])
                with rename_cols[0]:
                    new_name = st.text_input("New name:", value=project_name, key=f'rename_input_project_{project_id}', label_visibility="collapsed")
                with rename_cols[1]:
                    col1, col2 = st.columns(2)
                    with col1:
                        if st.button('✅', key=f'confirm_rename_project_{project_id}', help="Confirm rename"):
                            if new_name and new_name.strip() != project_name:
                                try:
                                    rename_project(project_id, new_name.strip())
                                    # Update current project name if it's the selected one
                                    if st.session_state.get('current_project_id') == project_id:
                                        st.session_state['current_project_name'] = new_name.strip()
                                    st.session_state[f'renaming_project_{project_id}'] = False
                                    st.success(f"✅ Renamed to '{new_name.strip()}'!")
                                    st.rerun()
                                except Exception as e:
                                    st.error(f"❌ Error: {str(e)}")
                            else:
                                st.warning("Please enter a different name.")
                    with col2:
                        if st.button('❌', key=f'cancel_rename_project_{project_id}', help="Cancel rename"):
                            st.session_state[f'renaming_project_{project_id}'] = False
                            st.rerun()
            
            # Show experiments if project is expanded and selected
            if (project_expanded and st.session_state.get('current_project_id') == project_id):
                existing_experiments = get_project_experiments(project_id)
                if existing_experiments:
                    st.markdown("##### 🧪 Experiments")
                    for experiment in existing_experiments:
                        experiment_id, experiment_name, file_name, data_json, created_date = experiment
                        
                        # Experiment row with indentation and three dots
                        exp_cols = st.columns([0.1, 0.7, 0.2])
                        with exp_cols[0]:
                            st.markdown("&nbsp;&nbsp;📊", unsafe_allow_html=True)
                        
                        with exp_cols[1]:
                            if st.button(experiment_name, key=f'exp_select_{experiment_id}', use_container_width=True):
                                st.session_state['loaded_experiment'] = {
                                    'experiment_id': experiment_id,
                                    'experiment_name': experiment_name,
                                    'project_id': project_id,
                                    'experiment_data': json.loads(data_json)
                                }
                                st.rerun()
                        
                        with exp_cols[2]:
                            exp_menu_open = st.session_state.get(f'exp_menu_open_{experiment_id}', False)
                            if st.button('⋯', key=f'exp_menu_btn_{experiment_id}', help="Experiment options"):
                                for e2 in existing_experiments:
                                    st.session_state[f'exp_menu_open_{e2[0]}'] = False
                                st.session_state[f'exp_menu_open_{experiment_id}'] = not exp_menu_open
                                st.rerun()
                            if exp_menu_open:
                                with st.container():
                                    if st.button('Rename', key=f'exp_rename_{experiment_id}_menu', use_container_width=True):
                                        st.session_state[f'renaming_experiment_{experiment_id}'] = True
                                        st.session_state[f'exp_menu_open_{experiment_id}'] = False
                                        st.rerun()
                                    if st.button('Delete', key=f'exp_delete_{experiment_id}_menu', use_container_width=True):
                                        st.session_state['confirm_delete_experiment'] = (experiment_id, experiment_name)
                                        st.session_state[f'exp_menu_open_{experiment_id}'] = False
                                        st.rerun()
                        
                        # Inline rename for experiment
                        if st.session_state.get(f'renaming_experiment_{experiment_id}', False):
                            exp_rename_cols = st.columns([0.1, 0.7, 0.2])
                            with exp_rename_cols[1]:
                                new_exp_name = st.text_input("New name:", value=experiment_name, key=f'rename_input_experiment_{experiment_id}', label_visibility="collapsed")
                            with exp_rename_cols[2]:
                                col1, col2 = st.columns(2)
                                with col1:
                                    if st.button('✅', key=f'confirm_rename_experiment_{experiment_id}', help="Confirm rename"):
                                        if new_exp_name and new_exp_name.strip() != experiment_name:
                                            try:
                                                rename_experiment(experiment_id, new_exp_name.strip())
                                                # Update loaded experiment name if it's the selected one
                                                if (st.session_state.get('loaded_experiment') and 
                                                    st.session_state['loaded_experiment'].get('experiment_id') == experiment_id):
                                                    st.session_state['loaded_experiment']['experiment_name'] = new_exp_name.strip()
                                                st.session_state[f'renaming_experiment_{experiment_id}'] = False
                                                st.success(f"✅ Renamed to '{new_exp_name.strip()}'!")
                                                st.rerun()
                                            except Exception as e:
                                                st.error(f"❌ Error: {str(e)}")
                                        else:
                                            st.warning("Please enter a different name.")
                                with col2:
                                    if st.button('❌', key=f'cancel_rename_experiment_{experiment_id}', help="Cancel rename"):
                                        st.session_state[f'renaming_experiment_{experiment_id}'] = False
                                        st.rerun()
                else:
                    st.markdown("&nbsp;&nbsp;&nbsp;&nbsp;*No experiments in this project*", unsafe_allow_html=True)
    else:
        st.info("No projects found. Create your first project below.")

    # Create new project
    with st.expander("➕ Create New Project", expanded=False):
        new_project_name = st.text_input("Project Name", key="new_project_name")
        new_project_description = st.text_area("Description (optional)", key="new_project_description")
        if st.button("Create Project", key="create_project_btn"):
            if new_project_name:
                project_id = create_project(TEST_USER_ID, new_project_name, new_project_description)
                st.success(f"Project '{new_project_name}' created successfully!")
                st.rerun()
            else:
                st.error("Please enter a project name.")

    st.markdown("---")
    st.markdown("#### 📁 Project Contents")
    # (Optional: Show currently loaded experiment status here)

    st.markdown("#### ℹ️ Quick Start")
    st.markdown("1. **Create or select a project** above")
    st.markdown("2. **Go to Cell Inputs tab** to upload data or edit experiments") 
    st.markdown("3. **View results** in Summary, Plots, and Export tabs")

# --- Dropdown CSS and Global State Management ---
# Add CSS for Gmail-style popup dropdown
st.markdown(
    """
    <style>
    /* Container for the three dots and popup */
    .dropdown-container {
        position: relative;
        display: inline-block;
    }
    
    /* Three dots button styling */
    .three-dots-btn {
        background: none;
        border: none;
        font-size: 18px;
        cursor: pointer;
        padding: 4px 8px;
        border-radius: 4px;
        color: #666;
        transition: background-color 0.2s;
    }
    
    .three-dots-btn:hover {
        background-color: #f0f0f0;
    }
    
    /* Gmail-style popup menu */
    .popup-menu {
        position: absolute;
        right: 0;
        top: 100%;
        background-color: white;
        border: 1px solid #ddd;
        border-radius: 8px;
        box-shadow: 0 4px 20px rgba(0,0,0,0.15);
        z-index: 1000;
        min-width: 120px;
        margin-top: 4px;
        padding: 4px 0;
        display: none;
    }
    
    /* Show popup when active */
    .popup-menu.show {
        display: block;
    }
    
    /* Popup menu items */
    .popup-item {
        display: block;
        width: 100%;
        padding: 8px 16px;
        text-align: left;
        background: none;
        border: none;
        cursor: pointer;
        font-size: 14px;
        color: #333;
        transition: background-color 0.2s;
        white-space: nowrap;
    }
    
    .popup-item:hover {
        background-color: #f5f5f5;
    }
    
    .popup-item.delete {
        color: #d32f2f;
    }
    
    .popup-item.delete:hover {
        background-color: #ffebee;
    }
    
    /* Arrow pointing up to the three dots */
    .popup-menu::before {
        content: '';
        position: absolute;
        top: -6px;
        right: 12px;
        width: 0;
        height: 0;
        border-left: 6px solid transparent;
        border-right: 6px solid transparent;
        border-bottom: 6px solid white;
        z-index: 1001;
    }
    
    .popup-menu::after {
        content: '';
        position: absolute;
        top: -7px;
        right: 12px;
        width: 0;
        height: 0;
        border-left: 6px solid transparent;
        border-right: 6px solid transparent;
        border-bottom: 6px solid #ddd;
        z-index: 1000;
    }
    </style>
    """,
    unsafe_allow_html=True
)



# --- Dialog Functions for Delete Confirmation ---

@st.dialog("⚠️ Confirm Project Deletion")
def confirm_delete_project():
    project_id = st.session_state['confirm_delete_project']
    user_projects = get_user_projects(TEST_USER_ID)
    project_name = next((p[1] for p in user_projects if p[0] == project_id), "this project")
    
    st.markdown(f"Are you sure you want to **permanently delete** the project **{project_name}**?")
    st.markdown("⚠️ This action cannot be undone and will delete all experiments in this project.")
    
    col1, col2 = st.columns(2)
    with col1:
        if st.button("🗑️ Yes, Delete", type="primary", use_container_width=True):
            try:
                delete_project(project_id)
                # Clear current project if it was deleted
                if st.session_state.get('current_project_id') == project_id:
                    st.session_state['current_project_id'] = None
                    st.session_state['current_project_name'] = None
                st.session_state['confirm_delete_project'] = None
                st.success(f"✅ Successfully deleted project '{project_name}'!")
                st.rerun()
            except Exception as e:
                st.error(f"❌ Error deleting project: {str(e)}")
    
    with col2:
        if st.button("Cancel", use_container_width=True):
            st.session_state['confirm_delete_project'] = None
            st.rerun()

@st.dialog("⚠️ Confirm Experiment Deletion")
def confirm_delete_experiment():
    experiment_id, experiment_name = st.session_state['confirm_delete_experiment']
    
    st.markdown(f"Are you sure you want to **permanently delete** the experiment **{experiment_name}**?")
    st.markdown("⚠️ This action cannot be undone.")
    
    col1, col2 = st.columns(2)
    with col1:
        if st.button("🗑️ Yes, Delete", type="primary", use_container_width=True):
            try:
                delete_cell_experiment(experiment_id)
                # Clear loaded experiment if it was deleted
                if (st.session_state.get('loaded_experiment') and 
                    st.session_state['loaded_experiment'].get('experiment_id') == experiment_id):
                    st.session_state['loaded_experiment'] = None
                st.session_state['confirm_delete_experiment'] = None
                st.success(f"✅ Successfully deleted experiment '{experiment_name}'!")
                st.rerun()
            except Exception as e:
                st.error(f"❌ Error deleting experiment: {str(e)}")
    
    with col2:
        if st.button("Cancel", use_container_width=True):
            st.session_state['confirm_delete_experiment'] = None
            st.rerun()

# Show delete confirmation dialogs when triggered
if st.session_state.get('confirm_delete_project'):
    confirm_delete_project()

if st.session_state.get('confirm_delete_experiment'):
    confirm_delete_experiment()

# --- Main Area: Tabs Layout ---
if 'datasets' not in st.session_state:
    st.session_state['datasets'] = []
datasets = st.session_state.get('datasets', [])
disc_diameter_mm = st.session_state.get('disc_diameter_mm', 15)
experiment_date = st.session_state.get('experiment_date', date.today())
# Ensure experiment_name is always defined
experiment_name = st.session_state.get('sidebar_experiment_name', '') or ''

# Tab selection state
if 'active_main_tab' not in st.session_state:
    st.session_state['active_main_tab'] = 0

# Remove unsupported arguments from st.tabs
# If 'show_cell_inputs_prompt' is set, show a message at the top
if 'show_cell_inputs_prompt' not in st.session_state:
    st.session_state['show_cell_inputs_prompt'] = False

if st.session_state.get('show_cell_inputs_prompt'):
    st.warning('Please click the "🧪 Cell Inputs" tab above to start your new experiment.')
    st.session_state['show_cell_inputs_prompt'] = False

# Create tabs - include Master Table tab if a project is selected
current_project_id = st.session_state.get('current_project_id')
if current_project_id:
    tab_inputs, tab1, tab2, tab3, tab_master = st.tabs(["🧪 Cell Inputs", "📊 Summary", "📈 Plots", "📤 Export", "📋 Master Table"])
else:
    tab_inputs, tab1, tab2, tab3 = st.tabs(["🧪 Cell Inputs", "📊 Summary", "📈 Plots", "📤 Export"])
    tab_master = None

# --- Cell Inputs Tab ---
with tab_inputs:
    # If user started a new experiment, clear cell input state
    if st.session_state.get('start_new_experiment'):
        st.session_state['datasets'] = []
        st.session_state['current_experiment_name'] = ''
        st.session_state['current_experiment_date'] = date.today()
        st.session_state['current_disc_diameter_mm'] = 15
        st.session_state['current_group_assignments'] = None
        st.session_state['current_group_names'] = ["Group A", "Group B", "Group C"]
        st.session_state['start_new_experiment'] = False
    st.header("🧪 Cell Inputs & Experiment Setup")
    st.markdown("---")
    
    # Check if we have a loaded experiment
    loaded_experiment = st.session_state.get('loaded_experiment')
    
    # If a new experiment is being started, always allow editing
    is_new_experiment = not loaded_experiment and st.session_state.get('current_project_id')
    
    if loaded_experiment:
        st.info(f"📊 Editing experiment: **{loaded_experiment['experiment_name']}**")
        experiment_data = loaded_experiment['experiment_data']
        
        # Load existing values from the experiment
        current_experiment_name = loaded_experiment['experiment_name']
        current_experiment_date = experiment_data.get('experiment_date')
        if isinstance(current_experiment_date, str):
            try:
                current_experiment_date = datetime.fromisoformat(current_experiment_date).date()
            except:
                current_experiment_date = date.today()
        elif current_experiment_date is None:
            current_experiment_date = date.today()
        
        current_disc_diameter = experiment_data.get('disc_diameter_mm', 15)
        current_group_assignments = experiment_data.get('group_assignments')
        current_group_names = experiment_data.get('group_names', ["Group A", "Group B", "Group C"])
        
        # Convert loaded cells data back to datasets format for editing
        cells_data = experiment_data.get('cells', [])
        loaded_datasets = []
        for cell_data in cells_data:
            # Create a mock file object for display purposes
            mock_file = type('MockFile', (), {
                'name': cell_data.get('file_name', 'loaded_data.csv'),
                'type': 'text/csv'
            })()
            
            loaded_datasets.append({
                'file': mock_file,
                'loading': cell_data.get('loading', 20.0),
                'active': cell_data.get('active_material', 90.0),
                'testnum': cell_data.get('test_number', cell_data.get('cell_name', '')),
                'formation_cycles': cell_data.get('formation_cycles', 4),
                'electrolyte': cell_data.get('electrolyte', '1M LiPF6 1:1:1'),
                'formulation': cell_data.get('formulation', [])
            })
        
        # Use loaded datasets as starting point
        datasets = loaded_datasets
        st.session_state['datasets'] = datasets
        
    elif is_new_experiment:
        st.info(f"📝 Creating a new experiment in project: **{st.session_state['current_project_name']}**")
        current_experiment_name = ""
        current_experiment_date = date.today()
        current_disc_diameter = 15
        current_group_assignments = None
        current_group_names = ["Group A", "Group B", "Group C"]
    else:
        st.info("📝 Create a new experiment or load an existing one from the sidebar")
        current_experiment_name = ""
        current_experiment_date = date.today()
        current_disc_diameter = 15
        current_group_assignments = None
        current_group_names = ["Group A", "Group B", "Group C"]
    
    # Experiment metadata inputs
    col1, col2 = st.columns(2)
    with col1:
        experiment_name_input = st.text_input(
            'Experiment Name', 
            value=current_experiment_name if loaded_experiment else experiment_name,
            placeholder='Enter experiment name for file naming and summary tab',
            key="main_experiment_name"
        )
    
    with col2:
        experiment_date_input = st.date_input(
            "Experiment Date", 
            value=current_experiment_date,
            help="Date associated with this experiment"
        )
    
    disc_diameter_input = st.number_input(
        'Disc Diameter (mm) for Areal Capacity Calculation', 
        min_value=1, 
        max_value=50, 
        value=current_disc_diameter, 
        step=1
    )
    
    st.markdown("---")
    
    # Cell inputs section
    if loaded_experiment:
        # For loaded experiments, show the cell input fields for editing
        if len(datasets) > 1:
            with st.expander(f'Cell 1: {datasets[0]["testnum"] or "Cell 1"}', expanded=False):
                col1, col2 = st.columns(2)
                with col1:
                    loading_0 = st.number_input(
                        f'Disc loading (mg) for Cell 1', 
                        min_value=0.0, 
                        step=1.0, 
                        value=float(datasets[0]["loading"]),
                        key=f'edit_loading_0'
                    )
                    formation_cycles_0 = st.number_input(
                        f'Formation Cycles for Cell 1', 
                        min_value=0, 
                        step=1, 
                        value=int(datasets[0]["formation_cycles"]),
                        key=f'edit_formation_0'
                    )
                with col2:
                    active_material_0 = st.number_input(
                        f'% Active material for Cell 1', 
                        min_value=0.0, 
                        max_value=100.0, 
                        step=1.0, 
                        value=float(datasets[0]["active"]),
                        key=f'edit_active_0'
                    )
                    test_number_0 = st.text_input(
                        f'Test Number for Cell 1', 
                        value=datasets[0]["testnum"] or "Cell 1",
                        key=f'edit_testnum_0'
                    )
                
                # Electrolyte selection
                electrolyte_options = ['1M LiPF6 1:1:1', '1M LiTFSI 3:7 +10% FEC']
                electrolyte_0 = st.selectbox(
                    f'Electrolyte for Cell 1', 
                    electrolyte_options,
                    index=electrolyte_options.index(datasets[0]["electrolyte"]) if datasets[0]["electrolyte"] in electrolyte_options else 0,
                    key=f'edit_electrolyte_0'
                )
                
                # Formulation table
                st.markdown("**Formulation:**")
                from ui_components import render_formulation_table
                # Initialize formulation data if needed
                formulation_key = f'formulation_data_edit_0_loaded'
                if formulation_key not in st.session_state:
                    st.session_state[formulation_key] = datasets[0]["formulation"] if datasets[0]["formulation"] else [{'Component': '', 'Dry Mass Fraction (%)': 0.0}]
                formulation_0 = render_formulation_table(f'edit_0_loaded', project_id, get_project_components)
                
                assign_all = st.checkbox('Assign values to all cells', key='assign_all_cells_loaded')
            # Update all datasets with new values
            edited_datasets = []
            for i, dataset in enumerate(datasets):
                if i == 0:
                    # First cell: preserve original file object and update other fields
                    edited_dataset = {
                        'file': dataset['file'],  # Always preserve original file object
                        'loading': loading_0,
                        'active': active_material_0,
                        'testnum': test_number_0,
                        'formation_cycles': formation_cycles_0,
                        'electrolyte': electrolyte_0,
                        'formulation': formulation_0
                    }
                else:
                    with st.expander(f'Cell {i+1}: {dataset["testnum"] or f"Cell {i+1}"}', expanded=False):
                        col1, col2 = st.columns(2)
                        if assign_all:
                            loading = loading_0
                            formation_cycles = formation_cycles_0
                            active_material = active_material_0
                            electrolyte = electrolyte_0
                            formulation = formulation_0
                            # Test number should remain individual (not assigned to all)
                            test_number = dataset['testnum'] or f'Cell {i+1}'
                        else:
                            with col1:
                                loading = st.number_input(
                                    f'Disc loading (mg) for Cell {i+1}', 
                                    min_value=0.0, 
                                    step=1.0, 
                                    value=float(dataset['loading']),
                                    key=f'edit_loading_{i}'
                                )
                                formation_cycles = st.number_input(
                                    f'Formation Cycles for Cell {i+1}', 
                                    min_value=0, 
                                    step=1, 
                                    value=int(dataset['formation_cycles']),
                                    key=f'edit_formation_{i}'
                                )
                            with col2:
                                active_material = st.number_input(
                                    f'% Active material for Cell {i+1}', 
                                    min_value=0.0, 
                                    max_value=100.0, 
                                    step=1.0, 
                                    value=float(dataset['active']),
                                    key=f'edit_active_{i}'
                                )
                                test_number = st.text_input(
                                    f'Test Number for Cell {i+1}', 
                                    value=dataset['testnum'] or f'Cell {i+1}',
                                    key=f'edit_testnum_{i}'
                                )
                            
                            # Electrolyte selection
                            electrolyte_options = ['1M LiPF6 1:1:1', '1M LiTFSI 3:7 +10% FEC']
                            electrolyte = st.selectbox(
                                f'Electrolyte for Cell {i+1}', 
                                electrolyte_options,
                                index=electrolyte_options.index(dataset['electrolyte']) if dataset['electrolyte'] in electrolyte_options else 0,
                                key=f'edit_electrolyte_{i}'
                            )
                            
                            # Formulation table
                            st.markdown("**Formulation:**")
                            from ui_components import render_formulation_table
                            # Initialize formulation data if needed
                            formulation_key = f'formulation_data_edit_{i}_loaded'
                            if formulation_key not in st.session_state:
                                st.session_state[formulation_key] = dataset['formulation'] if dataset['formulation'] else [{'Component': '', 'Dry Mass Fraction (%)': 0.0}]
                            formulation = render_formulation_table(f'edit_{i}_loaded', project_id, get_project_components)
                            
                        with col2:
                            # Always preserve original file object, only update other fields
                            edited_dataset = {
                                'file': dataset['file'],  # Always preserve original file object
                                'loading': loading,
                                'active': active_material,
                                'testnum': test_number,
                                'formation_cycles': formation_cycles,
                                'electrolyte': electrolyte,
                                'formulation': formulation
                            }
                edited_datasets.append(edited_dataset)
            datasets = edited_datasets
            st.session_state['datasets'] = datasets
        else:
            # Only one cell
            for i, dataset in enumerate(datasets):
                with st.expander(f'Cell {i+1}: {dataset["testnum"] or f"Cell {i+1}"}', expanded=False):
                    col1, col2 = st.columns(2)
                    with col1:
                        loading = st.number_input(
                            f'Disc loading (mg) for Cell {i+1}', 
                            min_value=0.0, 
                            step=1.0, 
                            value=float(dataset['loading']),
                            key=f'edit_loading_{i}'
                        )
                        formation_cycles = st.number_input(
                            f'Formation Cycles for Cell {i+1}', 
                            min_value=0, 
                            step=1, 
                            value=int(dataset['formation_cycles']),
                            key=f'edit_formation_{i}'
                        )
                    with col2:
                        active_material = st.number_input(
                            f'% Active material for Cell {i+1}', 
                            min_value=0.0, 
                            max_value=100.0, 
                            step=1.0, 
                            value=float(dataset['active']),
                            key=f'edit_active_{i}'
                        )
                        test_number = st.text_input(
                            f'Test Number for Cell {i+1}', 
                            value=dataset['testnum'] or f'Cell {i+1}',
                            key=f'edit_testnum_{i}'
                        )
                    
                    # Electrolyte selection
                    electrolyte_options = ['1M LiPF6 1:1:1', '1M LiTFSI 3:7 +10% FEC']
                    electrolyte = st.selectbox(
                        f'Electrolyte for Cell {i+1}', 
                        electrolyte_options,
                        index=electrolyte_options.index(dataset['electrolyte']) if dataset['electrolyte'] in electrolyte_options else 0,
                        key=f'edit_single_electrolyte_{i}'
                    )
                    
                    # Formulation table
                    st.markdown("**Formulation:**")
                    from ui_components import render_formulation_table
                    # Initialize formulation data if needed
                    formulation_key = f'formulation_data_edit_single_{i}_loaded'
                    if formulation_key not in st.session_state:
                        st.session_state[formulation_key] = dataset['formulation'] if dataset['formulation'] else [{'Component': '', 'Dry Mass Fraction (%)': 0.0}]
                    formulation = render_formulation_table(f'edit_single_{i}_loaded', project_id, get_project_components)
                    
                    # Always preserve original file object, only update other fields
                    edited_dataset = {
                        'file': dataset['file'],  # Always preserve original file object
                        'loading': loading,
                        'active': active_material,
                        'testnum': test_number,
                        'formation_cycles': formation_cycles,
                        'electrolyte': electrolyte,
                        'formulation': formulation
                    }
                dataset.update(edited_dataset)
    else:
        # New experiment flow - use unified render_cell_inputs
        st.markdown("#### 📁 Upload Cell Data Files")
        current_project_id = st.session_state.get('current_project_id')
        datasets = render_cell_inputs(context_key='main_cell_inputs', project_id=current_project_id, get_components_func=get_project_components)
        st.session_state['datasets'] = datasets
        # Store original uploaded files separately to prevent loss
        if datasets:
            st.session_state['original_uploaded_files'] = [ds['file'] for ds in datasets if ds.get('file')]
    
    # Group assignment section (if multiple cells)
    enable_grouping = False
    show_averages = False
    group_assignments = current_group_assignments
    group_names = current_group_names
    
    if datasets and len([d for d in datasets if d.get('file') or loaded_experiment or is_new_experiment]) > 1:
        st.markdown("---")
        st.markdown("#### 👥 Group Assignment (Optional)")
        enable_grouping = st.checkbox('Assign Cells into Groups?', value=bool(current_group_assignments))
        
        if enable_grouping:
            col1, col2, col3 = st.columns(3)
            with col1:
                group_names[0] = st.text_input('Group A Name', value=group_names[0], key='main_group_name_a')
            with col2:
                group_names[1] = st.text_input('Group B Name', value=group_names[1], key='main_group_name_b')
            with col3:
                group_names[2] = st.text_input('Group C Name', value=group_names[2], key='main_group_name_c')
            
            st.markdown("**Assign each cell to a group:**")
            group_assignments = []
            for i, cell in enumerate(datasets):
                if cell.get('file') or loaded_experiment or is_new_experiment:
                    cell_name = cell['testnum'] or f'Cell {i+1}'
                    default_group = current_group_assignments[i] if (current_group_assignments and i < len(current_group_assignments)) else group_names[0]
                    group = st.radio(
                        f"Assign {cell_name} to group:",
                        [group_names[0], group_names[1], group_names[2], "Exclude"],
                        index=[group_names[0], group_names[1], group_names[2], "Exclude"].index(default_group) if default_group in [group_names[0], group_names[1], group_names[2], "Exclude"] else 0,
                        key=f"main_group_assignment_{i}",
                        horizontal=True
                    )
                    group_assignments.append(group)
            
            show_averages = st.checkbox("Show Group Averages", value=True)
    
    # Update session state with current values
    st.session_state['current_experiment_name'] = experiment_name_input
    st.session_state['current_experiment_date'] = experiment_date_input
    st.session_state['current_disc_diameter_mm'] = disc_diameter_input
    st.session_state['current_group_assignments'] = group_assignments
    st.session_state['current_group_names'] = group_names
    
    # Save/Update experiment button
    st.markdown("---")
    if loaded_experiment:
        if st.button("💾 Update Experiment", type="primary", use_container_width=True):
            # Update the loaded experiment with new values
            experiment_id = loaded_experiment['experiment_id']
            project_id = loaded_experiment['project_id']
            
            # Prepare updated cells data
            updated_cells_data = []
            for i, dataset in enumerate(datasets):
                if i < len(experiment_data.get('cells', [])):
                    original_cell = experiment_data['cells'][i]
                    updated_cell = original_cell.copy()
                    updated_cell.update({
                        'loading': dataset['loading'],
                        'active_material': dataset['active'],
                        'formation_cycles': dataset['formation_cycles'],
                        'test_number': dataset['testnum'],
                        'cell_name': dataset['testnum'],
                        'electrolyte': dataset.get('electrolyte', '1M LiPF6 1:1:1'),
                        'formulation': dataset.get('formulation', [])
                    })
                    updated_cells_data.append(updated_cell)
            
            try:
                update_experiment(
                    experiment_id=experiment_id,
                    project_id=project_id,
                    experiment_name=experiment_name_input,
                    experiment_date=experiment_date_input,
                    disc_diameter_mm=disc_diameter_input,
                    group_assignments=group_assignments,
                    group_names=group_names,
                    cells_data=updated_cells_data
                )
                
                # Update the loaded experiment in session state
                st.session_state['loaded_experiment']['experiment_name'] = experiment_name_input
                st.session_state['loaded_experiment']['experiment_data'].update({
                    'experiment_date': experiment_date_input.isoformat(),
                    'disc_diameter_mm': disc_diameter_input,
                    'group_assignments': group_assignments,
                    'group_names': group_names,
                    'cells': updated_cells_data
                })
                
                st.success("✅ Experiment updated successfully!")
                st.rerun()
            except Exception as e:
                st.error(f"❌ Error updating experiment: {str(e)}")
    
    elif is_new_experiment:
        # Save new experiment (only if we have valid data and a selected project)
        valid_datasets = [ds for ds in datasets if ds.get('file') and ds.get('loading', 0) > 0 and 0 < ds.get('active', 0) <= 100]
        
        if valid_datasets and st.session_state.get('current_project_id'):
            if st.button("💾 Save New Experiment", type="primary", use_container_width=True):
                current_project_id = st.session_state['current_project_id']
                current_project_name = st.session_state['current_project_name']
                
                # Use experiment name from input or generate one
                exp_name = experiment_name_input if experiment_name_input else f"Experiment_{datetime.now().strftime('%Y%m%d_%H%M%S')}"
                
                # Prepare cells data
                cells_data = []
                for i, ds in enumerate(valid_datasets):
                    cell_name = ds['testnum'] if ds['testnum'] else f'Cell {i+1}'
                    file_name = ds['file'].name if ds['file'] else f'cell_{i+1}.csv'
                    
                    try:
                        # Process the data to get DataFrame
                        temp_dfs = load_and_preprocess_data([ds])
                        if temp_dfs and len(temp_dfs) > 0:
                            df = temp_dfs[0]['df']
                            
                            cells_data.append({
                                'cell_name': cell_name,
                                'file_name': file_name,
                                'loading': ds['loading'],
                                'active_material': ds['active'],
                                'formation_cycles': ds['formation_cycles'],
                                'test_number': ds['testnum'],
                                'electrolyte': ds.get('electrolyte', '1M LiPF6 1:1:1'),
                                'formulation': ds.get('formulation', []),
                                'data_json': df.to_json()
                            })
                        else:
                            st.warning(f"⚠️ Failed to process data for {cell_name}. Skipping this cell.")
                    except Exception as e:
                        st.error(f"❌ Error processing {cell_name}: {str(e)}")
                        continue
                
                # Save the experiment
                if cells_data:
                    try:
                        if check_experiment_name_exists(current_project_id, exp_name):
                            experiment_id = get_experiment_by_name(current_project_id, exp_name)
                            update_experiment(
                                experiment_id=experiment_id,
                                project_id=current_project_id,
                                experiment_name=exp_name,
                                experiment_date=experiment_date_input,
                                disc_diameter_mm=disc_diameter_input,
                                group_assignments=group_assignments,
                                group_names=group_names,
                                cells_data=cells_data
                            )
                            st.success(f"🔄 Updated experiment '{exp_name}' in project '{current_project_name}'!")
                        else:
                            save_experiment(
                                project_id=current_project_id,
                                experiment_name=exp_name,
                                experiment_date=experiment_date_input,
                                disc_diameter_mm=disc_diameter_input,
                                group_assignments=group_assignments,
                                group_names=group_names,
                                cells_data=cells_data
                            )
                            st.success(f"✅ Saved experiment '{exp_name}' with {len(cells_data)} cells in project '{current_project_name}'!")
                        
                        st.rerun()
                    except Exception as e:
                        st.error(f"❌ Failed to save experiment: {str(e)}")
                else:
                    st.error("❌ No valid cell data to save. Please check your files and try again.")
        
        elif valid_datasets and not st.session_state.get('current_project_id'):
            st.warning("⚠️ Please select a project in the sidebar before saving the experiment.")
        elif not valid_datasets:
            st.info("ℹ️ Upload cell data files and enter valid parameters to save an experiment.")

# --- Data Preprocessing Section ---
# Data processing now happens in the Cell Inputs tab

# Check if we have loaded experiment data to display
loaded_experiment = st.session_state.get('loaded_experiment')
if loaded_experiment:
    st.markdown("---")
    st.markdown(f"### 📊 Loaded Experiment: {loaded_experiment['experiment_name']}")
    
    # Convert saved JSON data back to DataFrames
    loaded_dfs = []
    experiment_data = loaded_experiment['experiment_data']
    cells_data = experiment_data.get('cells', [])
    
    for cell_data in cells_data:
        cell_name = cell_data.get('cell_name', 'Unknown')
        try:
            df = pd.read_json(cell_data['data_json'])
            loaded_dfs.append({
                'df': df,
                'testnum': cell_data.get('test_number'),
                'loading': cell_data.get('loading'),
                'active': cell_data.get('active_material'),
                'formation_cycles': cell_data.get('formation_cycles')
            })
        except Exception as e:
            st.error(f"Error loading data for {cell_name}: {str(e)}")
    
    if loaded_dfs:
        # Use loaded data for analysis
        dfs = loaded_dfs
        ready = True
        st.success(f"✅ Loaded {len(loaded_dfs)} cell(s) from saved experiment")
        
        # Display experiment metadata
        if experiment_data.get('experiment_date'):
            st.info(f"📅 Experiment Date: {experiment_data['experiment_date']}")
        if experiment_data.get('disc_diameter_mm'):
            st.info(f"🔘 Disc Diameter: {experiment_data['disc_diameter_mm']} mm")
    else:
        st.error("❌ Failed to load experiment data")
        ready = False

# Determine if we have data ready for analysis
if loaded_experiment:
    ready = len(loaded_dfs) > 0
else:
    # For new experiments, check if we have valid uploaded data
    datasets = st.session_state.get('datasets', [])
    # Only include datasets with a real uploaded file for processing
    valid_datasets = []
    for ds in datasets:
        file_obj = ds.get('file')
        if (file_obj and 
            hasattr(file_obj, 'read') and 
            hasattr(file_obj, 'name') and 
            hasattr(file_obj, 'type') and
            ds.get('loading', 0) > 0 and 
            0 < ds.get('active', 0) <= 100):
            # Additional check: ensure it's a real Streamlit UploadedFile, not a mock object
            try:
                # Try to access the file's size property (real uploaded files have this)
                if hasattr(file_obj, 'size') and file_obj.size is not None and file_obj.size > 0:
                    valid_datasets.append(ds)
            except (AttributeError, TypeError):
                # Skip files that don't have proper size attribute or other issues
                continue
    
    # Process uploaded data if we have valid datasets
    if valid_datasets:
        # Create a cache key based on file names and parameters
        cache_key = []
        for ds in valid_datasets:
            if ds.get('file'):
                file_info = f"{ds['file'].name}_{ds['loading']}_{ds['active']}_{ds['formation_cycles']}"
                cache_key.append(file_info)
        cache_key_str = "_".join(cache_key)
        
        # Check if we have cached processed data
        if ('processed_data_cache' in st.session_state and 
            st.session_state.get('cache_key') == cache_key_str):
            dfs = st.session_state['processed_data_cache']
        else:
            # Process data and cache it
            # Final safety check before processing
            safe_datasets = []
            for ds in valid_datasets:
                try:
                    # Test that we can actually read from the file
                    file_obj = ds['file']
                    current_pos = file_obj.tell() if hasattr(file_obj, 'tell') else 0
                    file_obj.seek(0)
                    # Try to read a small sample to verify it's readable
                    sample = file_obj.read(10)
                    file_obj.seek(current_pos)  # Reset position
                    if sample:  # File has content
                        safe_datasets.append(ds)
                except Exception as e:
                    st.warning(f"⚠️ Skipping invalid file: {ds.get('file', {}).get('name', 'Unknown')} - {str(e)}")
                    continue
            
            if safe_datasets:
                dfs = load_and_preprocess_data(safe_datasets)
                # After loading and preprocessing, re-attach the latest formation_cycles to each dfs entry
                for i, d in enumerate(dfs):
                    if i < len(safe_datasets):
                        d['formation_cycles'] = safe_datasets[i]['formation_cycles']
            else:
                dfs = []
                st.error("❌ No valid files found for processing.")
            
            # Cache the processed data
            st.session_state['processed_data_cache'] = dfs
            st.session_state['cache_key'] = cache_key_str
        
        ready = len(dfs) > 0
    else:
        ready = False

if ready:
    # Use values from Cell Inputs tab
    disc_diameter_mm = st.session_state.get('current_disc_diameter_mm', 15)
    experiment_name = st.session_state.get('current_experiment_name', '')
    group_assignments = st.session_state.get('current_group_assignments')
    group_names = st.session_state.get('current_group_names', ["Group A", "Group B", "Group C"])
    enable_grouping = bool(group_assignments)
    show_averages = enable_grouping
    datasets = st.session_state.get('datasets', [])
    disc_area_cm2 = np.pi * (disc_diameter_mm / 2 / 10) ** 2
    
    # --- Group Average Curve Calculation for Plotting ---
    group_curves = []
    if enable_grouping and group_assignments is not None:
        group_dfs = [[], [], []]
        for idx, name in enumerate(group_names):
            group_dfs[idx] = [df for df, g in zip(dfs, group_assignments) if g == name]
        def compute_group_avg_curve(group_dfs):
            if not group_dfs:
                return None, None, None, None
            dfs_trimmed = [d['df'] for d in group_dfs]
            x_col = dfs_trimmed[0].columns[0]
            common_cycles = set(dfs_trimmed[0][x_col])
            for df in dfs_trimmed[1:]:
                common_cycles = common_cycles & set(df[x_col])
            common_cycles = sorted(list(common_cycles))
            if not common_cycles:
                return None, None, None, None
            avg_qdis = []
            avg_qchg = []
            avg_eff = []
            for cycle in common_cycles:
                qdis_vals = []
                qchg_vals = []
                eff_vals = []
                for df in dfs_trimmed:
                    row = df[df[x_col] == cycle]
                    if not row.empty:
                        if 'Q Dis (mAh/g)' in row:
                            qdis_vals.append(row['Q Dis (mAh/g)'].values[0])
                        if 'Q Chg (mAh/g)' in row:
                            qchg_vals.append(row['Q Chg (mAh/g)'].values[0])
                        if 'Efficiency (-)' in row and not pd.isnull(row['Efficiency (-)'].values[0]):
                            eff_vals.append(row['Efficiency (-)'].values[0] * 100)
                avg_qdis.append(sum(qdis_vals)/len(qdis_vals) if qdis_vals else None)
                avg_qchg.append(sum(qchg_vals)/len(qchg_vals) if qchg_vals else None)
                avg_eff.append(sum(eff_vals)/len(eff_vals) if eff_vals else None)
            return common_cycles, avg_qdis, avg_qchg, avg_eff
        group_curves = [compute_group_avg_curve(group_dfs[idx]) for idx in range(3)]
    # --- Main Tabs Content ---
    with tab1:
        st.header("📊 Summary Tables")
        st.markdown("---")
        # Add toggle for showing average column
        show_average_col = False
        if len(dfs) > 1:
            show_average_col = st.toggle("Show average column", value=True, key="show_average_col_toggle")
        # Only one summary table should be rendered:
        from ui_components import display_summary_stats
        display_summary_stats(dfs, disc_area_cm2, show_average_col, group_assignments, group_names)
    with tab2:
        st.header("📈 Main Plot")
        st.markdown("---")
        show_lines, show_efficiency_lines, remove_last_cycle, show_graph_title, show_average_performance, avg_line_toggles, remove_markers, hide_legend, group_plot_toggles = render_toggle_section(dfs, enable_grouping=enable_grouping)
        fig = plot_capacity_graph(
            dfs, show_lines, show_efficiency_lines, remove_last_cycle, show_graph_title, experiment_name,
            show_average_performance, avg_line_toggles, remove_markers, hide_legend,
            group_a_curve=(group_curves[0][0], group_curves[0][1]) if enable_grouping and group_curves and group_curves[0][0] and group_curves[0][1] and group_plot_toggles.get("Group Q Dis", False) else None,
            group_b_curve=(group_curves[1][0], group_curves[1][1]) if enable_grouping and group_curves and group_curves[1][0] and group_curves[1][1] and group_plot_toggles.get("Group Q Dis", False) else None,
            group_c_curve=(group_curves[2][0], group_curves[2][1]) if enable_grouping and group_curves and group_curves[2][0] and group_curves[2][1] and group_plot_toggles.get("Group Q Dis", False) else None,
            group_a_qchg=(group_curves[0][0], group_curves[0][2]) if enable_grouping and group_curves and group_curves[0][0] and group_curves[0][2] and group_plot_toggles.get("Group Q Chg", False) else None,
            group_b_qchg=(group_curves[1][0], group_curves[1][2]) if enable_grouping and group_curves and group_curves[1][0] and group_curves[1][2] and group_plot_toggles.get("Group Q Chg", False) else None,
            group_c_qchg=(group_curves[2][0], group_curves[2][2]) if enable_grouping and group_curves and group_curves[2][0] and group_curves[2][2] and group_plot_toggles.get("Group Q Chg", False) else None,
            group_a_eff=(group_curves[0][0], group_curves[0][3]) if enable_grouping and group_curves and group_curves[0][0] and group_curves[0][3] and group_plot_toggles.get("Group Efficiency", False) else None,
            group_b_eff=(group_curves[1][0], group_curves[1][3]) if enable_grouping and group_curves and group_curves[1][0] and group_curves[1][3] and group_plot_toggles.get("Group Efficiency", False) else None,
            group_c_eff=(group_curves[2][0], group_curves[2][3]) if enable_grouping and group_curves and group_curves[2][0] and group_curves[2][3] and group_plot_toggles.get("Group Efficiency", False) else None,
            group_names=group_names
        )
        st.pyplot(fig)
    with tab3:
        st.header("📤 Export & Download")
        st.markdown("---")
        # PowerPoint and Excel export logic (move from main area)
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN
        from pptx.dml.color import RGBColor
        import tempfile
        # Only show export buttons if data is ready
        if ready:
            # Create PowerPoint
            prs = Presentation()
            
            if len(dfs) == 1:
                # Single cell - use bullet points format
                df_cell = dfs[0]['df']
                testnum = dfs[0]['testnum'] if dfs[0]['testnum'] else 'Cell 1'
                
                # Calculate summary values
                first_three_qdis = df_cell['Q Dis (mAh/g)'].head(3).tolist()
                max_qdis = max(first_three_qdis) if first_three_qdis else None
                if 'Efficiency (-)' in df_cell.columns and not df_cell['Efficiency (-)'].empty:
                    first_cycle_eff = df_cell['Efficiency (-)'].iloc[0]
                    try:
                        eff_pct = float(first_cycle_eff) * 100
                    except (ValueError, TypeError):
                        eff_pct = None
                else:
                    eff_pct = None
                
                # Cycle Life (80%) calculation
                qdis_series = get_qdis_series(df_cell)
                cycle_index_series = df_cell[df_cell.columns[0]].iloc[qdis_series.index]
                cycle_life_80 = calculate_cycle_life_80(qdis_series, cycle_index_series)
                
                # Safe formatting for summary values
                qdis_str = f"{max_qdis:.1f}" if isinstance(max_qdis, (int, float)) else "N/A"
                eff_str = f"{eff_pct:.1f}%" if isinstance(eff_pct, (int, float)) else "N/A"
                cycle_life_str = str(cycle_life_80) if isinstance(cycle_life_80, (int, float)) else "N/A"
                
                # Create slide with bullet points
                slide = prs.slides.add_slide(prs.slide_layouts[1])
                title_shape = slide.shapes.title
                if title_shape is not None:
                    if experiment_name:
                        title_shape.text = f"{experiment_name} - {testnum} Summary"
                    else:
                        title_shape.text = f"{testnum} Summary"
                if hasattr(title_shape, 'text_frame') and getattr(title_shape, 'text_frame', None) is not None:
                    title_shape.text_frame.paragraphs[0].font.size = Pt(30)
                
                # Echem header
                echem_header_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(3.5), Inches(0.5))
                echem_header_tf = echem_header_box.text_frame
                echem_header_tf.clear()
                echem_header_p = echem_header_tf.add_paragraph()
                echem_header_p.text = "Echem"
                echem_header_p.level = 0
                echem_header_p.font.size = Pt(24)
                echem_header_p.font.bold = True
                echem_header_p.font.color.rgb = RGBColor(0, 0, 0)
                echem_header_p.font.italic = False
                
                # Summary bullets
                content = slide.placeholders[1]
                if hasattr(content, 'text'):
                    content.text = ""
                if hasattr(content, 'text_frame') and getattr(content, 'text_frame', None) is not None:
                    p1 = content.text_frame.add_paragraph()
                    p1.text = f"1st Cycle Discharge Capacity: {qdis_str} mAh/g"
                    p1.level = 0
                    p1.font.size = Pt(18)
                    p2 = content.text_frame.add_paragraph()
                    p2.text = f"First Cycle Efficiency: {eff_str}"
                    p2.level = 0
                    p2.font.size = Pt(18)
                    p3 = content.text_frame.add_paragraph()
                    p3.text = f"Cycle Life (80%): {cycle_life_str}"
                    p3.level = 0
                    p3.font.size = Pt(18)
                
                # Create graph for single cell
                # Use plot_capacity_graph with show_graph_title as set by the user
                fig = plot_capacity_graph(
                    [dfs[0]], show_lines, show_efficiency_lines, remove_last_cycle, show_graph_title, experiment_name,
                    show_average_performance, avg_line_toggles, remove_markers, hide_legend
                )
                img_bytes = io.BytesIO()
                fig.savefig(img_bytes, format='png', bbox_inches='tight')
                plt.close(fig)
                img_bytes.seek(0)
                with tempfile.NamedTemporaryFile(delete=False, suffix='.png') as tmp_img:
                    tmp_img.write(img_bytes.read())
                    img_path = tmp_img.name
                slide.shapes.add_picture(img_path, Inches(5.5), Inches(1.2), width=Inches(4.5))
                
                # Download button
                pptx_bytes = io.BytesIO()
                prs.save(pptx_bytes)
                pptx_bytes.seek(0)
                # Update PowerPoint file name to include experiment name
                if experiment_name:
                    pptx_file_name = f'{experiment_name} {testnum} Summary.pptx'
                else:
                    pptx_file_name = f'{testnum} Summary.pptx'
                st.download_button(f'Download PowerPoint: {pptx_file_name}', data=pptx_bytes, file_name=pptx_file_name, mime='application/vnd.openxmlformats-officedocument.presentationml.presentation', key='download_pptx_single')
                
            else:
                # Multiple cells - use table format
                slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title and Content layout
                title_shape = slide.shapes.title
                if title_shape is not None:
                    if experiment_name:
                        title_shape.text = f"{experiment_name} - Cell Comparison Summary"
                    else:
                        title_shape.text = "Cell Comparison Summary"
                if hasattr(title_shape, 'text_frame') and getattr(title_shape, 'text_frame', None) is not None:
                    title_shape.text_frame.paragraphs[0].font.size = Pt(30)
                
                # Echem header
                echem_header_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(3.5), Inches(0.5))
                echem_header_tf = echem_header_box.text_frame
                echem_header_tf.clear()
                echem_header_p = echem_header_tf.add_paragraph()
                echem_header_p.text = "Echem"
                echem_header_p.level = 0
                echem_header_p.font.size = Pt(24)
                echem_header_p.font.bold = True
                echem_header_p.font.color.rgb = RGBColor(0, 0, 0)
                echem_header_p.font.italic = False
                
                # Calculate summary data for all cells
                table_data = []
                headers = ["Cell", "1st Cycle Discharge Capacity (mAh/g)", "First Cycle Efficiency (%)", "Cycle Life (80%)", "Initial Areal Capacity (mAh/cm²)", "Reversible Capacity (mAh/g)", "Coulombic Efficiency (post-formation)"]
                table_data.append(headers)
                
                # Individual cells
                for i, d in enumerate(dfs):
                    df_cell = d['df']
                    cell_name = d['testnum'] if d['testnum'] else f'Cell {i+1}'
                    # Calculate values
                    first_three_qdis = df_cell['Q Dis (mAh/g)'].head(3).tolist()
                    max_qdis = max(first_three_qdis) if first_three_qdis else None
                    qdis_str = f"{max_qdis:.1f}" if isinstance(max_qdis, (int, float)) else "N/A"
                    if 'Efficiency (-)' in df_cell.columns and not df_cell['Efficiency (-)'].empty:
                        first_cycle_eff = df_cell['Efficiency (-)'].iloc[0]
                        try:
                            eff_pct = float(first_cycle_eff) * 100
                            eff_str = f"{eff_pct:.1f}%"
                        except (ValueError, TypeError):
                            eff_str = "N/A"
                    else:
                        eff_str = "N/A"
                    # Cycle Life (80%)
                    qdis_series = get_qdis_series(df_cell)
                    cycle_index_series = df_cell[df_cell.columns[0]].iloc[qdis_series.index]
                    cycle_life_80 = calculate_cycle_life_80(qdis_series, cycle_index_series)
                    cycle_life_str = str(cycle_life_80) if isinstance(cycle_life_80, (int, float)) else "N/A"
                    # Initial Areal Capacity (mAh/cm²)
                    areal_capacity, chosen_cycle, diff_pct, eff_val = get_initial_areal_capacity(df_cell, disc_area_cm2)
                    if areal_capacity is None:
                        areal_str = "N/A"
                    else:
                        areal_str = f"{areal_capacity:.3f}"
                    # Reversible Capacity (mAh/g)
                    formation_cycles = d.get('formation_cycles', 4)
                    if len(df_cell) > formation_cycles:
                        reversible_capacity = df_cell['Q Dis (mAh/g)'].iloc[formation_cycles]
                        reversible_str = f"{reversible_capacity:.1f}"
                    else:
                        reversible_str = "N/A"
                    # Coulombic Efficiency (post-formation, %)
                    eff_col = 'Efficiency (-)'
                    qdis_col = 'Q Dis (mAh/g)'
                    n_cycles = len(df_cell)
                    ceff_values = []
                    if eff_col in df_cell.columns and qdis_col in df_cell.columns and n_cycles > formation_cycles+1:
                        prev_qdis = df_cell[qdis_col].iloc[formation_cycles]
                        prev_eff = df_cell[eff_col].iloc[formation_cycles]
                        for j in range(formation_cycles+1, n_cycles):
                            curr_qdis = df_cell[qdis_col].iloc[j]
                            curr_eff = df_cell[eff_col].iloc[j]
                            try:
                                pq = float(prev_qdis)
                                cq = float(curr_qdis)
                                pe = float(prev_eff)
                                ce = float(curr_eff)
                                if pq > 0 and (cq < 0.95 * pq or ce < 0.95 * pe):
                                    break
                                ceff_values.append(ce)
                                prev_qdis = cq
                                prev_eff = ce
                            except (ValueError, TypeError):
                                # Skip this cycle if any value is not numeric
                                continue
                    if ceff_values:
                        ceff_str = f"{sum(ceff_values)/len(ceff_values)*100:.2f}%"
                    else:
                        ceff_str = "N/A"
                    table_data.append([cell_name, qdis_str, eff_str, cycle_life_str, areal_str, reversible_str, ceff_str])
                # Prepare group summary rows (A, B, C) to insert before the average row
                group_summary_rows = []
                if enable_grouping and group_assignments is not None:
                    for group_idx, group_name in enumerate(group_names):
                        group_dfs = [df for df, g in zip(dfs, group_assignments) if g == group_name]
                        if len(group_dfs) > 1:
                            avg_qdis_values = []
                            avg_eff_values = []
                            avg_cycle_life_values = []
                            avg_areal_capacity_values = []
                            avg_reversible_capacities = []
                            avg_ceff_values = []
                            for d in group_dfs:
                                df_cell = d['df']
                                first_three_qdis = df_cell['Q Dis (mAh/g)'].head(3).tolist()
                                max_qdis = max(first_three_qdis) if first_three_qdis else None
                                if max_qdis is not None:
                                    avg_qdis_values.append(max_qdis)
                                if 'Efficiency (-)' in df_cell.columns and not df_cell['Efficiency (-)'].empty:
                                    first_cycle_eff = df_cell['Efficiency (-)'].iloc[0]
                                    try:
                                        eff_pct = float(first_cycle_eff) * 100
                                        avg_eff_values.append(eff_pct)
                                    except (ValueError, TypeError):
                                        avg_eff_values.append(None) # Handle non-numeric efficiency
                                else:
                                    avg_eff_values.append(None) # Handle missing efficiency
                                try:
                                    qdis_series = get_qdis_series(df_cell)
                                    cycle_index_series = df_cell[df_cell.columns[0]].iloc[qdis_series.index]
                                    cycle_life_80 = calculate_cycle_life_80(qdis_series, cycle_index_series)
                                    avg_cycle_life_values.append(cycle_life_80)
                                except Exception:
                                    pass
                                areal_capacity, chosen_cycle, diff_pct, eff_val = get_initial_areal_capacity(df_cell, disc_area_cm2)
                                if areal_capacity is not None:
                                    avg_areal_capacity_values.append(areal_capacity)
                                formation_cycles = d.get('formation_cycles', 4)
                                if len(df_cell) > formation_cycles:
                                    reversible_capacity = df_cell['Q Dis (mAh/g)'].iloc[formation_cycles]
                                    avg_reversible_capacities.append(reversible_capacity)
                                # Coulombic Efficiency (post-formation, %)
                                eff_col = 'Efficiency (-)'
                                qdis_col = 'Q Dis (mAh/g)'
                                n_cycles = len(df_cell)
                                ceff_values = []
                                if eff_col in df_cell.columns and qdis_col in df_cell.columns and n_cycles > formation_cycles+1:
                                    prev_qdis = df_cell[qdis_col].iloc[formation_cycles]
                                    prev_eff = df_cell[eff_col].iloc[formation_cycles]
                                    for j in range(formation_cycles+1, n_cycles):
                                        curr_qdis = df_cell[qdis_col].iloc[j]
                                        curr_eff = df_cell[eff_col].iloc[j]
                                        if prev_qdis > 0 and (curr_qdis < 0.95 * prev_qdis or curr_eff < 0.95 * prev_eff):
                                            break
                                        ceff_values.append(curr_eff)
                                        prev_qdis = curr_qdis
                                        prev_eff = curr_eff
                                if ceff_values:
                                    avg_ceff_values.append(sum(ceff_values)/len(ceff_values)*100)
                            avg_qdis = sum(avg_qdis_values) / len(avg_qdis_values) if avg_qdis_values else 0
                            # Filter out None values for average calculation
                            valid_avg_eff_values = [v for v in avg_eff_values if v is not None]
                            avg_eff = sum(valid_avg_eff_values) / len(valid_avg_eff_values) if valid_avg_eff_values else None
                            avg_cycle_life = sum(avg_cycle_life_values) / len(avg_cycle_life_values) if avg_cycle_life_values else 0
                            avg_areal = sum(avg_areal_capacity_values) / len(avg_areal_capacity_values) if avg_areal_capacity_values else 0
                            avg_reversible = sum(avg_reversible_capacities) / len(avg_reversible_capacities) if avg_reversible_capacities else None
                            avg_ceff = sum(avg_ceff_values) / len(avg_ceff_values) if avg_ceff_values else None
                            group_summary_rows.append([
                                group_name + " (Group Avg)",
                                f"{avg_qdis:.1f}",
                                f"{avg_eff:.1f}%" if avg_eff is not None else "N/A",
                                f"{avg_cycle_life:.0f}",
                                f"{avg_areal:.3f}",
                                f"{avg_reversible:.1f}" if avg_reversible is not None else "N/A",
                                f"{avg_ceff:.2f}%" if avg_ceff is not None else "N/A"
                            ])
                # Add average row for all cells at the end of the PowerPoint table
                # Calculate averages for each parameter
                avg_qdis_values = []
                avg_eff_values = []
                avg_cycle_life_values = []
                avg_areal_capacity_values = []
                avg_reversible_capacities = []
                avg_ceff_values = []
                for d in dfs:
                    df_cell = d['df']
                    first_three_qdis = df_cell['Q Dis (mAh/g)'].head(3).tolist()
                    max_qdis = max(first_three_qdis) if first_three_qdis else None
                    if max_qdis is not None:
                        avg_qdis_values.append(max_qdis)
                    if 'Efficiency (-)' in df_cell.columns and not df_cell['Efficiency (-)'].empty:
                        first_cycle_eff = df_cell['Efficiency (-)'].iloc[0]
                        try:
                            eff_pct = float(first_cycle_eff) * 100
                            avg_eff_values.append(eff_pct)
                        except (ValueError, TypeError):
                            avg_eff_values.append(None) # Handle non-numeric efficiency
                    else:
                        avg_eff_values.append(None) # Handle missing efficiency
                    qdis_series = get_qdis_series(df_cell)
                    cycle_index_series = df_cell[df_cell.columns[0]].iloc[qdis_series.index]
                    cycle_life_80 = calculate_cycle_life_80(qdis_series, cycle_index_series)
                    if cycle_life_80 is not None:
                        avg_cycle_life_values.append(cycle_life_80)
                    areal_capacity, chosen_cycle, diff_pct, eff_val = get_initial_areal_capacity(df_cell, disc_area_cm2)
                    if areal_capacity is not None:
                        avg_areal_capacity_values.append(areal_capacity)
                    formation_cycles = d.get('formation_cycles', 4)
                    if len(df_cell) > formation_cycles:
                        reversible_capacity = df_cell['Q Dis (mAh/g)'].iloc[formation_cycles]
                        avg_reversible_capacities.append(reversible_capacity)
                    # Coulombic Efficiency (post-formation, %)
                    eff_col = 'Efficiency (-)'
                    qdis_col = 'Q Dis (mAh/g)'
                    n_cycles = len(df_cell)
                    ceff_values = []
                    if eff_col in df_cell.columns and qdis_col in df_cell.columns and n_cycles > formation_cycles+1:
                        prev_qdis = df_cell[qdis_col].iloc[formation_cycles]
                        prev_eff = df_cell[eff_col].iloc[formation_cycles]
                        for i in range(formation_cycles+1, n_cycles):
                            curr_qdis = df_cell[qdis_col].iloc[i]
                            curr_eff = df_cell[eff_col].iloc[i]
                            try:
                                pq = float(prev_qdis)
                                cq = float(curr_qdis)
                                pe = float(prev_eff)
                                ce = float(curr_eff)
                                if pq > 0 and (cq < 0.95 * pq or ce < 0.95 * pe):
                                    break
                                ceff_values.append(ce)
                                prev_qdis = cq
                                prev_eff = ce
                            except (ValueError, TypeError):
                                continue
                    if ceff_values:
                        avg_ceff = sum(ceff_values) / len(ceff_values) * 100
                        avg_ceff_values.append(avg_ceff)
                avg_qdis = sum(avg_qdis_values) / len(avg_qdis_values) if avg_qdis_values else 0
                # Filter out None values for average calculation
                valid_avg_eff_values = [v for v in avg_eff_values if v is not None]
                avg_eff = sum(valid_avg_eff_values) / len(valid_avg_eff_values) if valid_avg_eff_values else None
                avg_cycle_life = sum(avg_cycle_life_values) / len(avg_cycle_life_values) if avg_cycle_life_values else 0
                avg_areal = sum(avg_areal_capacity_values) / len(avg_areal_capacity_values) if avg_areal_capacity_values else 0
                avg_reversible = sum(avg_reversible_capacities) / len(avg_reversible_capacities) if avg_reversible_capacities else None
                avg_ceff = sum(avg_ceff_values) / len(avg_ceff_values) if avg_ceff_values else None
                table_data.append([
                    "AVERAGE",
                    f"{avg_qdis:.1f}",
                    f"{avg_eff:.1f}%" if avg_eff is not None else "N/A",
                    f"{avg_cycle_life:.0f}",
                    f"{avg_areal:.3f}",
                    f"{avg_reversible:.1f}" if avg_reversible is not None else "N/A",
                    f"{avg_ceff:.2f}%" if avg_ceff is not None else "N/A"
                ])
                
                # Create table - centered on slide
                rows, cols = len(table_data), len(table_data[0])
                table_width = 7  # inches
                table_height = 2.5  # inches
                # Center the table: (slide width - table width) / 2 = (10 - 7) / 2 = 1.5 inches from left
                table_left = 1.5
                table = slide.shapes.add_table(rows, cols, Inches(table_left), Inches(2.0), Inches(table_width), Inches(table_height)).table
                
                # Populate table
                for i, row_data in enumerate(table_data):
                    for j, cell_data in enumerate(row_data):
                        cell = table.cell(i, j)
                        cell.text = str(cell_data)
                        
                        # Format header row
                        if i == 0:
                            cell.fill.solid()
                            cell.fill.fore_color.rgb = RGBColor(68, 114, 196)
                            for paragraph in cell.text_frame.paragraphs:
                                paragraph.font.color.rgb = RGBColor(255, 255, 255)
                                paragraph.font.bold = True
                                paragraph.font.size = Pt(12)
                        # Format average row
                        elif i == len(table_data) - 1:
                            cell.fill.solid()
                            cell.fill.fore_color.rgb = RGBColor(146, 208, 80)
                            for paragraph in cell.text_frame.paragraphs:
                                paragraph.font.bold = True
                                paragraph.font.size = Pt(12)
                        else:
                            for paragraph in cell.text_frame.paragraphs:
                                paragraph.font.size = Pt(11)
                
                # Create comparison graph
                # Use the same plot as the main display, including average line if toggled, with show_graph_title as set by the user
                fig = plot_capacity_graph(
                    dfs, show_lines, show_efficiency_lines, remove_last_cycle, show_graph_title, experiment_name,
                    show_average_performance, avg_line_toggles, remove_markers, hide_legend,
                    group_a_curve=(group_curves[0][0], group_curves[0][1]) if enable_grouping and group_curves and group_curves[0][0] and group_curves[0][1] else None,
                    group_b_curve=(group_curves[1][0], group_curves[1][1]) if enable_grouping and group_curves and group_curves[1][0] and group_curves[1][1] else None,
                    group_c_curve=(group_curves[2][0], group_curves[2][1]) if enable_grouping and group_curves and group_curves[2][0] and group_curves[2][1] else None,
                    group_a_qchg=(group_curves[0][0], group_curves[0][2]) if enable_grouping and group_curves and group_curves[0][0] and group_curves[0][2] else None,
                    group_b_qchg=(group_curves[1][0], group_curves[1][2]) if enable_grouping and group_curves and group_curves[1][0] and group_curves[1][2] else None,
                    group_c_qchg=(group_curves[2][0], group_curves[2][2]) if enable_grouping and group_curves and group_curves[2][0] and group_curves[2][2] else None,
                    group_a_eff=(group_curves[0][0], group_curves[0][3]) if enable_grouping and group_curves and group_curves[0][0] and group_curves[0][3] else None,
                    group_b_eff=(group_curves[1][0], group_curves[1][3]) if enable_grouping and group_curves and group_curves[1][0] and group_curves[1][3] else None,
                    group_c_eff=(group_curves[2][0], group_curves[2][3]) if enable_grouping and group_curves and group_curves[2][0] and group_curves[2][3] else None,
                    group_names=group_names
                )
                img_bytes = io.BytesIO()
                fig.savefig(img_bytes, format='png', bbox_inches='tight', dpi=300)
                plt.close(fig)
                img_bytes.seek(0)
                with tempfile.NamedTemporaryFile(delete=False, suffix='.png') as tmp_img:
                    tmp_img.write(img_bytes.read())
                    img_path = tmp_img.name
                # Set graph width to 11cm and calculate height to maintain aspect ratio
                slide_width_cm = 11  # cm
                slide_width_inches = slide_width_cm / 2.54  # convert cm to inches
                aspect_ratio = 8/5  # width/height from figsize
                slide_height_inches = slide_width_inches / aspect_ratio
                slide.shapes.add_picture(img_path, Inches(14/2.54), Inches(12/2.54), width=Inches(slide_width_inches), height=Inches(slide_height_inches))
                
                # Download button
                pptx_bytes = io.BytesIO()
                prs.save(pptx_bytes)
                pptx_bytes.seek(0)
                # Update PowerPoint file name to include experiment name
                if experiment_name:
                    pptx_file_name = f'{experiment_name} Cell Comparison Summary.pptx'
                else:
                    pptx_file_name = 'Cell Comparison Summary.pptx'
                st.download_button(f'Download PowerPoint: {pptx_file_name}', data=pptx_bytes, file_name=pptx_file_name, mime='application/vnd.openxmlformats-officedocument.presentationml.presentation', key='download_pptx_multiple')

            # --- Excel Export ---
            output = io.BytesIO()
            with pd.ExcelWriter(output, engine='openpyxl') as writer:
                for d in dfs:
                    df_cell = d['df']
                    sheet_name = d['testnum'] if d['testnum'] else f'Cell {dfs.index(d)+1}'
                    df_cell.to_excel(writer, index=False, startrow=4, sheet_name=sheet_name)
            output.seek(0)

            # Add summary values and native Excel chart to each sheet
            from openpyxl import load_workbook
            from openpyxl.chart import LineChart, Reference
            output.seek(0)
            wb = load_workbook(output)
            for d in dfs:
                ws = wb[d['testnum'] if d['testnum'] else f'Cell {dfs.index(d)+1}']
                df_cell = d['df']
                ws['A1'] = 'Total loading (mg)'
                ws['B1'] = d['loading']
                ws['A2'] = 'Active material loading (mg)'
                ws['B2'] = d['loading'] * (d['active'] / 100)
                first_three_qdis = df_cell['Q Dis (mAh/g)'].head(3).tolist()
                max_qdis = max(first_three_qdis) if first_three_qdis else None
                if 'Efficiency (-)' in df_cell.columns and not df_cell['Efficiency (-)'].empty:
                    first_cycle_eff = df_cell['Efficiency (-)'].iloc[0]
                    try:
                        eff_pct = float(first_cycle_eff) * 100
                    except (ValueError, TypeError):
                        eff_pct = None
                else:
                    eff_pct = None
                qdis_series = get_qdis_series(df_cell)
                cycle_index_series = df_cell[df_cell.columns[0]].iloc[qdis_series.index]
                cycle_life_80 = calculate_cycle_life_80(qdis_series, cycle_index_series)
                ws['P1'] = '1st Cycle Discharge Capacity (mAh/g)'
                ws['Q1'] = f"{max_qdis:.1f}" if isinstance(max_qdis, (int, float)) and max_qdis is not None else ''
                ws['P2'] = 'First Cycle Efficiency (%)'
                ws['Q2'] = f"{eff_pct:.1f}" if isinstance(eff_pct, (int, float)) and eff_pct is not None else ''
                ws['P3'] = 'Cycle Life (80%)'
                ws['Q3'] = cycle_life_80 if cycle_life_80 is not None else ''
                data_start_row = 6
                data_end_row = data_start_row + len(df_cell) - 1
                headers = list(df_cell.columns)
                x_col_idx = 1
                q_dis_col = headers.index('Q Dis (mAh/g)') + 1
                q_chg_col = headers.index('Q Chg (mAh/g)') + 1
                chart = LineChart()
                chart.title = 'Gravimetric Capacity vs. ' + df_cell.columns[0]
                chart.y_axis.title = 'Capacity (mAh/g)'
                chart.x_axis.title = df_cell.columns[0]
                chart.x_axis.majorTickMark = 'in'
                chart.y_axis.majorTickMark = 'in'
                chart.width = 20
                chart.height = 12
                data = Reference(ws, min_col=min(q_dis_col, q_chg_col), max_col=max(q_dis_col, q_chg_col), min_row=data_start_row-1, max_row=data_end_row)
                chart.add_data(data, titles_from_data=True)
                cats = Reference(ws, min_col=x_col_idx, min_row=data_start_row, max_row=data_end_row)
                chart.set_categories(cats)
                ws.add_chart(chart, 'S5')
            output2 = io.BytesIO()
            wb.save(output2)
            output2.seek(0)

            # Update file name to include experiment name
            if experiment_name:
                if len(dfs) > 1:
                    file_name = f'{experiment_name} Comparison Cycling data.xlsx'
                else:
                    file_name = f'{experiment_name} {datasets[0]["testnum"]} Cycling data.xlsx' if datasets[0]["testnum"] else f'{experiment_name} Cycling data.xlsx'
            else:
                file_name = 'Comparison Cycling data.xlsx' if len(dfs) > 1 else f'{datasets[0]["testnum"]} Cycling data.xlsx' if datasets[0]["testnum"] else 'Cycling data.xlsx'
            
            st.success('Processing complete!')

            # Add summary worksheet if multiple cells
            if len(dfs) > 1:
                summary_sheet_name = f"{experiment_name} Summary" if experiment_name else "Summary"
                summary_sheet_name = summary_sheet_name[:31]  # Excel sheet names limited to 31 characters
                ws_summary = wb.create_sheet(summary_sheet_name)
                # Move summary sheet to the first position (leftmost tab)
                wb._sheets.insert(0, wb._sheets.pop(wb._sheets.index(ws_summary)))
                # Write headers
                ws_summary['A1'] = 'Cell'
                ws_summary['B1'] = '1st Cycle Discharge Capacity (mAh/g)'
                ws_summary['C1'] = 'First Cycle Efficiency (%)'
                ws_summary['D1'] = 'Cycle Life (80%)'
                ws_summary['E1'] = 'Initial Areal Capacity (mAh/cm²)'
                ws_summary['F1'] = 'Reversible Capacity (mAh/g)'
                ws_summary['G1'] = 'Coulombic Efficiency (post-formation)'
                # Write data for each cell
                for i, d in enumerate(dfs):
                    df_cell = d['df']
                    cell_name = d['testnum'] if d['testnum'] else f'Cell {i+1}'
                    # Calculate summary values (same as in your summary table)
                    first_three_qdis = df_cell['Q Dis (mAh/g)'].head(3).tolist()
                    max_qdis = max(first_three_qdis) if first_three_qdis else None
                    qdis_str = f"{max_qdis:.1f}" if isinstance(max_qdis, (int, float)) else "N/A"
                    if 'Efficiency (-)' in df_cell.columns and not df_cell['Efficiency (-)'].empty:
                        first_cycle_eff = df_cell['Efficiency (-)'].iloc[0]
                        try:
                            eff_pct = float(first_cycle_eff) * 100
                            eff_str = f"{eff_pct:.1f}%"
                        except (ValueError, TypeError):
                            eff_str = "N/A"
                    else:
                        eff_str = "N/A"
                    # Cycle Life (80%)
                    qdis_series = get_qdis_series(df_cell)
                    cycle_index_series = df_cell[df_cell.columns[0]].iloc[qdis_series.index]
                    cycle_life_80 = calculate_cycle_life_80(qdis_series, cycle_index_series)
                    cycle_life_str = str(cycle_life_80) if isinstance(cycle_life_80, (int, float)) else "N/A"
                    # Initial Areal Capacity (mAh/cm²)
                    areal_capacity, chosen_cycle, diff_pct, eff_val = get_initial_areal_capacity(df_cell, disc_area_cm2)
                    if areal_capacity is None:
                        areal_str = "N/A"
                    else:
                        areal_str = f"{areal_capacity:.3f}"
                    # Reversible Capacity (mAh/g)
                    formation_cycles = d.get('formation_cycles', 4)
                    if len(df_cell) > formation_cycles:
                        reversible_capacity = df_cell['Q Dis (mAh/g)'].iloc[formation_cycles]
                        reversible_str = f"{reversible_capacity:.1f}"
                    else:
                        reversible_str = "N/A"
                    # Coulombic Efficiency (post-formation, %)
                    eff_col = 'Efficiency (-)'
                    qdis_col = 'Q Dis (mAh/g)'
                    n_cycles = len(df_cell)
                    ceff_values = []
                    if eff_col in df_cell.columns and qdis_col in df_cell.columns and n_cycles > formation_cycles+1:
                        prev_qdis = df_cell[qdis_col].iloc[formation_cycles]
                        prev_eff = df_cell[eff_col].iloc[formation_cycles]
                        for j in range(formation_cycles+1, n_cycles):
                            curr_qdis = df_cell[qdis_col].iloc[j]
                            curr_eff = df_cell[eff_col].iloc[j]
                            try:
                                pq = float(prev_qdis)
                                cq = float(curr_qdis)
                                pe = float(prev_eff)
                                ce = float(curr_eff)
                                if pq > 0 and (cq < 0.95 * pq or ce < 0.95 * pe):
                                    break
                                ceff_values.append(ce)
                                prev_qdis = cq
                                prev_eff = ce
                            except (ValueError, TypeError):
                                # Skip this cycle if any value is not numeric
                                continue
                    if ceff_values:
                        ceff_str = f"{sum(ceff_values)/len(ceff_values)*100:.2f}%"
                    else:
                        ceff_str = "N/A"
                    ws_summary[f'A{i+2}'] = cell_name
                    ws_summary[f'B{i+2}'] = qdis_str
                    ws_summary[f'C{i+2}'] = eff_str
                    ws_summary[f'D{i+2}'] = cycle_life_str
                    ws_summary[f'E{i+2}'] = areal_str
                    ws_summary[f'F{i+2}'] = reversible_str
                    ws_summary[f'G{i+2}'] = ceff_str
                # Optionally, add group/average rows as well
                if enable_grouping and group_assignments is not None:
                    for group_idx, group_name in enumerate(group_names):
                        group_dfs = [df for df, g in zip(dfs, group_assignments) if g == group_name]
                        if len(group_dfs) > 1:
                            avg_qdis_values = []
                            avg_eff_values = []
                            avg_cycle_life_values = []
                            avg_areal_capacity_values = []
                            avg_reversible_capacities = []
                            avg_ceff_values = []
                            for d in group_dfs:
                                df_cell = d['df']
                                first_three_qdis = df_cell['Q Dis (mAh/g)'].head(3).tolist()
                                max_qdis = max(first_three_qdis) if first_three_qdis else None
                                if max_qdis is not None:
                                    avg_qdis_values.append(max_qdis)
                                if 'Efficiency (-)' in df_cell.columns and not df_cell['Efficiency (-)'].empty:
                                    first_cycle_eff = df_cell['Efficiency (-)'].iloc[0]
                                    try:
                                        eff_pct = float(first_cycle_eff) * 100
                                        avg_eff_values.append(eff_pct)
                                    except (ValueError, TypeError):
                                        avg_eff_values.append(None) # Handle non-numeric efficiency
                                else:
                                    avg_eff_values.append(None) # Handle missing efficiency
                                try:
                                    qdis_series = get_qdis_series(df_cell)
                                    cycle_index_series = df_cell[df_cell.columns[0]].iloc[qdis_series.index]
                                    cycle_life_80 = calculate_cycle_life_80(qdis_series, cycle_index_series)
                                    avg_cycle_life_values.append(cycle_life_80)
                                except Exception:
                                    pass
                                areal_capacity, chosen_cycle, diff_pct, eff_val = get_initial_areal_capacity(df_cell, disc_area_cm2)
                                if areal_capacity is not None:
                                    avg_areal_capacity_values.append(areal_capacity)
                                formation_cycles = d.get('formation_cycles', 4)
                                if len(df_cell) > formation_cycles:
                                    reversible_capacity = df_cell['Q Dis (mAh/g)'].iloc[formation_cycles]
                                    avg_reversible_capacities.append(reversible_capacity)
                                # Coulombic Efficiency (post-formation, %)
                                eff_col = 'Efficiency (-)'
                                qdis_col = 'Q Dis (mAh/g)'
                                n_cycles = len(df_cell)
                                ceff_values = []
                                if eff_col in df_cell.columns and qdis_col in df_cell.columns and n_cycles > formation_cycles+1:
                                    prev_qdis = df_cell[qdis_col].iloc[formation_cycles]
                                    prev_eff = df_cell[eff_col].iloc[formation_cycles]
                                    for j in range(formation_cycles+1, n_cycles):
                                        curr_qdis = df_cell[qdis_col].iloc[j]
                                        curr_eff = df_cell[eff_col].iloc[j]
                                        if prev_qdis > 0 and (curr_qdis < 0.95 * prev_qdis or curr_eff < 0.95 * prev_eff):
                                            break
                                        ceff_values.append(curr_eff)
                                        prev_qdis = curr_qdis
                                        prev_eff = curr_eff
                                if ceff_values:
                                    avg_ceff_values.append(sum(ceff_values)/len(ceff_values)*100)
                            avg_qdis = sum(avg_qdis_values) / len(avg_qdis_values) if avg_qdis_values else 0
                            # Filter out None values for average calculation
                            valid_avg_eff_values = [v for v in avg_eff_values if v is not None]
                            avg_eff = sum(valid_avg_eff_values) / len(valid_avg_eff_values) if valid_avg_eff_values else None
                            avg_cycle_life = sum(avg_cycle_life_values) / len(avg_cycle_life_values) if avg_cycle_life_values else 0
                            avg_areal = sum(avg_areal_capacity_values) / len(avg_areal_capacity_values) if avg_areal_capacity_values else 0
                            avg_reversible = sum(avg_reversible_capacities) / len(avg_reversible_capacities) if avg_reversible_capacities else None
                            avg_ceff = sum(avg_ceff_values) / len(avg_ceff_values) if avg_ceff_values else None
                            ws_summary.append([
                                group_name + " (Group Avg)",
                                f"{avg_qdis:.1f}",
                                f"{avg_eff:.1f}%" if avg_eff is not None else "N/A",
                                f"{avg_cycle_life:.0f}",
                                f"{avg_areal:.3f}",
                                f"{avg_reversible:.1f}" if avg_reversible is not None else "N/A",
                                f"{avg_ceff:.2f}%" if avg_ceff is not None else "N/A"
                            ])
            # Save the updated workbook
            wb.save(output2)
            output2.seek(0)

            # Update file name to include experiment name
            if experiment_name:
                if len(dfs) > 1:
                    file_name = f'{experiment_name} Comparison Cycling data.xlsx'
                else:
                    file_name = f'{experiment_name} {datasets[0]["testnum"]} Cycling data.xlsx' if datasets[0]["testnum"] else f'{experiment_name} Cycling data.xlsx'
            else:
                file_name = 'Comparison Cycling data.xlsx' if len(dfs) > 1 else f'{datasets[0]["testnum"]} Cycling data.xlsx' if datasets[0]["testnum"] else 'Cycling data.xlsx'
            
            st.success('Processing complete!')
            st.download_button(f'Download XLSX: {file_name}', data=output2, file_name=file_name, mime='application/vnd.openxmlformats-officedocument.spreadsheetml.sheet', key='download_xlsx_final')

else:
    st.info('Please upload a file and enter valid disc loading and % active material.')

# --- Master Table Tab ---
if tab_master and current_project_id:
    with tab_master:
        st.header("📋 Master Table")
        current_project_name = st.session_state.get('current_project_name', 'Selected Project')
        st.markdown(f"**Project:** {current_project_name}")
        st.markdown("---")
        
        # Get all experiments data for this project
        all_experiments_data = get_all_project_experiments_data(current_project_id)
        
        if not all_experiments_data:
            st.info("📊 No experiments found in this project. Create experiments to see master table data.")
        else:
            # Process experiment data
            experiment_summaries = []
            individual_cells = []
            
            for exp_data in all_experiments_data:
                exp_id, exp_name, file_name, loading, active_material, formation_cycles, test_number, electrolyte, formulation_json, data_json, created_date = exp_data
                
                try:
                    parsed_data = json.loads(data_json)
                    
                    # Check if this is a multi-cell experiment or single cell
                    if 'cells' in parsed_data:
                        # Multi-cell experiment
                        cells_data = parsed_data['cells']
                        disc_diameter = parsed_data.get('disc_diameter_mm', 15)
                        disc_area_cm2 = np.pi * (disc_diameter / 2 / 10) ** 2
                        
                        experiment_cells = []
                        for cell_data in cells_data:
                            try:
                                df = pd.read_json(cell_data['data_json'])
                                cell_summary = calculate_cell_summary(df, cell_data, disc_area_cm2)
                                cell_summary['experiment_name'] = exp_name
                                cell_summary['experiment_date'] = parsed_data.get('experiment_date', created_date)
                                experiment_cells.append(cell_summary)
                                individual_cells.append(cell_summary)
                            except Exception as e:
                                continue
                        
                        # Calculate experiment average
                        if experiment_cells:
                            exp_summary = calculate_experiment_average(experiment_cells, exp_name, parsed_data.get('experiment_date', created_date))
                            experiment_summaries.append(exp_summary)
                    
                    else:
                        # Legacy single cell experiment
                        df = pd.read_json(data_json)
                        cell_summary = calculate_cell_summary(df, {
                            'cell_name': test_number or exp_name,
                            'loading': loading,
                            'active_material': active_material,
                            'formation_cycles': formation_cycles,
                            'test_number': test_number
                        }, np.pi * (15 / 2 / 10) ** 2)  # Default disc size
                        cell_summary['experiment_name'] = exp_name
                        cell_summary['experiment_date'] = created_date
                        individual_cells.append(cell_summary)
                        
                        # Also add as experiment summary (since it's a single cell)
                        exp_summary = cell_summary.copy()
                        exp_summary['cell_name'] = f"{exp_name} (Single Cell)"
                        experiment_summaries.append(exp_summary)
                        
                except Exception as e:
                    st.error(f"Error processing experiment {exp_name}: {str(e)}")
                    continue
            
            # Section 1: Average Cell Data per Experiment
            with st.expander("### 📊 Section 1: Average Cell Data per Experiment", expanded=True):
                if experiment_summaries:
                    display_experiment_summaries_table(experiment_summaries)
                else:
                    st.info("No experiment summary data available.")
            
            st.markdown("---")
            
            # Section 2: All Individual Cells Data
            with st.expander("### 🧪 Section 2: All Individual Cells Data", expanded=False):
                if individual_cells:
                    display_individual_cells_table(individual_cells)
                else:
                    st.info("No individual cell data available.")
            
            st.markdown("---")
            
            # Section 3: Best Performing Cells Analysis
            with st.expander("### 🏅 Section 3: Best Performing Cells Analysis", expanded=True):
                display_best_performers_analysis(individual_cells)

# --- Data Preprocessing Section ---